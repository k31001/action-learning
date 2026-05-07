"""
핵심전략 11개 카테고리화 발표자료 (6슬라이드)

작업 2 산출물 — `presentation/core-strategies-categorization.pptx`

11개 전략을 4개 축으로 분류해 한 슬라이드씩 시각화 + SE-3 deep dive:
  슬라이드 1  시나리오별        — 5개 시나리오에서 어떤 전략이 작동하는지 매트릭스
  슬라이드 2  기대효과별        — 매출/마진/신시장/헤지 4개 카테고리
  슬라이드 3  실행부서별        — 책임 부서 매핑
  슬라이드 4  시간축별          — 단기(0~6M) / 중기(6~18M) / 장기(18M+) 타임라인
  슬라이드 5  부록: 현황 요약   — 11개 전략 정성·정량 현황 + 다음 마일스톤
  슬라이드 6  SE-3 Deep Dive    — AI 가치사슬 마진 분포 + 3-Tier 수직 진출 헷지 모델 (신규)

기존 `generate_pptx.py`의 디자인 시스템(THEME, helper 함수)을 재사용한다.
"""

import os
import sys

# 동일 디렉토리의 generate_pptx.py에서 디자인 시스템과 helper 함수 import
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from generate_pptx import (
    THEME, FONT_KO, FONT_EN, SLIDE_W, SLIDE_H, TEMPLATE,
    add_text, add_rect, add_line, add_footer, add_header, add_so_what,
    remove_existing_slides,
)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
OUTPUT = os.path.join(ROOT, 'presentation', 'core-strategies-categorization.pptx')

TOTAL = 6

# ============================================================
# 핵심전략 10개 데이터 (단일 소스)
# ============================================================

STRATEGIES = [
    # (id, name, type, color_hint)
    ('MB-4', '커스텀 AI 메모리 솔루션',           'main',  'samsung_blue'),
    ('RS-3', '고객특화·전환비용 (CMX/SCADA/FDP)', 'main',  'samsung_blue'),
    ('RS-6', '공정 리더십 통합',                  'main',  'samsung_blue'),
    ('MB-2', '동서 균형 공급망',                  'main',  'samsung_blue'),
    ('SD-1', 'HBM 조직 독립 P&L',                'main',  'samsung_blue'),
    ('RS-5', '재무 규율 + 초과이익 재투자',       'main',  'samsung_blue'),
    ('SA-2', '일본 R&D 허브 (EUV 우회 NIL)',     'side',  'amber'),
    ('SD-2', '산업용 AI 메모리 (자동차·의료)',   'side',  'amber'),
    ('SE-1', '3D DRAM + IMEC + M&A',             'side',  'amber'),
    ('SE-2', 'CXL SIG 표준 주도',                 'side',  'amber'),
    ('SE-3', 'AI 인프라 수직 진출 (Vertical Ascent)', 'side', 'amber'),
]


def _strategy_color(strat_type):
    return THEME['samsung_blue'] if strat_type == 'main' else THEME['amber']


# ============================================================
# 슬라이드 1: 시나리오별 카테고리화
# ============================================================

def build_slide_1_scenario(prs):
    """시나리오 A~E 각각에서 어느 전략이 어떤 강도로 작동하는지 매트릭스."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '카테고리 1 · 시나리오별 — 어떤 미래에 어떤 전략이 작동하는가',
               '5개 시나리오 × 10개 전략 매트릭스. ●●● 핵심 가치 / ●● 의미 있는 가치 / ● 부분 가치 / − 제한적')

    # 매트릭스: 행=10개 전략, 열=5개 시나리오
    # 강도: '●●●' = 3, '●●' = 2, '●' = 1, '−' = 0
    intensity = {
        # (id, [A, B, C, D, E])
        'MB-4': ['●●', '●●●', '●', '●', '●●●'],
        'RS-3': ['●●●', '●●●', '●●', '●●●', '●●●'],
        'RS-6': ['●●●', '●●●', '●●●', '●●●', '●●'],
        'MB-2': ['●●', '●●●', '●', '●●●', '●●'],
        'SD-1': ['●●', '●●', '●●●', '●●●', '●●'],
        'RS-5': ['●●●', '●●', '●●●', '●●●', '●●●'],
        'SA-2': ['●●●', '●●', '●●●', '●●', '●●●'],
        'SD-2': ['●', '●●', '●●', '●●●', '●●●'],
        'SE-1': ['●●', '●●', '●●●', '●●●', '●●●'],
        'SE-2': ['●●', '●●●', '●●', '●●●', '●●●'],
        'SE-3': ['●●●', '●●●', '●', '●●', '●●●'],
    }

    scenarios = [
        ('A', '황금 요새', 'AI 지속+디커플링'),
        ('B ★', 'AI 르네상스', 'AI 지속+공존 (메인)'),
        ('C', '기술 냉전', 'AI 붕괴+디커플링'),
        ('D', '조용한 재편', 'AI 붕괴+공존'),
        ('E', '패러다임 전환', 'HBM 대체'),
    ]

    # 표 구조
    table_x = Inches(0.4)
    table_y = Inches(1.85)
    name_col_w = Inches(3.4)
    scen_col_w = Inches(1.85)

    # 헤더 행
    add_rect(slide, table_x, table_y, name_col_w, Inches(0.7), fill=THEME['samsung_blue'])
    add_text(slide, table_x + Inches(0.15), table_y + Inches(0.1), name_col_w - Inches(0.3), Inches(0.5),
             '핵심전략 (메인 6 + 사이드 4)',
             font=FONT_KO, size=10.5, bold=True, color=THEME['white'])

    for i, (code, name, desc) in enumerate(scenarios):
        x = table_x + name_col_w + Inches(i * 1.85)
        bg = THEME['amber'] if '★' in code else THEME['samsung_blue']
        add_rect(slide, x, table_y, scen_col_w, Inches(0.7), fill=bg)
        add_text(slide, x + Inches(0.05), table_y + Inches(0.05), scen_col_w - Inches(0.1), Inches(0.3),
                 f'{code} · {name}',
                 font=FONT_KO, size=9.5, bold=True, color=THEME['white'], align='center')
        add_text(slide, x + Inches(0.05), table_y + Inches(0.36), scen_col_w - Inches(0.1), Inches(0.3),
                 desc, font=FONT_KO, size=8, italic=True, color=THEME['white'], align='center')

    # 데이터 행
    row_h = Inches(0.39)
    for i, (sid, name, stype, _) in enumerate(STRATEGIES):
        y = table_y + Inches(0.7) + row_h * i
        bg = THEME['amber_light'] if stype == 'main' else THEME['soft_blue_bg']
        # 전략명 셀
        add_rect(slide, table_x, y, name_col_w, row_h, fill=bg, line=THEME['light_gray'])
        # 좌측 색상 띠
        add_rect(slide, table_x, y, Inches(0.08), row_h, fill=_strategy_color(stype))
        # 코드
        add_text(slide, table_x + Inches(0.15), y + Inches(0.04), Inches(0.6), Inches(0.32),
                 sid, font=FONT_EN, size=10, bold=True, color=_strategy_color(stype))
        # 이름
        add_text(slide, table_x + Inches(0.75), y + Inches(0.04), name_col_w - Inches(0.85), Inches(0.32),
                 name, font=FONT_KO, size=9, color=THEME['dark_text'])

        # 시나리오 셀
        for j, mark in enumerate(intensity[sid]):
            cx = table_x + name_col_w + Inches(j * 1.85)
            add_rect(slide, cx, y, scen_col_w, row_h, fill=THEME['white'], line=THEME['light_gray'])
            color = THEME['samsung_blue'] if mark == '●●●' else (
                THEME['amber'] if mark == '●●' else (
                    THEME['gray_caption'] if mark == '●' else THEME['light_gray']))
            add_text(slide, cx, y + Inches(0.05), scen_col_w, Inches(0.3),
                     mark, font=FONT_EN, size=14, bold=True, color=color, align='center')

    add_so_what(slide, Inches(6.95),
                'RS-3·RS-5·RS-6·SE-3은 거의 모든 시나리오에서 ●● 이상 — 중심축. SE-3는 가치사슬 마진 헷지를 추가.')
    add_footer(slide, 1, TOTAL, '카테고리 1 · 시나리오')
    return slide


# ============================================================
# 슬라이드 2: 기대효과별 카테고리화
# ============================================================

def build_slide_2_impact(prs):
    """4개 카테고리: 매출/점유율, 비용/마진, 신시장, 리스크 헤지."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '카테고리 2 · 기대효과별 — 어떤 가치를 만드는가',
               '매출 견인 / 마진 방어 / 신시장 개척 / 리스크 헤지 — 10개 전략의 1차 효과 분류')

    # 4개 카테고리 카드 (2x2)
    categories = [
        {
            'title': '매출 증대 · 점유율 회복',
            'subtitle': 'Revenue & Market Share',
            'color': THEME['samsung_blue'],
            'items': [
                ('MB-4', '커스텀 AI 메모리 솔루션',
                 '하이퍼스케일러 ASIC 통합, 매출 비중 20%+ (2029)'),
                ('RS-3', '고객특화·전환비용',
                 'NVIDIA 3대 통합, 통합 매출 $8~9B/년 (2030)'),
                ('MB-2', '동서 균형 공급망',
                 '신흥시장 매출 3배+ (2028, vs 2025)'),
                ('SD-1', 'HBM 조직 독립 P&L',
                 'HBM 사업부 영업이익률 35%+ 회복'),
            ],
        },
        {
            'title': '비용 절감 · 마진 확대',
            'subtitle': 'Cost & Margin Defense',
            'color': THEME['green_pos'],
            'items': [
                ('RS-6', '공정 리더십 통합',
                 '1c nm 원가 -30%, NAND capex 회피 1.5~2조 원'),
                ('RS-5', '재무 규율 + 재투자',
                 '다운턴 흑자 구조, HBM 가격 60% 하락에도 흑자'),
            ],
        },
        {
            'title': '신규 시장 개척',
            'subtitle': 'New Market Expansion',
            'color': THEME['amber'],
            'items': [
                ('SA-2', '일본 R&D 허브 (EUV 우회 NIL)',
                 'ASML 의존 탈피, 한일 소재·장비 생태계 신시장'),
                ('SD-2', '산업용 AI 메모리 (자동차·의료)',
                 'AEC-Q100 등급, 의료·자율주행·로봇 — 사이클 안정 신시장'),
            ],
        },
        {
            'title': '리스크 헤지 · 패러다임 대비',
            'subtitle': 'Risk Hedge & Paradigm Shift',
            'color': THEME['red_alert'],
            'items': [
                ('SE-1', '3D DRAM + IMEC + M&A',
                 'HBM 패러다임 전환 시 피벗 자원 (5천억 펀드)'),
                ('SE-2', 'CXL SIG 표준 주도',
                 '차세대 메모리 패브릭 표준 결정권'),
                ('SE-3', 'AI 인프라 수직 진출 (Vertical Ascent)',
                 '가치사슬 마진 헷지 — 3-Tier 모델 (지분/AI Factory/자체 캠퍼스)'),
            ],
        },
    ]

    card_w = Inches(6.3)
    card_h = Inches(2.35)
    gap = Inches(0.13)
    base_x = Inches(0.4)
    base_y = Inches(1.85)

    for i, cat in enumerate(categories):
        col = i % 2
        row = i // 2
        x = base_x + (card_w + gap) * col
        y = base_y + (card_h + gap) * row

        # 카드 배경
        add_rect(slide, x, y, card_w, card_h, fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        # 좌측 색상 띠
        add_rect(slide, x, y, Inches(0.10), card_h, fill=cat['color'])
        # 제목
        add_text(slide, x + Inches(0.25), y + Inches(0.08), card_w - Inches(0.35), Inches(0.35),
                 cat['title'], font=FONT_KO, size=12, bold=True, color=cat['color'])
        # 영문 부제
        add_text(slide, x + Inches(0.25), y + Inches(0.42), card_w - Inches(0.35), Inches(0.25),
                 cat['subtitle'], font=FONT_EN, size=8.5, italic=True,
                 color=THEME['gray_caption'])

        # 항목
        for j, (sid, name, desc) in enumerate(cat['items']):
            iy = y + Inches(0.75) + Inches(j * 0.38)
            add_text(slide, x + Inches(0.25), iy, Inches(0.7), Inches(0.3),
                     sid, font=FONT_EN, size=10, bold=True, color=cat['color'])
            add_text(slide, x + Inches(0.95), iy, Inches(2.6), Inches(0.3),
                     name, font=FONT_KO, size=9.5, bold=True, color=THEME['dark_text'])
            add_text(slide, x + Inches(0.95), iy + Inches(0.18), card_w - Inches(1.10), Inches(0.2),
                     desc, font=FONT_KO, size=8, italic=True,
                     color=THEME['gray_caption'])

    add_so_what(slide, Inches(6.55),
                '메인벳 6개는 매출+마진에 집중, 사이드벳 4개는 신시장+패러다임에 분산 — 자원 배분 균형')
    add_footer(slide, 2, TOTAL, '카테고리 2 · 기대효과')
    return slide


# ============================================================
# 슬라이드 3: 실행부서별 카테고리화
# ============================================================

def build_slide_3_owner(prs):
    """책임 부서/조직 매핑."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '카테고리 3 · 실행부서별 — 누가 책임지는가',
               '메모리사업부 R&D / 영업·전략 / 재무·이사회 / 일본 R&D / 표준·M&A — 5개 책임 그룹 매핑')

    # 5개 부서 그룹 (수직 배치)
    departments = [
        {
            'name': '메모리사업부 R&D',
            'subtitle': 'DS Memory R&D · 평택·기흥',
            'color': THEME['samsung_blue'],
            'strategies': [
                ('MB-4', '커스텀 AI 메모리', 'HBM 베이스다이 커스텀 로직, CXL 모듈 R&D'),
                ('RS-3', 'NVIDIA 통합 (CMX/SCADA/FDP)', 'AI SSD 컨트롤러 펌웨어 + Co-Validation'),
                ('RS-6', '공정 리더십 통합', '1c nm DRAM + NAND 4트랙 R&D + Hybrid bonding IP'),
                ('SD-1', 'HBM 조직 독립 P&L', 'HBM 전용 R&D 분리, 패키징 100인+ 채용'),
            ],
        },
        {
            'name': '영업 · 전략기획 · 글로벌 사업',
            'subtitle': 'Sales · Strategy · Global Business',
            'color': THEME['amber'],
            'strategies': [
                ('MB-2', '동서 균형 공급망',
                 '신흥시장 영업팀 신설 (사우디·UAE·인도), 시안 라이선스 협상'),
                ('SD-2', '산업용 AI 메모리 (자동차·의료)',
                 '자동차·의료 OEM 신규 영업, AEC-Q100 인증'),
            ],
        },
        {
            'name': '재무 · 전략 · 이사회',
            'subtitle': 'CFO · Strategy · Board',
            'color': THEME['green_pos'],
            'strategies': [
                ('RS-5', '재무 규율 + 초과이익 재투자',
                 '재고일수 상한, FCF 기준, 다운사이클 capex 하한 4조/년 이사회 결의'),
            ],
        },
        {
            'name': '일본 R&D 허브 (요코하마)',
            'subtitle': 'Japan R&D Hub',
            'color': RGBColor(0xE3, 0x4A, 0x4A),
            'strategies': [
                ('SA-2', '일본 R&D 허브 (EUV 우회 NIL)',
                 'JSR·신에쓰화학·캐논 파트너십, 나노임프린트 NIL 공동 개발 (1조 원)'),
            ],
        },
        {
            'name': '표준 컨소시엄 · M&A · IP',
            'subtitle': 'Standards · M&A · IP Strategy',
            'color': RGBColor(0x6B, 0x46, 0xC1),
            'strategies': [
                ('SE-1', '3D DRAM + IMEC + 스타트업 M&A',
                 'IMEC 공동 연구, 3D DRAM 스타트업 2~3개 M&A (5천억 펀드)'),
                ('SE-2', 'CXL SIG 표준 주도',
                 'CXL 4.0+ 워킹그룹 인력 10→25명 확대, 자사 아키텍처 표준 반영'),
            ],
        },
        {
            'name': '그룹 컨소시엄 (수직 진출)',
            'subtitle': 'Group Consortium · DS+C&T+SDS+Heavy',
            'color': RGBColor(0x0A, 0x7E, 0x8C),
            'strategies': [
                ('SE-3', 'AI 인프라 수직 진출 (Vertical Ascent)',
                 'Stargate Korea LOI→DA 본계약, 4사 운영 협약, 부유식 DC, GPUaaS 확장, neocloud equity-for-supply swap'),
            ],
        },
    ]

    base_y = Inches(1.85)
    row_h = Inches(0.78)
    x = Inches(0.4)
    w = Inches(12.5)

    for i, dept in enumerate(departments):
        y = base_y + row_h * i
        # 카드 배경
        add_rect(slide, x, y, w, row_h - Inches(0.05), fill=THEME['soft_white_bg'],
                 line=THEME['light_gray'])
        # 좌측 색상 띠
        add_rect(slide, x, y, Inches(0.10), row_h - Inches(0.05), fill=dept['color'])
        # 부서명
        add_text(slide, x + Inches(0.25), y + Inches(0.08), Inches(3.5), Inches(0.32),
                 dept['name'], font=FONT_KO, size=11, bold=True, color=dept['color'])
        # 영문
        add_text(slide, x + Inches(0.25), y + Inches(0.42), Inches(3.5), Inches(0.25),
                 dept['subtitle'], font=FONT_EN, size=8.5, italic=True,
                 color=THEME['gray_caption'])

        # 전략 항목 (가로 배치)
        n = len(dept['strategies'])
        item_w = Inches(8.7 / max(n, 1))
        for j, (sid, name, action) in enumerate(dept['strategies']):
            ix = x + Inches(3.85) + item_w * j
            iy = y + Inches(0.08)
            # 코드 + 이름
            add_text(slide, ix, iy, item_w - Inches(0.15), Inches(0.28),
                     f'{sid} · {name}', font=FONT_KO, size=9.5, bold=True,
                     color=dept['color'])
            # 행동
            add_text(slide, ix, iy + Inches(0.30), item_w - Inches(0.15), Inches(0.45),
                     action, font=FONT_KO, size=8.5,
                     color=THEME['dark_text'], line_spacing=1.3)

    add_so_what(slide, Inches(6.65),
                '메모리사업부 R&D가 4개 전략의 책임 — AI 자동화로 잉여 자원 확보 없이는 동시 실행 불가')
    add_footer(slide, 3, TOTAL, '카테고리 3 · 실행부서')
    return slide


# ============================================================
# 슬라이드 4: 시간축별 카테고리화
# ============================================================

def build_slide_4_timeline(prs):
    """단기 / 중기 / 장기 타임라인 시각화."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '카테고리 4 · 시간축별 — 언제 시작하고 언제 효과를 보는가',
               '단기 (0~6M, 2026 H2) / 중기 (6~18M, 2027) / 장기 (18M+, 2028+) — 착수 시점 + 효과 발현 시점')

    # 3개 시간 구간 (수직 배치, 각 구간은 가로 카드)
    phases = [
        {
            'name': '단기',
            'period': '0~6M',
            'window': '2026 H2 ~ 2027 H1',
            'desc': '즉시 착수 — 거버넌스·조직·계약 구조',
            'color': THEME['red_alert'],
            'strategies': [
                ('RS-5', '재무 규율 이사회 결의 + 다운사이클 capex 하한 4조 명문화'),
                ('SD-1', 'HBM 사업부 독립 P&L 분리 + 패키징 인재 채용 착수'),
                ('MB-2', '신흥시장 영업팀 신설 (사우디·UAE·인도) + 시안 라이선스 협상'),
                ('SE-2', 'CXL SIG 워킹그룹 인력 10→25명 즉시 확대'),
                ('RS-3', 'NVIDIA CMX·SCADA 공동 개발 MOU 체결, FDP 파일럿'),
                ('SE-3', 'Stargate Korea LOI→DA 본 계약 + 4사(Electronics·C&T·SDS·Heavy) 컨소시엄 운영 협약'),
            ],
        },
        {
            'name': '중기',
            'period': '6~18M',
            'window': '2027 H1 ~ 2028 H1',
            'desc': '본격 실행 — R&D·인프라·시장 진입',
            'color': THEME['amber'],
            'strategies': [
                ('MB-4', '구글·아마존·MS와 메모리 공동 로드맵 협의체 가동, CXL 모듈 양산'),
                ('RS-6', '1c nm 80% yield 달성, V10 hybrid bonding 양산, 자체 IP 200건 출원'),
                ('SA-2', '일본 R&D 허브 1조 원 투자 + 나노임프린트 NIL 공동 개발'),
                ('SE-1', 'IMEC 공동 연구 협약 + 3D DRAM 스타트업 2~3개 M&A 집행'),
                ('SD-2', '자동차·의료 AI 메모리 AEC-Q100 인증 + OEM 파일럿 공급'),
                ('SE-3', 'Tier 1 neocloud equity-for-supply swap 1건 + 첫 BTO/BOO 컨소시엄 수주 (OpenAI/Anthropic 한국 capacity)'),
            ],
        },
        {
            'name': '장기',
            'period': '18M+',
            'window': '2028 H2 ~ 2030+',
            'desc': '효과 발현 — 매출·마진·생존력',
            'color': THEME['samsung_blue'],
            'strategies': [
                ('RS-6', 'V11 자체 IP 비중 70%+ 달성, NAND 공정 주기 24개월+ 연장'),
                ('MB-4', '커스텀 AI 메모리 매출 비중 20%+ (2029), CMX SSD 점유 40%+ (2028)'),
                ('RS-3', 'FDP 검증 SSD 매출 비중 30% (2030), 통합 매출 $8~9B/년'),
                ('SE-1', '3D DRAM 시제품 가동, 패러다임 전환 시 즉시 양산 피벗'),
                ('SD-2', '산업용 메모리 매출 안정 성장 — 사이클 충격 흡수'),
                ('SE-3', 'Hwaseong/Pyeongtaek 인접 자체 AI DC 캠퍼스 1단계 50~100MW + Sovereign AI 매출 $3B+ (2032)'),
            ],
        },
    ]

    x = Inches(0.4)
    base_y = Inches(1.85)
    phase_h = Inches(1.6)
    gap = Inches(0.1)
    w = Inches(12.5)

    for i, phase in enumerate(phases):
        y = base_y + (phase_h + gap) * i
        # 카드 배경
        add_rect(slide, x, y, w, phase_h, fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        # 좌측 색상 띠 (시간 구간)
        add_rect(slide, x, y, Inches(0.10), phase_h, fill=phase['color'])
        # 좌측 영역: 구간 정보
        add_text(slide, x + Inches(0.25), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 phase['name'], font=FONT_KO, size=18, bold=True, color=phase['color'])
        add_text(slide, x + Inches(0.25), y + Inches(0.50), Inches(2.5), Inches(0.30),
                 phase['period'], font=FONT_EN, size=14, bold=True, color=phase['color'])
        add_text(slide, x + Inches(0.25), y + Inches(0.82), Inches(2.5), Inches(0.25),
                 phase['window'], font=FONT_EN, size=9, italic=True,
                 color=THEME['gray_caption'])
        add_text(slide, x + Inches(0.25), y + Inches(1.10), Inches(2.5), Inches(0.30),
                 phase['desc'], font=FONT_KO, size=9, italic=True,
                 color=THEME['dark_text'])

        # 우측 영역: 전략 리스트
        list_x = x + Inches(2.95)
        for j, (sid, action) in enumerate(phase['strategies']):
            iy = y + Inches(0.08 + j * 0.245)
            # 코드
            add_text(slide, list_x, iy, Inches(0.7), Inches(0.23),
                     sid, font=FONT_EN, size=9.5, bold=True, color=phase['color'])
            # 행동
            add_text(slide, list_x + Inches(0.7), iy, Inches(8.8), Inches(0.23),
                     action, font=FONT_KO, size=8.5, color=THEME['dark_text'])

    add_so_what(slide, Inches(6.55),
                '단기는 거버넌스·조직(돈 안 듦), 중기는 R&D·인프라(자원 집중), 장기는 매출·마진(수확) — 순서 지켜야 함')
    add_footer(slide, 4, TOTAL, '카테고리 4 · 시간축')
    return slide


# ============================================================
# 슬라이드 5: 부록 — 10개 전략 현황 요약 매트릭스 (작업 3)
# ============================================================

def build_slide_5_appendix(prs):
    """현재 위치 + 다음 마일스톤 + 신뢰도 한 눈 매트릭스."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '부록 · 11개 전략 현황 요약 — 우리는 지금 어디에 있는가',
               '외부 공개 자료(IR·산업 리포트·뉴스)만으로 정성·정량 분석. ✅ 확정 / 🔵 추정 / ⚠️ 정보 공백')

    # 매트릭스 — 행 11개, 열 4개 (전략·현재 위치·다음 마일스톤·신뢰도)
    rows = [
        ('MB-4', '커스텀 AI 메모리',
         'HBM 회복 (Q3 2025 35%), 베이스다이 커스텀 미가시',
         '2026 Tech Day SCADA + Pangea v3', '🔵'),
        ('RS-3', '고객특화·전환비용 (NVIDIA 통합)',
         'CMX 진입 (PM1753), SCADA 공개 로드맵 부재',
         '2026 SCADA AI SSD 로드맵 발표', '✅ / ⚠️'),
        ('RS-6', '공정 리더십 통합',
         '1c yield 60%, hybrid bonding IP 공백 (YMTC 의존 정황)',
         '2026 H2 V10 BV NAND 양산 + 80% yield', '🔵 / ⚠️'),
        ('MB-2', '동서 균형 공급망',
         '5거점 구축 (한·미·일·중·인도 계획), 시안 매년 갱신 리스크',
         '2026 Q4 시안 라이선스 + CHIPS 2.0', '✅'),
        ('SD-1', 'HBM 조직 독립 P&L',
         '메모리사업부 내 통합 운영, 분리 P&L 미공개',
         '2026 H1 P&L 분리 결정 (목표)', '⚠️'),
        ('RS-5', '재무 규율 + 초과이익 재투자',
         '현금 $63B 강점, 다운사이클 capex 하한 정책 명문화 부재',
         '2026 H1 이사회 결의 + IR 사전 공시', '✅ / ⚠️'),
        ('SA-2', '일본 R&D 허브 (EUV 우회 NIL)',
         'Canon NIL 양산 채택 사례 부재, R&D 단계 베팅',
         '2027 NIL 공동 개발 1차 결과', '⚠️'),
        ('SD-2', '산업용 AI 메모리 (자동차·의료)',
         'Tesla 다년 계약 ✅, AEC-Q100 양산 미공개, Micron에 후행',
         '2026~2028 AEC-Q100 인증 사이클', '🔵 / ⚠️'),
        ('SE-1', '3D DRAM + IMEC + M&A',
         'SK 30년 로드맵 대비 후행, Samsung 전담 조직 미공개',
         '2026~2027 전담 R&D 조직 + IMEC 협약', '⚠️'),
        ('SE-2', 'CXL SIG 표준 주도',
         'CMM-D 첫 제품 ✅, Pangea v3 (CXL 3.2) 2026 발표 예정',
         '2026 H1 Pangea v3 + 워킹그룹 인력 2배', '✅'),
        ('SE-3', 'AI 인프라 수직 진출 (Vertical Ascent)',
         'Stargate Korea LOI(2025-10) ✅, 그룹 통합 P&L·통합 의사결정 부재',
         '2026 H1 Stargate Korea 4사 컨소시엄 운영 협약 + DA 본계약', '✅ / ⚠️'),
    ]

    # 표 헤더
    table_x = Inches(0.4)
    table_y = Inches(1.85)
    col_widths = [Inches(0.7), Inches(2.6), Inches(4.4), Inches(3.6), Inches(1.2)]
    headers = ['ID', '전략', '현재 위치', '다음 마일스톤', '신뢰도']

    # 헤더 셀
    for i, h in enumerate(headers):
        x = table_x + Inches(sum(c.inches for c in col_widths[:i]))
        add_rect(slide, x, table_y, col_widths[i], Inches(0.45),
                 fill=THEME['samsung_blue'])
        add_text(slide, x + Inches(0.08), table_y + Inches(0.08),
                 col_widths[i] - Inches(0.16), Inches(0.30),
                 h, font=FONT_KO, size=10, bold=True, color=THEME['white'],
                 align='center' if i in [0, 4] else 'left')

    # 데이터 행
    row_h = Inches(0.40)
    for ri, (sid, name, current, milestone, confidence) in enumerate(rows):
        y = table_y + Inches(0.45) + row_h * ri
        # 메인/사이드 색상 구분
        is_main = sid in ('MB-4', 'RS-3', 'RS-6', 'MB-2', 'SD-1', 'RS-5')
        bg = THEME['amber_light'] if is_main else THEME['soft_blue_bg']
        accent = THEME['samsung_blue'] if is_main else THEME['amber']

        # 행 배경
        cells = [sid, name, current, milestone, confidence]
        for i, value in enumerate(cells):
            x = table_x + Inches(sum(c.inches for c in col_widths[:i]))
            add_rect(slide, x, y, col_widths[i], row_h,
                     fill=bg, line=THEME['light_gray'])

            if i == 0:
                # ID 셀 — 좌측 띠 + 코드
                add_rect(slide, x, y, Inches(0.06), row_h, fill=accent)
                add_text(slide, x + Inches(0.08), y + Inches(0.06),
                         col_widths[i] - Inches(0.12), Inches(0.28),
                         value, font=FONT_EN, size=9.5, bold=True,
                         color=accent, align='center')
            elif i == 1:
                add_text(slide, x + Inches(0.08), y + Inches(0.08),
                         col_widths[i] - Inches(0.16), Inches(0.28),
                         value, font=FONT_KO, size=8.5, bold=True,
                         color=THEME['dark_text'])
            elif i == 4:
                add_text(slide, x, y + Inches(0.08), col_widths[i], Inches(0.28),
                         value, font=FONT_EN, size=10, bold=True,
                         color=THEME['samsung_blue'], align='center')
            else:
                add_text(slide, x + Inches(0.10), y + Inches(0.08),
                         col_widths[i] - Inches(0.20), Inches(0.28),
                         value, font=FONT_KO, size=8,
                         color=THEME['dark_text'])

    # 핵심 인사이트 박스 (하단)
    insight_y = table_y + Inches(0.45) + row_h * 11 + Inches(0.12)
    add_rect(slide, table_x, insight_y, Inches(12.5), Inches(0.65),
             fill=THEME['amber'])
    add_text(slide, table_x + Inches(0.20), insight_y + Inches(0.08),
             Inches(12.1), Inches(0.27),
             '핵심 정보 공백 (⚠️ 7개): HBM P&L 분리 · hybrid bonding 자체 IP · SCADA 로드맵 · 3D DRAM 전담 조직 · AEC-Q100 양산 · NIL 협력 · 다운턴 capex 하한 명문화',
             font=FONT_KO, size=9, bold=True, color=THEME['white'])
    add_text(slide, table_x + Inches(0.20), insight_y + Inches(0.36),
             Inches(12.1), Inches(0.27),
             '→ 2026 H1 Tech Day + 이사회 결의 + IR 사전 공시로 정보 공백 해소가 다음 단계의 첫 행동',
             font=FONT_KO, size=8.5, italic=True, color=THEME['white'])

    add_footer(slide, 5, TOTAL, '부록 · 현황 요약')
    return slide


# ============================================================
# 슬라이드 6: SE-3 Deep Dive — AI 가치사슬 마진 분포 + 3-Tier 헷지 모델 (신규)
# ============================================================

def build_slide_6_se3_deepdive(prs):
    """SE-3 전략 deep dive — 가치사슬 마진 분포 + 3-Tier 단계 진입 모델."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'SE-3 Deep Dive · AI 가치사슬 수직 진출 — 메모리 위 70% 마진 헷지',
               'NVIDIA DC FY26 영업이익률 60.4% / AWS·Google Cloud 32.9% / 메모리 슈퍼사이클 ~25% — 3-Tier 단계 진입으로 가치사슬 상류 이동')

    # 좌측 상단: 가치사슬 마진 분포 (가로 바 차트형 표)
    margin_x = Inches(0.4)
    margin_y = Inches(1.85)
    margin_w = Inches(6.3)
    margin_h = Inches(2.55)

    add_rect(slide, margin_x, margin_y, margin_w, margin_h,
             fill=THEME['soft_white_bg'], line=THEME['light_gray'])
    add_rect(slide, margin_x, margin_y, Inches(0.10), margin_h, fill=THEME['red_alert'])
    add_text(slide, margin_x + Inches(0.20), margin_y + Inches(0.08),
             margin_w - Inches(0.30), Inches(0.32),
             '① 가치사슬 마진 분포 — 메모리는 "재료 공급자"',
             font=FONT_KO, size=11.5, bold=True, color=THEME['red_alert'])
    add_text(slide, margin_x + Inches(0.20), margin_y + Inches(0.40),
             margin_w - Inches(0.30), Inches(0.22),
             'Margin distribution across the AI value chain (FY2025/2026)',
             font=FONT_EN, size=8.5, italic=True, color=THEME['gray_caption'])

    # 마진 바 (시각화)
    bars = [
        ('NVIDIA Data Center', '$197B', '60.4%', 0.604, THEME['red_alert']),
        ('AWS / Google Cloud', '$115B / $48B', '32.9%', 0.329, THEME['amber']),
        ('Hyperscaler Avg', 'Big4 $388B capex', '~30%', 0.30, THEME['amber']),
        ('Memory (슈퍼사이클)', 'cyclical', '~25%', 0.25, THEME['samsung_blue']),
        ('Server OEM (Dell)', '$12.5B AI', '3~5%', 0.045, THEME['gray_caption']),
    ]
    bar_base_x = margin_x + Inches(0.20)
    bar_base_y = margin_y + Inches(0.75)
    label_w = Inches(2.0)
    rev_w = Inches(1.3)
    bar_max_w = Inches(2.4)
    bar_h = Inches(0.27)
    gap = Inches(0.08)

    for i, (label, rev, pct, frac, color) in enumerate(bars):
        by = bar_base_y + (bar_h + gap) * i
        # 라벨
        add_text(slide, bar_base_x, by, label_w, bar_h,
                 label, font=FONT_KO, size=9, bold=True, color=THEME['dark_text'])
        # 매출
        add_text(slide, bar_base_x + label_w, by, rev_w, bar_h,
                 rev, font=FONT_EN, size=8.5, color=THEME['gray_caption'])
        # 바 (영업이익률 상대 길이)
        bar_w = Inches(2.4 * frac)
        add_rect(slide, bar_base_x + label_w + rev_w, by + Inches(0.05),
                 bar_w, bar_h - Inches(0.10), fill=color)
        # %
        add_text(slide, bar_base_x + label_w + rev_w + bar_max_w + Inches(0.05), by,
                 Inches(0.6), bar_h,
                 pct, font=FONT_EN, size=9.5, bold=True, color=color)

    # 우측 상단: TAM (시장 규모)
    tam_x = Inches(6.85)
    tam_y = Inches(1.85)
    tam_w = Inches(6.3)
    tam_h = Inches(2.55)

    add_rect(slide, tam_x, tam_y, tam_w, tam_h,
             fill=THEME['soft_white_bg'], line=THEME['light_gray'])
    add_rect(slide, tam_x, tam_y, Inches(0.10), tam_h, fill=THEME['samsung_blue'])
    add_text(slide, tam_x + Inches(0.20), tam_y + Inches(0.08),
             tam_w - Inches(0.30), Inches(0.32),
             '② 진입 가능 TAM (2025 → 2030)',
             font=FONT_KO, size=11.5, bold=True, color=THEME['samsung_blue'])
    add_text(slide, tam_x + Inches(0.20), tam_y + Inches(0.40),
             tam_w - Inches(0.30), Inches(0.22),
             'Total addressable market for AI infra ascent',
             font=FONT_EN, size=8.5, italic=True, color=THEME['gray_caption'])

    tam_rows = [
        ('전 세계 DC 용량', '103 GW', '200 GW', '+14% CAGR'),
        ('AI DC 시장', '$147B', '$811B (2033)', '+24% CAGR'),
        ('AI 서버 시장', '$245B', '$524B', '+18% CAGR'),
        ('Big 4 Capex', '$388B', '$630B (2026)', '+62% YoY'),
        ('Neocloud', '$25B', '$180B', '+69% CAGR'),
    ]
    tr_x = tam_x + Inches(0.20)
    tr_y = tam_y + Inches(0.75)
    tr_w = tam_w - Inches(0.40)
    tr_h = Inches(0.32)

    # 헤더
    add_rect(slide, tr_x, tr_y, tr_w, tr_h, fill=THEME['samsung_blue'])
    add_text(slide, tr_x + Inches(0.10), tr_y + Inches(0.05),
             Inches(2.0), Inches(0.22),
             '지표', font=FONT_KO, size=9, bold=True, color=THEME['white'])
    add_text(slide, tr_x + Inches(2.10), tr_y + Inches(0.05),
             Inches(1.2), Inches(0.22),
             '2025', font=FONT_EN, size=9, bold=True, color=THEME['white'], align='center')
    add_text(slide, tr_x + Inches(3.30), tr_y + Inches(0.05),
             Inches(1.4), Inches(0.22),
             '2030', font=FONT_EN, size=9, bold=True, color=THEME['white'], align='center')
    add_text(slide, tr_x + Inches(4.70), tr_y + Inches(0.05),
             Inches(1.2), Inches(0.22),
             '성장', font=FONT_KO, size=9, bold=True, color=THEME['white'], align='center')

    for i, (name, v25, v30, cagr) in enumerate(tam_rows):
        ry = tr_y + tr_h + Inches(i * 0.32)
        bg = THEME['soft_blue_bg'] if i % 2 == 0 else THEME['white']
        add_rect(slide, tr_x, ry, tr_w, Inches(0.32),
                 fill=bg, line=THEME['light_gray'])
        add_text(slide, tr_x + Inches(0.10), ry + Inches(0.05),
                 Inches(2.0), Inches(0.22),
                 name, font=FONT_KO, size=8.5, bold=True, color=THEME['dark_text'])
        add_text(slide, tr_x + Inches(2.10), ry + Inches(0.05),
                 Inches(1.2), Inches(0.22),
                 v25, font=FONT_EN, size=8.5, color=THEME['dark_text'], align='center')
        add_text(slide, tr_x + Inches(3.30), ry + Inches(0.05),
                 Inches(1.4), Inches(0.22),
                 v30, font=FONT_EN, size=8.5, bold=True, color=THEME['samsung_blue'], align='center')
        add_text(slide, tr_x + Inches(4.70), ry + Inches(0.05),
                 Inches(1.2), Inches(0.22),
                 cagr, font=FONT_EN, size=8.5, bold=True, color=THEME['amber'], align='center')

    # 하단: 3-Tier 카드 (가로 3개)
    tier_base_y = Inches(4.55)
    tier_h = Inches(2.40)
    tier_w = Inches(4.18)
    tier_gap = Inches(0.10)
    tier_x = Inches(0.4)

    tiers = [
        {
            'name': 'Tier 1 · 지분 헷지',
            'subtitle': 'Equity Hedge (Light)',
            'color': THEME['green_pos'],
            'capital': '$5~10B (5y)',
            'risk': '낮음 — 운영 부담 없음',
            'actions': [
                'Catalyst Fund/계열 VC',
                '→ CoreWeave-급 neocloud 1~3% 지분',
                'HBM 장기공급 ↔ Equity swap',
                '하이퍼스케일러 Pref Equity',
            ],
            'kpi': '지분 IRR 15%+ / AI 이익 환수율 5%+',
        },
        {
            'name': 'Tier 2 · AI Factory 파트너',
            'subtitle': 'Already Started (Stargate Korea LOI)',
            'color': THEME['samsung_blue'],
            'capital': '$15~25B (5y)',
            'risk': '중 — 카니발리제이션 가능',
            'actions': [
                'Stargate Korea LOI → DA 본계약',
                '4사 컨소시엄 운영 협약',
                'BTO/BOO 모델',
                'Floating DC 상용화',
                'Samsung SDS GPUaaS 확장',
            ],
            'kpi': 'DC 매출 $5B (2030) / C&T 신규 수주 $10B/yr',
        },
        {
            'name': 'Tier 3 · 자체 AI 캠퍼스',
            'subtitle': 'Co-Owner (옵션)',
            'color': THEME['red_alert'],
            'capital': '$20~40B (5y)',
            'risk': '높음 — 메모리 capex 충돌',
            'actions': [
                'Hwaseong/Pyeongtaek 인접',
                '50~300MW AI DC 캠퍼스',
                'HBM 직공급 + 액침냉각',
                '자체 LNG/연료전지',
                'Sovereign AI (KR·SEA·중동)',
            ],
            'kpi': '자체 DC 200MW (2030) / Sovereign $3B (2032)',
        },
    ]

    for i, tier in enumerate(tiers):
        tx = tier_x + (tier_w + tier_gap) * i
        # 카드 배경
        add_rect(slide, tx, tier_base_y, tier_w, tier_h,
                 fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        # 좌측 색상 띠
        add_rect(slide, tx, tier_base_y, Inches(0.10), tier_h, fill=tier['color'])
        # 헤더: 이름
        add_text(slide, tx + Inches(0.20), tier_base_y + Inches(0.08),
                 tier_w - Inches(0.30), Inches(0.30),
                 tier['name'], font=FONT_KO, size=11, bold=True, color=tier['color'])
        # 부제
        add_text(slide, tx + Inches(0.20), tier_base_y + Inches(0.40),
                 tier_w - Inches(0.30), Inches(0.22),
                 tier['subtitle'], font=FONT_EN, size=8.5, italic=True,
                 color=THEME['gray_caption'])

        # 자본 / 리스크
        add_rect(slide, tx + Inches(0.20), tier_base_y + Inches(0.68),
                 tier_w - Inches(0.40), Inches(0.30),
                 fill=tier['color'], line=tier['color'])
        add_text(slide, tx + Inches(0.30), tier_base_y + Inches(0.71),
                 tier_w - Inches(0.50), Inches(0.24),
                 f"{tier['capital']}  ·  {tier['risk']}",
                 font=FONT_KO, size=8.5, bold=True, color=THEME['white'])

        # 액션
        for j, act in enumerate(tier['actions']):
            ay = tier_base_y + Inches(1.07) + Inches(j * 0.21)
            add_text(slide, tx + Inches(0.25), ay,
                     tier_w - Inches(0.40), Inches(0.20),
                     f'• {act}', font=FONT_KO, size=8.5, color=THEME['dark_text'])

        # KPI (하단)
        kpi_y = tier_base_y + tier_h - Inches(0.36)
        add_rect(slide, tx + Inches(0.20), kpi_y,
                 tier_w - Inches(0.40), Inches(0.28),
                 fill=THEME['amber_light'], line=THEME['light_gray'])
        add_text(slide, tx + Inches(0.30), kpi_y + Inches(0.04),
                 tier_w - Inches(0.50), Inches(0.20),
                 f"KPI · {tier['kpi']}",
                 font=FONT_KO, size=8, bold=True, color=THEME['samsung_blue'])

    add_so_what(slide, Inches(7.05),
                'AI 락인 시 메모리 위에서 마진 70%가 가치사슬 상류로 누적 — 3-Tier 동시 착수 (Tier 1·2 우선, Tier 3는 옵션). RS-5 capex 하한과 별도 자본 트랙으로 운용')
    add_footer(slide, 6, TOTAL, 'SE-3 Deep Dive · 가치사슬 헷지')
    return slide


# ============================================================
# Main
# ============================================================

def main():
    print(f'Loading template (theme only): {TEMPLATE}')
    prs = Presentation(TEMPLATE)
    remove_existing_slides(prs)
    print(f'After cleanup — slides: {len(prs.slides)}')

    builders = [
        build_slide_1_scenario,
        build_slide_2_impact,
        build_slide_3_owner,
        build_slide_4_timeline,
        build_slide_5_appendix,
        build_slide_6_se3_deepdive,
    ]
    for i, builder in enumerate(builders, 1):
        print(f'  Building slide {i}: {builder.__name__}...')
        builder(prs)

    print(f'Saving to: {OUTPUT}')
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f'Done. Output size: {size_kb:.1f} KB · {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
