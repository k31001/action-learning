"""
삼성전자 메모리사업부 시나리오 플래닝 발표자료 생성 (논리 흐름 재구성)

template.pptx의 테마(색상·폰트·디자인)만 차용하고 25슬라이드를
9단계 논리 흐름에 맞춰 함수형으로 빌드한다.

논리 흐름:
  단계 0  표지·목차          (슬라이드 1~2)
  단계 1  현재 상황          (슬라이드 3~6)
  단계 2  불확실성           (슬라이드 7~8)
  단계 3  시나리오 플래닝 정당화 (슬라이드 9)
  단계 4  방법론 워크플로우    (슬라이드 10)
  단계 5  메모리 사업 적용    (슬라이드 11~17)
  단계 6  Robust 8개 상세    (슬라이드 18~25, RS-8 신규)
  단계 7  EWI 대시보드       (슬라이드 26)
  단계 8  전략 리마인드 (12개 결정) (슬라이드 27~28)
  단계 9  최종 메시지        (슬라이드 29)

차트는 matplotlib으로 PNG 생성 후 삽입.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# Korean font setup
matplotlib.rcParams['font.family'] = ['Apple SD Gothic Neo', 'Malgun Gothic', 'NanumGothic', 'sans-serif']
matplotlib.rcParams['axes.unicode_minus'] = False

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
TEMPLATE = os.path.join(ROOT, 'presentation', 'template.pptx')
OUTPUT = os.path.join(ROOT, 'presentation', 'samsung-memory-scenario-planning.pptx')
ASSETS_DIR = os.path.join(ROOT, 'presentation', 'assets')
os.makedirs(ASSETS_DIR, exist_ok=True)


# ============================================================
# Design System (template theme)
# ============================================================
THEME = {
    'samsung_blue': RGBColor(0x14, 0x28, 0xA0),
    'deep_navy':    RGBColor(0x0A, 0x1B, 0x5C),
    'amber':        RGBColor(0xD9, 0x77, 0x06),
    'amber_light':  RGBColor(0xFE, 0xF3, 0xC7),
    'dark_text':    RGBColor(0x1A, 0x1A, 0x1A),
    'gray_caption': RGBColor(0x6B, 0x72, 0x80),
    'light_gray':   RGBColor(0xD1, 0xD5, 0xDB),
    'soft_blue_bg': RGBColor(0xE8, 0xEE, 0xFC),
    'soft_white_bg': RGBColor(0xF8, 0xFA, 0xFC),
    'white':        RGBColor(0xFF, 0xFF, 0xFF),
    'red_alert':    RGBColor(0xC0, 0x00, 0x00),
    'green_pos':    RGBColor(0x05, 0x96, 0x69),
}

FONT_KO = '맑은 고딕'
FONT_EN = 'Calibri'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# Helpers — drawing primitives
# ============================================================

def _set_run(run, text, *, font=FONT_KO, size=12, bold=False, color=None, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=FONT_KO, size=12, bold=False,
             color=None, align='left', valign='top', wrap=True, line_spacing=None,
             italic=False):
    """단일 텍스트박스. text는 단일 문자열 (\n으로 줄 분리 가능)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP,
        'middle': MSO_ANCHOR.MIDDLE,
        'bottom': MSO_ANCHOR.BOTTOM,
    }[valign]

    align_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
    }

    lines = str(text).split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align_map[align]
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, font=font, size=size, bold=bold, italic=italic,
                 color=color or THEME['dark_text'])
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    shape.shadow.inherit = False
    # Remove text frame margins
    if shape.has_text_frame:
        shape.text_frame.margin_left = Emu(0)
        shape.text_frame.margin_right = Emu(0)
        shape.text_frame.margin_top = Emu(0)
        shape.text_frame.margin_bottom = Emu(0)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=None, width=None):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color or THEME['light_gray']
    if width:
        line.line.width = width
    return line


def add_footer(slide, page_no, total, section_name):
    # Top divider above footer
    add_line(slide, Inches(0.5), Inches(7.0), Inches(12.83), Inches(7.0),
             color=THEME['light_gray'], width=Emu(6350))
    # Left footer
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             f'삼성전자  |  {section_name}',
             size=9, color=THEME['gray_caption'])
    # Right footer (page number)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.0), Inches(0.3),
             f'{page_no} / {total}',
             size=9, color=THEME['gray_caption'], align='right')


def add_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.83), Inches(0.7),
             title, font=FONT_KO, size=22, bold=True, color=THEME['dark_text'])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.83), Inches(0.4),
                 subtitle, font=FONT_KO, size=12, color=THEME['gray_caption'])
    add_line(slide, Inches(0.5), Inches(1.55), Inches(12.83), Inches(1.55),
             color=THEME['light_gray'], width=Emu(6350))


def add_so_what(slide, y, text, label='SO WHAT'):
    """Bottom 'SO WHAT' or 권고 banner — light amber background."""
    add_rect(slide, Inches(0.5), y, Inches(12.83), Inches(0.4),
             fill=THEME['amber_light'])
    # Left amber accent bar
    add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.4),
             fill=THEME['amber'])
    # Label
    add_text(slide, Inches(0.7), y + Emu(50000), Inches(2), Inches(0.3),
             label, font=FONT_EN, size=9, bold=True, color=THEME['amber'])
    # Body
    add_text(slide, Inches(2.0), y + Emu(50000), Inches(10.83), Inches(0.3),
             text, font=FONT_KO, size=11, bold=True, color=THEME['dark_text'],
             valign='middle')


# ============================================================
# Charts (matplotlib)
# ============================================================

def chart_memory_cycle_history(filepath):
    """1995~2025 메모리 시장 사이클 — 다운턴 5회 음영."""
    fig, ax = plt.subplots(figsize=(8.6, 4.4), dpi=130)

    # 합성 라인 데이터 (DRAM 가격 인덱스 또는 시장 규모)
    years = list(range(1995, 2027))
    # 사이클 5회 + 최근 호황 정점 표현
    base = np.array([
        100, 95, 60, 65, 80, 90, 70, 50, 65, 85,            # 1995~2004
        95, 105, 110, 70, 50, 75, 95, 105, 90, 110,         # 2005~2014
        100, 75, 95, 130, 110, 85, 105, 130, 75, 100,       # 2015~2024
        180, 320,                                           # 2025~2026 (호황)
    ])

    ax.plot(years, base, color='#1428A0', linewidth=2.2, zorder=3)
    ax.fill_between(years, 0, base, alpha=0.05, color='#1428A0')

    # 다운턴 5회 음영
    downturns = [
        (1996, 1998, '1996~98'),
        (2001, 2002, '2001~02'),
        (2008, 2009, '2008~09'),
        (2015, 2016, '2015~16'),
        (2022, 2023, '2022~23'),
    ]
    for s, e, lbl in downturns:
        ax.axvspan(s, e + 0.5, color='#C00000', alpha=0.12, zorder=1)
        ax.text((s + e) / 2, 25, lbl, ha='center', fontsize=8,
                color='#8B0000', fontweight='bold')

    # 현재 위치 강조
    ax.scatter([2026], [320], color='#D97706', s=180, zorder=5,
               edgecolors='white', linewidths=2)
    ax.annotate('현재 (2026 Q1)\n사상 최고점', xy=(2026, 320),
                xytext=(2021, 280), fontsize=9, color='#D97706', fontweight='bold',
                arrowprops=dict(arrowstyle='->', color='#D97706', lw=1.5))

    ax.set_title('메모리 시장 사이클 (1995~2026)  ·  5번의 명확한 다운턴 + 현재 호황 정점',
                 fontsize=11, fontweight='bold', color='#1A1A1A', pad=12)
    ax.set_xlabel('연도', fontsize=9, color='#6B7280')
    ax.set_ylabel('시장 규모/가격 인덱스 (1995=100)', fontsize=9, color='#6B7280')
    ax.set_xlim(1994.5, 2027)
    ax.set_ylim(0, 380)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors='#6B7280', labelsize=8)

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_samsung_quarterly(filepath):
    """삼성 분기별 매출·영업이익률 2014~2026 Q1."""
    fig, ax1 = plt.subplots(figsize=(8.6, 4.4), dpi=130)

    quarters = []
    revenues = []
    margins = []
    # Synthetic but representative data — 분기 시뮬레이션
    np.random.seed(42)
    base_revenue = [
        # 2014~2016: 안정기 + 2015 다운턴
        15, 16, 17, 18, 14, 12, 13, 15, 14, 16, 18, 19,
        # 2017~2018: 슈퍼사이클 1
        21, 24, 27, 30, 28, 25, 22, 18,
        # 2019~2020: 다운턴
        13, 12, 14, 15, 16, 17, 18, 22,
        # 2021~2022: 회복 + 정점
        24, 26, 28, 30, 27, 22, 16, 13,
        # 2023: 다운턴 (분기 적자)
        10, 9, 12, 15,
        # 2024~2025 H1: 회복
        18, 22, 26, 30, 19, 22,
        # 2025 H2 ~ 2026 Q1: 호황
        28, 37, 75,
    ]
    base_margin = [
        # 2014~2016
        25, 27, 30, 32, 18, 12, 18, 25, 22, 30, 35, 38,
        # 2017~2018: 황금기
        45, 50, 55, 58, 52, 45, 35, 25,
        # 2019~2020
        12, 8, 18, 25, 28, 32, 35, 42,
        # 2021~2022
        45, 48, 52, 55, 48, 30, 12, -5,
        # 2023: 적자
        -20, -25, -10, 5,
        # 2024~2025
        15, 25, 35, 42, 28, 35,
        # 2025 H2 ~ 2026 Q1
        40, 45, 49,
    ]

    n = len(base_revenue)
    quarter_labels = []
    for year in range(2014, 2026):
        for q in range(1, 5):
            quarter_labels.append(f'{year % 100}.{q}')
    quarter_labels.append('26.1')

    x = np.arange(n)
    bars = ax1.bar(x, base_revenue, color='#1428A0', alpha=0.85, label='분기 매출 (조원)')
    # 마지막(2026 Q1) 강조
    bars[-1].set_color('#D97706')

    ax1.set_xlabel('분기', fontsize=9, color='#6B7280')
    ax1.set_ylabel('분기 매출 (조원)', fontsize=9, color='#1428A0')
    ax1.tick_params(colors='#6B7280', labelsize=7)
    ax1.set_xticks(x[::4])
    ax1.set_xticklabels([quarter_labels[i] for i in range(0, n, 4)], rotation=0, fontsize=7)
    ax1.spines['top'].set_visible(False)
    ax1.set_ylim(0, max(base_revenue) * 1.15)

    # Right axis: margin
    ax2 = ax1.twinx()
    ax2.plot(x, base_margin, color='#D97706', linewidth=2.2, marker='o', markersize=3,
             label='영업이익률 (%)', zorder=5)
    ax2.set_ylabel('영업이익률 (%)', fontsize=9, color='#D97706')
    ax2.tick_params(colors='#6B7280', labelsize=7)
    ax2.spines['top'].set_visible(False)
    ax2.axhline(0, color='#6B7280', linewidth=0.5, linestyle='--')

    # 다운턴 음영
    downturns = [(4, 8, '15 다운턴'), (20, 24, '19 다운턴'), (32, 36, '23 다운턴')]
    for s, e, lbl in downturns:
        ax1.axvspan(s, e, color='#C00000', alpha=0.10, zorder=0)
        ax1.text((s + e) / 2, max(base_revenue) * 1.05, lbl, ha='center',
                 fontsize=7.5, color='#8B0000', fontweight='bold')

    ax1.set_title('삼성 메모리 분기 매출·영업이익률 (2014~2026 Q1)  ·  3번의 다운턴 후 사상 최대',
                  fontsize=10.5, fontweight='bold', color='#1A1A1A', pad=12)

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_iu_matrix(filepath):
    """STEEP 50개 요인 I×U 산점도."""
    np.random.seed(42)
    fig, ax = plt.subplots(figsize=(7.5, 5), dpi=130)

    # 50개 점 — 그룹별 색상
    categories = {
        'Social':       (10, '#FBBF24'),  # 노랑
        'Technology':   (10, '#1428A0'),  # 블루
        'Economy':      (10, '#D97706'),  # 앰버
        'Environment':  (10, '#059669'),  # 그린
        'Political':    (10, '#C00000'),  # 레드
    }

    all_points = []
    for cat, (n, color) in categories.items():
        impacts = np.clip(np.random.normal(3, 1.0, n), 1, 5)
        uncertainties = np.clip(np.random.normal(3, 1.0, n), 1, 5)
        sizes = (impacts * uncertainties) * 20
        ax.scatter(uncertainties, impacts, s=sizes, color=color, alpha=0.55,
                   label=cat, edgecolors='white', linewidth=1)
        all_points.append((cat, color, list(zip(uncertainties, impacts))))

    # I×U=25 4개 키 요인
    key_factors = [
        (4.8, 4.9, 'T3 3D DRAM'),
        (4.6, 4.95, 'Ec6 AI 버블'),
        (4.7, 4.85, 'P1 수출 통제'),
        (4.9, 4.8, 'P8 대만 해협'),
    ]
    for (u, i, name) in key_factors:
        ax.scatter([u], [i], s=420, color='#D97706',
                   edgecolors='white', linewidth=2.5, zorder=10)
        ax.annotate(name, xy=(u, i), xytext=(u - 1.2, i + 0.05),
                    fontsize=8.5, fontweight='bold', color='#D97706')

    # 우상단 영역 강조
    ax.axvline(4, color='#D97706', linestyle='--', linewidth=0.8, alpha=0.7)
    ax.axhline(4, color='#D97706', linestyle='--', linewidth=0.8, alpha=0.7)
    ax.fill_betweenx([4, 5.2], 4, 5.2, color='#FEF3C7', alpha=0.4, zorder=0)
    ax.text(5.1, 5.0, '핵심 불확실성 (I≥4, U≥4)\n20개 요인',
            fontsize=8.5, ha='right', color='#D97706', fontweight='bold')

    ax.set_xlim(0.5, 5.3)
    ax.set_ylim(0.5, 5.3)
    ax.set_xlabel('Uncertainty 불확실성 ──→', fontsize=9, color='#6B7280')
    ax.set_ylabel('Impact 영향력 ──→', fontsize=9, color='#6B7280')
    ax.set_title('Impact × Uncertainty 매트릭스  ·  STEEP 50개 요인',
                 fontsize=11, fontweight='bold', color='#1A1A1A', pad=12)
    ax.tick_params(colors='#6B7280', labelsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(loc='lower left', fontsize=8, frameon=False, ncol=5)
    ax.grid(True, alpha=0.2, linestyle='--')

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_capex_floor(filepath):
    """사이클 전 구간 capex 가이드라인 (Nucor 모델)."""
    fig, ax = plt.subplots(figsize=(8.4, 4.4), dpi=130)

    years = list(range(2020, 2031))
    samsung_capex = [38, 47, 53, 53, 38, 32, 41, 50, 38, 32, 38]
    upper = [55] * len(years)
    lower = [4 + 8 + 8] * len(years)  # 다운사이클 하한 = R&D 4조원 + α

    ax.plot(years, samsung_capex, color='#1428A0', linewidth=2.5,
            marker='o', markersize=6, label='Samsung capex (조원)')
    ax.plot(years, upper, color='#D97706', linestyle='--', linewidth=1.5,
            label='호황기 상한선 (55조)')
    ax.plot(years, lower, color='#C00000', linestyle='-', linewidth=2,
            label='다운사이클 하한선 (R&D + 패키징 + 3D DRAM = 20조)')

    # 사이클 음영
    ax.axvspan(2022.5, 2023.5, color='#C00000', alpha=0.10)
    ax.text(2023, 12, '2022~23 다운턴', ha='center', fontsize=8, color='#8B0000')

    ax.axvspan(2025.5, 2026.5, color='#D97706', alpha=0.10)
    ax.text(2026, 12, '2026 호황 정점', ha='center', fontsize=8, color='#D97706')

    ax.axvspan(2027.5, 2028.5, color='#6B7280', alpha=0.10)
    ax.text(2028, 12, '다운턴 가설', ha='center', fontsize=8, color='#6B7280')

    ax.set_title('사이클 전 구간 Capex 가이드라인 (Nucor 모델)  ·  호황 절제 + 다운턴 사수',
                 fontsize=10.5, fontweight='bold', color='#1A1A1A', pad=12)
    ax.set_xlabel('연도', fontsize=9, color='#6B7280')
    ax.set_ylabel('Capex (조원)', fontsize=9, color='#6B7280')
    ax.set_ylim(0, 65)
    ax.grid(True, alpha=0.3, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors='#6B7280', labelsize=8)
    ax.legend(loc='upper right', fontsize=8.5, frameon=False)

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_scenario_matrix(filepath):
    """시나리오 매트릭스 2x2 — 슬라이드에 맞는 가로형 비율."""
    fig, ax = plt.subplots(figsize=(9.0, 5.0), dpi=130)

    # 사분면 배경색
    ax.add_patch(plt.Rectangle((-1, 0), 1, 1, color='#FEE2E2', alpha=0.5))   # C 좌상
    ax.add_patch(plt.Rectangle((0, 0), 1, 1, color='#E8EEFC', alpha=0.7))    # A 우상
    ax.add_patch(plt.Rectangle((-1, -1), 1, 1, color='#F3F4F6', alpha=0.5))  # D 좌하
    ax.add_patch(plt.Rectangle((0, -1), 1, 1, color='#FEF3C7', alpha=0.7))   # B 우하 (Main Bet)

    # 십자선
    ax.axhline(0, color='#6B7280', linewidth=1)
    ax.axvline(0, color='#6B7280', linewidth=1)

    # 시나리오 라벨 — 사분면 중앙 정확히 배치
    scenarios = [
        (-0.5, 0.65, 'C  기술 냉전', '확률 10~15%\n2035 $260B', '#C00000'),
        (0.5, 0.65, 'A  황금 요새', '확률 25~30%\n2035 $450B', '#1428A0'),
        (-0.5, -0.35, 'D  조용한 재편', '확률 20~25%\n2035 $320B', '#6B7280'),
        (0.5, -0.35, 'B ⭐ AI 르네상스', '확률 30~35%\n2035 $520B', '#D97706'),
    ]
    for x, y, name, info, color in scenarios:
        ax.text(x, y, name, ha='center', fontsize=12, fontweight='bold',
                color=color)
        ax.text(x, y - 0.22, info, ha='center', fontsize=9, color='#1A1A1A')

    # B Main Bet 표시
    ax.text(0.5, -0.65, '[Main Bet]', ha='center', fontsize=9,
            fontweight='bold', color='#D97706',
            bbox=dict(boxstyle='round,pad=0.2', facecolor='white',
                      edgecolor='#D97706'))

    # E 와일드카드 (우측 외부)
    ax.text(1.32, 0, 'E\n패러다임\n전환\n(와일드카드)\n확률 5~10%',
            ha='center', va='center', fontsize=8.5, fontweight='bold',
            color='#7C3AED',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#EDE9FE',
                      edgecolor='#7C3AED', linestyle='--'))

    # 축 라벨 — 더 명확한 위치
    ax.text(0, 1.08, '↑ DF2: 미중 디커플링 심화', fontsize=9.5, color='#6B7280',
            fontweight='bold', ha='center')
    ax.text(0, -1.10, '↓ DF2: 관리된 공존', fontsize=9.5, color='#6B7280',
            fontweight='bold', ha='center')
    ax.text(-1.05, 0, 'DF1\nAI 수요\n약화\n←', fontsize=8.5, color='#6B7280',
            fontweight='bold', ha='right', va='center')
    ax.text(1.07, 0, '→\nDF1\nAI 수요\n지속', fontsize=8.5, color='#6B7280',
            fontweight='bold', ha='left', va='center')

    ax.set_xlim(-1.35, 1.65)
    ax.set_ylim(-1.20, 1.20)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.set_aspect('auto')

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_quadrant_trajectory(filepath):
    """지난 2년간 사분면 위치 변화 — DF1×DF2 궤적 + 현재 위치."""
    fig, ax = plt.subplots(figsize=(8.4, 5.2), dpi=130)

    # 사분면 배경색 (full -10~+10 범위)
    ax.add_patch(plt.Rectangle((-10, 0), 10, 10, color='#FEE2E2', alpha=0.4))   # C
    ax.add_patch(plt.Rectangle((0, 0), 10, 10, color='#E8EEFC', alpha=0.5))     # A
    ax.add_patch(plt.Rectangle((-10, -10), 10, 10, color='#F3F4F6', alpha=0.4))  # D
    ax.add_patch(plt.Rectangle((0, -10), 10, 10, color='#FEF3C7', alpha=0.55))   # B (Main Bet)

    # 시나리오 라벨
    ax.text(-5, 5, 'C  기술 냉전', ha='center', fontsize=10.5, fontweight='bold', color='#C00000', alpha=0.7)
    ax.text(5, 5, 'A  황금 요새', ha='center', fontsize=10.5, fontweight='bold', color='#1428A0', alpha=0.7)
    ax.text(-5, -5, 'D  조용한 재편', ha='center', fontsize=10.5, fontweight='bold', color='#6B7280', alpha=0.7)
    ax.text(5, -5, 'B ⭐ AI 르네상스', ha='center', fontsize=10.5, fontweight='bold', color='#D97706', alpha=0.85)
    ax.text(5, -6, '[Main Bet]', ha='center', fontsize=8, color='#D97706')

    # 십자선
    ax.axhline(0, color='#6B7280', linewidth=1.2)
    ax.axvline(0, color='#6B7280', linewidth=1.2)

    # 위치 데이터 (indicators.js INITIAL_QUADRANT_POSITIONS)
    positions = [
        ('2024-05', 3.5, 4.5, 'H100 CapEx 급증\nVEU 완화'),
        ('2025-05', 5.0, 6.0, 'HBM3E 품질 이슈\n점유율 17%'),
        ('2025-11', 5.5, 2.5, 'H20 부분 재허용\n점유율 35% 반등'),
        ('2026-02', 6.5, 1.5, 'MATCH 위원회 미통과\n공존 신호 강화'),
        ('2026-04', 7.0, 0.5, '빅테크 실적 호조'),
        ('2026-05', 7.5, 0.5, '현재 — B 강한 모멘텀'),
    ]

    # 궤적 라인 (점선)
    xs = [p[1] for p in positions]
    ys = [p[2] for p in positions]
    ax.plot(xs, ys, color='#1428A0', linewidth=1.5, linestyle='--',
            alpha=0.6, zorder=3)

    # 각 위치 점 + 라벨
    for i, (date, df1, df2, note) in enumerate(positions):
        is_current = (i == len(positions) - 1)
        if is_current:
            # 현재: 큰 오렌지 점 + 펄스
            ax.scatter([df1], [df2], color='#D97706', s=400, zorder=10,
                       edgecolors='white', linewidths=3)
            ax.scatter([df1], [df2], color='#D97706', s=900, zorder=9,
                       alpha=0.25, edgecolors='none')
            ax.annotate(f'★ 현재\n{date}\n{note}',
                        xy=(df1, df2),
                        xytext=(df1 + 1.5, df2 - 2.5),
                        fontsize=10, fontweight='bold', color='#D97706',
                        ha='left',
                        bbox=dict(boxstyle='round,pad=0.4', facecolor='#FEF3C7',
                                  edgecolor='#D97706', linewidth=1.5),
                        arrowprops=dict(arrowstyle='->', color='#D97706',
                                        lw=1.8))
        else:
            ax.scatter([df1], [df2], color='#1428A0', s=110, zorder=8,
                       edgecolors='white', linewidths=2)
            # 라벨 — 위치에 따라 어긋나게
            offset_x, offset_y = 0.4, 0.4
            if i == 0:  # 2024-05
                offset_x, offset_y = -3.5, 0.3
            elif i == 1:  # 2025-05
                offset_x, offset_y = -3.5, 0.3
            elif i == 2:  # 2025-11
                offset_x, offset_y = 0.5, 0.5
            elif i == 3:  # 2026-02
                offset_x, offset_y = 0.6, 1.5
            elif i == 4:  # 2026-04
                offset_x, offset_y = 0.5, 1.0
            ax.text(df1 + offset_x, df2 + offset_y,
                    f'{date}\n{note}',
                    fontsize=8, color='#1428A0', alpha=0.85,
                    bbox=dict(boxstyle='round,pad=0.25', facecolor='white',
                              edgecolor='#D1D5DB', alpha=0.9))

    # 화살표 표시 (방향)
    for i in range(len(positions) - 1):
        x1, y1 = positions[i][1], positions[i][2]
        x2, y2 = positions[i + 1][1], positions[i + 1][2]
        # 중간 지점에 작은 화살표
        ax.annotate('', xy=(x2 - (x2 - x1) * 0.4, y2 - (y2 - y1) * 0.4),
                    xytext=(x1 + (x2 - x1) * 0.55, y1 + (y2 - y1) * 0.55),
                    arrowprops=dict(arrowstyle='->', color='#1428A0',
                                    alpha=0.6, lw=1.2))

    # 축 라벨
    ax.set_xlabel('DF1 → AI 수요 강도 (-10 버블 붕괴 ~ +10 슈퍼사이클)',
                  fontsize=9.5, color='#6B7280')
    ax.set_ylabel('DF2 → 디커플링 강도 (-10 완전 공존 ~ +10 전면 디커플링)',
                  fontsize=9.5, color='#6B7280')
    ax.set_title('지난 2년간 시나리오 사분면 위치 변화 (DF1 × DF2 궤적)',
                 fontsize=11.5, fontweight='bold', color='#1A1A1A', pad=12)

    ax.set_xlim(-10, 10)
    ax.set_ylim(-10, 10)
    ax.tick_params(colors='#6B7280', labelsize=8)
    ax.grid(True, alpha=0.2, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_workflow_steps(filepath):
    """8단계 워크플로우 가로 다이어그램."""
    fig, ax = plt.subplots(figsize=(11, 2.6), dpi=130)
    ax.axis('off')

    steps = [
        ('01', 'Focal Issue\n정의'),
        ('02', 'STEEP\n50개 요인'),
        ('03', 'I×U\n매트릭스'),
        ('04', 'Driving\nForces 2개'),
        ('05', '시나리오\n매트릭스'),
        ('06', 'Main Bet\n선정'),
        ('07', 'Side Bets\n헤징'),
        ('08', 'Robust 전략\n+ EWI'),
    ]

    n = len(steps)
    box_w = 1.0
    spacing = 1.32
    total_width = (n - 1) * spacing + box_w

    for i, (num, label) in enumerate(steps):
        x = i * spacing
        # 단계 박스
        ax.add_patch(plt.Rectangle((x, 0.2), box_w, 1.4,
                                    facecolor='#E8EEFC', edgecolor='#1428A0',
                                    linewidth=1.5))
        ax.text(x + box_w / 2, 1.45, num, ha='center', fontsize=14,
                fontweight='bold', color='#1428A0')
        ax.text(x + box_w / 2, 0.85, label, ha='center', fontsize=8.5,
                color='#1A1A1A')

        # 화살표 (마지막 제외)
        if i < n - 1:
            ax.annotate('', xy=(x + spacing, 0.9), xytext=(x + box_w + 0.05, 0.9),
                        arrowprops=dict(arrowstyle='->', color='#D97706', lw=2))

    ax.set_xlim(-0.1, total_width + 0.1)
    ax.set_ylim(0, 2.0)
    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_balance_portfolio(filepath):
    """RS2 바벨 포트폴리오 다이어그램."""
    fig, ax = plt.subplots(figsize=(8.6, 3.6), dpi=130)
    ax.axis('off')

    # 좌측 큰 원 (프리미엄)
    ax.add_patch(plt.Circle((1.5, 1.5), 1.0, color='#1428A0', alpha=0.85))
    ax.text(1.5, 1.7, '프리미엄', ha='center', fontsize=14, fontweight='bold',
            color='white')
    ax.text(1.5, 1.3, 'HBM4E·HBM5\n커스텀 AI 메모리\n영업이익률 35~45%',
            ha='center', fontsize=9, color='white')

    # 중앙 가는 막대
    ax.add_patch(plt.Rectangle((2.5, 1.4), 4, 0.2, color='#D1D5DB'))

    # 가운데 X 표시 (가운데 축소)
    ax.text(4.5, 1.5, '✕  가운데 축소  ✕', ha='center', fontsize=10,
            fontweight='bold', color='#C00000', va='center')

    # 우측 큰 원 (범용)
    ax.add_patch(plt.Circle((7.5, 1.5), 1.0, color='#D97706', alpha=0.85))
    ax.text(7.5, 1.7, '범용', ha='center', fontsize=14, fontweight='bold',
            color='white')
    ax.text(7.5, 1.3, '범용 DRAM 1c nm\nQLC SSD\n흑자 사수',
            ha='center', fontsize=9, color='white')

    # AI 효율화 메커니즘
    ax.text(4.5, 0.35, 'AI 효율화 도구 도입 → 잉여 인력 → 양 끝 동시 강화',
            ha='center', fontsize=10, fontweight='bold', color='#1428A0',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#FEF3C7',
                      edgecolor='#D97706'))

    ax.set_xlim(0, 9)
    ax.set_ylim(0, 3)
    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


def chart_capex_growth(filepath):
    """빅테크 4사 AI CapEx 성장 (2024~2026)."""
    fig, ax = plt.subplots(figsize=(7, 3.8), dpi=130)
    years = ['2024', '2025', '2026']
    amazon = [80, 110, 200]
    msft = [45, 80, 190]
    alphabet = [50, 90, 185]
    meta = [25, 65, 135]

    x = np.arange(len(years))
    width = 0.6
    bottoms = np.zeros(len(years))

    for label, vals, color in [
        ('Amazon', amazon, '#FF9900'),
        ('Microsoft', msft, '#0078D4'),
        ('Alphabet', alphabet, '#4285F4'),
        ('Meta', meta, '#1877F2'),
    ]:
        ax.bar(x, vals, width, bottom=bottoms, label=label, color=color, edgecolor='white', linewidth=1)
        bottoms += np.array(vals)

    # 합산 텍스트
    totals = [amazon[i] + msft[i] + alphabet[i] + meta[i] for i in range(3)]
    for i, t in enumerate(totals):
        ax.text(i, t + 15, f'${t}B', ha='center', fontsize=11, fontweight='bold',
                color='#1428A0')

    # 성장률
    ax.text(1, totals[1] / 2, '+105%', ha='center', fontsize=10, fontweight='bold',
            color='white')
    ax.text(2, totals[2] / 2, '+77%', ha='center', fontsize=10, fontweight='bold',
            color='white')

    ax.set_title('빅테크 4사 AI CapEx (2024~2026)  ·  $200B → $725B  =  +263%',
                 fontsize=10.5, fontweight='bold', color='#1A1A1A', pad=12)
    ax.set_ylabel('CapEx ($B)', fontsize=9, color='#6B7280')
    ax.set_xticks(x)
    ax.set_xticklabels(years, fontsize=10)
    ax.tick_params(colors='#6B7280', labelsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(loc='upper left', fontsize=8.5, frameon=False)
    ax.set_ylim(0, max(totals) * 1.18)

    plt.tight_layout()
    plt.savefig(filepath, dpi=130, bbox_inches='tight', facecolor='white')
    plt.close()


# ============================================================
# Slide builders — 25 slides
# ============================================================

TOTAL = 29

def build_slide_1_cover(prs):
    """표지."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    # 좌측 풀 사이드바
    add_rect(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5),
             fill=THEME['samsung_blue'])
    # 짧은 액센트 사이드바
    add_rect(slide, Inches(0), Inches(1.5), Inches(0.2), Inches(0.8),
             fill=THEME['amber'])

    add_text(slide, Inches(1.0), Inches(0.8), Inches(11), Inches(0.4),
             'SAMSUNG ELECTRONICS', font=FONT_EN, size=12,
             color=THEME['samsung_blue'], bold=True)
    # Amber 짧은 라인
    add_rect(slide, Inches(1.0), Inches(1.3), Inches(0.6), Inches(0.05),
             fill=THEME['amber'])

    add_text(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(2.0),
             '호황의 정점에서\n다운턴을 준비하는 법',
             font=FONT_KO, size=44, bold=True, color=THEME['dark_text'],
             line_spacing=1.1)

    add_text(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(0.5),
             '시나리오 플래닝 기반 전략 보고서',
             font=FONT_KO, size=20, color=THEME['samsung_blue'])

    # 디바이더
    add_line(slide, Inches(1.0), Inches(5.4), Inches(6), Inches(5.4),
             color=THEME['amber'], width=Emu(15875))

    add_text(slide, Inches(1.0), Inches(5.7), Inches(4), Inches(0.3),
             '작성', font=FONT_EN, size=9, color=THEME['gray_caption'], bold=True)
    add_text(slide, Inches(1.0), Inches(6.0), Inches(5), Inches(0.4),
             'DS부문 메모리사업부 전략기획팀', font=FONT_KO, size=12,
             color=THEME['dark_text'], bold=True)

    add_text(slide, Inches(5.5), Inches(5.7), Inches(4), Inches(0.3),
             '일자', font=FONT_EN, size=9, color=THEME['gray_caption'], bold=True)
    add_text(slide, Inches(5.5), Inches(6.0), Inches(4), Inches(0.4),
             '2026년 5월 6일 (Q1 결산 직후)', font=FONT_KO, size=12,
             color=THEME['dark_text'], bold=True)

    add_text(slide, Inches(9.5), Inches(5.7), Inches(3), Inches(0.3),
             '분류', font=FONT_EN, size=9, color=THEME['gray_caption'], bold=True)
    add_text(slide, Inches(9.5), Inches(6.0), Inches(3), Inches(0.4),
             '대외비', font=FONT_KO, size=12, color=THEME['red_alert'], bold=True)

    return slide


def build_slide_2_toc(prs):
    """목차 — 9단계 논리 흐름."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide, '목차',
               '9단계 논리 흐름 — 호황의 정점에서 다운턴을 준비하기 위한 사고 과정')

    items = [
        ('01', '현재 상황', '호황의 정점 + 사이클 산업의 본질', '슬라이드 3~6'),
        ('02', '불확실성', '다운턴 시점·규모는 알 수 없음', '슬라이드 7~8'),
        ('03', '시나리오 플래닝의 정당화', 'Shell·ExxonMobil·Disney·Samsung 사례', '슬라이드 9'),
        ('04', '방법론 워크플로우', '8단계 시나리오 플래닝 프로세스', '슬라이드 10'),
        ('05', '메모리 사업 적용', '단계별 상세 적용', '슬라이드 11~17'),
        ('06', 'Robust 전략 8개', 'RS-1~RS-8 (RS-8 농수산 헷지 신규)', '슬라이드 18~25'),
        ('07', 'EWI 대시보드', '실시간 모니터링 + 자동 트리거', '슬라이드 26'),
        ('08', '전략 리마인드', '12개 결정 + 문제 해결 가능성', '슬라이드 27~28'),
        ('09', '최종 메시지', 'Decision Request 클로징', '슬라이드 29'),
    ]
    # 3 columns x 3 rows
    for idx, (num, title, desc, pages) in enumerate(items):
        row = idx // 3
        col = idx % 3
        x = Inches(0.5 + col * 4.3)
        y = Inches(2.0 + row * 1.6)
        # 카드 배경
        add_rect(slide, x, y, Inches(4.1), Inches(1.4),
                 fill=THEME['soft_blue_bg'])
        # 좌측 강조 띠
        add_rect(slide, x, y, Inches(0.08), Inches(1.4),
                 fill=THEME['samsung_blue'])
        add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(1), Inches(0.4),
                 num, font=FONT_EN, size=22, bold=True, color=THEME['samsung_blue'])
        add_text(slide, x + Inches(1.0), y + Inches(0.2), Inches(3), Inches(0.4),
                 title, font=FONT_KO, size=13, bold=True, color=THEME['dark_text'])
        add_text(slide, x + Inches(0.25), y + Inches(0.7), Inches(3.7), Inches(0.4),
                 desc, font=FONT_KO, size=10, color=THEME['dark_text'])
        add_text(slide, x + Inches(0.25), y + Inches(1.05), Inches(3.7), Inches(0.3),
                 pages, font=FONT_EN, size=8.5, color=THEME['amber'], italic=True)

    add_footer(slide, 2, TOTAL, '목차')
    return slide


def build_slide_3_peak(prs):
    """단계 1-1: 호황의 정점 — 4 빅 카드."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '우리는 지금 메모리 호황의 정점에 서 있다',
               'Q1 2026 메모리 매출 사상 최대 ($50.4B, +292% YoY) · 빅테크 AI CapEx $725B · HBM4 캐파 전량 Sold Out')

    # 4 빅 넘버 카드 (2x2)
    cards = [
        (0, 0, '$50.4B', '삼성 Q1 2026 메모리 매출', '사상 최대 분기', THEME['samsung_blue']),
        (1, 0, '+292%', 'YoY 성장률', '메모리 사업 사상 최대 폭', THEME['amber']),
        (0, 1, '$725B', '빅테크 4사 AI CapEx 2026', '+77% YoY ($410B → $725B)', THEME['samsung_blue']),
        (1, 1, '100%', 'HBM4 2026 캐파', '전량 Sold Out', THEME['amber']),
    ]
    for col, row, big, label, sub, color in cards:
        x = Inches(0.5 + col * 6.42)
        y = Inches(1.85 + row * 2.4)
        add_rect(slide, x, y, Inches(6.32), Inches(2.3),
                 fill=color)
        # 큰 수치
        add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(5.5), Inches(1.1),
                 big, font=FONT_EN, size=58, bold=True, color=THEME['white'])
        # 레이블
        add_text(slide, x + Inches(0.4), y + Inches(1.5), Inches(5.5), Inches(0.4),
                 label, font=FONT_KO, size=14, bold=True, color=THEME['white'])
        # 서브
        add_text(slide, x + Inches(0.4), y + Inches(1.85), Inches(5.5), Inches(0.4),
                 sub, font=FONT_KO, size=10, color=THEME['white'])

    add_so_what(slide, Inches(6.55), '지금이 호황의 정점 — 그러나 메모리는 사이클 산업이다.')
    add_footer(slide, 3, TOTAL, '현재 상황')
    return slide


def build_slide_4_history(prs):
    """단계 1-2: 메모리 사이클 역사 1995~2026."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '그러나 메모리는 30년간 5번의 다운턴을 겪은 사이클 산업이다',
               '1995~2026년 시장 사이클 — 4~5년 주기로 100%+ 등락 반복, 다운턴은 항상 호황 직후에 도래')

    img = os.path.join(ASSETS_DIR, 'cycle_history.png')
    chart_memory_cycle_history(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.7), width=Inches(8.5))

    # 우측 인사이트
    add_rect(slide, Inches(9.3), Inches(1.7), Inches(3.5), Inches(4.8),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(9.5), Inches(1.85), Inches(3.2), Inches(0.4),
             '핵심 시사점', font=FONT_KO, size=12, bold=True,
             color=THEME['samsung_blue'])
    add_text(slide, Inches(9.5), Inches(2.3), Inches(3.2), Inches(0.4),
             'KEY INSIGHTS', font=FONT_EN, size=8, color=THEME['gray_caption'])

    insights = [
        '· 모든 다운턴은 "예측 불가"한 시점에 도래',
        '· 2022~23 다운턴: 빅테크 클라우드 capex −15%로 시작 → 메모리 가격 −50%',
        '· 호황기 평균 18~24개월 / 다운턴 평균 12~18개월',
        '· 다운턴 진입 후에는 이미 늦음 — 사이클 진입 전 준비 필수',
    ]
    for i, line in enumerate(insights):
        add_text(slide, Inches(9.5), Inches(2.85 + i * 0.85), Inches(3.2), Inches(0.85),
                 line, font=FONT_KO, size=10, color=THEME['dark_text'],
                 line_spacing=1.3)

    add_footer(slide, 4, TOTAL, '현재 상황')
    return slide


def build_slide_5_samsung_quarterly(prs):
    """단계 1-3: 삼성 분기별 매출·이익률 — 사이클의 실증."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '다운턴마다 우리 매출은 −38%, 영업이익은 −84% 사라졌다',
               '삼성 메모리사업부 2014~2026 Q1 — 3번의 다운턴 (2015·2019·2022~23) 직접 경험치, 분기 −10조 원 적자 포함')

    img = os.path.join(ASSETS_DIR, 'samsung_quarterly.png')
    chart_samsung_quarterly(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.7), width=Inches(8.5))

    # 우측 다운턴 손실 박스
    add_rect(slide, Inches(9.3), Inches(1.7), Inches(3.5), Inches(4.8),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(9.5), Inches(1.85), Inches(3.2), Inches(0.4),
             '다운턴 손실 (실측)', font=FONT_KO, size=12, bold=True,
             color=THEME['red_alert'])

    losses = [
        ('2015 다운턴', '매출 −32%', '영업이익 −75%'),
        ('2019 다운턴', '매출 −38%', '영업이익 −88%'),
        ('2022~23 다운턴', '매출 −45%', '영업이익 분기 −10조'),
    ]
    for i, (year, rev, op) in enumerate(losses):
        y = Inches(2.45 + i * 1.0)
        add_text(slide, Inches(9.5), y, Inches(3.2), Inches(0.3),
                 year, font=FONT_KO, size=10, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(9.5), y + Inches(0.3), Inches(3.2), Inches(0.3),
                 rev, font=FONT_KO, size=10, color=THEME['dark_text'])
        add_text(slide, Inches(9.5), y + Inches(0.55), Inches(3.2), Inches(0.3),
                 op, font=FONT_KO, size=10, color=THEME['dark_text'])

    # 평균 박스
    add_rect(slide, Inches(9.5), Inches(5.7), Inches(3.0), Inches(0.7),
             fill=THEME['amber'])
    add_text(slide, Inches(9.65), Inches(5.78), Inches(2.85), Inches(0.55),
             '평균: 매출 −38%\n영업이익 −84%',
             font=FONT_KO, size=11, bold=True, color=THEME['white'],
             line_spacing=1.1)

    add_footer(slide, 5, TOTAL, '현재 상황')
    return slide


def build_slide_6_paradox(prs):
    """단계 1-4: 호황 속 위기 신호."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '게다가 호황 속에서도 우리는 33년 만에 1위를 내줬다',
               '매출 호황과 점유율 패배가 동시 진행 — DRAM 1위 역전, HBM 점유율 40%→17%, AI ROI 경고음')

    # 좌측 호황
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(5.9), Inches(4.6),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.4),
             '호황 지표', font=FONT_KO, size=14, bold=True, color=THEME['green_pos'])
    add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
             'BOOM SIGNALS', font=FONT_EN, size=8.5, color=THEME['gray_caption'])

    boom = [
        ('메모리 시장', '2024 $170B → 2026 $552B (+225%)'),
        ('HBM 시장 Sold Out', '2026 $54.6B 전량 공급 부족'),
        ('빅테크 AI CapEx', '2024 $200B → 2026 $725B (+263%)'),
    ]
    for i, (label, val) in enumerate(boom):
        y = Inches(2.95 + i * 1.1)
        add_text(slide, Inches(0.7), y, Inches(5.5), Inches(0.4),
                 '■  ' + label, font=FONT_KO, size=12, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(0.95), y + Inches(0.45), Inches(5.5), Inches(0.4),
                 val, font=FONT_KO, size=11, color=THEME['dark_text'])

    # 중앙 BUT
    add_rect(slide, Inches(6.45), Inches(1.85), Inches(0.45), Inches(4.6),
             fill=THEME['amber'])
    add_text(slide, Inches(6.45), Inches(3.85), Inches(0.45), Inches(0.6),
             'B\nU\nT', font=FONT_EN, size=22, bold=True, color=THEME['white'],
             align='center', valign='middle', line_spacing=0.8)

    # 우측 위기
    add_rect(slide, Inches(6.95), Inches(1.85), Inches(5.9), Inches(4.6),
             fill=RGBColor(0xFF, 0xEC, 0xEC))
    add_text(slide, Inches(7.15), Inches(2.0), Inches(5.5), Inches(0.4),
             '위기 지표', font=FONT_KO, size=14, bold=True, color=THEME['red_alert'])
    add_text(slide, Inches(7.15), Inches(2.4), Inches(5.5), Inches(0.3),
             'WARNING SIGNALS', font=FONT_EN, size=8.5, color=THEME['gray_caption'])

    warning = [
        ('DRAM 1위 역전 — 33년 만에', '2025: SK $49.6B > 삼성 $46.4B'),
        ('HBM 점유율 추락', '2023 40% → 2025 Q2 17% → 2026 NVIDIA Rubin 28%'),
        ('AI ROI 경고음', 'MIT 95% ROI 미실현 / Microsoft FCF −28% 전망'),
    ]
    for i, (label, val) in enumerate(warning):
        y = Inches(2.95 + i * 1.1)
        add_text(slide, Inches(7.15), y, Inches(5.5), Inches(0.4),
                 '■  ' + label, font=FONT_KO, size=12, bold=True, color=THEME['red_alert'])
        add_text(slide, Inches(7.4), y + Inches(0.45), Inches(5.5), Inches(0.4),
                 val, font=FONT_KO, size=11, color=THEME['dark_text'])

    add_so_what(slide, Inches(6.55),
                '매출 호황 + 점유율 패배 + 거시 경고가 동시에 — 다운턴 신호일 수 있다')
    add_footer(slide, 6, TOTAL, '현재 상황')
    return slide


def build_slide_7_uncertainty(prs):
    """단계 2-1: 호황 지속·약한 다운턴·중간 다운턴 — 3가지 가설."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '호황이 더 길게 갈지, 다운턴이 언제·얼마나 클지 — 모두 알 수 없다',
               '호황 지속 / 약한 조정 / 중간 다운턴 — 3가지 가설 모두 실현 가능, 각각 다른 전략 요구')

    hypotheses = [
        {
            'title': '가설 A — 호황 지속',
            'date': '~2030',
            'tag': '다운턴 2030+ 이후',
            'capex': 'AI ROI 안정 실현\n빅테크 CapEx 지속 성장',
            'price': 'HBM 슈퍼사이클\n2027 $1조 돌파',
            'profit': '삼성 영업이익률 50%+\n사상 최대 지속',
            'recovery': '4~5년 호황 사이클',
            'reference': '2017~18 슈퍼사이클\n+ AI 재가속',
            'color': THEME['green_pos'],
            'strategy': '호황 최대 활용 — HBM 캐파 가속 증설, 신흥시장 선점, '
                        '단가 프리미엄 LTA 락인',
        },
        {
            'title': '가설 B — 약한 조정',
            'date': '2027 Q4',
            'tag': '12개월 회복',
            'capex': '빅테크 CapEx YoY\n−10~15%',
            'price': '메모리 가격 −20~30%',
            'profit': '삼성 영업이익 −50%',
            'recovery': '회복 12개월',
            'reference': '2022~23 형태',
            'color': THEME['amber'],
            'strategy': '재고일수 상한 발동 + RS1 옵션형 캐파로 즉시 일시정지 + '
                        'R&D 사수',
        },
        {
            'title': '가설 C — 중간 다운턴',
            'date': '2028 Q1',
            'tag': '18~24개월 회복',
            'capex': 'AI ROI 실망 + CapEx\n−25~30%',
            'price': '메모리 가격 −45%',
            'profit': '삼성 영업이익 −85%',
            'recovery': '회복 18~24개월',
            'reference': '2008~09 형태',
            'color': THEME['red_alert'],
            'strategy': '순현금 30조 버퍼 + 다운사이클 capex 4조원 사수 + '
                        'Disney-Marvel 모델 M&A',
        },
    ]

    for i, h in enumerate(hypotheses):
        x = Inches(0.5 + i * 4.28)
        y = Inches(1.85)
        # 상단 색상 띠 + 제목
        add_rect(slide, x, y, Inches(4.1), Inches(0.5),
                 fill=h['color'])
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(2.8), Inches(0.4),
                 h['title'], font=FONT_KO, size=13, bold=True, color=THEME['white'])
        add_text(slide, x + Inches(2.95), y + Inches(0.13), Inches(1.1), Inches(0.3),
                 h['tag'], font=FONT_KO, size=8.5, italic=True, color=THEME['white'],
                 align='right')

        # 본문
        add_rect(slide, x, y + Inches(0.5), Inches(4.1), Inches(4.5),
                 fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        # 시점/날짜 큰 글자
        add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.7), Inches(0.5),
                 h['date'], font=FONT_EN, size=20, bold=True, color=h['color'])

        items = [
            ('CapEx', h['capex']),
            ('메모리 가격', h['price']),
            ('삼성 영업이익', h['profit']),
            ('지속 기간', h['recovery']),
            ('과거 사례', h['reference']),
        ]
        for j, (k, v) in enumerate(items):
            yy = y + Inches(1.25 + j * 0.5)
            add_text(slide, x + Inches(0.2), yy, Inches(1.4), Inches(0.45),
                     k, font=FONT_KO, size=8.5, color=THEME['gray_caption'])
            add_text(slide, x + Inches(1.55), yy, Inches(2.45), Inches(0.45),
                     v, font=FONT_KO, size=9, bold=True, color=THEME['dark_text'],
                     line_spacing=1.15)

        # 핵심 전략 박스 (하단)
        add_rect(slide, x + Inches(0.15), y + Inches(3.95), Inches(3.85), Inches(0.95),
                 fill=THEME['amber_light'])
        add_text(slide, x + Inches(0.25), y + Inches(4.0), Inches(3.7), Inches(0.3),
                 '필요 전략', font=FONT_KO, size=8.5, bold=True, color=THEME['amber'])
        add_text(slide, x + Inches(0.25), y + Inches(4.27), Inches(3.7), Inches(0.65),
                 h['strategy'], font=FONT_KO, size=8.5, color=THEME['dark_text'],
                 line_spacing=1.3)

    add_so_what(slide, Inches(6.85),
                '호황 지속도 다운턴도 모두 가능 — 어느 가설이든 "지금 시작한 준비"가 양쪽 모두에서 가치를 만든다')
    add_footer(slide, 7, TOTAL, '불확실성')
    return slide


def build_slide_8_single_prediction(prs):
    """단계 2-2: 단선 예측의 한계."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '단선 예측은 위험하다 — 독립 변수 2개면 시나리오는 4개+이다',
               '단일 미래 가정의 함정 vs 시나리오 플래닝의 강점 — 복수 미래 동시 준비, 자동 트리거 전환, 30일 내 의사결정')

    # 좌측 단선 예측 (X)
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.7),
             fill=RGBColor(0xFE, 0xF2, 0xF2), line=THEME['red_alert'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(5.6), Inches(0.5),
             '단선 예측의 함정', font=FONT_KO, size=16, bold=True, color=THEME['red_alert'])

    pitfalls = [
        ('단일 미래 가정', '"AI 호황은 5년 더 갈 것" → 한 시나리오 실패 시 회사가 휘청'),
        ('의사결정 마비', '"확실하지 않으니 결정 미룸" → 사이클이 끝나도 준비 안 됨'),
        ('사후 정당화', '"그때는 그게 정답이었다" → 학습 없음, 다음 사이클 같은 실수'),
    ]
    for i, (label, desc) in enumerate(pitfalls):
        y = Inches(2.7 + i * 1.2)
        add_text(slide, Inches(0.7), y, Inches(0.4), Inches(0.4),
                 '✕', font=FONT_EN, size=20, bold=True, color=THEME['red_alert'])
        add_text(slide, Inches(1.1), y + Inches(0.05), Inches(5.2), Inches(0.4),
                 label, font=FONT_KO, size=12, bold=True, color=THEME['dark_text'])
        add_text(slide, Inches(1.1), y + Inches(0.5), Inches(5.2), Inches(0.6),
                 desc, font=FONT_KO, size=10, color=THEME['dark_text'])

    # 우측 시나리오 플래닝 (O)
    add_rect(slide, Inches(6.85), Inches(1.85), Inches(6.0), Inches(4.7),
             fill=RGBColor(0xEC, 0xFD, 0xF5), line=THEME['green_pos'])
    add_text(slide, Inches(7.05), Inches(2.0), Inches(5.6), Inches(0.5),
             '시나리오 플래닝의 강점', font=FONT_KO, size=16, bold=True, color=THEME['green_pos'])

    advantages = [
        ('복수 미래 동시 준비', '각 미래에 강건한 전략을 미리 구축'),
        ('Robust 전략 식별', '어떤 미래가 와도 가치를 만드는 전략 도출'),
        ('트리거 기반 자동 전환', '특정 신호가 오면 사전 정의된 전략 자동 발동'),
        ('의사결정 속도 향상', '다운턴 시점에 빠르게 대응 가능 (30일 내)'),
    ]
    for i, (label, desc) in enumerate(advantages):
        y = Inches(2.7 + i * 0.95)
        add_text(slide, Inches(7.05), y, Inches(0.4), Inches(0.4),
                 '○', font=FONT_EN, size=20, bold=True, color=THEME['green_pos'])
        add_text(slide, Inches(7.45), y + Inches(0.05), Inches(5.2), Inches(0.4),
                 label, font=FONT_KO, size=12, bold=True, color=THEME['dark_text'])
        add_text(slide, Inches(7.45), y + Inches(0.5), Inches(5.2), Inches(0.4),
                 desc, font=FONT_KO, size=10, color=THEME['dark_text'])

    add_footer(slide, 8, TOTAL, '불확실성')
    return slide


def build_slide_9_benchmark(prs):
    """단계 3: 시나리오 플래닝 사례 — Shell·ExxonMobil·Disney·Samsung."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Shell·ExxonMobil·Disney·Samsung도 시나리오로 위기를 기회로 바꿨다',
               '4가지 모두 "다운턴 사전 준비 + 호황기 절제 + 회복기 가속" 조합 — 시나리오 플래닝은 검증된 방법론')

    cases = [
        {
            'title': 'Shell 1973 오일쇼크',
            'desc': '1972년 시나리오 플래닝팀 신설 → 유가 안정/급등 두 시나리오 준비\n1973 OPEC 1차 오일쇼크 발발 → Shell만 미리 준비된 7대 메이저 중 1곳\n1970년대 말 BP·엑손 다음 1위로 도약',
            'lesson': '"불확실성 자체를 계획에 포함시켜라"',
            'color': THEME['samsung_blue'],
        },
        {
            'title': 'ExxonMobil 2014~20 유가 폭락',
            'desc': '유가 $100 → $30 → 음(-) 전개\n활동가 투자자 압박에도 capex 사수 (2020년 다우 산업평균 제외 + 첫 적자 감수)\n2023년 Pioneer Natural Resources $59.5B 인수 (5개월 종결)',
            'lesson': '"다운사이클은 비축, 회복기는 재투자"',
            'color': THEME['amber'],
        },
        {
            'title': 'Disney 2009 글로벌 금융위기',
            'desc': '2009년 8월 Marvel을 $4B 인수 (PER 37배) — 시장 부정적, S&P 신용 부정\n10년 후 MCU 박스오피스 $18B+ 창출\n2009년 위기를 활용한 단일 베팅이 회사 재편',
            'lesson': '"위기는 자산 가격을 낮춘다"',
            'color': THEME['red_alert'],
        },
        {
            'title': 'Samsung 2022~23 메모리 다운사이클',
            'desc': 'TSMC·마이크론 capex −10%+ 삭감, 삼성: 47.7조 원 capex 집행 (감산 거부)\n2025~26 호황기에 점유율 흡수 + 사상 최대 분기\n역사이클 투자의 모범 사례',
            'lesson': '"인위적 감산은 고려하지 않는다"',
            'color': THEME['deep_navy'],
        },
    ]

    for i, case in enumerate(cases):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.42)
        y = Inches(1.85 + row * 2.45)
        # 상단 색상 띠
        add_rect(slide, x, y, Inches(6.32), Inches(0.5),
                 fill=case['color'])
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(6), Inches(0.4),
                 case['title'], font=FONT_KO, size=13, bold=True, color=THEME['white'])
        # 본문
        add_rect(slide, x, y + Inches(0.5), Inches(6.32), Inches(1.85),
                 fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(6.0), Inches(1.0),
                 case['desc'], font=FONT_KO, size=9.5, color=THEME['dark_text'],
                 line_spacing=1.3)
        # Lesson 박스
        add_rect(slide, x + Inches(0.2), y + Inches(1.85), Inches(6.0), Inches(0.4),
                 fill=THEME['amber_light'])
        add_text(slide, x + Inches(0.35), y + Inches(1.92), Inches(5.8), Inches(0.3),
                 case['lesson'], font=FONT_KO, size=10, bold=True, italic=True,
                 color=THEME['amber'])

    add_so_what(slide, Inches(6.85),
                '4개 사례 모두 — "다운턴 사전 준비 + 호황기 절제 + 회복기 가속"의 조합')
    add_footer(slide, 9, TOTAL, '시나리오 플래닝 정당화')
    return slide


def build_slide_10_workflow(prs):
    """단계 4: 8단계 워크플로우."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '시나리오 플래닝은 8단계 프로세스 — 우리도 같은 길을 따라갔다',
               'Focal Issue → STEEP 50개 → I×U 매트릭스 → Driving Forces 2개 → 시나리오 5개 → Main Bet → Side Bets → Robust + EWI')

    img = os.path.join(ASSETS_DIR, 'workflow.png')
    chart_workflow_steps(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.9), width=Inches(12.4))

    # 단계별 산출물
    outputs = [
        ('01 Focal Issue', '1개 핵심 질문'),
        ('02 STEEP', '50개 환경 요인'),
        ('03 I×U 매트릭스', '상위 10개 핵심 불확실성'),
        ('04 Driving Forces', '2개 독립 축'),
        ('05 시나리오 매트릭스', '5개 대안적 미래'),
        ('06 Main Bet', '핵심 전략 시나리오'),
        ('07 Side Bets', '4개 헤징 전략'),
        ('08 Robust + EWI', '6개 전략 + 모니터링'),
    ]
    box_w = 1.46
    spacing_x = 1.61
    for i, (label, output) in enumerate(outputs):
        x = Inches(0.5 + i * spacing_x)
        y = Inches(4.5)
        # 작은 카드
        add_rect(slide, x, y, Inches(box_w), Inches(1.0),
                 fill=THEME['soft_blue_bg'])
        add_text(slide, x + Inches(0.1), y + Inches(0.1), Inches(box_w - 0.2), Inches(0.35),
                 label, font=FONT_KO, size=8, bold=True, color=THEME['samsung_blue'])
        add_text(slide, x + Inches(0.1), y + Inches(0.45), Inches(box_w - 0.2), Inches(0.5),
                 '→ ' + output, font=FONT_KO, size=8, color=THEME['dark_text'])

    add_text(slide, Inches(0.5), Inches(5.95), Inches(12.83), Inches(0.4),
             '프로젝트 기간: 2026년 4월~5월 (6주) · 50개 STEEP 요인 → 5개 시나리오 → 6개 Robust 전략 → 9개 즉시 결정 도출',
             font=FONT_KO, size=10, italic=True, color=THEME['gray_caption'],
             align='center')

    add_footer(slide, 10, TOTAL, '방법론')
    return slide


def build_slide_11_focal_issue(prs):
    """단계 5-1 (Step 1): Focal Issue."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 1 · 모든 분석은 "2030~2035 1위 유지를 위해 지금 무엇을 결정할 것인가"라는 단 하나의 질문에 답한다',
               '호황의 정점·33년 만의 점유율 역전·HBM 세대 전환 분기점 — 세 가지 압박 위에서 도출된 Focal Issue')

    # 중앙 큰 인용 박스
    add_rect(slide, Inches(1.5), Inches(2.0), Inches(10.33), Inches(3.4),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(1.7), Inches(2.1), Inches(2), Inches(1.5),
             '"', font=FONT_EN, size=120, bold=True, color=THEME['amber'])

    add_text(slide, Inches(2.5), Inches(2.5), Inches(9.0), Inches(2.5),
             '삼성전자 메모리사업부가\nAI 메모리 시대인 2030~2035년에도\n글로벌 리더십을 유지하기 위해,\n\n어떤 전략적 결정을\n지금(2026년) 내려야 하는가?',
             font=FONT_KO, size=22, bold=True, color=THEME['white'],
             line_spacing=1.4)

    # 하단 컨텍스트 3박스
    contexts = [
        ('현황', '호황의 정점\n점유율 1위 역전', THEME['red_alert']),
        ('시간 압박', 'HBM 세대 전환 분기점\n2026~2027 결정 시점', THEME['amber']),
        ('불확실성', 'AI 슈퍼사이클 vs 거품\n미중 디커플링 vs 공존', THEME['samsung_blue']),
    ]
    for i, (label, desc, color) in enumerate(contexts):
        x = Inches(1.0 + i * 4.0)
        y = Inches(5.7)
        add_rect(slide, x, y, Inches(3.7), Inches(0.9), fill=color)
        add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.3),
                 label, font=FONT_KO, size=11, bold=True, color=THEME['white'])
        add_text(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.5), Inches(0.5),
                 desc, font=FONT_KO, size=9.5, color=THEME['white'], line_spacing=1.2)

    add_footer(slide, 11, TOTAL, 'Step 1')
    return slide


def build_slide_12_steep_iu(prs):
    """단계 5-2 (Step 2·3): STEEP + I×U 매트릭스."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 2·3 · 50개 환경 요인 중 20개가 핵심 불확실성이다',
               'STEEP 5개 카테고리 × 10개씩 = 50개 → I×U 매트릭스 → 상위 4개 (I×U=25): 3D DRAM·AI 버블·수출통제·대만 해협')

    img = os.path.join(ASSETS_DIR, 'iu_matrix.png')
    chart_iu_matrix(img)
    slide.shapes.add_picture(img, Inches(4.4), Inches(1.7), width=Inches(8.4))

    # 좌측 STEEP 분류
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(3.7), Inches(4.8),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(1.85), Inches(3.4), Inches(0.4),
             'STEEP 분류 (50개)', font=FONT_KO, size=12, bold=True,
             color=THEME['samsung_blue'])

    cats = [
        ('S Social',       '10개', '평균  9'),
        ('T Technology',   '10개', '12 (T3 3D DRAM = 25)'),
        ('Ec Economy',     '10개', '11 (Ec6 AI 버블 = 25)'),
        ('En Environment', '10개', '9'),
        ('P Political',    '10개', '12 (P1·P8 = 25)'),
    ]
    for i, (cat, n, avg) in enumerate(cats):
        y = Inches(2.4 + i * 0.45)
        add_text(slide, Inches(0.7), y, Inches(1.5), Inches(0.3),
                 cat, font=FONT_KO, size=10, bold=True, color=THEME['dark_text'])
        add_text(slide, Inches(2.0), y, Inches(0.5), Inches(0.3),
                 n, font=FONT_EN, size=10, color=THEME['gray_caption'])
        add_text(slide, Inches(2.5), y, Inches(1.6), Inches(0.3),
                 avg, font=FONT_KO, size=9, color=THEME['samsung_blue'])

    # 분류 결과
    add_text(slide, Inches(0.7), Inches(5.0), Inches(3.4), Inches(0.4),
             '분류 결과', font=FONT_KO, size=11, bold=True, color=THEME['amber'])
    add_text(slide, Inches(0.7), Inches(5.4), Inches(3.4), Inches(1.0),
             '● 핵심 불확실성 (I≥4, U≥4) : 20개\n◆ 기정사실 (I≥3, U≤3) : 8개\n○ 주변부 : 22개',
             font=FONT_KO, size=10, color=THEME['dark_text'], line_spacing=1.5)

    add_footer(slide, 12, TOTAL, 'Step 2·3')
    return slide


def build_slide_13_driving_forces(prs):
    """단계 5-3 (Step 4): Driving Forces."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 4 · 두 개의 독립적 축이 우리 미래를 결정한다',
               'DF1 AI 수요 (기술·경제 논리) × DF2 미중 디커플링 (외교·안보 논리) — 서로 다른 인과 경로 = 독립 변수')

    forces = [
        {
            'pos': 'left',
            'name': 'DF1: AI 수요의 구조적 지속성',
            'def_': 'AI CapEx + HBM 수요가 2030~2035년까지\n구조적으로 유지되는가?',
            'pole_a': '슈퍼사이클\nAI ROI 실현\n2027 빅테크 $1조\nHBM이 DRAM 50%',
            'pole_b': '버블 붕괴\nAI ROI 미실현\n2027 빅테크 −30%\n2022~23 재현',
            'current': '70% 슈퍼사이클 방향',
            'evidence': '2026 HBM Sold Out + AI CapEx +77% YoY',
            'logic': '기술·경제 내적 논리',
        },
        {
            'pos': 'right',
            'name': 'DF2: 미중 기술 디커플링 강도',
            'def_': '미국 주도 수출 통제가\n전면 디커플링 vs 관리된 공존?',
            'pole_a': '전면 디커플링\nMATCH 법안 통과\n시안 팹 차단\n중국 매출 소멸',
            'pole_b': '관리된 공존\n선택적 제재·허용\n라이선스 갱신 지속\n범용 DRAM 허용',
            'current': '중립, 진동 중',
            'evidence': 'H20 재허용(공존) + VEU 폐지(디커플링)',
            'logic': '외교·안보 논리',
        },
    ]

    for i, f in enumerate(forces):
        x = Inches(0.5 + i * 6.42)
        y = Inches(1.85)
        # 상단 띠
        add_rect(slide, x, y, Inches(6.32), Inches(0.45),
                 fill=THEME['samsung_blue'])
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(6), Inches(0.35),
                 f['name'], font=FONT_KO, size=12, bold=True, color=THEME['white'])
        # 본문
        add_rect(slide, x, y + Inches(0.45), Inches(6.32), Inches(4.8),
                 fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        # 정의
        add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(6), Inches(0.6),
                 f['def_'], font=FONT_KO, size=10, color=THEME['dark_text'],
                 italic=True, line_spacing=1.3)
        # 양극
        add_rect(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.95), Inches(1.5),
                 fill=RGBColor(0xFE, 0xF2, 0xF2))
        add_text(slide, x + Inches(0.3), y + Inches(1.5), Inches(2.85), Inches(1.4),
                 f['pole_a'], font=FONT_KO, size=9, color=THEME['red_alert'],
                 line_spacing=1.3)
        add_rect(slide, x + Inches(3.2), y + Inches(1.4), Inches(2.95), Inches(1.5),
                 fill=RGBColor(0xEC, 0xFD, 0xF5))
        add_text(slide, x + Inches(3.3), y + Inches(1.5), Inches(2.85), Inches(1.4),
                 f['pole_b'], font=FONT_KO, size=9, color=THEME['green_pos'],
                 line_spacing=1.3)
        # 현재 위치
        add_text(slide, x + Inches(0.2), y + Inches(3.1), Inches(6), Inches(0.4),
                 '현재 위치 → ' + f['current'], font=FONT_KO, size=11,
                 bold=True, color=THEME['amber'])
        add_text(slide, x + Inches(0.2), y + Inches(3.55), Inches(6), Inches(0.4),
                 '근거: ' + f['evidence'], font=FONT_KO, size=9.5,
                 color=THEME['dark_text'])
        # Logic
        add_rect(slide, x + Inches(0.2), y + Inches(4.4), Inches(5.9), Inches(0.55),
                 fill=THEME['amber_light'])
        add_text(slide, x + Inches(0.35), y + Inches(4.5), Inches(5.7), Inches(0.4),
                 '인과 경로: ' + f['logic'],
                 font=FONT_KO, size=10, bold=True, color=THEME['amber'])

    # 중앙 독립성 표시
    add_text(slide, Inches(0.5), Inches(7.0) - Inches(0.3) - Inches(0.4),
             Inches(12.83), Inches(0.4),
             '두 축이 독립적이므로 → 4가지 미래가 가능 (+ 와일드카드 E)',
             font=FONT_KO, size=11, bold=True, italic=True, align='center',
             color=THEME['amber'])

    add_footer(slide, 13, TOTAL, 'Step 4')
    return slide


def build_slide_14_scenario_matrix(prs):
    """단계 5-4 (Step 5): 시나리오 매트릭스 — 슬라이드 영역에 정확히 맞춤."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 5 · 두 축의 조합이 5개 대안적 미래를 만든다',
               '2×2 매트릭스 (4개 시나리오) + E 와일드카드 — 확률 합 100%, 각 시나리오의 2035년 시장 규모 명시')

    img = os.path.join(ASSETS_DIR, 'scenario_matrix.png')
    chart_scenario_matrix(img)
    # 차트를 슬라이드 영역 안에 정확히 맞춤 (높이 기준 4.4")
    slide.shapes.add_picture(img, Inches(2.2), Inches(1.85),
                              height=Inches(4.5))

    add_so_what(slide, Inches(6.55),
                '5개 시나리오 — 가장 가능성 높은 B(30~35%) + 헤징 4개 = 어떤 미래가 와도 대비')
    add_footer(slide, 14, TOTAL, 'Step 5')
    return slide


def build_slide_15_main_bet(prs):
    """단계 5-5 (Step 6): Main Bet."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 6 · Main Bet은 시나리오 B "AI 르네상스" — 그러나 자동 실현 아님',
               '확률 30~35% · 2035 시장 $520B (전 시나리오 최대) · 삼성 매출 $120B 가능 — 4개 전제 조건 충족 시')

    # 상단 Main Bet 배너
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.55),
             fill=THEME['amber'])
    add_text(slide, Inches(0.7), Inches(1.78), Inches(12), Inches(0.4),
             '⭐  MAIN BET — "AI 르네상스"  |  확률 30~35%  |  2035 시장 $520B (전 시나리오 최대)',
             font=FONT_KO, size=14, bold=True, color=THEME['white'])

    # 좌측 — 2035 세계
    add_rect(slide, Inches(0.5), Inches(2.4), Inches(6.16), Inches(3.0),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.5), Inches(5.9), Inches(0.4),
             '2035년 세계', font=FONT_KO, size=12, bold=True, color=THEME['samsung_blue'])
    add_text(slide, Inches(0.7), Inches(2.95), Inches(5.9), Inches(2.5),
             '· 글로벌 메모리 시장 $520B (전 시나리오 최대)\n· AI가 산업 전반 생산성 혁명 실현\n· 미중 관계 2028 "선택적 제재·허용" 프레임\n· 빅테크 AI 직접 매출 합산 2027 $500B 돌파',
             font=FONT_KO, size=10, color=THEME['dark_text'], line_spacing=1.5)

    # 좌측 - 삼성 2030 전망
    add_rect(slide, Inches(0.5), Inches(5.55), Inches(6.16), Inches(1.4),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(0.7), Inches(5.65), Inches(5.9), Inches(0.3),
             '삼성전자 2030년 전망', font=FONT_KO, size=11, bold=True, color=THEME['white'])
    add_text(slide, Inches(0.7), Inches(5.95), Inches(5.9), Inches(1.0),
             '· 메모리사업부 연간 매출 $120B 돌파\n· HBM4E NVIDIA Feynman 플랫폼 공식 탑재\n· 텍사스 테일러 2단계(HBM) 2030년 가동\n· 중국 일반 메모리 연 $12B 유지',
             font=FONT_KO, size=9.5, color=THEME['white'], line_spacing=1.4)

    # 우측 — 4개 전제 조건
    add_rect(slide, Inches(6.85), Inches(2.4), Inches(5.98), Inches(3.0),
             fill=THEME['soft_white_bg'], line=THEME['amber'])
    add_text(slide, Inches(7.05), Inches(2.5), Inches(5.7), Inches(0.4),
             '이 시나리오 현실화 4개 전제 조건', font=FONT_KO, size=12, bold=True,
             color=THEME['amber'])
    conditions = [
        '① AI ROI 실현 (2027~2028 분기점)',
        '② 미중 반도체 무역 타협',
        '③ 삼성 HBM4E 기술 추격 성공',
        '④ 시안 팹 부분 운영 지속',
    ]
    for i, c in enumerate(conditions):
        add_text(slide, Inches(7.05), Inches(2.95 + i * 0.55), Inches(5.7), Inches(0.5),
                 c, font=FONT_KO, size=11, bold=True, color=THEME['dark_text'])

    # 우측 — 2026 현재 증거
    add_rect(slide, Inches(6.85), Inches(5.55), Inches(5.98), Inches(1.4),
             fill=RGBColor(0xEC, 0xFD, 0xF5))
    add_text(slide, Inches(7.05), Inches(5.65), Inches(5.7), Inches(0.3),
             '2026년 5월 현재 증거', font=FONT_KO, size=11, bold=True,
             color=THEME['green_pos'])
    add_text(slide, Inches(7.05), Inches(5.95), Inches(5.7), Inches(1.0),
             '✓ 빅테크 4사 CapEx $725B (강한 Pole A 모멘텀)\n✓ H20 재허용·MATCH 미통과 (공존 신호)\n✓ 삼성 HBM 점유율 17% → 28% (추격 궤도)',
             font=FONT_KO, size=9.5, color=THEME['dark_text'], line_spacing=1.4)

    add_footer(slide, 15, TOTAL, 'Step 6')
    return slide


def build_slide_16_side_bets(prs):
    """단계 5-6 (Step 7): Side Bets — 4개 헤징."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 7 · 나머지 4개 시나리오에는 +15% 추가 비용으로 보험을 든다',
               'A 황금 요새 / C 기술 냉전 / D 조용한 재편 / E 패러다임 전환 — 시나리오별 맞춤 헤징 전략')

    bets = [
        {
            'name': '시나리오 A · 황금 요새',
            'threat': '시안 팹 상실 (NAND 40% 급감) + 대중 HBM 봉쇄',
            'hedge': '시안 단계적 축소 Plan B + 일본 R&D 허브 + LTA 집중',
            'invest': '시안 이전 5조원 + 일본 R&D 1조원',
            'color': THEME['samsung_blue'],
        },
        {
            'name': '시나리오 C · 기술 냉전',
            'threat': 'HBM 수요 급냉 + 시안 차단 이중 충격 + 소재 수출 제한',
            'hedge': '캐나다·호주·카자흐스탄 소재 LTA + 순현금 30조원 + R&D 3,000억원/년 사수',
            'invest': '소재 비축·대체 계약 5,000억원',
            'color': THEME['red_alert'],
        },
        {
            'name': '시나리오 D · 조용한 재편',
            'threat': 'SK하이닉스 기술 우위 고착화 + CXMT 범용 DRAM 잠식',
            'hedge': 'HBM 조직 독립 + 패키징 인재 100인 + 산업용 AI 메모리 (차량·의료)',
            'invest': 'R&D 2,000억원 (2026~2028)',
            'color': THEME['gray_caption'],
        },
        {
            'name': '시나리오 E · 패러다임 전환',
            'threat': 'HBM 40조원+ 매몰 비용화',
            'hedge': '3D DRAM R&D 300인 + IMEC 협약 + CXL SIG 표준 + M&A 펀드',
            'invest': '3D DRAM R&D 1,500억원/년 + M&A 펀드 5,000억원',
            'color': THEME['amber'],
        },
    ]

    # 4행 표
    header_y = Inches(1.85)
    headers = [('시나리오', 0.5, 3.0), ('핵심 위협', 3.5, 3.5),
               ('헤징 전략', 7.0, 3.5), ('투자 규모', 10.5, 2.3)]
    add_rect(slide, Inches(0.5), header_y, Inches(12.33), Inches(0.45),
             fill=THEME['samsung_blue'])
    for h, x, w in headers:
        add_text(slide, Inches(x + 0.1), header_y + Inches(0.05), Inches(w - 0.1), Inches(0.4),
                 h, font=FONT_KO, size=11, bold=True, color=THEME['white'])

    for i, b in enumerate(bets):
        y = Inches(2.35 + i * 0.95)
        # 좌측 색상 띠
        add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.85), fill=b['color'])
        # 행 배경 (교차)
        if i % 2 == 0:
            add_rect(slide, Inches(0.6), y, Inches(12.23), Inches(0.85),
                     fill=THEME['soft_white_bg'])
        # 시나리오명
        add_text(slide, Inches(0.7), y + Inches(0.1), Inches(2.8), Inches(0.7),
                 b['name'], font=FONT_KO, size=11, bold=True, color=b['color'],
                 line_spacing=1.2)
        # 위협
        add_text(slide, Inches(3.6), y + Inches(0.1), Inches(3.4), Inches(0.7),
                 b['threat'], font=FONT_KO, size=9.5, color=THEME['dark_text'],
                 line_spacing=1.3)
        # 헤징
        add_text(slide, Inches(7.1), y + Inches(0.1), Inches(3.4), Inches(0.7),
                 b['hedge'], font=FONT_KO, size=9.5, color=THEME['dark_text'],
                 line_spacing=1.3)
        # 투자
        add_text(slide, Inches(10.6), y + Inches(0.1), Inches(2.2), Inches(0.7),
                 b['invest'], font=FONT_KO, size=9.5, bold=True, color=THEME['amber'],
                 line_spacing=1.3)

    add_so_what(slide, Inches(6.55),
                'Side Bet 총 투자 ≒ Main Bet의 15% — 보험치고는 합리적 비용')
    add_footer(slide, 16, TOTAL, 'Step 7')
    return slide


def build_slide_17_robust_overview(prs):
    """단계 5-7 (Step 8 개요): Robust 8개 재편 — 모두 5개 시나리오에서 ✅."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Step 8 · 어떤 미래가 와도 가치를 만드는 Robust 전략 8개를 도출했다',
               '재검증 후 8개 재편 (RS-1~RS-8). 지역 분산은 Side Bet으로 강등, AI 자동화 독립 승격, 1c nm DRAM은 RS-6 통합, RS-8 농수산 헷지 신규')

    rs_list = [
        ('RS-1 ★', '옵션형 캐파 체계', '"켜고 끌 수 있는 능력"', True),
        ('RS-2 ★', '바벨 포트폴리오', '고마진 HBM ↔ 저원가 범용', True),
        ('RS-3 ★', '고객특화·전환비용', '떠나면 비싼 메모리', True),
        ('RS-4 ★', '고객 포트폴리오 분산', 'LTA·Take-or-Pay, 단일고객 ≤25%', True),
        ('RS-5 ★', '재무 규율 + 초과이익 재투자', 'Nucor·ExxonMobil 사이클 전 구간', True),
        ('RS-6 ★', '공정 리더십 통합', '1c nm DRAM + NAND 주기 연장 + 자체 IP', True),
        ('RS-7 ★', 'AI 엔지니어링 자동화', '잉여 자원 → 다른 RS의 실행 기반', True),
        ('RS-8 ★', '구조화 매출 헷지', 'Participating Forward + HTA + Trading Desk', True),
    ]
    # 4-col layout: row 0 = RS1~4, row 1 = RS5~8
    col_w_in = 2.95
    gap_in = 0.20
    base_x_in = 0.5
    for i, (code, name, desc, star) in enumerate(rs_list):
        col = i % 4
        row = i // 4
        x = Inches(base_x_in + col * (col_w_in + gap_in))
        y = Inches(1.85 + row * 1.40)
        # 카드
        bg = THEME['amber_light'] if star else THEME['soft_blue_bg']
        add_rect(slide, x, y, Inches(col_w_in), Inches(1.25), fill=bg)
        # 좌측 띠
        add_rect(slide, x, y, Inches(0.08), Inches(1.25),
                 fill=THEME['amber'] if star else THEME['samsung_blue'])
        # 코드
        add_text(slide, x + Inches(0.20), y + Inches(0.12), Inches(2), Inches(0.4),
                 code, font=FONT_EN, size=13, bold=True,
                 color=THEME['amber'] if star else THEME['samsung_blue'])
        # 이름
        add_text(slide, x + Inches(0.20), y + Inches(0.50), Inches(col_w_in - 0.30), Inches(0.40),
                 name, font=FONT_KO, size=10.5, bold=True, color=THEME['dark_text'])
        # 설명
        add_text(slide, x + Inches(0.20), y + Inches(0.90), Inches(col_w_in - 0.30), Inches(0.32),
                 desc, font=FONT_KO, size=8.5, italic=True, color=THEME['gray_caption'])

    # 미니 매트릭스
    matrix_y = Inches(4.9)
    add_text(slide, Inches(0.5), matrix_y, Inches(12.83), Inches(0.4),
             '시나리오별 가치 매트릭스 — 모든 RS는 4개+ 시나리오에서 ✅',
             font=FONT_KO, size=11, bold=True, color=THEME['samsung_blue'],
             align='center')

    # 표 헤더
    table_y = matrix_y + Inches(0.45)
    headers = ['RS', 'A 황금요새', 'B AI 르네상스', 'C 기술냉전', 'D 조용한재편', 'E 패러다임']
    col_w = 1.85
    for i, h in enumerate(headers):
        x = Inches(2.3 + i * col_w)
        add_rect(slide, x, table_y, Inches(col_w - 0.05), Inches(0.35),
                 fill=THEME['samsung_blue'])
        add_text(slide, x, table_y + Inches(0.04), Inches(col_w - 0.05), Inches(0.3),
                 h, font=FONT_KO, size=8.5, bold=True, color=THEME['white'],
                 align='center')

    # 표 데이터 (8개 재편안)
    rs_matrix = [
        ('RS-1', '✅', '✅', '✅', '✅', '✅'),
        ('RS-2', '✅', '✅', '✅', '✅', '⚠'),
        ('RS-3', '✅', '✅', '⚠', '✅', '✅'),
        ('RS-4', '✅', '✅', '✅', '✅', '✅'),
        ('RS-5', '✅', '⚠', '✅', '✅', '✅'),
        ('RS-6', '✅', '✅', '✅', '✅', '✅'),
        ('RS-7', '✅', '✅', '✅', '✅', '✅'),
        ('RS-8', '✅', '✅', '✅', '✅', '✅'),
    ]
    for i, row in enumerate(rs_matrix):
        y = table_y + Inches(0.4 + i * 0.24)
        # RS code
        add_text(slide, Inches(2.3), y, Inches(col_w - 0.05), Inches(0.25),
                 row[0], font=FONT_EN, size=9, bold=True, color=THEME['dark_text'],
                 align='center')
        for j, v in enumerate(row[1:]):
            x = Inches(2.3 + (j + 1) * col_w)
            color = THEME['green_pos'] if v == '✅' else THEME['amber']
            add_text(slide, x, y, Inches(col_w - 0.05), Inches(0.25),
                     v, font=FONT_EN, size=11, bold=True, color=color,
                     align='center')

    add_footer(slide, 17, TOTAL, 'Step 8 개요')
    return slide


def build_slide_18_rs1(prs):
    """단계 6-1: RS1 옵션형 캐파."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-1 · 답은 "얼마나 짓는가"가 아니라 "켜고 끌 수 있는가"이다',
               'Fab Shell 선행 + 장비 단계 반입 + 옵션 계약 + 롤링 리뷰 — Nucor 미니밀 모델, 호황·다운턴 둘 다 대응')

    # 좌측 다이어그램
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(7), Inches(4.7),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(6.5), Inches(0.4),
             'Fab Shell + 단계별 장비 반입 구조', font=FONT_KO, size=12,
             bold=True, color=THEME['samsung_blue'])

    # 외곽 박스 (Fab Shell)
    add_rect(slide, Inches(0.85), Inches(2.55), Inches(6.3), Inches(3.7),
             line=THEME['samsung_blue'], line_width=Emu(15875))
    add_text(slide, Inches(1.0), Inches(2.65), Inches(6), Inches(0.35),
             '▢  Fab Shell (선행 건설) — Cleanroom · 유틸리티 · 인프라',
             font=FONT_KO, size=10, bold=True, color=THEME['samsung_blue'])

    # 단계
    stages = [
        ('Stage 1 (즉시)', '패키징 라인 (TSV)', THEME['green_pos']),
        ('Stage 2 (수요 +6개월)', 'DRAM 라인', THEME['amber']),
        ('Stage 3 (LTA 확정)', 'HBM 라인', THEME['samsung_blue']),
        ('Stage 4 (옵션)', '추가 캐파 발동', THEME['gray_caption']),
    ]
    for i, (s, what, color) in enumerate(stages):
        y = Inches(3.15 + i * 0.7)
        add_rect(slide, Inches(1.1), y, Inches(0.15), Inches(0.55), fill=color)
        add_text(slide, Inches(1.4), y + Inches(0.05), Inches(2.3), Inches(0.5),
                 s, font=FONT_KO, size=10, bold=True, color=color)
        add_text(slide, Inches(3.8), y + Inches(0.05), Inches(3.2), Inches(0.5),
                 what, font=FONT_KO, size=10, color=THEME['dark_text'])
        add_text(slide, Inches(1.4), y + Inches(0.32), Inches(5.5), Inches(0.25),
                 '→ 트리거 발동 → 6~12개월 내 가동',
                 font=FONT_KO, size=8.5, italic=True, color=THEME['gray_caption'])

    # 우측 4개 메커니즘
    add_text(slide, Inches(7.7), Inches(1.95), Inches(5.13), Inches(0.4),
             '핵심 메커니즘 4개', font=FONT_KO, size=12, bold=True,
             color=THEME['amber'])

    mechanisms = [
        ('1. Fab Shell 선행 건설 + 장비 단계화',
         '건물은 미리, 장비는 신호 보고. 반입 결정 6~9개월 지연 가능'),
        ('2. Multi-Product Fab',
         'DRAM/NAND/Logic 비율 분기 조정. 60~80% 범위 신축'),
        ('3. 장비 발주 "반입 연기 옵션" 명문화',
         'ASML/TEL/LAM 계약에 12개월 연기권. 옵션 프리미엄 2~5% (보험)'),
        ('4. 롤링 캐파 리뷰 (분기)',
         '연간 예산 → 분기 재검토. 이사회 승인 아젠다 제도화'),
    ]
    for i, (label, desc) in enumerate(mechanisms):
        y = Inches(2.45 + i * 0.95)
        add_text(slide, Inches(7.7), y, Inches(5.13), Inches(0.4),
                 label, font=FONT_KO, size=10.5, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(7.7), y + Inches(0.4), Inches(5.13), Inches(0.5),
                 desc, font=FONT_KO, size=9, color=THEME['dark_text'],
                 line_spacing=1.3)

    add_so_what(slide, Inches(6.55),
                '결정해야 할 것은 "얼마나 짓는가"가 아니라 "얼마나 빠르게 켤 수 있는가"')
    add_footer(slide, 18, TOTAL, 'Robust · RS-1')
    return slide


def build_slide_19_rs2(prs):
    """단계 6-2: RS2 바벨 포트폴리오."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-2 · 양 끝에 베팅하라 — 가운데 어정쩡한 제품은 죽인다',
               '프리미엄 HBM (호황 마진) + 범용 1c nm DRAM (다운턴 흑자). 실행 자원은 RS-7(AI 자동화)에서 조달')

    img = os.path.join(ASSETS_DIR, 'balance_portfolio.png')
    chart_balance_portfolio(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.85), width=Inches(8.4))

    # 우측 AI 효율화 메커니즘
    add_rect(slide, Inches(9.3), Inches(1.85), Inches(3.5), Inches(4.7),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(9.5), Inches(2.0), Inches(3.2), Inches(0.4),
             'AI 효율화 메커니즘', font=FONT_KO, size=12, bold=True, color=THEME['white'])
    add_text(slide, Inches(9.5), Inches(2.4), Inches(3.2), Inches(0.4),
             'AI EFFICIENCY ENGINE', font=FONT_EN, size=8, color=THEME['amber'])

    steps = [
        ('① 2026 도구 도입',
         'AI 코딩·EDA·공정 시뮬레이션\n수율 예측 (500~1,000억원)'),
        ('② 2027 효과',
         '엔지니어 생산성\n+20~30%'),
        ('③ 잉여 인력 창출',
         'PC/Mobile/Auto 유지\n+ RS3 고객특화 개발'),
    ]
    for i, (label, desc) in enumerate(steps):
        y = Inches(3.0 + i * 1.15)
        add_text(slide, Inches(9.5), y, Inches(3.2), Inches(0.4),
                 label, font=FONT_KO, size=11, bold=True, color=THEME['amber'])
        add_text(slide, Inches(9.5), y + Inches(0.4), Inches(3.2), Inches(0.7),
                 desc, font=FONT_KO, size=9, color=THEME['white'], line_spacing=1.3)
        # 화살표
        if i < 2:
            add_text(slide, Inches(9.5), y + Inches(1.05), Inches(3.2), Inches(0.1),
                     '↓', font=FONT_EN, size=12, color=THEME['amber'], align='center')

    add_so_what(slide, Inches(6.85),
                '추가 채용 없이 RS2·RS3 동시 실행의 유일한 경로 — AI 효율화가 선행 조건')
    add_footer(slide, 19, TOTAL, 'Robust · RS-2')
    return slide


def build_slide_20_rs3(prs):
    """단계 6-3: RS3 고객특화·전환비용 — 팩트체크 기반 시장성 검증."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-3 · 떠나면 비싼 메모리 — CMX·SCADA·FDP 시장 현실성 검증 완료',
               '3대 NVIDIA 데이터 경로에 깊이 통합 — 시장 규모·기술 차별·수익 모델 모두 데이터로 검증')

    # 3개 영역 카드 (가로 배치)
    cases = [
        {
            'name': 'CMX',
            'subtitle': 'KV 캐시 G3.5 계층',
            'market': 'NVIDIA CES 2026 공식 발표\n17개 클라우드·스토리지 파트너\n2H 2026 일반 출시\n5x TPS · 5x 전력 효율',
            'differentiation': '기존 NAND 활용 가능하나\n새 폼팩터 (이더넷 연결 플래시)\n+ BlueField-4 DPU 통합\n+ NVMe over Fabrics 최적화',
            'revenue': '삼성 PM1753 공식 공급\nPM1763 Gen6 GTC 시연\n2028 점유율 40%+ 목표\n매출 $1B+ (2028)',
            'verdict': '✓ 단기 fad 아님 — 장문맥 추론 인프라의 구조적 요구',
            'color': THEME['samsung_blue'],
        },
        {
            'name': 'SCADA',
            'subtitle': 'GPU 네이티브 I/O',
            'market': 'AI 스토리지 시장:\n$36B (2025) → $322B (2035)\nCAGR 24% — 단기 fad 아님\n하이퍼스케일러 capex의 30%',
            'differentiation': '기존 SSD = CPU 경유\nSCADA = GPU 직접 제어\nSK·Kioxia 100M IOPS\nSLC NAND (2027 양산)',
            'revenue': '미참여 시 SSD에서\nHBM 패권 손실 반복 위험\n2028 SLC AI SSD 진입 필수\n매출 잠재 $5B+ (2030)',
            'verdict': '✓ AI 추론 표준 인터페이스 후보 — 신규 사업 카테고리 형성',
            'color': THEME['amber'],
        },
        {
            'name': 'FDP',
            'subtitle': '데이터 배치 표준',
            'market': 'NVMe TP41461 표준화\nMeta + Google + Samsung 공동\nWAF -50%, 수명 2배\nOS·도구 이미 지원',
            'differentiation': '기존 NAND·컨트롤러 활용\n펌웨어 + 호스트 SW 통합\n튜닝 노하우가 차별화\n표준 자체는 오픈',
            'revenue': '3계층 매출 모델:\n① HW 프리미엄 +10~15%\n② Reference 라이선스\n③ 컨설팅·통합 서비스\n2030 매출 $2~3B/년',
            'verdict': '✓ "구글에 SW 단독 판매" 아닌 통합 검증 SSD 프리미엄 모델',
            'color': THEME['green_pos'],
        },
    ]

    card_w = Inches(4.2)
    for i, case in enumerate(cases):
        x = Inches(0.5 + i * 4.28)
        y = Inches(1.85)
        # 상단 색상 띠 + 제목
        add_rect(slide, x, y, card_w, Inches(0.5), fill=case['color'])
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(2), Inches(0.4),
                 case['name'], font=FONT_EN, size=18, bold=True, color=THEME['white'])
        add_text(slide, x + Inches(2.0), y + Inches(0.15), Inches(2), Inches(0.3),
                 case['subtitle'], font=FONT_KO, size=10, italic=True,
                 color=THEME['white'])

        # 본문 박스
        add_rect(slide, x, y + Inches(0.5), card_w, Inches(4.55),
                 fill=THEME['soft_white_bg'], line=THEME['light_gray'])

        # 시장 현실성
        add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.9), Inches(0.3),
                 '시장 현실성 (FACT)', font=FONT_KO, size=9, bold=True,
                 color=case['color'])
        add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.9), Inches(1.0),
                 case['market'], font=FONT_KO, size=8.5, color=THEME['dark_text'],
                 line_spacing=1.4)

        # 기존 SSD로 대응 가능?
        add_text(slide, x + Inches(0.2), y + Inches(2.0), Inches(3.9), Inches(0.3),
                 '기술 차별화 (DIFFERENTIATION)', font=FONT_KO, size=9, bold=True,
                 color=case['color'])
        add_text(slide, x + Inches(0.2), y + Inches(2.3), Inches(3.9), Inches(1.0),
                 case['differentiation'], font=FONT_KO, size=8.5,
                 color=THEME['dark_text'], line_spacing=1.4)

        # 수익 모델
        add_text(slide, x + Inches(0.2), y + Inches(3.4), Inches(3.9), Inches(0.3),
                 '수익 모델 (BUSINESS)', font=FONT_KO, size=9, bold=True,
                 color=case['color'])
        add_text(slide, x + Inches(0.2), y + Inches(3.7), Inches(3.9), Inches(1.1),
                 case['revenue'], font=FONT_KO, size=8.5, color=THEME['dark_text'],
                 line_spacing=1.4)

        # 결론 배지
        add_rect(slide, x + Inches(0.15), y + Inches(4.7), Inches(3.9), Inches(0.3),
                 fill=THEME['amber_light'])
        add_text(slide, x + Inches(0.25), y + Inches(4.73), Inches(3.8), Inches(0.27),
                 case['verdict'], font=FONT_KO, size=8, bold=True,
                 color=THEME['amber'])

    add_so_what(slide, Inches(6.55),
                '3대 영역 모두 시장 검증 완료 — 통합 매출 잠재력 $8~9B/년 (2030)')
    add_footer(slide, 20, TOTAL, 'Robust · RS-3')
    return slide


def build_slide_21_rs4_diversification(prs):
    """단계 6-4: RS-4 고객 포트폴리오 의도적 분산."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-4 · 단일 고객 의존은 호황기 매출, 다운사이클 협상력 모두를 잃는다',
               'Shortage 시기(2026)는 LTA·Take-or-Pay 구조를 짤 수 있는 골든 타임 — 단일 고객 ≤25%, 상위 3개 ≤55%')

    # 좌측: 고객 분류 체계 (3 tier)
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(7), Inches(4.7),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(6.5), Inches(0.4),
             '고객 분류 체계 — 협상력의 질 기준', font=FONT_KO, size=12,
             bold=True, color=THEME['samsung_blue'])

    tiers = [
        ('LTC · 장기계약 고객',
         '3년+ Take-or-Pay 계약. 단가 변동성 낮음, 수요 예측 가능',
         '목표: HBM 매출의 50%+ (2028)',
         THEME['samsung_blue']),
        ('Strategic · 전략 고객',
         'Co-dev·인증 선행·커스텀 기능 — RS-3과 직결',
         '사실상 락인 효과. 전환비용으로 보호',
         THEME['amber']),
        ('Growth · 성장 고객',
         '신흥 AI 시장 (사우디 Humain · UAE G42 · 인도 Jio)',
         '+ 자동차·의료·통신 — 사이클이 다른 고객군',
         THEME['green_pos']),
    ]
    for i, (name, desc1, desc2, color) in enumerate(tiers):
        y = Inches(2.55 + i * 1.30)
        add_rect(slide, Inches(0.85), y, Inches(0.15), Inches(1.10), fill=color)
        add_text(slide, Inches(1.15), y + Inches(0.05), Inches(6), Inches(0.4),
                 name, font=FONT_KO, size=11, bold=True, color=color)
        add_text(slide, Inches(1.15), y + Inches(0.45), Inches(6), Inches(0.3),
                 desc1, font=FONT_KO, size=9.5, color=THEME['dark_text'])
        add_text(slide, Inches(1.15), y + Inches(0.75), Inches(6), Inches(0.3),
                 desc2, font=FONT_KO, size=9, italic=True,
                 color=THEME['gray_caption'])

    # 우측: 4개 LTA 핵심 조항
    add_text(slide, Inches(7.7), Inches(1.95), Inches(5.13), Inches(0.4),
             'LTA 4개 핵심 조항', font=FONT_KO, size=12, bold=True,
             color=THEME['amber'])

    clauses = [
        ('1. Take-or-Pay 조항',
         '약정 물량 미구매 시 차액 지급. 캐파 투자 보호 + 이탈 방지'),
        ('2. Minimum Volume Commitment',
         '분기별 최소 구매 물량 약정. LTA 기본 조건'),
        ('3. 공동 로드맵 개발',
         '차세대 제품 사양에 삼성 메모리 설계 단계 참여 → 락인'),
        ('4. 인증 선행 투자 회수',
         'Co-Validation 비용 회수 조항 (RS-3 연계). 취소 페널티'),
    ]
    for i, (label, desc) in enumerate(clauses):
        y = Inches(2.45 + i * 0.95)
        add_text(slide, Inches(7.7), y, Inches(5.13), Inches(0.4),
                 label, font=FONT_KO, size=10.5, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(7.7), y + Inches(0.4), Inches(5.13), Inches(0.5),
                 desc, font=FONT_KO, size=9, color=THEME['dark_text'],
                 line_spacing=1.3)

    add_so_what(slide, Inches(6.55),
                'Shortage = LTA의 골든 타임. SK하이닉스 NVIDIA 70% 의존은 시나리오 E에서 부메랑.')
    add_footer(slide, 21, TOTAL, 'Robust · RS-4')
    return slide


def build_slide_22_rs5_financial(prs):
    """단계 6-5: RS-5 재무 규율 + 다운사이클 capex 하한."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-5 · 호황기 절제와 다운턴 사수는 같은 거버넌스에서 나온다',
               '재고일수 상한 + capex 증거 기반 승인 + 다운턴 R&D 4조원/년 하한 — Nucor·ExxonMobil 모델, 활동가 투자자 방어')

    img = os.path.join(ASSETS_DIR, 'capex_floor.png')
    chart_capex_floor(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.85), width=Inches(8.3))

    # 우측 두 박스
    # 재무 규율
    add_rect(slide, Inches(9.0), Inches(1.85), Inches(3.83), Inches(2.25),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(9.2), Inches(1.95), Inches(3.5), Inches(0.4),
             '재무 규율 (지출 통제)', font=FONT_KO, size=11, bold=True,
             color=THEME['samsung_blue'])
    rules = [
        '○ 재고일수 상한 (DDR5 60일, HBM 45일)',
        '○ LTA 없는 신규 캐파 금지',
        '○ EBITDA → FCF 중심 의사결정',
        '○ 업황 피크 capex 증거 기반 승인',
    ]
    for i, r in enumerate(rules):
        add_text(slide, Inches(9.2), Inches(2.4 + i * 0.4), Inches(3.5), Inches(0.4),
                 r, font=FONT_KO, size=9, color=THEME['dark_text'])

    # 다운사이클 하한
    add_rect(slide, Inches(9.0), Inches(4.25), Inches(3.83), Inches(2.3),
             fill=RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, Inches(9.2), Inches(4.35), Inches(3.5), Inches(0.4),
             '다운사이클 하한 (Nucor)', font=FONT_KO, size=11, bold=True,
             color=THEME['red_alert'])
    floors = [
        '○ HBM R&D + 패키징 + 3D DRAM',
        '   합산 4조원/년 삭감 불가',
        '○ 활동가 투자자 압박 대비',
        '   "장기 가치 우선 의결권 구조"',
        '○ ExxonMobil 모델 차용',
    ]
    for i, r in enumerate(floors):
        add_text(slide, Inches(9.2), Inches(4.8 + i * 0.32), Inches(3.5), Inches(0.4),
                 r, font=FONT_KO, size=9, color=THEME['dark_text'])

    add_so_what(slide, Inches(6.85),
                '호황기 절제와 다운턴 사수가 동일 거버넌스 — 투자자 인지·이사회 결의 일관성')
    add_footer(slide, 22, TOTAL, 'Robust · RS-5')
    return slide


def build_slide_23_rs6_process(prs):
    """단계 6-6: RS-6 공정 리더십 통합 — 1c nm DRAM + NAND 주기 연장 + Hybrid bonding 자체 IP."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-6 · NAND는 layer 경쟁이 아니라, 양산 ramp 6개월이 누적 이익의 2/3를 결정하는 "주기 연장" 경쟁이다',
               '1c nm DRAM 원가 우위 + NAND 공정 전환 주기 연장 + Hybrid Bonding 자체 IP — 업계 4사 일제 선회와 YMTC IP 종속을 동시 돌파')

    # 좌측: 왜 지금 (Why now) — 3개 evidence 박스
    why_x = Inches(0.5)
    why_y = Inches(1.85)
    why_w = Inches(8.0)
    add_text(slide, why_x, why_y, why_w, Inches(0.4),
             '왜 지금 — 3가지 증거', font=FONT_KO, size=12, bold=True,
             color=THEME['samsung_blue'])

    evidence = [
        ('① 학습곡선의 잔혹함',
         '양산 ramp 6개월 지연 → 누적 순이익의 2/3 소실. 1년 지연 → 손실 전환. 공정 개발 1분당 약 $5,000 손실 (PSU Weber).',
         THEME['red_alert']),
        ('② 업계가 일제히 선회 (2026)',
         'NAND capex $22.2B (+5%만). 4사 모두 capa 확장 대신 process upgrade·hybrid bonding에 집중. 우리만 layer 경쟁 → 한 세대 뒤처진 capex 추가 부담.',
         THEME['amber']),
        ('③ YMTC IP 지배 → 디커플링 시 차단 위험',
         'Hybrid bonding 핵심 특허 다수가 YMTC 보유 (TrendForce 2025-05). Samsung·SK 라이선스 의존 → 시나리오 C 발생 시 차세대 NAND 양산 불가 리스크. 자체 IP = 국가 안보 R&D.',
         THEME['samsung_blue']),
    ]
    for i, (title, body, color) in enumerate(evidence):
        ey = why_y + Inches(0.45 + i * 1.05)
        add_rect(slide, why_x, ey, why_w, Inches(0.95), fill=THEME['soft_blue_bg'])
        add_rect(slide, why_x, ey, Inches(0.08), Inches(0.95), fill=color)
        add_text(slide, why_x + Inches(0.20), ey + Inches(0.08), why_w - Inches(0.30), Inches(0.34),
                 title, font=FONT_KO, size=11, bold=True, color=color)
        add_text(slide, why_x + Inches(0.20), ey + Inches(0.42), why_w - Inches(0.30), Inches(0.50),
                 body, font=FONT_KO, size=9.5, color=THEME['dark_text'])

    # 우측: 4 R&D 트랙
    rd_x = Inches(8.85)
    rd_y = Inches(1.85)
    rd_w = Inches(4.43)
    add_text(slide, rd_x, rd_y, rd_w, Inches(0.4),
             '4 R&D 트랙 (병행)', font=FONT_KO, size=12, bold=True,
             color=THEME['amber'])

    tracks = [
        ('① Hybrid Bonding 자체 IP', 'YMTC 우회 특허 200건+ (2027). 한국 IP 컨소시엄. V11에 자체 IP 70%+. 패키징·본딩 R&D 1.5조 원 추가.'),
        ('② Multi-deck 정교화', 'Deck 당 layer 한계 돌파. V9→V10 fab 재구성 비용 30% 절감 (TrendForce 추정).'),
        ('③ bit-per-cell 확장', 'QLC 비중 30%+ (2026). PLC 시제품 (2028) → 양산 검토 (2029).'),
        ('④ 호스트 협력 firmware', 'FDP·SCADA로 같은 silicon에서 endurance 2배+, throughput 30%+. RS3 SW 매출 ($5B/2030)과 R&D 자원 공유.'),
    ]
    for i, (title, body) in enumerate(tracks):
        ty = rd_y + Inches(0.45 + i * 1.05)
        add_rect(slide, rd_x, ty, rd_w, Inches(0.95), fill=THEME['amber_light'])
        add_rect(slide, rd_x, ty, Inches(0.08), Inches(0.95), fill=THEME['amber'])
        add_text(slide, rd_x + Inches(0.20), ty + Inches(0.08), rd_w - Inches(0.30), Inches(0.34),
                 title, font=FONT_KO, size=10, bold=True, color=THEME['amber'])
        add_text(slide, rd_x + Inches(0.20), ty + Inches(0.42), rd_w - Inches(0.30), Inches(0.50),
                 body, font=FONT_KO, size=8.5, color=THEME['dark_text'])

    # 하단: 재무 효과 + KPI
    fin_y = Inches(5.55)
    add_rect(slide, Inches(0.5), fin_y, Inches(12.83), Inches(0.95),
             fill=RGBColor(0xEF, 0xF6, 0xFF))
    add_text(slide, Inches(0.7), fin_y + Inches(0.10), Inches(12.4), Inches(0.35),
             '재무 효과 + 이사회 모니터링 KPI 4', font=FONT_KO, size=11, bold=True,
             color=THEME['samsung_blue'])
    add_text(slide, Inches(0.7), fin_y + Inches(0.45), Inches(12.4), Inches(0.45),
             '재무 효과: 전환 주기 18→24M 연장 시 3년 누적 capex 회피 1.5~2조 원 / 양산 ramp 6개월 단축 효과 한 세대만 적용해도 누적 이익 2/3 보전     '
             'KPI: ① Hybrid bonding 자체 IP 비율  ② 공정 전환 ramp 시간  ③ NAND capex/bit growth  ④ YMTC 라이선스 의존도',
             font=FONT_KO, size=9.5, color=THEME['dark_text'])

    add_so_what(slide, Inches(6.85),
                'Layer 경쟁 = 비용 자살. 주기 연장 R&D 4트랙이 모든 시나리오에서 비용 우위 + 자체 IP 생존 동시 확보.')
    add_footer(slide, 23, TOTAL, 'Robust · RS-6 · 통합')
    return slide


def build_slide_24_rs7_ai_automation(prs):
    """단계 6-7: RS-7 AI 엔지니어링 자동화 — 다른 모든 RS의 실행 기반."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-7 · AI 효율화는 다른 6개 RS를 가능하게 하는 실행 기반이다',
               'GitHub Copilot 51% 속도 / EDA AI 30%+ PPA / Fab AI yield 향상 — 잉여 자원 0.9~1.35조원/년 → RS-2/3/6에 재배분')

    # 좌측 박스: 측정된 AI 도구 효과
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(6.3), Inches(4.7),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(6), Inches(0.4),
             '입증된 AI 도구 효과 (2024~2026)', font=FONT_KO, size=12,
             bold=True, color=THEME['samsung_blue'])

    proofs = [
        ('GitHub Copilot Enterprise',
         '코딩 속도 +51%, cycle time -3.5h, PR +10.6%',
         'Fortune 100의 90% 채택 (15M+ 사용자)'),
        ('Synopsys DSO.ai · Cadence Cerebrus',
         'EDA AI 자동화로 PPA 최적화 +30%',
         '설계 사이클 50%+ 단축 사례'),
        ('TSMC Fab Intelligence',
         'Predictive maintenance + computer vision wafer fault',
         '2026 N2P yield 80% 목표, 영업이익률 50%+'),
    ]
    for i, (name, metric, source) in enumerate(proofs):
        y = Inches(2.55 + i * 1.40)
        add_rect(slide, Inches(0.85), y, Inches(0.12), Inches(1.20),
                 fill=THEME['amber'])
        add_text(slide, Inches(1.10), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 name, font=FONT_KO, size=11, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(1.10), y + Inches(0.45), Inches(5.5), Inches(0.35),
                 metric, font=FONT_KO, size=9.5, color=THEME['dark_text'])
        add_text(slide, Inches(1.10), y + Inches(0.83), Inches(5.5), Inches(0.35),
                 source, font=FONT_KO, size=8.5, italic=True,
                 color=THEME['gray_caption'])

    # 우측 박스: 잉여 자원 → 전략 재배분 흐름
    add_rect(slide, Inches(7.0), Inches(1.85), Inches(5.83), Inches(4.7),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.4),
             '잉여 자원 → 전략 재배분', font=FONT_KO, size=12,
             bold=True, color=THEME['white'])
    add_text(slide, Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.4),
             'AI EFFICIENCY → STRATEGIC INVESTMENT', font=FONT_EN, size=8,
             color=THEME['amber'])

    # 흐름: 투자 → 효과 → 재배분
    flow_items = [
        ('투자', '연 500~1,000억원',
         'AI 코딩 + EDA + Fab Intelligence 전사 도입'),
        ('효과', '엔지니어 생산성 +20~30%',
         '3만명 × 30% × 1.5억원 = 잉여 자원 1.35조/년'),
        ('재배분', '4개 우선순위 항목',
         'RS-2 양 끝 운영 / RS-3 Co-Validation / RS-6 1c+NAND R&D / E 시나리오 피벗'),
    ]
    for i, (stage, headline, detail) in enumerate(flow_items):
        y = Inches(2.95 + i * 1.15)
        # 단계 라벨
        add_text(slide, Inches(7.2), y, Inches(0.8), Inches(0.4),
                 stage, font=FONT_KO, size=11, bold=True, color=THEME['amber'])
        # 헤드라인
        add_text(slide, Inches(8.0), y, Inches(4.7), Inches(0.4),
                 headline, font=FONT_KO, size=10.5, bold=True, color=THEME['white'])
        # 상세
        add_text(slide, Inches(8.0), y + Inches(0.4), Inches(4.7), Inches(0.7),
                 detail, font=FONT_KO, size=9, color=THEME['white'],
                 line_spacing=1.3)
        # 화살표
        if i < 2:
            add_text(slide, Inches(7.2), y + Inches(1.05), Inches(5.5), Inches(0.1),
                     '↓', font=FONT_EN, size=12, color=THEME['amber'], align='center')

    # ROI 강조
    add_rect(slide, Inches(7.2), Inches(6.10), Inches(5.5), Inches(0.35),
             fill=THEME['amber'])
    add_text(slide, Inches(7.3), Inches(6.13), Inches(5.3), Inches(0.3),
             'ROI 9~27배 — 다른 모든 RS의 prerequisite, 시나리오에 둔감',
             font=FONT_KO, size=9.5, bold=True, color=THEME['white'])

    add_so_what(slide, Inches(6.55),
                'AI 효율화 없이는 RS-2 양 끝, RS-3 락인, RS-6 3축 R&D가 인력 부담으로 무너진다.')
    add_footer(slide, 24, TOTAL, 'Robust · RS-7 · 신규')
    return slide


def build_slide_25_rs8_structured_hedging(prs):
    """단계 6-8: RS-8 구조화 매출 헷지 — 농수산·외환 헷지 100년 노하우의 메모리 첫 도입."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'RS-8 · 단순 LTA는 사이클당 50~100% 매출 기회를 포기한다',
               '농수산 100년 + 외환 50년 헷지 노하우의 메모리 첫 도입 — Participating Forward + HTA + Tiered + Trading Desk = 매출 변동성 ±25%→±12%')

    # 좌측 박스: 왜 지금 (3가지 증거)
    add_rect(slide, Inches(0.5), Inches(1.85), Inches(6.3), Inches(4.7),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(0.7), Inches(2.0), Inches(6), Inches(0.4),
             '왜 지금 — 단순 LTA의 한계', font=FONT_KO, size=12,
             bold=True, color=THEME['samsung_blue'])

    evidence = [
        ('① 297% 가격 점프 사례 (2025년 9~12월)',
         '16Gb DDR5 contract price $6.84 → $27.20 (3개월 +297%)',
         '9월 단순 고정가 LTA 시 297% 상방 모두 포기 — 고객만 차익 향유'),
        ('② 메모리 변동성 σ ≈ 60~120%',
         '원유 30%·S&P 18% 대비 2~6배 — "농산물 같은 IT 부품"',
         '농가도 작황의 30~50%만 forward 헷지 — 100% 락인은 표준 아님'),
        ('③ Participating Forward 메모리 미도입',
         '외환·통화시장 50년 표준 도구 — 메모리 산업 거의 도입 0건',
         '첫 도입자 표준 효과 = Fish Pool 모델, 후발주자 진입장벽 형성'),
    ]
    for i, (name, metric, source) in enumerate(evidence):
        y = Inches(2.55 + i * 1.40)
        add_rect(slide, Inches(0.85), y, Inches(0.12), Inches(1.20),
                 fill=THEME['amber'])
        add_text(slide, Inches(1.10), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 name, font=FONT_KO, size=11, bold=True, color=THEME['samsung_blue'])
        add_text(slide, Inches(1.10), y + Inches(0.45), Inches(5.5), Inches(0.35),
                 metric, font=FONT_KO, size=9.5, color=THEME['dark_text'])
        add_text(slide, Inches(1.10), y + Inches(0.83), Inches(5.5), Inches(0.35),
                 source, font=FONT_KO, size=8.5, italic=True,
                 color=THEME['gray_caption'])

    # 우측 박스: 4트랙 헷지 구조
    add_rect(slide, Inches(7.0), Inches(1.85), Inches(5.83), Inches(4.7),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.4),
             '4트랙 헷지 구조 (RS-4 LTA 위에 적층)', font=FONT_KO, size=12,
             bold=True, color=THEME['white'])
    add_text(slide, Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.4),
             'STRUCTURED REVENUE HEDGING', font=FONT_EN, size=8,
             color=THEME['amber'])

    tracks = [
        ('트랙 1', 'Participating Forward 시범',
         'Floor + 50% 상방 참여 / NVIDIA·하이퍼스케일러 1~2사 / HBM4E 일부 물량 (연 10~20%)'),
        ('트랙 2', 'Wafer Slot Reservation HTA',
         '슬롯/가격 분리 결정 + Take-or-Pay / Nvidia $5.4~7.7억 선급금 산업 표준화'),
        ('트랙 3', 'Tiered Pricing',
         '시장 구간별 차등 share / 자동차·산업용·신흥 AI 시장 (SD-2 연계)'),
        ('트랙 4', 'Memory Trading Desk',
         '골드만/JP모건 트레이더 20~30인 / DRAMeXchange 기반 OTC 스왑 시장 조성'),
    ]
    for i, (label, headline, detail) in enumerate(tracks):
        y = Inches(2.95 + i * 0.88)
        add_text(slide, Inches(7.2), y, Inches(0.8), Inches(0.4),
                 label, font=FONT_KO, size=11, bold=True, color=THEME['amber'])
        add_text(slide, Inches(8.0), y, Inches(4.7), Inches(0.4),
                 headline, font=FONT_KO, size=10.5, bold=True, color=THEME['white'])
        add_text(slide, Inches(8.0), y + Inches(0.38), Inches(4.7), Inches(0.5),
                 detail, font=FONT_KO, size=9, color=THEME['white'],
                 line_spacing=1.25)

    # 금지 사항 강조 (셰일 Three-way Collar 함정)
    add_rect(slide, Inches(7.2), Inches(6.10), Inches(5.5), Inches(0.35),
             fill=THEME['red_alert'])
    add_text(slide, Inches(7.3), Inches(6.13), Inches(5.3), Inches(0.3),
             '금지: Three-way Collar(sub-put 매도형) — 셰일 Pioneer 2014/2020 손실 가속',
             font=FONT_KO, size=9.5, bold=True, color=THEME['white'])

    add_so_what(slide, Inches(6.55),
                '계약 *물량*(RS-4)이 아니라 계약 *구조*(RS-8)를 혁신 — 매출 변동성 절반 + 사이클당 +50~100% 회수.')
    add_footer(slide, 25, TOTAL, 'Robust · RS-8 · 신규')
    return slide


def build_slide_26_dashboard(prs):
    """단계 7: EWI 대시보드 — 사분면 위치 변화 궤적 + 트리거."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'EWI 대시보드의 사분면 궤적과 12개 자동 트리거가 다운턴 신호를 30일 내 실행 전환으로 바꾼다',
               '지난 2년 궤적은 B 르네상스 모멘텀이지만 자동 실현은 아님 — 12개 결정 D-day와 12개 트리거가 신호를 즉시 행동으로 변환')

    # 좌측 (큰 영역) — 사분면 궤적 차트
    img = os.path.join(ASSETS_DIR, 'quadrant_trajectory.png')
    chart_quadrant_trajectory(img)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.85), width=Inches(8.4))

    # 우측 — 12개 결정 D-day + 12개 트리거 컴팩트
    # 12개 결정
    add_rect(slide, Inches(9.1), Inches(1.85), Inches(3.73), Inches(2.4),
             fill=THEME['soft_blue_bg'])
    add_text(slide, Inches(9.3), Inches(1.95), Inches(3.4), Inches(0.3),
             '12개 결정 D-day', font=FONT_KO, size=11, bold=True,
             color=THEME['samsung_blue'])
    decisions = [
        ('D1', 'HBM4 NVIDIA 점유율', 'D-150', True),
        ('D2', '소재 비중국 다각화', 'D-150', True),
        ('D6', 'RS1·RS4·RS6 정책화', 'D-150', True),
        ('D3', '3D DRAM + IMEC', 'D-240', False),
        ('D4', '텍사스 1·2단계', 'D-240', False),
        ('D5', 'AI 효율화', 'D-240', False),
        ('D8', '텍사스 추가 보조금', 'D-240', False),
        ('D9', 'M&A 펀드', 'D-240', False),
        ('D10', 'NAND 주기 연장 R&D', 'D-240', False),
        ('D11', 'Stargate Korea 본계약', 'D-240', False),
        ('D12', 'RS-8 구조화 헷지 시범', 'D-240', False),
        ('D7', '잉여 인력 전환', 'D-330', False),
    ]
    for i, (d, name, dday, urgent) in enumerate(decisions):
        y = Inches(2.30 + i * 0.16)
        add_text(slide, Inches(9.3), y, Inches(0.5), Inches(0.18),
                 d, font=FONT_EN, size=8, bold=True, color=THEME['amber'])
        add_text(slide, Inches(9.7), y, Inches(2.3), Inches(0.18),
                 name, font=FONT_KO, size=7.5, color=THEME['dark_text'])
        ddc = THEME['red_alert'] if urgent else THEME['gray_caption']
        add_text(slide, Inches(12.0), y, Inches(0.7), Inches(0.18),
                 dday, font=FONT_EN, size=7.5, bold=True, color=ddc, align='right')

    # 12개 트리거 — 카테고리별 압축
    add_rect(slide, Inches(9.1), Inches(4.35), Inches(3.73), Inches(2.2),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(9.3), Inches(4.45), Inches(3.4), Inches(0.3),
             '12개 자동 트리거', font=FONT_KO, size=11, bold=True,
             color=THEME['white'])
    add_text(slide, Inches(9.3), Inches(4.7), Inches(3.4), Inches(0.25),
             '발동 시 30일 내 전략 전환', font=FONT_KO, size=8,
             italic=True, color=THEME['amber'])

    trigger_groups = [
        ('다운턴 신호 (5)', '빅테크 CapEx -25% / MS CapEx -20% / HBM 가격 -30% / CXMT HBM4 / Hsinx 라이선스'),
        ('지정학 (2)', 'MATCH 통과 / 텍사스 2단계 미발표'),
        ('가속 신호 (2)', 'AI 매출 $500B / 3D DRAM 50% PoC'),
        ('전략 위기 (3)', 'HBM4 미인증 / 단일 고객 >25% / M&A 타깃 EV<5x'),
    ]
    for i, (cat, items) in enumerate(trigger_groups):
        y = Inches(5.0 + i * 0.38)
        add_text(slide, Inches(9.3), y, Inches(3.4), Inches(0.2),
                 cat, font=FONT_KO, size=8.5, bold=True, color=THEME['amber'])
        add_text(slide, Inches(9.3), y + Inches(0.18), Inches(3.4), Inches(0.2),
                 items, font=FONT_KO, size=7, color=THEME['white'])

    add_so_what(slide, Inches(6.85),
                '현재 위치 (DF1 +7.5, DF2 +0.5) — B 시나리오 강한 모멘텀, 그러나 자동 실현 아님')
    add_footer(slide, 26, TOTAL, '대시보드')
    return slide


def build_slide_27_decisions_summary(prs):
    """단계 8-1: 12개 결정 한 페이지 요약."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               '12개 결정은 묶음이다 — 분리해서 처리하면 효과가 사라진다',
               'D1·D5·D7 직렬 / D4·D8 동시 / D6이 D9·D12 거버넌스 기반 / D10·D11·D12 신성장축 — 모든 결정 시한 2026 Q4')

    decisions = [
        {
            'id': 'D1', 'urgent': True,
            'title': 'HBM4 NVIDIA 점유율',
            'spec': '수율 90%+ Q3 2026 / 점유율 28% → 40%+',
            'owner': 'D-150 · 마감 임박',
        },
        {
            'id': 'D2', 'urgent': False,
            'title': '소재 비중국 다각화',
            'spec': 'Q3 2026 LTA 체결 / 6개월 비축',
            'owner': '캐나다·호주·카자흐',
        },
        {
            'id': 'D3', 'urgent': False,
            'title': '3D DRAM R&D + IMEC',
            'spec': 'Q4 2026 협약 / $200M / 200~300인',
            'owner': '마이크론·SK 추격 차단',
        },
        {
            'id': 'D4', 'urgent': False,
            'title': '텍사스 1·2단계',
            'spec': '1단계 가동 + 2단계 발표 / Q4 2026',
            'owner': '$4.745B 보조금 활용',
        },
        {
            'id': 'D5', 'urgent': False,
            'title': 'AI 효율화 전사 도입',
            'spec': 'Q4 파일럿 → Q1 전사 / RS2·RS3 선행',
            'owner': '500~1,000억 원',
        },
        {
            'id': 'D6', 'urgent': False,
            'title': 'RS1·RS4·RS6 정책화',
            'spec': '다운사이클 capex 4조원 명문화',
            'owner': '활동가 투자자 방어',
        },
        {
            'id': 'D7', 'urgent': False,
            'title': '잉여 인력 전환 배치',
            'spec': '2027 Q1 가시화 / 분기별 사업부장 보고',
            'owner': '추가 채용 없이 RS2·RS3',
        },
        {
            'id': 'D8', 'urgent': False,
            'title': '텍사스 추가 보조금',
            'spec': 'CHIPS 추가 협상 + LTA / Tesla 외 빅테크',
            'owner': 'JV 모델 대안',
        },
        {
            'id': 'D9', 'urgent': False,
            'title': '다운사이클 M&A 펀드',
            'spec': '5,000억 원 사전 적립 / 자동 트리거',
            'owner': 'Disney-Marvel 모델',
        },
        {
            'id': 'D10', 'urgent': False,
            'title': 'NAND 주기 연장 R&D 4트랙',
            'spec': 'Hybrid Bonding 자체 IP / Multi-deck / QLC·PLC / FDP·SCADA firmware',
            'owner': '1.5조 원/3년 · YMTC 우회',
        },
        {
            'id': 'D11', 'urgent': False,
            'title': 'Stargate Korea 본 계약',
            'spec': 'LOI → DA / Electronics+C&T+SDS+Heavy 컨소시엄 / Neocloud 1~3% 지분',
            'owner': 'AI 가치사슬 마진 헷지',
        },
        {
            'id': 'D12', 'urgent': False,
            'title': 'RS-8 구조화 매출 헷지 시범',
            'spec': 'Participating Forward + HTA + Tiered + Trading Desk / Three-way 금지',
            'owner': '농수산·외환 100년 헷지 이식',
        },
    ]

    # 3x4 grid (3 columns × 4 rows)
    row_height = 1.18
    for i, d in enumerate(decisions):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.28)
        y = Inches(1.80 + row * row_height)
        bg = RGBColor(0xFE, 0xF2, 0xF2) if d['urgent'] else THEME['soft_blue_bg']
        border_color = THEME['red_alert'] if d['urgent'] else THEME['samsung_blue']
        add_rect(slide, x, y, Inches(4.1), Inches(row_height - 0.13),
                 fill=bg, line=border_color, line_width=Emu(15875) if d['urgent'] else None)
        # 좌측 띠
        add_rect(slide, x, y, Inches(0.1), Inches(row_height - 0.13),
                 fill=THEME['red_alert'] if d['urgent'] else THEME['samsung_blue'])
        # ID
        add_text(slide, x + Inches(0.20), y + Inches(0.05), Inches(0.85), Inches(0.35),
                 d['id'], font=FONT_EN, size=15, bold=True,
                 color=THEME['red_alert'] if d['urgent'] else THEME['samsung_blue'])
        # 제목
        add_text(slide, x + Inches(1.05), y + Inches(0.10), Inches(2.95), Inches(0.30),
                 d['title'], font=FONT_KO, size=10.5, bold=True, color=THEME['dark_text'])
        # 스펙
        add_text(slide, x + Inches(0.22), y + Inches(0.45), Inches(3.78), Inches(0.42),
                 d['spec'], font=FONT_KO, size=8.5, color=THEME['dark_text'],
                 line_spacing=1.2)
        # Owner / context
        add_text(slide, x + Inches(0.22), y + Inches(0.83), Inches(3.78), Inches(0.22),
                 d['owner'], font=FONT_KO, size=8, italic=True,
                 color=THEME['amber'])

    # 하단 강조 박스 (4행 그리드 종료점 = 1.80 + 4×1.18 = 6.52)
    add_rect(slide, Inches(0.5), Inches(6.58), Inches(12.33), Inches(0.45),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(0.7), Inches(6.63), Inches(12), Inches(0.35),
             '12개는 묶음이다 — D1·D5·D7 직렬 / D4·D8 동시 / D6이 D9·D12 거버넌스 / D10·D11·D12 신성장축',
             font=FONT_KO, size=10.5, bold=True, color=THEME['white'], align='center')

    add_footer(slide, 27, TOTAL, '결정 요약')
    return slide


def build_slide_28_problem_solution(prs):
    """단계 8-2: 핵심 전략 ↔ 해결 문제 매핑 — 3열 레이아웃."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    add_header(slide,
               'Main·Side·Robust 9개 핵심 전략의 묶음이 다운턴 흑자·1위 회복·자산 폭락 시 M&A를 동시에 해결한다',
               '각 영역에서 가장 중요한 3개씩 — 단일 전략으로는 한 문제만 풀리지만, 묶음으로 실행하면 모든 문제가 동시에 해결된다')

    # 3열 헤더
    column_titles = [
        ('MAIN BET', '시나리오 B 베팅', THEME['amber']),
        ('SIDE BETS', '나머지 시나리오 보험', THEME['samsung_blue']),
        ('ROBUST', '시나리오 무관 가치', THEME['green_pos']),
    ]
    for i, (label, sub, color) in enumerate(column_titles):
        x = Inches(0.5 + i * 4.28)
        y = Inches(1.85)
        add_rect(slide, x, y, Inches(4.1), Inches(0.55), fill=color)
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.9), Inches(0.3),
                 label, font=FONT_EN, size=12, bold=True, color=THEME['white'])
        add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.9), Inches(0.25),
                 sub, font=FONT_KO, size=10, italic=True, color=THEME['white'])

    # Main Bet — 가장 중요한 3개
    main_bets = [
        ('HBM4 점유율 회복', 'NVIDIA Rubin 28% → 40%+',
         '시나리오 B에서 1위 자리 미확보 시 Main Bet 무의미'),
        ('동서 균형 공급망', '서방 + 신흥 + 중국 3원',
         '한 지역 의존이 다운턴·지정학 충격에 취약'),
        ('1c nm 원가 우위', 'CXMT 공세 방어 기반',
         '범용 DRAM 마진 잠식이 다운턴 손실 키움'),
    ]
    # Side Bet — 가장 중요한 3개
    side_bets = [
        ('순현금 30조 버퍼', '시나리오 C 생존 기반',
         '다운턴 자금 부족 시 R&D 삭감 강요됨'),
        ('소재 공급망 다각화', '캐나다·호주·카자흐 LTA',
         '중국 소재 수출 통제 시 원가 30% 급등'),
        ('3D DRAM + IMEC 협약', '시나리오 E 헤징',
         'HBM 패러다임 붕괴 시 40조 매몰 비용'),
    ]
    # Robust — 가장 중요한 3개
    robust = [
        ('옵션형 캐파', '켜고 끌 수 있는 능력',
         '잘못된 타이밍 고정 증설이 다운턴 손실 핵심'),
        ('바벨 + AI 효율화', '양 끝 + 잉여 인력 전환',
         '추가 채용 없이 RS3 동시 실행의 유일 경로'),
        ('재무 규율 + capex 하한', '호황 절제 + 다운턴 사수',
         '활동가 투자자 압박에 R&D 무너지면 회복 불가'),
    ]

    columns_data = [main_bets, side_bets, robust]
    column_colors = [THEME['amber'], THEME['samsung_blue'], THEME['green_pos']]

    for col_idx, items in enumerate(columns_data):
        x = Inches(0.5 + col_idx * 4.28)
        color = column_colors[col_idx]
        for j, (title, desc, problem) in enumerate(items):
            y = Inches(2.5 + j * 1.45)
            # 카드 배경
            add_rect(slide, x, y, Inches(4.1), Inches(1.35),
                     fill=THEME['soft_white_bg'], line=THEME['light_gray'])
            # 좌측 띠
            add_rect(slide, x, y, Inches(0.1), Inches(1.35), fill=color)
            # 번호
            add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.4), Inches(0.4),
                     f'#{j+1}', font=FONT_EN, size=11, bold=True, color=color)
            # 핵심 3-4단어
            add_text(slide, x + Inches(0.6), y + Inches(0.1), Inches(3.4), Inches(0.4),
                     title, font=FONT_KO, size=12, bold=True, color=THEME['dark_text'])
            # 한 줄 설명
            add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.8), Inches(0.3),
                     desc, font=FONT_KO, size=9, italic=True, color=color)
            # 해결하는 문제
            add_rect(slide, x + Inches(0.15), y + Inches(0.85), Inches(3.85), Inches(0.45),
                     fill=THEME['amber_light'])
            add_text(slide, x + Inches(0.25), y + Inches(0.88), Inches(3.7), Inches(0.4),
                     '문제: ' + problem, font=FONT_KO, size=8.5,
                     color=THEME['dark_text'], line_spacing=1.2)

    # 하단 강조 박스 — 모든 문제 해결
    add_rect(slide, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.45),
             fill=THEME['samsung_blue'])
    add_text(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.35),
             '9개 핵심 전략의 묶음 = 다운턴이 와도 흑자 유지 + 회복기 1번 자리 + 자산 폭락 시 M&A',
             font=FONT_KO, size=11, bold=True, color=THEME['white'], align='center')

    add_footer(slide, 28, TOTAL, '전략 리마인드')
    return slide


def build_slide_29_closing(prs):
    """단계 9: 최종 메시지."""
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    # 좌측 풀 사이드바
    add_rect(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5),
             fill=THEME['samsung_blue'])
    add_rect(slide, Inches(0), Inches(1.2), Inches(0.2), Inches(0.7),
             fill=THEME['amber'])

    # DECISION REQUEST
    add_text(slide, Inches(1.0), Inches(0.7), Inches(11), Inches(0.4),
             'DECISION REQUEST', font=FONT_EN, size=12, bold=True,
             color=THEME['samsung_blue'])
    add_text(slide, Inches(1.0), Inches(1.1), Inches(11), Inches(0.5),
             '12개 결정의 묶음 의결로 호황의 정점을 다운턴 준비의 마지막 기회로 잡는다',
             font=FONT_KO, size=20, bold=True, color=THEME['dark_text'])

    # Amber 짧은 라인
    add_rect(slide, Inches(1.0), Inches(1.85), Inches(0.6), Inches(0.05),
             fill=THEME['amber'])

    # 큰 인용
    add_text(slide, Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.5),
             '"호황의 정점이\n 다운턴 준비의 마지막 기회다."',
             font=FONT_KO, size=36, bold=True, color=THEME['samsung_blue'],
             line_spacing=1.2)

    add_text(slide, Inches(1.0), Inches(4.0), Inches(11.5), Inches(0.4),
             '─ 12개 결정의 묶음 의결을 요청드립니다 ─',
             font=FONT_KO, size=14, italic=True, color=THEME['gray_caption'])

    # 3개 결의 요청 카드
    asks = [
        ('01', '전략 방향 승인',
         '시나리오 B Main Bet + Side Bets + Robust 8개 전략 채택\n2030~2035 메모리 1위 회복 로드맵 의결'),
        ('02', '12개 즉시 결정 패키지',
         'HBM4·소재·3D DRAM·텍사스·AI 효율화·정책화·잉여 인력·\n추가 보조금·M&A 펀드·NAND 4트랙·Stargate Korea·RS-8 헷지 — 시한: 2026 Q4'),
        ('03', '거버넌스 + 트리거 승인',
         '30일 내 의사결정 + 자동 트리거 12개 + 다운사이클 capex 하한 4조 원/년\n+ 활동가 투자자 방어 의결권 구조 + 단일 Floor 원칙(RS-8)'),
    ]
    for i, (num, title, desc) in enumerate(asks):
        x = Inches(1.0 + i * 4.0)
        y = Inches(4.7)
        add_rect(slide, x, y, Inches(3.85), Inches(2.0),
                 fill=THEME['soft_blue_bg'])
        add_rect(slide, x, y, Inches(0.08), Inches(2.0), fill=THEME['amber'])
        add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(1), Inches(0.5),
                 num, font=FONT_EN, size=22, bold=True, color=THEME['amber'])
        add_text(slide, x + Inches(0.25), y + Inches(0.65), Inches(3.5), Inches(0.4),
                 title, font=FONT_KO, size=12, bold=True, color=THEME['dark_text'])
        add_text(slide, x + Inches(0.25), y + Inches(1.05), Inches(3.5), Inches(0.9),
                 desc, font=FONT_KO, size=9.5, color=THEME['dark_text'],
                 line_spacing=1.3)

    # 푸터
    add_text(slide, Inches(1.0), Inches(7.05), Inches(11.5), Inches(0.3),
             '삼성전자  ·  DS부문 메모리사업부 전략기획팀  ·  대외비  ·  2026.05.06',
             font=FONT_KO, size=10, color=THEME['gray_caption'])
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1), Inches(0.3),
             '29 / 29', font=FONT_EN, size=10, color=THEME['gray_caption'],
             align='right')

    return slide


# ============================================================
# Main
# ============================================================

def remove_existing_slides(prs):
    """템플릿의 모든 슬라이드 제거 — 마스터/테마는 유지."""
    sldIdLst = prs.slides._sldIdLst
    rIds = []
    for sldId in list(sldIdLst):
        rIds.append(sldId.get(qn('r:id')))
        sldIdLst.remove(sldId)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def main():
    print(f'Loading template (theme only): {TEMPLATE}')
    prs = Presentation(TEMPLATE)
    remove_existing_slides(prs)
    print(f'After cleanup — slides: {len(prs.slides)}')

    builders = [
        build_slide_1_cover,
        build_slide_2_toc,
        build_slide_3_peak,
        build_slide_4_history,
        build_slide_5_samsung_quarterly,
        build_slide_6_paradox,
        build_slide_7_uncertainty,
        build_slide_8_single_prediction,
        build_slide_9_benchmark,
        build_slide_10_workflow,
        build_slide_11_focal_issue,
        build_slide_12_steep_iu,
        build_slide_13_driving_forces,
        build_slide_14_scenario_matrix,
        build_slide_15_main_bet,
        build_slide_16_side_bets,
        build_slide_17_robust_overview,
        build_slide_18_rs1,
        build_slide_19_rs2,
        build_slide_20_rs3,
        build_slide_21_rs4_diversification,
        build_slide_22_rs5_financial,
        build_slide_23_rs6_process,
        build_slide_24_rs7_ai_automation,
        build_slide_25_rs8_structured_hedging,
        build_slide_26_dashboard,
        build_slide_27_decisions_summary,
        build_slide_28_problem_solution,
        build_slide_29_closing,
    ]

    for i, builder in enumerate(builders, start=1):
        print(f'  Building slide {i}: {builder.__name__}...')
        builder(prs)

    print(f'Saving to: {OUTPUT}')
    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT) / 1024
    print(f'Done. Output size: {size:.1f} KB · {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
