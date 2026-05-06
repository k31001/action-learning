"""
삼성전자 메모리사업부 시나리오 플래닝 — Figma 하이브리드 발표자료

10-slide hero deck:
  Slides 1~3: Figma에서 디자인한 PNG를 full-bleed 배경으로 사용
              (presentation/assets/figma/slide-01..03.png)
  Slides 4~10: python-pptx로 동일 디자인 시스템(딥 네이비 + 앰버) 네이티브 렌더링
              (Figma Starter 플랜 rate limit으로 PNG export 불가)

Output: presentation/samsung-memory-scenario-planning-figma.pptx (16:9, 13.33×7.5 inch)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
OUTPUT = os.path.join(ROOT, 'presentation', 'samsung-memory-scenario-planning-figma.pptx')
FIGMA_DIR = os.path.join(ROOT, 'presentation', 'assets', 'figma')

W_IN, H_IN = 13.333, 7.5

# Design system — matches Figma file
NAVY = RGBColor(0x0A, 0x1B, 0x5C)
NAVY2 = RGBColor(0x14, 0x28, 0xA0)
NAVY3 = RGBColor(0x1E, 0x3A, 0xC8)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
PANEL = RGBColor(0xE8, 0xEE, 0xFC)
BORDER = RGBColor(0xCB, 0xD5, 0xE1)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xC0, 0x00, 0x00)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)


# ─── helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.0, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w if line_w else 0.75)
    else:
        shape.line.fill.background()
    if radius is not None:
        # rounded corner adjust value (0~0.5)
        try:
            shape.adjustments[0] = min(max(radius, 0.0), 0.5)
        except Exception:
            pass
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=TEXT,
             font='Arial Black' if False else 'Calibri',
             font_face=None,
             align='left', valign='top',
             line_spacing=None, char_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_face or font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if char_spacing:
        from pptx.oxml.ns import qn as _qn
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(char_spacing * 100)))
    if valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    return box


def add_oval(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


# =============================================================================
# Slides 1~3: Figma full-bleed PNG
# =============================================================================
def add_figma_slide(prs, png_path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Full-bleed image
    slide.shapes.add_picture(png_path, 0, 0, width=Inches(W_IN), height=Inches(H_IN))
    return slide


# =============================================================================
# Slide 4 — CYCLE HISTORY (dramatic line chart with downturns)
# =============================================================================
def build_slide_4_cycle(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, NAVY)
    add_rect(slide, 0, 0, W_IN, 0.08, AMBER)

    add_text(slide, 0.55, 0.55, 8, 0.3, "01  ·  CYCLE HISTORY",
             size=12, bold=True, color=AMBER, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12, 0.7,
             "메모리는 30년간 5번의 다운턴을 겪은 사이클 산업이다",
             size=32, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, 0.55, 1.40, 12, 0.4,
             "1995~2026  ·  4~5년 주기로 100%+ 등락 반복  ·  다운턴은 항상 호황 직후",
             size=14, color=BORDER)

    # Chart frame
    cx, cy, cw, ch = 0.55, 1.95, 12.2, 3.8
    add_rect(slide, cx, cy, cw, ch, NAVY2, line=NAVY3)

    # Downturn shaded regions
    downturns = [
        ('1996~98', 0.05, 0.07, '−65%'),
        ('2001~02', 0.20, 0.05, '−55%'),
        ('2008~09', 0.42, 0.05, '−50%'),
        ('2015~16', 0.65, 0.05, '−40%'),
        ('2022~23', 0.85, 0.05, '−45%'),
    ]
    for name, t, w, peak in downturns:
        rx = cx + t * cw
        rw = w * cw
        s = add_rect(slide, rx, cy + 0.2, rw, ch - 0.4, RED)
        s.fill.transparency = 0  # alpha not directly settable easily, use overlay alt
        try:
            s.fill.fore_color.rgb = RGBColor(0xFC, 0xA5, 0xA5)
        except Exception:
            pass
        add_text(slide, rx + 0.02, cy + 0.25, rw + 0.5, 0.25,
                 name, size=9, bold=True, color=AMBER)
        add_text(slide, rx + 0.02, cy + 0.45, rw + 0.5, 0.25,
                 peak, size=11, bold=True, color=AMBER)

    # Synthesized cycle line via connectors
    points = [
        (0.00, 0.55), (0.04, 0.85), (0.08, 0.30), (0.14, 0.60), (0.18, 0.90),
        (0.22, 0.35), (0.30, 0.55), (0.40, 0.95), (0.45, 0.42), (0.55, 0.65),
        (0.60, 0.92), (0.66, 0.50), (0.74, 0.65), (0.82, 0.90), (0.87, 0.45),
        (0.94, 0.85), (1.00, 0.10),
    ]
    prev = None
    for tp, vp in points:
        x = cx + tp * cw
        y = cy + 0.3 + (1 - vp) * (ch - 0.6)
        if prev is not None:
            conn = slide.shapes.add_connector(1, Inches(prev[0]), Inches(prev[1]),
                                              Inches(x), Inches(y))
            conn.line.color.rgb = AMBER
            conn.line.width = Pt(2.5)
        prev = (x, y)

    # Highlight current point
    lx, ly = prev
    add_oval(slide, lx - 0.12, ly - 0.12, 0.24, 0.24, AMBER, line=WHITE)
    add_text(slide, lx - 0.7, ly - 0.45, 1.4, 0.25,
             "현재  Q1 2026", size=10, bold=True, color=AMBER)
    add_text(slide, lx - 0.55, ly - 0.22, 1.1, 0.22,
             "사상 최고점", size=9, bold=True, color=WHITE)

    # X-axis labels
    for i, yr in enumerate(['1995','2000','2005','2010','2015','2020','2025']):
        add_text(slide, cx + (i / 6) * cw - 0.2, cy + ch - 0.25, 0.5, 0.2,
                 yr, size=9, color=BORDER, font_face='Calibri')

    # Bottom takeaway
    add_text(slide, 0.55, 6.05, 0.5, 0.5, "✱", size=24, bold=True, color=AMBER)
    add_text(slide, 1.0, 6.05, 11.5, 0.4,
             "호황기 평균 18~24개월  ·  다운턴 평균 12~18개월",
             size=18, bold=True, color=WHITE)
    add_text(slide, 0.55, 6.55, 12.2, 0.5,
             "다운턴마다 우리 매출은 −38%, 영업이익은 −84% 사라졌다 (분기 −10조 원 적자 포함)",
             size=14, color=BORDER)
    add_text(slide, 12.4, 7.15, 0.8, 0.25, "04 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 5 — SCENARIO MATRIX (2×2 quadrant)
# =============================================================================
def build_slide_5_matrix(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, LIGHT)
    add_rect(slide, 0, 0, W_IN, 0.08, NAVY)

    add_text(slide, 0.55, 0.55, 8, 0.3, "02  ·  SCENARIO MATRIX",
             size=12, bold=True, color=AMBER, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12, 0.7,
             "두 축의 조합이 5개의 대안적 미래를 만든다",
             size=32, bold=True, color=NAVY, font_face='Arial Black')
    add_text(slide, 0.55, 1.40, 12, 0.4,
             "DF1 (AI 수요)  ×  DF2 (미중 디커플링)  +  DF3 와일드카드",
             size=14, color=TEXT_GRAY)

    # 2×2 matrix area
    mx, my, mw, mh = 1.95, 1.95, 8.0, 4.7
    half_w, half_h = mw/2, mh/2
    # Quadrant tints
    add_rect(slide, mx,         my,         half_w, half_h, RGBColor(0xFE, 0xE2, 0xE2))  # C
    add_rect(slide, mx+half_w,  my,         half_w, half_h, RGBColor(0xDB, 0xEA, 0xFE))  # A
    add_rect(slide, mx,         my+half_h,  half_w, half_h, RGBColor(0xFE, 0xF3, 0xC7))  # D
    add_rect(slide, mx+half_w,  my+half_h,  half_w, half_h, RGBColor(0xDC, 0xFC, 0xE7))  # B

    # Frame
    add_rect(slide, mx, my, mw, mh, None, line=BORDER, line_w=1)

    # Cross axes
    add_rect(slide, mx + half_w - 0.005, my + 0.15, 0.01, mh - 0.3, TEXT)
    add_rect(slide, mx + 0.15, my + half_h - 0.005, mw - 0.3, 0.01, TEXT)

    # Axis labels
    add_text(slide, mx + half_w + 1.5, my + 0.10, 1.5, 0.25,
             "디커플링 심화", size=10, bold=True, color=TEXT)
    add_text(slide, mx + half_w + 1.5, my + mh - 0.30, 1.5, 0.25,
             "관리된 공존", size=10, bold=True, color=TEXT)
    add_text(slide, mx + 0.15, my + half_h + 0.10, 1.4, 0.25,
             "AI 거품 붕괴", size=10, bold=True, color=TEXT)
    add_text(slide, mx + mw - 1.5, my + half_h + 0.10, 1.4, 0.25,
             "AI 수요 지속", size=10, bold=True, color=TEXT, align='right')

    # Quadrant scenarios
    scenarios = [
        ('C', '기술 냉전',     '10~15%', mx + 0.20,        my + 0.40,        RED,    '이중 충격', False),
        ('A', '황금 요새',     '25~30%', mx + half_w + 0.20, my + 0.40,      NAVY2,  '서방 듀오폴리', False),
        ('D', '조용한 재편',   '20~25%', mx + 0.20,        my + half_h + 0.40, AMBER, '사이클 재현', False),
        ('B', 'AI 르네상스',   '30~35%', mx + half_w + 0.20, my + half_h + 0.40, GREEN, '★ MAIN BET ★', True),
    ]
    for sid, name, prob, sx, sy, color, bullet, main in scenarios:
        circle = 0.55 if main else 0.42
        add_oval(slide, sx + 0.15, sy + 0.10, circle, circle, color)
        add_text(slide, sx + 0.15, sy + (0.16 if main else 0.18),
                 circle, circle - 0.1,
                 sid, size=24 if main else 18, bold=True,
                 color=WHITE, align='center', font_face='Arial Black')
        # Probability tag
        tag_x = sx + circle + 0.30
        add_rect(slide, tag_x, sy + 0.14, 1.0, 0.32, color, radius=0.5)
        add_text(slide, tag_x, sy + 0.18, 1.0, 0.28,
                 prob, size=10, bold=True, color=WHITE, align='center')
        # Name
        add_text(slide, sx + 0.10, sy + circle + 0.25, 3.5, 0.4,
                 name, size=18 if main else 16, bold=True, color=TEXT,
                 font_face='Arial Black')
        # Bullet
        add_text(slide, sx + 0.10, sy + circle + 0.65, 3.5, 0.3,
                 bullet, size=12, bold=True, color=color)

    # Wildcard E circle in center
    add_oval(slide, mx + half_w - 0.40, my + half_h - 0.40, 0.80, 0.80, PURPLE)
    add_text(slide, mx + half_w - 0.30, my + half_h - 0.20, 0.60, 0.40,
             "E", size=22, bold=True, color=WHITE, align='center',
             font_face='Arial Black')
    add_text(slide, mx + half_w - 1.0, my + half_h + 0.45, 2.0, 0.25,
             "패러다임 전환  5~10%  🃏", size=10, bold=True, color=PURPLE, align='center')

    # Right legend
    lx = 10.10
    add_text(slide, lx, 1.95, 3, 0.4, "5개 대안적 미래",
             size=16, bold=True, color=NAVY, font_face='Arial Black')
    legend = [
        ("A", "황금 요새", "AI 지속 + 디커플링", NAVY2),
        ("B", "AI 르네상스 ⭐", "AI 지속 + 공존", GREEN),
        ("C", "기술 냉전", "AI 붕괴 + 디커플링", RED),
        ("D", "조용한 재편", "AI 붕괴 + 공존", AMBER),
        ("E", "패러다임 전환", "HBM 대체 (와일드카드)", PURPLE),
    ]
    for i, (sid, nm, desc, col) in enumerate(legend):
        ly = 2.55 + i * 0.65
        add_oval(slide, lx, ly + 0.05, 0.30, 0.30, col)
        add_text(slide, lx + 0.03, ly + 0.10, 0.30, 0.25,
                 sid, size=12, bold=True, color=WHITE, align='center')
        add_text(slide, lx + 0.45, ly, 2.7, 0.30,
                 nm, size=12, bold=True, color=TEXT)
        add_text(slide, lx + 0.45, ly + 0.30, 2.7, 0.30,
                 desc, size=10, color=TEXT_GRAY)

    add_text(slide, 12.4, 7.15, 0.8, 0.25, "05 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 6 — MAIN BET (Scenario B + 5 initiatives)
# =============================================================================
def build_slide_6_mainbet(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, WHITE)
    add_rect(slide, 0, 0, W_IN, 0.08, GREEN)

    add_text(slide, 0.55, 0.55, 8, 0.3, "03  ·  MAIN BET  ·  SCENARIO B",
             size=12, bold=True, color=GREEN, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12, 0.7,
             "Main Bet은 시나리오 B 「AI 르네상스」",
             size=32, bold=True, color=NAVY, font_face='Arial Black')
    add_text(slide, 0.55, 1.40, 12, 0.4,
             "확률 30~35%  ·  자동 실현되지 않는다 — 9개의 즉시 결정 묶음이 전제 조건",
             size=14, color=TEXT_GRAY)

    # Hero block
    hx, hy, hw, hh = 0.55, 1.95, 4.2, 5.0
    add_rect(slide, hx, hy, hw, hh, NAVY, radius=0.05)
    add_rect(slide, hx, hy, hw, 0.06, AMBER)

    add_text(slide, hx + 0.30, hy + 0.30, 3.5, 0.30,
             "SCENARIO", size=10, bold=True, color=AMBER, char_spacing=4)
    add_text(slide, hx + 0.20, hy + 0.55, 3.8, 2.5,
             "B", size=180, bold=True, color=AMBER,
             font_face='Arial Black', align='left')

    # Probability bar
    add_rect(slide, hx + 0.30, hy + 2.85, hw - 0.6, 0.10, NAVY3, radius=0.5)
    add_rect(slide, hx + 0.30, hy + 2.85, (hw - 0.6) * 0.325, 0.10, AMBER, radius=0.5)
    add_text(slide, hx + 0.30, hy + 3.05, hw - 0.6, 0.6,
             "30~35%", size=44, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, hx + 0.30, hy + 3.65, hw - 0.6, 0.30,
             "5개 시나리오 중 1순위 확률", size=12, color=BORDER)

    add_rect(slide, hx + 0.30, hy + 4.05, 0.6, 0.04, AMBER)
    add_text(slide, hx + 0.30, hy + 4.15, hw - 0.6, 0.40,
             "AI 르네상스", size=24, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, hx + 0.30, hy + 4.55, hw - 0.6, 0.30,
             "AI 수요 지속 + 미중 관리된 공존", size=12, color=BORDER)

    # Right side — 5 initiatives
    ix, iy = 5.0, 1.95
    add_text(slide, ix, iy, 8, 0.35,
             "5대 핵심 이니셔티브 (MB-1 ~ MB-5)",
             size=18, bold=True, color=NAVY, font_face='Arial Black')
    add_text(slide, ix, iy + 0.32, 8, 0.30,
             "2030~2035 글로벌 1위 회복 로드맵 — 단일 결정이 아닌 결정의 묶음",
             size=11, color=TEXT_GRAY, italic=True)

    inits = [
        ('MB-1', 'HBM4E·HBM5 기술 1위 탈환', 'Q3 2026 NVIDIA Rubin 공급 점유 30%+ → HBM5 리더십 2027년 확보', NAVY2),
        ('MB-2', '동서 균형 공급망 4거점', '한국·미국·일본·신흥시장 — 디커플링·공존 모두 대응', GREEN),
        ('MB-3', '1c nm 공정 전환 + 원가 우위', '2027년 수율 80%+ 도달 + 비트당 원가 −25%', AMBER),
        ('MB-4', '커스텀 AI 메모리 솔루션', '하이퍼스케일러 co-design, FDP/SCADA 호스트 협력', PURPLE),
        ('MB-5', '텍사스 테일러 2단계 발표', 'CHIPS Act 2.0 추가 보조금 협상 + 2030년 가동', RED),
    ]
    cw, ch = 7.7, 0.85
    gap = 0.10
    for i, (num, title, body, col) in enumerate(inits):
        y = iy + 0.75 + i * (ch + gap)
        add_rect(slide, ix, y, cw, ch, LIGHT, radius=0.05)
        add_rect(slide, ix, y, 0.08, ch, col)
        # Number tag
        add_rect(slide, ix + 0.20, y + 0.15, 0.65, 0.30, col, radius=0.1)
        add_text(slide, ix + 0.20, y + 0.18, 0.65, 0.25,
                 num, size=11, bold=True, color=WHITE, align='center')
        add_text(slide, ix + 1.0, y + 0.13, cw - 1.1, 0.32,
                 title, size=14, bold=True, color=TEXT)
        add_text(slide, ix + 1.0, y + 0.45, cw - 1.1, 0.40,
                 body, size=10, color=TEXT_GRAY)

    add_text(slide, 12.4, 7.15, 0.8, 0.25, "06 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 7 — ROBUST 7 STRATEGIES (4×2 grid, 7 cards + 1 explanation)
# =============================================================================
def build_slide_7_robust(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, LIGHT)
    add_rect(slide, 0, 0, W_IN, 0.08, AMBER)

    add_text(slide, 0.55, 0.55, 8, 0.3, "04  ·  ROBUST STRATEGIES",
             size=12, bold=True, color=AMBER, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12, 0.7,
             "어떤 미래가 와도 가치를 만드는 Robust 7개",
             size=32, bold=True, color=NAVY, font_face='Arial Black')
    add_text(slide, 0.55, 1.40, 12, 0.4,
             "5개 시나리오 모두 ✅ 4개+ — 시나리오 불확실성과 무관하게 즉시 실행",
             size=14, color=TEXT_GRAY)

    items = [
        ('RS1', '옵션형 캐파 체계',           '"켜고 끌 수 있는 능력" — Nucor 미니밀',     NAVY2,  True,  False),
        ('RS2', '바벨 포트폴리오 + AI 효율화',  '양 끝의 강함 — 중간 어정쩡 제거',         NAVY2,  True,  False),
        ('RS3', '고객특화·전환비용 극대화',     'Marriott Asset-Light SW $5B/2030',     NAVY2,  True,  False),
        ('RS4', '고객 포트폴리오 분산',         'LTA·Take-or-Pay 협상력 tier화',          AMBER,  False, False),
        ('RS5', '정책 리스크 지역 분산',        '한국·미국·일본·인도 4거점',                AMBER,  False, False),
        ('RS6', '재무 규율 + capex 하한',      'Nucor·ExxonMobil — R&D 4조원/년',       NAVY2,  True,  False),
        ('RS7', 'NAND 공정 전환 주기 연장 R&D','신규 — Hybrid Bonding 자체 IP, 4 트랙', GREEN,  True,  True),
    ]
    cw, ch = 2.95, 2.10
    sx, sy = 0.55, 1.95
    gx, gy = 0.12, 0.12
    for i, (code, name, body, col, star, is_new) in enumerate(items):
        col_idx = i % 4
        row_idx = i // 4
        x = sx + col_idx * (cw + gx)
        y = sy + row_idx * (ch + gy)
        add_rect(slide, x, y, cw, ch, WHITE, line=BORDER, line_w=1, radius=0.05)
        add_rect(slide, x, y, 0.06, ch, col)
        if is_new:
            add_rect(slide, x + cw - 0.6, y + 0.13, 0.5, 0.22, GREEN, radius=0.5)
            add_text(slide, x + cw - 0.6, y + 0.16, 0.5, 0.18,
                     "NEW", size=8, bold=True, color=WHITE, align='center', char_spacing=2)
        # Code tag
        add_rect(slide, x + 0.20, y + 0.20, 0.55, 0.25, col, radius=0.1)
        add_text(slide, x + 0.20, y + 0.22, 0.55, 0.22,
                 code, size=10, bold=True, color=WHITE, align='center')
        if star:
            add_text(slide, x + cw - 0.35, y + 0.15, 0.30, 0.30,
                     "★", size=18, bold=True, color=AMBER)
        # Name
        add_text(slide, x + 0.20, y + 0.62, cw - 0.40, 0.65,
                 name, size=14, bold=True, color=TEXT)
        # Body
        add_text(slide, x + 0.20, y + 1.30, cw - 0.40, 0.55,
                 body, size=10, color=TEXT_GRAY)
        # Bottom scenario coverage
        add_text(slide, x + 0.20, y + ch - 0.30, cw - 0.40, 0.25,
                 "✅ 5/5 시나리오 (A·B·C·D·E)",
                 size=9, bold=True, color=GREEN)

    # 8th cell — explanation
    x8 = sx + 3 * (cw + gx)
    y8 = sy + 1 * (ch + gy)
    add_rect(slide, x8, y8, cw, ch, NAVY, radius=0.05)
    add_text(slide, x8 + 0.20, y8 + 0.20, cw - 0.40, 0.4,
             "★ 5개 핵심 (RS1·RS2·RS3·RS6·RS7)",
             size=12, bold=True, color=AMBER, font_face='Arial Black')
    add_text(slide, x8 + 0.20, y8 + 0.65, cw - 0.40, 0.30,
             "각 슬라이드에서 상세 시각화", size=10, color=BORDER)
    add_rect(slide, x8 + 0.20, y8 + 1.05, 0.6, 0.03, AMBER)
    add_text(slide, x8 + 0.20, y8 + 1.15, cw - 0.40, 0.85,
             "RS7 NAND 공정 전환 주기 연장 R&D는 2026-05-06 신규 추가 (Weber/PSU 학습곡선 + YMTC IP 종속 회피)",
             size=9, color=WHITE)

    add_text(slide, 12.4, 7.15, 0.8, 0.25, "07 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 8 — RS7 NAND CYCLE EXTENSION (Why Now + 4 R&D tracks)
# =============================================================================
def build_slide_8_rs7(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, WHITE)
    add_rect(slide, 0, 0, W_IN, 0.08, GREEN)

    add_text(slide, 0.55, 0.55, 8, 0.3, "05  ·  RS7  ·  NEW",
             size=12, bold=True, color=GREEN, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12.5, 0.7,
             'NAND는 layer 경쟁이 아니라 "주기 연장" 경쟁이다',
             size=28, bold=True, color=NAVY, font_face='Arial Black')
    add_text(slide, 0.55, 1.40, 12.5, 0.4,
             "양산 ramp 6개월 지연 = 누적 이익의 2/3 소실 (Weber/PSU). 4사 모두 process upgrade 선회. YMTC IP 종속 회피 = 디커플링 시 생존",
             size=11, color=TEXT_GRAY)

    # Left — 3 evidence boxes
    ex, ey, ew = 0.55, 1.95, 6.1
    add_text(slide, ex, ey, ew, 0.30,
             "WHY NOW · 3가지 증거",
             size=12, bold=True, color=AMBER, char_spacing=3)
    evidence = [
        ('①', '학습곡선의 잔혹함', '−2/3', 'ramp 6개월 지연 → 누적 이익 손실',
         '1년 지연 → 손실 전환 / 1분당 약 $5,000 손실 (PSU Weber)', RED),
        ('②', '업계 일제 선회 (2026)', '+5%', '2026 NAND capex $22.2B 증가율만',
         '4사 모두 capa 확장 대신 process upgrade·hybrid bonding (TrendForce)', AMBER),
        ('③', 'YMTC IP 지배 — 디커플링 시 차단', '70%+', 'YMTC가 hybrid bonding 핵심 IP 보유',
         '시나리오 C 발생 시 양산 불가. 자체 IP = 국가 안보 R&D', NAVY2),
    ]
    for i, (icon, title, kpi, kpi_sub, body, col) in enumerate(evidence):
        y = ey + 0.40 + i * 1.55
        add_rect(slide, ex, y, ew, 1.43, LIGHT, radius=0.05)
        add_rect(slide, ex, y, 0.06, 1.43, col)
        add_text(slide, ex + 0.18, y + 0.10, 0.6, 0.6,
                 icon, size=28, bold=True, color=col, font_face='Arial Black')
        add_text(slide, ex + 0.85, y + 0.15, ew - 1.0, 0.35,
                 title, size=14, bold=True, color=TEXT)
        add_text(slide, ex + 0.85, y + 0.50, 1.5, 0.6,
                 kpi, size=28, bold=True, color=col, font_face='Arial Black')
        add_text(slide, ex + 2.30, y + 0.55, ew - 2.5, 0.30,
                 kpi_sub, size=10, bold=True, color=TEXT_GRAY)
        add_text(slide, ex + 0.85, y + 1.10, ew - 1.0, 0.30,
                 body, size=10, color=TEXT)

    # Right — 4 R&D tracks
    tx, ty, tw = 6.85, 1.95, 5.85
    add_text(slide, tx, ty, tw, 0.30,
             "4 R&D 트랙 (병행)",
             size=12, bold=True, color=GREEN, char_spacing=3)
    tracks = [
        ('01', 'Hybrid Bonding 자체 IP',
         'YMTC 우회 특허 200건+ (2027). 한국 IP 컨소시엄. V11에 자체 IP 70%+'),
        ('02', 'Multi-deck 정교화',
         'Deck 당 layer 한계 돌파. V9→V10 fab 재구성 비용 30% 절감'),
        ('03', 'bit-per-cell 확장 (TLC→QLC→PLC)',
         'QLC 비중 30%+ (2026). PLC 시제품 (2028) → 양산 검토 (2029)'),
        ('04', '호스트 협력 firmware (FDP·SCADA)',
         'endurance 2배+, throughput 30%+. RS3 SW 매출과 R&D 자원 공유'),
    ]
    for i, (num, title, body) in enumerate(tracks):
        y = ty + 0.40 + i * 1.10
        add_rect(slide, tx, ty + 0.40 + i * 1.10, tw, 1.00, AMBER_BG, radius=0.05)
        add_rect(slide, tx, y, 0.06, 1.00, GREEN)
        add_rect(slide, tx + 0.20, y + 0.18, 0.45, 0.30, GREEN, radius=0.1)
        add_text(slide, tx + 0.20, y + 0.21, 0.45, 0.25,
                 num, size=11, bold=True, color=WHITE, align='center')
        add_text(slide, tx + 0.80, y + 0.15, tw - 0.95, 0.35,
                 title, size=13, bold=True, color=TEXT)
        add_text(slide, tx + 0.80, y + 0.50, tw - 0.95, 0.45,
                 body, size=10, color=TEXT)

    # Footer KPI strip
    add_rect(slide, 0.55, 6.70, W_IN - 1.1, 0.45, NAVY, radius=0.1)
    add_text(slide, 0.75, 6.78, W_IN - 1.4, 0.30,
             "✱  3년 누적 capex 회피 1.5~2조 원  ·  KPI 4: 자체 IP 비율 / ramp 시간 / capex/bit / YMTC 의존도",
             size=11, bold=True, color=AMBER)
    add_text(slide, 12.4, 7.15, 0.8, 0.25, "08 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 9 — EWI DASHBOARD MOCKUP
# =============================================================================
def build_slide_9_ewi(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, NAVY)
    add_rect(slide, 0, 0, W_IN, 0.08, AMBER)

    add_text(slide, 0.55, 0.55, 8, 0.3, "06  ·  EWI DASHBOARD",
             size=12, bold=True, color=AMBER, char_spacing=4)
    add_text(slide, 0.55, 0.80, 12, 0.7,
             "9개 결정 + 28개 지표 + 12개 자동 트리거",
             size=28, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, 0.55, 1.30, 12, 0.4,
             "Vercel 자동 배포  ·  GitHub git integration  ·  데이터 변경 시 양쪽 갈래 동시 갱신",
             size=11, color=BORDER)

    # Top KPI row (6 cards)
    tk = [
        ('DECISIONS',  '10',  '즉시 결정 추적',  AMBER),
        ('INDICATORS', '32',  'EWI 모니터링',   GREEN),
        ('TRIGGERS',   '12',  '자동 트리거',    RED),
        ('SCENARIOS',  '5',   '대안적 미래',    PURPLE),
        ('COMMITS',    '30+', 'GitHub 추적',    AMBER),
        ('ROBUST',     '7',   '즉시 실행 대상', GREEN),
    ]
    kw, kh = 1.95, 0.95
    kx, ky = 0.55, 1.85
    kg = 0.10
    for i, (label, val, sub, col) in enumerate(tk):
        x = kx + i * (kw + kg)
        add_rect(slide, x, ky, kw, kh, NAVY2, radius=0.05)
        add_rect(slide, x, ky, kw, 0.04, col)
        add_text(slide, x + 0.15, ky + 0.10, kw - 0.30, 0.20,
                 label, size=8, bold=True, color=col, char_spacing=2)
        add_text(slide, x + 0.15, ky + 0.30, kw - 0.30, 0.45,
                 val, size=24, bold=True, color=WHITE, font_face='Arial Black')
        add_text(slide, x + 0.15, ky + 0.75, kw - 0.30, 0.20,
                 sub, size=9, color=BORDER)

    # Quadrant trajectory
    qx, qy, qw, qh = 0.55, 3.00, 6.0, 3.4
    add_rect(slide, qx, qy, qw, qh, NAVY2, radius=0.05)
    add_text(slide, qx + 0.15, qy + 0.12, qw - 0.30, 0.30,
             "QUADRANT TRAJECTORY · 2년 사분면 위치 변화",
             size=10, bold=True, color=AMBER, char_spacing=2)
    add_text(slide, qx + 0.15, qy + 0.40, qw - 0.30, 0.30,
             "DF1 (AI 수요)  ×  DF2 (디커플링)  —  현재 → B 시나리오 강한 모멘텀",
             size=8, color=BORDER)
    # Cross axes inside quadrant
    cx_ax = qx + qw / 2
    cy_ax = qy + qh / 2 + 0.15
    add_rect(slide, qx + 0.4, cy_ax - 0.005, qw - 0.8, 0.01, BORDER)
    add_rect(slide, cx_ax - 0.005, qy + 0.7, 0.01, qh - 1.1, BORDER)

    traj = [
        ('24.05', 0.35, 0.45, RGBColor(0x6B, 0x6B, 0x6B), False),
        ('25.05', 0.50, 0.60, RGBColor(0x80, 0x80, 0x80), False),
        ('25.11', 0.55, 0.25, RGBColor(0x9A, 0x9A, 0x9A), False),
        ('26.02', 0.65, 0.15, RGBColor(0xB5, 0xB5, 0xB5), False),
        ('26.04', 0.70, 0.05, AMBER, False),
        ('26.05', 0.75, 0.05, AMBER, True),
    ]
    def px(v): return qx + 0.5 + v * (qw - 1.0)
    def py(v): return qy + 0.7 + (1 - v) * (qh - 1.1)
    prev = None
    for date, df1, df2, col, current in traj:
        x, y = px(df1), py(df2)
        if prev:
            conn = slide.shapes.add_connector(1, Inches(prev[0]), Inches(prev[1]),
                                              Inches(x), Inches(y))
            conn.line.color.rgb = AMBER
            conn.line.width = Pt(1.5)
        prev = (x, y)
    for date, df1, df2, col, current in traj:
        x, y = px(df1), py(df2)
        r = 0.10 if current else 0.06
        add_oval(slide, x - r, y - r, r * 2, r * 2, col,
                 line=WHITE if current else None)
        if current:
            add_text(slide, x - 0.20, y - 0.35, 0.45, 0.20,
                     "현재", size=8, bold=True, color=AMBER, align='center')
        add_text(slide, x - 0.25, y + r + 0.02, 0.55, 0.18,
                 date, size=7, bold=True, color=WHITE, align='center')

    # Quadrant labels
    add_text(slide, qx + 0.4,  qy + 0.85, 1.5, 0.2, "C 기술냉전",   size=8, bold=True, color=RED)
    add_text(slide, qx + qw - 1.6, qy + 0.85, 1.4, 0.2, "A 황금 요새", size=8, bold=True, color=NAVY3, align='right')
    add_text(slide, qx + 0.4,  qy + qh - 0.5, 1.5, 0.2, "D 조용한 재편", size=8, bold=True, color=AMBER)
    add_text(slide, qx + qw - 1.6, qy + qh - 0.5, 1.4, 0.2, "B AI 르네상스 ★", size=8, bold=True, color=GREEN, align='right')

    # Right side decisions
    dx, dy, dw, dh = 6.70, 3.00, 6.10, 3.4
    add_rect(slide, dx, dy, dw, dh, NAVY2, radius=0.05)
    add_text(slide, dx + 0.15, dy + 0.12, dw - 0.30, 0.30,
             "DECISIONS · 9개 즉시 결정 진척",
             size=10, bold=True, color=AMBER, char_spacing=2)
    decisions = [
        ('D1',  'HBM4 NVIDIA 점유율 회복',         '2026.09', 60, AMBER, False),
        ('D2',  '소재 공급망 비중국 다각화',         '2026.09', 40, AMBER, False),
        ('D3',  '3D DRAM R&D + IMEC 협약',         '2026.12', 20, RED,   False),
        ('D5',  'AI 개발 효율화 도구 전사 도입',     '2026.12', 30, AMBER, False),
        ('D6',  '이사회 정책 패키지 의결',            '2026.06', 75, GREEN, False),
        ('D9',  '다운사이클 M&A 펀드 적립',          '2026.12', 0,  RED,   False),
        ('D10', 'NAND 공정 전환 주기 R&D (RS7)',    '2026.12', 5,  GREEN, True),
    ]
    rh = 0.40
    for i, (did, title, deadline, pct, col, is_new) in enumerate(decisions):
        y = dy + 0.50 + i * rh
        add_text(slide, dx + 0.15, y + 0.05, 0.4, 0.25,
                 did, size=9, bold=True, color=col)
        add_text(slide, dx + 0.50, y + 0.05, 2.6, 0.25,
                 title, size=9, color=WHITE)
        add_text(slide, dx + 3.10, y + 0.05, 0.7, 0.25,
                 deadline, size=8, color=BORDER)
        # progress bar
        add_rect(slide, dx + 3.85, y + 0.10, 1.5, 0.10, NAVY3, radius=0.5)
        add_rect(slide, dx + 3.85, y + 0.10, 1.5 * (pct / 100.0), 0.10, col, radius=0.5)
        add_text(slide, dx + 5.40, y + 0.05, 0.4, 0.25,
                 f"{pct}%", size=9, bold=True, color=col, align='right')
        if is_new:
            add_rect(slide, dx + 5.85, y + 0.07, 0.20, 0.18, GREEN, radius=0.5)

    add_text(slide, 0.55, 6.85, 12.2, 0.3,
             "✱  데이터 변경 한 번 → 양쪽 갈래(① PPTX + ② 대시보드) 동시 갱신 → git push origin main → Vercel 자동 배포",
             size=10, italic=True, color=AMBER)
    add_text(slide, 12.4, 7.15, 0.8, 0.25, "09 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Slide 10 — CLOSING (Decision Request)
# =============================================================================
def build_slide_10_closing(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W_IN, H_IN, NAVY)
    add_rect(slide, 0, 0, 0.10, H_IN, AMBER)

    add_text(slide, 0.60, 0.55, 8, 0.3, "DECISION REQUEST",
             size=14, bold=True, color=AMBER, char_spacing=5)

    add_text(slide, 0.60, 1.55, 12, 1.0,
             "호황의 정점이",
             size=64, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, 0.60, 2.50, 12, 1.0,
             "다운턴 준비의",
             size=64, bold=True, color=WHITE, font_face='Arial Black')
    add_text(slide, 0.60, 3.45, 12, 1.0,
             "마지막 기회다.",
             size=64, bold=True, color=AMBER, font_face='Arial Black')

    add_rect(slide, 0.60, 4.50, 1.6, 0.04, AMBER)
    add_text(slide, 0.60, 4.65, 12, 0.4,
             "─ 9개 결정의 묶음 의결을 요청드립니다 ─",
             size=14, italic=True, color=WHITE)

    asks = [
        ('01', '전략 방향 승인',
         '시나리오 B Main Bet + Side Bets +\nRobust 7개 전략 채택 의결'),
        ('02', '9개 즉시 결정 패키지',
         'HBM4·소재·3D DRAM·텍사스·AI효율화·이사회정책·\n잉여인력·보조금·M&A·NAND R&D'),
        ('03', '거버넌스 + 트리거 승인',
         '30일 내 의사결정 + 트리거 12개 +\ncapex 하한 4조원/년 + 의결권 구조'),
    ]
    aw, ah = 4.00, 1.50
    ax, ay = 0.60, 5.30
    ag = 0.20
    for i, (num, title, body) in enumerate(asks):
        x = ax + i * (aw + ag)
        add_rect(slide, x, ay, aw, ah, NAVY2, radius=0.08)
        add_rect(slide, x, ay, 0.06, ah, AMBER)
        add_text(slide, x + 0.20, ay + 0.15, 1.0, 0.55,
                 num, size=32, bold=True, color=AMBER, font_face='Arial Black')
        add_text(slide, x + 0.20, ay + 0.70, aw - 0.40, 0.30,
                 title, size=14, bold=True, color=WHITE)
        add_text(slide, x + 0.20, ay + 1.00, aw - 0.40, 0.50,
                 body, size=9, color=BORDER, line_spacing=1.3)

    add_text(slide, 0.60, 7.10, 9.5, 0.25,
             "삼성전자  ·  DS부문 메모리사업부 전략기획팀  ·  대외비  ·  2026.05.06",
             size=9, color=BORDER)
    add_text(slide, 12.4, 7.10, 0.8, 0.25, "10 / 10",
             size=10, bold=True, color=AMBER, align='right')


# =============================================================================
# Main
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    # Slides 1~3: Figma full-bleed PNGs
    for i in (1, 2, 3):
        png = os.path.join(FIGMA_DIR, f'slide-{i:02d}.png')
        if not os.path.exists(png):
            raise SystemExit(f'Missing Figma PNG: {png}')
        print(f'  Building slide {i}: figma full-bleed ({os.path.basename(png)})...')
        add_figma_slide(prs, png)

    # Slides 4~10: native python-pptx (figma-style)
    builders = [
        (4,  build_slide_4_cycle),
        (5,  build_slide_5_matrix),
        (6,  build_slide_6_mainbet),
        (7,  build_slide_7_robust),
        (8,  build_slide_8_rs7),
        (9,  build_slide_9_ewi),
        (10, build_slide_10_closing),
    ]
    for i, fn in builders:
        print(f'  Building slide {i}: {fn.__name__}...')
        fn(prs)

    print(f'Saving to: {OUTPUT}')
    prs.save(OUTPUT)
    sz = os.path.getsize(OUTPUT) / 1024
    print(f'Done. Output size: {sz:.1f} KB · {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
