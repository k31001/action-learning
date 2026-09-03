#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
내부 인터뷰 핵심 인용문 1장 슬라이드 (내부 인터뷰 대상자 3인, 발언 1개씩).

단일 소스: dashboard/src/data/interviews.js 의 keyQuotes (내부 인터뷰 3건)
           sources/raw-notes/{choi-jangseok,lee-changsoo,song-yongho}-*.md
출력:     outputs/presentation/internal-interview-quotes.pptx
재생성:   .venv/bin/python outputs/presentation/scripts/generate_interview_quotes_slide.py

디자인 규율: samsung-memory-ppt-design-skill (화이트 배경, Samsung Blue 단일 액센트,
             액션 타이틀, em-dash 금지, 문서등급 플레이스홀더).
포맷:      인용문 전용 슬라이드 문법 (v2). 레포 기존 인용 포맷을 승계했다:
             generate_pptx.py Focal Issue(대형 인용부호 + 인용문 블록)·클로징(큰 인용문 + 출처 한 줄),
             generate_dev_transformation_summary.cjs(인용 카드 + "이름, 출처(일자)" 출처 라인).
             인용문이 주인공: 3단 컬럼마다 대형 인용부호 → 인용문(대형 볼드) → 헤어라인 → 화자.
             프레임 라벨·해설·과제 반영 문단은 제거하고 발표자 노트로 이동.
인용문은 interviews.js keyQuotes 원문에서 em-dash만 문장부호로 치환했다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 토큰 ────────────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x14, 0x28, 0xA0)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
G_MID  = RGBColor(0x8A, 0x8A, 0x8E)
G_LINE = RGBColor(0xD9, 0xD9, 0xD9)
G_BG   = RGBColor(0xF6, 0xF6, 0xF7)
B_TINT = RGBColor(0xE9, 0xED, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = 'SamsungOneKorean'

def _font(run, size, bold=False, color=INK, name=FONT, italic=False):
    f = run.font
    f.name, f.size, f.bold, f.italic = name, Pt(size), bold, italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set('typeface', name)

def txt(slide, x, y, w, h, parts, *, size=10, align=PP_ALIGN.LEFT,
        leading=1.15, anchor=MSO_ANCHOR.TOP, space_after=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(parts, str):
        parts = [[(parts, {})]]
    elif parts and isinstance(parts[0], tuple):
        parts = [parts]
    for i, line in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(line, str):
            line = [(line, {})]
        for text, kw in line:
            r = p.add_run()
            r.text = text
            _font(r, kw.get('size', size), kw.get('bold', False), kw.get('color', INK),
                  italic=kw.get('italic', False))
    return tb

def box(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, rnd=False, radius=0.04):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if rnd:
        s.adjustments[0] = radius
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def hline(slide, x, y, w, color=G_LINE, weight=0.5):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False

# ── 데이터 (interviews.js keyQuotes 원문, em-dash만 치환) ─────────────────────
QUOTES = [
    dict(
        step='01', frame='다운사이드의 실체',
        name='최장석 상무', role='메모리사업부 상품기획팀장', date='2026-07-29 인터뷰',
        quote='HBM 하나 만들면 DDR을 4~5개 더 만들 수 있는데 그걸 희생하면서 HBM을 만드는 거다. '
              '디맨드 대비 서플라이가 60~70% 수준인데, HBM이 무너지는 순간 그 갭이 확 줄어들어 '
              '쇼티지가 바로 오버서플라이가 될 수 있다.',
        context='HBM 다운사이드의 실체: 캐파 상쇄 메커니즘',
        takeaway='"수요 감소"를 가속기 → HBM → DDR의 제품 단위 인과로 정의. 쇼티지가 오버서플라이로 반전되는 지점이 과제의 문제 정의다.',
    ),
    dict(
        step='02', frame='1차 저지선과 과제의 자리',
        name='이창수 부사장', role='메모리사업부 영업팀장', date='2026-08-03 인터뷰',
        quote='선수금을 받아놓고 구매 의무를 저버리면 그만큼 차감한다. 메모리 최초의 take-or-pay '
              '바인딩에 NTB 가격 하한까지. 이게 1차 방어선이고, 여러분들의 이 과제가 2차 방어선이다.',
        context='멀티이어 계약의 구조와 과제의 위치',
        takeaway='계약(LTA·take-or-pay)이 다운턴 물량의 바닥을 만드는 1차 저지선. 계약이 뚫리는 극한 국면을 대비하는 2차 저지선이 이 과제의 정체성이다.',
    ),
    dict(
        step='03', frame='2차 저지선의 기본',
        name='송용호 부사장', role='AX/PI센터장 (前 솔루션개발실장)', date='2026-09-03 인터뷰',
        quote='다음 다운턴이 시작됐다고 가정해 보자. 가격이 빠지고 주문이 슬로우해지고 재고가 쌓이면 '
              '무슨 일이 벌어질까. 첫 번째는 대규모 RMA다. 그리고 대규모 RMA가 시작될 때 빌미는 품질이다.',
        context='다음 다운턴 시뮬레이션: 과제에 빠져 있던 키워드 "품질"',
        takeaway='고객은 재고 대신 반품으로 대응한다. 새로운 것(전환 TAT·라인 일체화·본딩)에 앞서 품질·원가라는 기본을 2차 저지선의 별도 축으로 세운다.',
    ),
]

TITLE = '내부 임원 3인의 증언: 다운사이드는 캐파 반전으로 오고, 계약은 1차 저지선일 뿐, 2차 저지선의 기본은 품질이다'
LEAD  = ('상품기획·영업·AX/PI 관점의 내부 인터뷰 3건(2026-07~09)에서 핵심 발언 1개씩. '
         '발언의 순서가 곧 과제의 스토리라인(문제 정의 → 1차 저지선 → 2차 저지선)이다.')
FOOTER = ('출처: 사내 인터뷰 녹취(Clova Note 자동 전사) 정리본. sources/raw-notes/'
          'choi-jangseok-product-planning-interview-2026-07-29.md · lee-changsoo-memory-sales-interview-2026-08-03.md · '
          'song-yongho-ax-pi-interview-2026-09-03.md. 인용문은 대시보드 인터뷰 메뉴 keyQuotes 원문(문장부호만 정리). '
          '수치는 발언자 표현 그대로이며 별도 검증치가 아니다.')
NOTES = """내부 인터뷰 핵심 인용문 1장.

세 발언을 이 순서로 놓은 이유: 최장석(상품기획)은 다운턴이 어떤 메커니즘으로 오는지(HBM 캐파 상쇄가 반전되는 지점), 이창수(영업)는 그 다운턴에서 회사가 이미 세운 1차 저지선(take-or-pay 멀티이어 계약)과 이 과제의 자리(2차 저지선), 송용호(AX/PI)는 2차 저지선에 빠져 있던 기본(품질·원가)을 말한다. 타이틀만 읽어도 과제의 3단 논증이 이어진다.

외부 전문가(신문섭·크리스 밀러·Sachin Katti) 발언은 이 장에서 제외했다. 내부 인터뷰 대상자 3인만 담는다.

인용문은 dashboard/src/data/interviews.js 의 keyQuotes 원문이며 em-dash만 마침표·쉼표로 치환했다. 수치(4~5개, 60~70%)는 발언자 표현 그대로.

각 발언의 과제 반영(장표에서는 제외, 구두 설명용):
- 1. 
- 2. 
- 3. 

문서등급은 플레이스홀더. 사내 인물 실명·직책이 들어가므로 배포 범위 확인 후 등급 표기."""

# ── 슬라이드 (인용문 전용 포맷) ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

MX, MW = 0.7, 13.333 - 1.4
# 타이틀 (액션 타이틀) + 문서등급
txt(slide, MX, 0.5, MW - 1.7, 0.9, [[(TITLE, {'size': 20, 'bold': True, 'color': INK})]], leading=1.15)
txt(slide, MX + MW - 1.5, 0.5, 1.5, 0.2, [[('[문서등급 표기]', {'size': 9, 'color': GRAY})]], align=PP_ALIGN.RIGHT)
txt(slide, MX, 1.42, MW, 0.3, [[('내부 인터뷰 3건(2026-07~09) · 발언 순서가 곧 과제의 논증 순서: 문제 정의 → 1차 저지선 → 2차 저지선', {'size': 11, 'color': GRAY})]])

# 3단: 인용부호 → 인용문 → 헤어라인 → 화자
TOP = 2.15
GAP = 0.55
CW = (MW - 2 * GAP) / 3
QH = 2.45                       # 인용문 영역 높이
for i, q in enumerate(QUOTES):
    x = MX + i * (CW + GAP)
    # 대형 인용부호 (Samsung Blue, 액센트 1개)
    txt(slide, x - 0.04, TOP - 0.12, 1.0, 1.0, [[('“', {'size': 72, 'bold': True, 'color': BLUE})]], leading=0.85)
    # 인용문 (주인공)
    txt(slide, x, TOP + 0.85, CW, QH, [[(q['quote'], {'size': 16, 'bold': True, 'color': INK})]], leading=1.32)
    # 헤어라인 + 화자
    Y = TOP + 0.85 + QH + 0.12
    hline(slide, x, Y, CW, color=G_LINE, weight=0.75)
    # 이니셜 원 (아바타 대체, 텍스트 없는 사진 슬롯 대신)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(Y + 0.18), Inches(0.46), Inches(0.46))
    circ.fill.solid(); circ.fill.fore_color.rgb = B_TINT
    circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); r.text = q['name'][0]; _font(r, 13, True, BLUE)
    txt(slide, x + 0.6, Y + 0.16, CW - 0.6, 0.28,
        [[(q['name'], {'size': 12.5, 'bold': True, 'color': INK})]])
    txt(slide, x + 0.6, Y + 0.45, CW - 0.6, 0.24,
        [[(q['role'], {'size': 9.5, 'color': GRAY})]])
    txt(slide, x + 0.6, Y + 0.69, CW - 0.6, 0.24,
        [[(q['date'], {'size': 9.5, 'color': G_MID})]])
    # 컬럼 구분 세로 헤어라인
    if i < 2:
        ln = slide.shapes.add_connector(1, Inches(x + CW + GAP / 2), Inches(TOP + 0.1),
                                        Inches(x + CW + GAP / 2), Inches(Y + 0.95))
        ln.line.color.rgb = G_LINE; ln.line.width = Pt(0.5); ln.shadow.inherit = False

# 각주
hline(slide, MX, 7.5 - 0.62, MW, color=G_LINE, weight=0.5)
txt(slide, MX, 7.5 - 0.55, MW - 0.6, 0.4, [[(FOOTER, {'size': 8.5, 'color': G_MID})]], leading=1.15)
txt(slide, MX + MW - 0.5, 7.5 - 0.55, 0.5, 0.2, [[('1', {'size': 9, 'color': G_MID})]], align=PP_ALIGN.RIGHT)

# 발표자 노트: 과제 반영 문장을 노트로 이동
notes = NOTES
for i, q in enumerate(QUOTES):
    notes = notes.replace(f"- {i+1}. ", f"- {i+1}. {q['name']} ({q['frame']}): {q['takeaway']}")
slide.notes_slide.notes_text_frame.text = notes
OUT = 'outputs/presentation/internal-interview-quotes.pptx'
prs.save(OUT)
print('saved:', OUT)
