# -*- coding: utf-8 -*-
"""향후 계획 2장 덱 생성 (전략 3 FDP 4주 계획 + 4대 전략 종합 타임라인).

디자인 시스템: generate_ssd_strategy_pptx.py 승계
  20 x 11.25 in 캔버스 / Arial 단일 폰트 / Samsung Blue #1428A0 단일 액센트
  헤더(조직명·문서등급·킥커·33pt 액션 타이틀·21pt 리드·헤어라인) / 푸터(출처·페이지)
  틴트 카드 #F4F6FC(무테) · 아웃라인 카드 흰색+#D9D9D9 0.75pt · 다크 블루 정리 밴드 · 직각 사각형

실행: .venv/bin/python outputs/presentation/scripts/generate_future_plan_pptx.py
출력: outputs/presentation/future-plan.pptx
콘텐츠 소스: outputs/presentation/future-plan-outline.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ---- 디자인 토큰 ----
BLUE = RGBColor(0x14, 0x28, 0xA0)
BLUE_T2 = RGBColor(0xAA, 0xB8, 0xE8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GRAY_L = RGBColor(0x8A, 0x8A, 0x8A)
LINE = RGBColor(0xD9, 0xD9, 0xD9)
TINT = RGBColor(0xF4, 0xF6, 0xFC)
SHADE = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

MX = 0.79
CW = 18.42
RIGHT = MX + CW
GRADE = "[문서등급 표기]"
TOTAL = 2

OUT = os.path.join(os.path.dirname(__file__), "..", "future-plan.pptx")

prs = Presentation()
prs.slide_width = Emu(18288000)   # 20.00 in
prs.slide_height = Emu(10287000)  # 11.25 in
BLANK = prs.slide_layouts[6]


def _font(run, size, bold, color):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True,
       spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, size, bold, color) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = t
        _font(r, size, bold, color)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE,
         dash=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
        if dash:
            sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return sp


def header(slide, kicker, title, lead):
    tb(slide, MX, 0.46, 6.0, 0.34, [("삼성전자 메모리사업부", 18, True, BLUE)])
    tb(slide, 13.21, 0.46, 6.0, 0.34, [(GRADE, 18, False, GRAY)], align=PP_ALIGN.RIGHT)
    tb(slide, MX, 0.93, CW, 0.34, [(kicker, 18, True, BLUE)])
    tb(slide, MX, 1.33, CW, 0.62, [(title, 33, True, INK)])
    tb(slide, MX, 2.04, CW, 0.44, [(lead, 21, False, GRAY)])
    rect(slide, MX, 2.62, CW, 0.014, fill=LINE)


def footer(slide, source, no):
    tb(slide, MX, 10.49, 15.6, 0.34, [(source, 18, False, GRAY)])
    tb(slide, 17.55, 10.49, 1.66, 0.34, [(f"{no:02d} / {TOTAL:02d}", 18, False, GRAY)],
       align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---- 공통 타임라인 그리드 ----
LW = 4.30                      # 좌측 라벨 열
TX = MX + LW + 0.16            # 타임라인 시작 x
WKW = 2.30                     # 주 열 폭
GAPW = 0.46                    # 추석 연휴 열 폭
WEEKS = [("1주", "9/7"), ("2주", "9/14"), None, ("3주", "9/29"), ("4주", "10/6")]
OX_GAP = 0.16
OUT_W = RIGHT - (TX + WKW * 4 + GAPW + OX_GAP)   # 산출물 열 폭
HDR_Y = 2.86


def col_x():
    """각 열의 (x, w) 리스트. None 자리는 추석 연휴 열."""
    xs, x = [], TX
    for wk in WEEKS:
        w = GAPW if wk is None else WKW
        xs.append((x, w))
        x += w
    return xs, x + OX_GAP


def timeline_header(slide, out_label):
    cols, ox = col_x()
    for (x, w), wk in zip(cols, WEEKS):
        if wk is None:
            tb(slide, x, HDR_Y, w, 0.40, [("추석", 15, True, GRAY_L)], align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)
        else:
            tb(slide, x + 0.06, HDR_Y, w - 0.12, 0.40,
               [(f"{wk[0]} · {wk[1]}", 18, True, INK)], anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, ox, HDR_Y, OUT_W, 0.40, [(out_label, 18, True, INK)], anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, MX, HDR_Y + 0.46, CW, 0.014, fill=LINE)
    return cols, ox


def shade_holiday(slide, y, h):
    cols, _ = col_x()
    x, w = cols[2]
    rect(slide, x, y, w, h, fill=SHADE)


def cell(slide, x, w, y, h, text, hot=False, size=16.5, pad=0.14):
    if hot:
        rect(slide, x + 0.06, y, w - 0.12, h, fill=BLUE)
        color = WHITE
    else:
        rect(slide, x + 0.06, y, w - 0.12, h, fill=TINT)
        color = INK
    tb(slide, x + 0.06 + pad, y + pad * 0.8, w - 0.12 - 2 * pad, h - pad * 1.4,
       [(text, size, hot, color)])


# =====================================================================
# 슬라이드 1 · 전략 3 FDP 향후 계획 (세 검토 축 × 4주)
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "삼성 SSD 전략적 방향성 · 향후 계획 (전략 3 FDP)",
       "남은 4주, 세 가지 검토로 FDP 전략을 실행안으로 끌어올립니다",
       "현실 인식은 인터뷰로, AI 워크로드는 조사와 기술검토로, 실행 전략은 전략구상으로 확인해 보고서 v1.2로 수렴합니다")
cols, ox = timeline_header(s, "4주 후 산출물")

RY0, RH, RG = 3.42, 1.80, 0.16
shade_holiday(s, RY0, RH * 3 + RG * 2)

axes = [
    ("① 현실 인식 · 인터뷰",
     "우리 FDP 기술과 고객 관계는 지금 어디에 있나",
     "보고서의 사내 주장과 추정치를 확인된 사실로 교체합니다",
     ["인터뷰 가이드 설계, 확인 대장(사내 확인·추정치 전수)",
      "◆ 인터뷰 1차\n현황 확인 + ②·③ 사내 질문 동반",
      "현황 대장(기술·고객·조직) + 주장 대 실제 갭 표",
      "◆ 인터뷰 2차\n②·③ 결론 검증(서면 대체 가능)"],
     [True if i in (1, 3) else False for i in range(4)],
     ("현황 대장 + 갭 표", "사내 확인·추정치 표기 0건이 목표")),
    ("② AI 워크로드 · 기술검토",
     "FDP가 AI 워크로드에 통하는가, 통하면 얼마나 어려운가",
     "KV Cache 오프로드 등에서 효과·난이도·차별화 여지를 판정합니다",
     ["AI 워크로드 분류(KV Cache·체크포인트·데이터 로더), 조사 스코프",
      "스택 조사(vLLM·LMCache·Dynamo·커널 6.16), 코드 검토, 공개 실측 대조",
      "효과·난이도 판정표, 기술 숙제 3분류(우리·생태계·고객)",
      "차별화 여지 결론\n(따라오기 쉬운 층 / 관계 때문에 어려운 층)"],
     [False, False, False, False],
     ("워크로드별 효과·난이도 판정표", "기술 숙제 3분류 + 차별화 여지")),
    ("③ 실행 전략 · 전략구상",
     "고객에게 어떻게 다가가고 내부는 어떻게 준비하나",
     "고객별 제안·협력 구조와 개발실 조직·개발 체계 준비안을 만듭니다",
     ["고객 분류 틀(캡티브 단계 × FDP 채택), 내부 준비 항목 초안",
      "고객별 주고받기 초안(주는 것·받는 것)",
      "고객 접근 단계·계약 구조 / 개발실 조직·개발 체계 준비안",
      "준비안 확정, 부록 D 6과제 재배치(오너·일정·첫 액션)"],
     [False, False, False, False],
     ("고객 접근 플레이북", "개발실 준비안(조직·체계·로드맵)")),
]

for i, (ttl, q, purpose, tasks, hots, out) in enumerate(axes):
    y = RY0 + i * (RH + RG)
    # 좌측 라벨 카드
    rect(s, MX, y, LW, RH, fill=TINT)
    tb(s, MX + 0.28, y + 0.18, LW - 0.56, 0.34, [(ttl, 19.5, True, BLUE)])
    tb(s, MX + 0.28, y + 0.56, LW - 0.56, 0.62, [(q, 18, True, INK)])
    tb(s, MX + 0.28, y + 1.24, LW - 0.56, 0.52, [(purpose, 15.5, False, GRAY)])
    # 주별 셀
    wi = 0
    for (x, w), wk in zip(cols, WEEKS):
        if wk is None:
            continue
        cell(s, x, w, y + 0.14, RH - 0.28, tasks[wi], hot=hots[wi], size=16)
        wi += 1
    # 산출물 카드
    rect(s, ox, y, OUT_W, RH, fill=WHITE, line=LINE)
    tb(s, ox + 0.26, y + 0.24, OUT_W - 0.52, 0.72, [(out[0], 18, True, INK)])
    tb(s, ox + 0.26, y + 0.98, OUT_W - 0.52, 0.62, [(out[1], 15.5, False, GRAY)])

# 하단 정리 밴드
BY, BH = 9.32, 1.02
rect(s, MX, BY, CW, BH, fill=BLUE)
tb(s, MX + 0.50, BY + 0.16, 1.4, 0.32, [("수렴", 18.75, False, WHITE)])
tb(s, MX + 0.50, BY + 0.48, CW - 1.0, 0.46,
   [("10월 초 보고서 v1.2와 덱 갱신으로 세 검토를 묶고, 10/12부터 발표자료 준비에 들어갑니다", 23, True, WHITE)])
tb(s, MX + 2.2, BY + 0.16, CW - 2.7, 0.32,
   [("인터뷰가 한 번만 가능하면 2주차에 ②·③ 질문을 함께 넣고, 4주차 검증은 서면으로 대체합니다", 16.5, False, BLUE_T2)])
footer(s, "일정 기준: 9/24-9/28 추석 연휴 · 10/5 대체공휴일 · 10/9 한글날 제외 · 출처: SSD 전략 보고서 v1.1 부록 D", 1)
notes(s, "남은 6주 중 추석 연휴와 발표자료 준비 기간을 빼면 검토에 쓸 수 있는 시간은 4주입니다. "
         "세 검토 축은 순서가 있습니다. 현실 인식(인터뷰)이 출발선이고, 그 위에서 AI 워크로드 검토가 "
         "FDP가 통하는지와 난이도를 판정하며, 실행 전략은 두 결과를 받아 고객 접근과 내부 준비안으로 "
         "바꿉니다. 인터뷰는 1차(현황)와 2차(결론 검증)로 나눠 세 축 모두에 씁니다. 산출물은 10월 초 "
         "보고서 v1.2 한 묶음으로 수렴하고, 10/12부터는 발표자료 준비 기간입니다.")

# =====================================================================
# 슬라이드 2 · 4대 전략 향후 계획 종합 (전략 1·2·4 작성 예정)
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "삼성 SSD 전략적 방향성 · 향후 계획 종합 (4대 전략)",
       "네 전략 모두 10월 초 검토를 마치고 10/12부터 발표자료로 수렴합니다",
       "전략 1·2·4의 계획은 담당별 작성 예정입니다. 전략 3 FDP는 현실 인식, AI 워크로드, 실행 전략 세 검토를 4주에 배치했습니다")
cols, ox = timeline_header(s, "산출물 → 10/12~ 발표자료 준비")

RY0, RG = 3.42, 0.16
BLANK_H, FDP_H = 1.12, 2.62
rows = [("전략 1", BLANK_H, None), ("전략 2", BLANK_H, None),
        ("전략 3", FDP_H, "fdp"), ("전략 4", BLANK_H, None)]
total_h = sum(h for _, h, _ in rows) + RG * (len(rows) - 1)
shade_holiday(s, RY0, total_h)

lanes = [
    ["① 가이드·확인 대장 설계", "◆ 인터뷰 1차\n현황 + ②·③ 질문", "현황 대장 · 갭 표", "◆ 인터뷰 2차\n결론 검증"],
    ["② 워크로드 분류·스코프", "스택·코드 조사, 실측 대조", "효과·난이도 판정표, 숙제 3분류", "차별화 여지 결론"],
    ["③ 고객 분류·준비 항목", "고객별 주고받기 초안", "접근 단계·계약 / 조직·체계안", "준비안 확정 · 부록 D 재배치"],
]
lane_hot = [[False, True, False, True], [False] * 4, [False] * 4]

y = RY0
for name, h, kind in rows:
    if kind is None:
        # 빈 전략 행: 점선 플레이스홀더
        rect(s, MX, y, LW, h, fill=WHITE, line=LINE, dash=True)
        tb(s, MX + 0.28, y + 0.18, LW - 0.56, 0.34, [(name, 19.5, True, GRAY_L)])
        tb(s, MX + 0.28, y + 0.56, LW - 0.56, 0.40, [("[전략명 · 목적 작성 예정]", 15.5, False, GRAY_L)])
        tl_x, tl_w = cols[0][0] + 0.06, (cols[-1][0] + cols[-1][1]) - cols[0][0] - 0.12
        rect(s, tl_x, y + 0.14, tl_w, h - 0.28, fill=WHITE, line=LINE, dash=True)
        tb(s, tl_x, y + 0.14, tl_w, h - 0.28, [("[담당별 4주 계획 작성 예정]", 16, False, GRAY_L)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, ox, y, OUT_W, h, fill=WHITE, line=LINE, dash=True)
        tb(s, ox, y, OUT_W, h, [("[산출물]", 16, False, GRAY_L)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        rect(s, MX, y, LW, h, fill=TINT)
        tb(s, MX + 0.28, y + 0.20, LW - 0.56, 0.34, [("전략 3 · FDP 플랫폼", 19.5, True, BLUE)])
        tb(s, MX + 0.28, y + 0.60, LW - 0.56, 0.62,
           [("자체 SSD 수요를 표준으로 흡수하는 실행안 확정", 18, True, INK)])
        tb(s, MX + 0.28, y + 1.30, LW - 0.56, 1.2,
           [("① 현실 인식: 인터뷰", 15.5, False, GRAY),
            ("② AI 워크로드: 조사·기술검토", 15.5, False, GRAY),
            ("③ 실행 전략: 전략구상", 15.5, False, GRAY)], spacing=1.25)
        LN = 3
        lane_h = (h - 0.14 * 2 - 0.10 * (LN - 1)) / LN
        for li in range(LN):
            ly = y + 0.14 + li * (lane_h + 0.10)
            wi = 0
            for (x, w), wk in zip(cols, WEEKS):
                if wk is None:
                    continue
                cell(s, x, w, ly, lane_h, lanes[li][wi], hot=lane_hot[li][wi], size=15.5, pad=0.11)
                wi += 1
        rect(s, ox, y, OUT_W, h, fill=WHITE, line=LINE)
        tb(s, ox + 0.26, y + 0.20, OUT_W - 0.52, 1.6,
           [("① 현황 대장 + 갭 표", 16.5, True, INK),
            ("② 효과·난이도 판정표 + 숙제 3분류", 16.5, True, INK),
            ("③ 고객 플레이북 + 개발실 준비안", 16.5, True, INK)], spacing=1.3)
        tb(s, ox + 0.26, y + h - 0.62, OUT_W - 0.52, 0.42,
           [("◆ 10월 초 보고서 v1.2 · 덱 갱신", 16.5, True, BLUE)])
    y += h + RG

tb(s, MX, y + 0.06, CW, 0.34,
   [("◆ 마일스톤 · 파란 셀은 담당임원 인터뷰(1차 현황, 2차 검증) · 점선은 작성 예정", 15.5, False, GRAY)])
footer(s, "일정 기준: 9/24-9/28 추석 연휴 · 10/5 대체공휴일 · 10/9 한글날 제외 · 출처: SSD 전략 보고서 v1.1 부록 D", 2)
notes(s, "네 전략의 향후 계획을 한 장에 모은 종합 슬라이드입니다. 전략 1·2·4는 담당별로 같은 4주 격자에 "
         "채워 넣으면 되고, 전략 3 FDP는 세 검토 축을 세 레인으로 보였습니다. 파란 셀은 담당임원 인터뷰 "
         "두 번이고, 네 전략 모두 10월 초 산출물을 내고 10/12부터 발표자료 준비로 합류합니다.")

prs.save(os.path.abspath(OUT))
print(f"생성 완료: {os.path.abspath(OUT)} ({len(prs.slides._sldIdLst)}장)")
