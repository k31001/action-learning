# -*- coding: utf-8 -*-
"""삼성 SSD 전략적 방향성 — 4장 요약 덱 생성.

디자인 시스템: outputs/presentation/storyline-overview.pptx 승계
  20 x 11.25 in 캔버스 / Arial 단일 폰트 / Samsung Blue #1428A0 단일 액센트
  헤더(조직명·문서등급·킥커·33pt 액션 타이틀·21pt 리드·헤어라인) / 푸터(출처·페이지)
  틴트 카드 #F4F6FC(무테) · 아웃라인 카드 흰색+#D9D9D9 0.75pt · 다크 블루 정리 밴드 · 직각 사각형

실행: .venv/bin/python outputs/presentation/scripts/generate_ssd_strategy_pptx.py
출력: outputs/presentation/ssd-strategy.pptx
콘텐츠 소스: outputs/report/ssd-strategy-report.md (본문 v1.0 PPT 압축 맵)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 디자인 토큰 (storyline-overview.pptx 실측) ----
BLUE = RGBColor(0x14, 0x28, 0xA0)
BLUE_T2 = RGBColor(0xAA, 0xB8, 0xE8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xD9, 0xD9, 0xD9)
TINT = RGBColor(0xF4, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD9, 0x30, 0x25)
FONT = "Arial"

MX = 0.79          # 좌우 마진
CW = 18.42         # 콘텐츠 폭 (우측 끝 19.21)
GRADE = "[문서등급 표기]"
TOTAL = 5

OUT = os.path.join(os.path.dirname(__file__), "..", "ssd-strategy.pptx")

prs = Presentation()
prs.slide_width = Emu(18288000)   # 20.00 in
prs.slide_height = Emu(10287000)  # 11.25 in
BLANK = prs.slide_layouts[6]


def _font(run, size, bold, color):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: [(text, size, bold, color)] — 항목당 문단 1개."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, size, bold, color) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = t
        _font(r, size, bold, color)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def header(slide, kicker, title, lead):
    tb(slide, MX, 0.46, 6.0, 0.34, [("삼성전자 메모리사업부", 18, True, BLUE)])
    tb(slide, 13.21, 0.46, 6.0, 0.34, [(GRADE, 18, False, GRAY)], align=PP_ALIGN.RIGHT)
    tb(slide, MX, 0.93, CW, 0.34, [(kicker, 18, True, BLUE)])
    tb(slide, MX, 1.33, CW, 0.62, [(title, 33, True, INK)])
    tb(slide, MX, 2.04, CW, 0.44, [(lead, 21, False, GRAY)])
    rect(slide, MX, 2.62, CW, 0.014, fill=LINE)


def footer(slide, source, no):
    tb(slide, MX, 10.49, 15.6, 0.34, [(source, 18, False, GRAY)])
    tb(slide, 17.55, 10.49, 1.66, 0.34, [(f"{no:02d} / {TOTAL:02d}", 18, False, GRAY)],
       align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, label, key, body, tint=False, label_color=BLUE,
         tag=None, tag_color=BLUE, key_h=0.42, pad=0.36):
    if tint:
        rect(slide, x, y, w, h, fill=TINT)
    else:
        rect(slide, x, y, w, h, fill=WHITE, line=LINE, line_w=0.75)
    ix, iw = x + pad, w - 2 * pad
    tb(slide, ix, y + 0.32, iw if tag is None else iw - 1.7, 0.36,
       [(label, 20.25, True, label_color)])
    if tag is not None:
        tb(slide, x + w - pad - 1.7, y + 0.32, 1.7, 0.36, [(tag, 18.75, True, tag_color)],
           align=PP_ALIGN.RIGHT)
    tb(slide, ix, y + 0.80, iw, key_h, [(key, 21.75, True, INK)])
    tb(slide, ix, y + 0.80 + key_h + 0.14, iw, h - (0.80 + key_h + 0.14) - 0.28,
       [(body, 18.75, False, GRAY)])


def band(slide, y, h, label, main, main_size=25.5, sub=None):
    rect(slide, MX, y, CW, h, fill=BLUE)
    tb(slide, MX + 0.50, y + 0.26, CW - 1.0, 0.32, [(label, 18.75, False, WHITE)])
    tb(slide, MX + 0.50, y + 0.66, CW - 1.0, 0.60, [(main, main_size, True, WHITE)])
    if sub:
        tb(slide, MX + 0.50, y + h - 0.86, CW - 1.0, 0.72, [(sub, 18.75, False, BLUE_T2)])


def stat(slide, x, y, w, num, label, num_color=BLUE):
    tb(slide, x, y, w, 0.46, [(num, 25.5, True, num_color)])
    tb(slide, x, y + 0.52, w, 0.66, [(label, 18.75, False, GRAY)])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# =========================================================
# S0. 요약 — 문제 · 원인 · 결론 · 실행 (한 장 논증)
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "삼성 SSD 전략적 방향성 · 요약",
       "서버 점유율을 위협하는 자체 SSD, 열린 원인은 워크로드이고 답은 FDP입니다",
       "문제와 원인, 결론, 실행 전략을 한 장으로 요약합니다. 상세 근거는 본편 4장에 있습니다.")

R0Y, R0H = 2.95, 3.42
# 문제 카드 — 좌 (가로 바 차트)
rect(s, MX, R0Y, 5.92, R0H, fill=WHITE, line=LINE, line_w=0.75)
p_ix = MX + 0.36
tb(s, p_ix, R0Y + 0.30, 5.20, 0.36, [("문제 · 자체 설계 SSD의 부상", 20.25, True, BLUE)])
tb(s, p_ix, R0Y + 0.74, 5.20, 0.80,
   [("서버 SSD 1위(38.2%)가 걸린 시장에서 고객이 경쟁자가 됩니다", 21.75, True, INK)])
BAR_X, BAR_MAX = p_ix + 1.55, 2.15
for (bi, (yr, val, frac, bfill)) in enumerate(
        [("2023", "5% 미만", 5 / 30, LINE), ("2026 전망", "30% 초과", 1.0, GRAY)]):
    by = R0Y + 1.72 + bi * 0.52
    tb(s, p_ix, by + 0.02, 1.45, 0.32, [(yr, 18, False, GRAY)])
    rect(s, BAR_X, by, BAR_MAX * frac, 0.32, fill=bfill)
    tb(s, BAR_X + BAR_MAX * frac + 0.12, by + 0.02, 1.55, 0.32, [(val, 18.75, True, INK)])
tb(s, p_ix, R0Y + 2.80, 5.20, 0.30,
   [("자체 설계 SSD 비중 전망 (가트너)", 18, False, GRAY)])
tb(s, p_ix, R0Y + 3.10, 5.20, 0.30,
   [("기업 저장 인프라 기준 · 웨이퍼 +246%", 18, False, GRAY)])

# 원인 카드 — 우 (표: 네 가지 원인 × 우리의 대응)
CX0, CW0 = 7.04, 12.17
rect(s, CX0, R0Y, CW0, R0H, fill=WHITE, line=LINE, line_w=0.75)
c_ix = CX0 + 0.36
c_iw = CW0 - 0.72
tb(s, c_ix, R0Y + 0.30, c_iw, 0.36, [("원인 · 고객이 자체 SSD를 만드는 네 가지 이유", 20.25, True, BLUE)])
tb(s, c_ix, R0Y + 0.76, c_iw, 0.42,
   [("셋은 계약의 몫이거나 우리 손 밖이고, 워크로드 최적화만 열려 있습니다", 21.75, True, INK)])
T1, T2, T3 = c_ix, c_ix + 5.15, c_ix + 7.65
tb(s, T1, R0Y + 1.32, 4.9, 0.30, [("원인", 18, False, GRAY)])
tb(s, T2, R0Y + 1.32, 2.3, 0.30, [("성격", 18, False, GRAY)])
tb(s, T3, R0Y + 1.32, 3.8, 0.30, [("우리의 대응", 18, False, GRAY)])
rect(s, c_ix, R0Y + 1.64, c_iw, 0.012, fill=LINE)
rows0 = [
    ("가격 (완제품 마진 회피)", "시황 의존", "장기 계약이 흡수 (기가동)", False),
    ("공급 안보", "구조적", "다년 계약·공급 신뢰", False),
    ("보안·수직 통합", "구조적", "대응 불가 영역 (양보)", False),
    ("워크로드 최적화 (WAF·수명·성능)", "구조적·확대 중", "FDP + 생태계 (본 전략)", True),
]
ry = R0Y + 1.74
for (r1, r2, r3, hot) in rows0:
    if hot:
        rect(s, c_ix - 0.10, ry - 0.05, c_iw + 0.20, 0.44, fill=TINT)
    tb(s, T1, ry, 4.9, 0.36, [(r1, 18.75, hot, BLUE if hot else INK)])
    tb(s, T2, ry, 2.3, 0.36, [(r2, 18.75, False, GRAY)])
    tb(s, T3, ry, 3.8, 0.36, [(r3, 18.75, hot, BLUE if hot else GRAY)])
    ry += 0.40

# 결론 스트립 (Samsung Blue)
rect(s, MX, 6.55, CW, 0.72, fill=BLUE)
tb(s, MX + 0.50, 6.55, 1.3, 0.72, [("결론", 18.75, False, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
tb(s, MX + 1.90, 6.55, 16.5, 0.72,
   [("커스텀 요구를 개별 대응이 아니라 표준(FDP)으로 흡수하고, 생태계로 완성하는 것이 최선입니다", 21.75, True, WHITE)],
   anchor=MSO_ANCHOR.MIDDLE)

# 실행 3카드
cards0 = [
    ("실행 1 · 전략적 제휴 패키지",
     "물량을 주고 워크로드를 받습니다",
     "장기 물량·가격 확약과 워크로드 협력을 한 계약으로 묶습니다. 1순위 구글, 선례는 마이크론과 앤트로픽입니다."),
    ("실행 2 · FDP 플랫폼",
     "제품은 하나, 최적화는 고객별로",
     "공통 펌웨어 위에 SDK·프로파일·공동 검증을 얹습니다. 기존 데이터센터 전담 조직의 강점을 그대로 씁니다."),
    ("실행 3 · 오픈소스 생태계",
     "LMCache 업스트림이 1호 과제입니다",
     "커널 6.16의 FDP 수용으로 열린 공백을 선점합니다. 성과는 실제 활성화된 FDP 용량으로 잽니다."),
]
for (lx, (label, key, body)) in zip((MX, 7.04, 13.29), cards0):
    card(s, lx, 7.50, 5.92, 2.68, label, key, body, tint=True)
footer(s, "출처: 삼성 SSD 전략적 방향성 보고서 v1.1 (2026년 8월) · 상세 근거는 본편 4장", 1)
notes(s, "한 장 논증 구조입니다. 문제: 자체 설계 SSD 확대가 서버 SSD 점유율(1위 38.2%)을 구조적으로 위협합니다. 원인: 가격·공급 안보·보안·워크로드 최적화 네 가지 중 앞의 셋은 계약이 흡수하거나 우리 손 밖이고, 워크로드 최적화만 우리가 풀 수 있는 표적입니다. 결론: 고객별 커스텀 대행은 발산하므로, 고객이 함께 설계한 표준 FDP로 흡수하고 시스템 소프트웨어 생태계로 완성하는 것이 최선입니다. 실행: 제휴 패키지(물량과 워크로드의 교환), FDP 플랫폼(공통 제품+SDK·검증), 오픈소스 생태계(LMCache FDP 백엔드 업스트림 1호 과제). 가트너 수치는 기업 저장 인프라 기준 참고이며, 하이퍼스케일러 캡티브 공식 통계는 없어 정황 지표(웨이퍼 +246%)를 병기했습니다.")

# =========================================================
# S1. 역사 — 세 번의 선택
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "역사",
       "지난 20년, 사이클의 진폭을 줄이는 같은 선택을 세 번 반복했습니다",
       "자기잠식을 무릅쓴 진출과 표준보다 빠른 실행, 데이터센터 집중이 오늘의 1위를 만들었습니다. 목적은 늘 하나였습니다.")

C3W, C3XS, C3Y, C3H = 5.92, (MX, 7.04, 13.29), 2.95, 4.52
cards1 = [
    ("첫 번째 선택 · 2005 진출",
     "HDD를 지키는 대신, HDD를 대체할 SSD를 먼저 만들었습니다",
     "2006년 세계 최초 SSD 양산으로 낸드의 초장기 수요처를 열었고, 2011년 HDD 사업을 약 14억 달러에 매각해 SSD로 일원화했습니다. 2013년 점유율 28.5%로 1위에 올랐습니다."),
    ("두 번째 선택 · 2013 표준 선행",
     "표준은 인텔이 주도했지만, 첫 제품은 삼성이 냈습니다",
     "업계 최초 NVMe SSD(XS1715)를 표준 주도자 인텔보다 약 1년 먼저 냈습니다. 빠른 결단과 실행의 동시 증명이며, 인텔은 이후 스토리지에서 퇴장했습니다."),
    ("세 번째 선택 · 2017 데이터센터 집중",
     "전담 조직과 기술 내재화로 서버 1위를 지키고 있습니다",
     "데이터센터 전용 기술 내재화와 전담 개발 조직의 결과로, 2026년 1분기 기업용 SSD 38.2% 1위와 엔비디아 인공지능 PC · GPU 서버 동시 채용으로 이어졌습니다."),
]
for (lx, (label, key, body)) in zip(C3XS, cards1):
    card(s, lx, C3Y, C3W, C3H, label, key, body, tint=True, key_h=1.02)

band(s, 7.79, 2.36, "공통 구조",
     "아직 아프지 않을 때 결단하고, 남보다 빨리 실행해, 부가가치를 한 계층 위로 옮겼습니다",
     sub="실증: 2019년 하강기에 반도체 영업이익이 69% 급감하는 동안, 스토리지 솔루션은 약 0.9조 원 적자(추정)에 그쳤습니다. 믹스가 완충재였습니다.")
footer(s, "출처: 스토리지 솔루션 사업사 리서치 · SSD 시장 전환기 리서치 · 기업용 SSD 2026년 1분기 집계 (2026년 8월)", 2)
notes(s, "삼성 SSD 20년은 우연이 아니라 같은 패턴의 세 번 반복입니다. 자기잠식을 무릅쓴 선제 결단, 표준 주도자보다 빠른 실행, 솔루션 부가가치. 세 선택 모두 목적은 메모리 사이클의 진폭을 줄이는 것이었고, 2019년 하강기에 완충 효과가 실측됐습니다.")

# =========================================================
# S2. 균열 — 호황의 역설
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "균열",
       "호황이 완충재를 줄이게 하고, 최대 고객을 경쟁자로 만들고 있습니다",
       "믹스는 서버로 쏠리고, 대형 클라우드 사업자는 낸드를 직접 사서 자체 SSD를 만듭니다. 하강기에는 두 균열이 동시에 작동합니다.")

C2Y, C2H, C2W = 2.95, 5.42, 9.04
# 좌 카드 — 믹스 쏠림
card(s, MX, C2Y, C2W, C2H,
     "균열 1 · 믹스 쏠림",
     "이익을 좇는 쏠림이 하강기의 완충재를 지웁니다",
     "기업용 SSD 계약가가 한 분기에 80% 오르자 산업 전체가 서버로 쏠리고, 마이크론은 소비자 시장에서 철수했습니다. 그런데 지금 줄이는 소비자 채널은 지난 하강기(2012년)에 심어 2013년 1위의 발판이 된 채널입니다.",
     tint=False)
stat(s, MX + 0.36, C2Y + 4.22, 3.9, "+80%", "기업용 SSD 계약가, 분기 상승률")
stat(s, MX + 4.60, C2Y + 4.22, 4.0, "철수", "마이크론 소비자 브랜드(크루셜)", num_color=INK)

# 우 카드 — Captive 4단계
RX = 10.17
rect(s, RX, C2Y, C2W, C2H, fill=WHITE, line=LINE, line_w=0.75)
tb(s, RX + 0.36, C2Y + 0.32, C2W - 0.72, 0.36,
   [("균열 2 · 자체 설계 SSD(Captive)의 확대", 20.25, True, BLUE)])
tb(s, RX + 0.36, C2Y + 0.80, C2W - 0.72, 0.42,
   [("통제권은 10년째 한 방향으로만 움직였습니다", 21.75, True, INK)])
steps = [
    ("2016년까지", "벤더 표준품을 그대로 구매했습니다"),
    ("2017년부터", "펌웨어를 직접 규정하기 시작했습니다"),
    ("2021년부터", "컨트롤러를 자체 설계했습니다 (아마존)"),
]
sy = C2Y + 1.40
for (yr, desc) in steps:
    tb(s, RX + 0.52, sy, 1.95, 0.36, [(yr, 18.75, True, BLUE)])
    tb(s, RX + 2.55, sy, C2W - 2.95, 0.36, [(desc, 18.75, False, GRAY)])
    sy += 0.58
rect(s, RX + 0.36, sy - 0.06, C2W - 0.72, 0.78, fill=TINT)
tb(s, RX + 0.52, sy + 0.04, 1.95, 0.36, [("2022년부터", 18.75, True, BLUE)])
tb(s, RX + 2.55, sy + 0.04, C2W - 3.05, 0.66,
   [("표준을 함께 만들고 웨이퍼를 직접 삽니다 (구글 · 메타)", 18.75, True, INK)])
stat(s, RX + 0.36, C2Y + 4.22, 3.9, "+246%", "낸드 웨이퍼 계약가, 2025년 1분기 대비")
stat(s, RX + 4.60, C2Y + 4.22, 4.1, "30% 초과", "자체 설계 SSD 비중 전망 (가트너 · 기업 저장 인프라 기준)")

band(s, 8.72, 1.42, "이중 리스크",
     "하강기에는 돌아갈 소비자 시장이 줄어 있고, 자체 설계로 간 물량은 돌아오지 않습니다")
footer(s, "출처: Captive SSD 리서치 · 가트너 재인용 검증, 기업 저장 인프라 기준 (2026년 8월)", 3)
notes(s, "쏠림 자체가 아니라 되돌릴 수 없는 방식의 쏠림이 문제이며, 소비자 채널은 외주 생산 활용으로 유지 비용을 낮추는 별도 트랙이 진행 중입니다(복귀 통로 유지). 통제권 상승은 10년간 한 방향이었지만, 되돌아오지 않는 것은 역량이고 물량 배분은 TCO가 결정하는 가역 영역입니다 — 우리의 싸움은 배분의 방어입니다. 가트너 전망의 원 수치는 2023년 5% 미만에서 2026년 30% 초과이며 기업(온프레미스) 저장 인프라 기준입니다. 대형 클라우드 자체 설계 비중의 공식 통계는 없어 정황 지표(웨이퍼 직구매 +246%, 아마존 Nitro, 구글 Titanium)를 주근거로 봅니다.")

# =========================================================
# S3. 네 번째 방향성 — 커스텀의 표준화 (FDP)
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "네 번째 방향성",
       "모든 커스텀을 받아주는 대신, 커스텀 요구를 표준으로 흡수합니다",
       "고객이 원하는 것은 자기 워크로드 최적화입니다. 그 요구를 표준(FDP)으로 받고, 시스템 소프트웨어 생태계로 완성합니다.")

# 좌: 두 갈래
card(s, MX, 2.95, 9.04, 2.75,
     "길 1 · 풀커스텀 대행",
     "고객마다 만들면 개발과 평가, 유지보수가 발산합니다",
     "고객 수만큼 펌웨어가 갈라지고 검증 부담이 곱으로 늘어납니다. 커스텀은 계약 체질부터 다른 사업이며, 보상 계약 없이 비용을 떠안은 수업료가 이미 있습니다(사내 확인).",
     tint=False, label_color=GRAY, tag="지속 불가", tag_color=RED)
card(s, MX, 5.90, 9.04, 2.75,
     "길 2 · 표준화된 유연성",
     "FDP 표준 제품 하나로 고객별 최적화를 받아냅니다",
     "펌웨어는 공통으로 지키고, 고객별 차이는 설정과 소프트웨어로 받습니다. FDP는 메타와 구글이 요구하고 삼성이 함께 완성해 2023년 비준한, 고객이 설계한 표준입니다.",
     tint=True, tag="선택", tag_color=BLUE)
tb(s, MX, 8.85, 9.04, 0.80,
   [("성립 조건: 리눅스와 오픈소스 응용 계층에 FDP 지원을 심는 시스템 소프트웨어 생태계 확산이 함께 가야 합니다.", 18.75, False, GRAY)])

# 우: 부가가치 계층 스택
tb(s, RX, 2.95, C2W, 0.36, [("부가가치는 표준 위 계층으로 옮겨 갑니다", 20.25, True, BLUE)])
layers = [
    ("시스템 소프트웨어 · 생태계", "부가가치의 새 자리, 우리가 만들 층", BLUE, WHITE, WHITE),
    ("FDP 표준", "고객과 공동 설계, 2023년 비준", BLUE_T2, INK, INK),
    ("SSD 하드웨어", "여러 공급사가 함께 서는 층", None, INK, GRAY),
    ("낸드 웨이퍼", "고객 직구매가 늘어나는 층", None, INK, GRAY),
]
LY, LH, LGAP, LW = 3.55, 1.02, 0.14, 7.55
for i, (t1, t2, fill, c1, c2) in enumerate(layers):
    y = LY + i * (LH + LGAP)
    if fill is None:
        rect(s, RX, y, LW, LH, fill=WHITE, line=LINE, line_w=0.75)
    else:
        rect(s, RX, y, LW, LH, fill=fill)
    tb(s, RX + 0.36, y + 0.14, LW - 0.72, 0.38, [(t1, 19.5, True, c1)])
    tb(s, RX + 0.36, y + 0.55, LW - 0.72, 0.34, [(t2, 18, False, c2)])
rect(s, RX + LW + 0.42, LY + 0.55, 0.52, 3.40, fill=BLUE_T2, shape=MSO_SHAPE.UP_ARROW)
tb(s, RX + LW + 0.10, LY + 4.15, 1.3, 0.70, [("부가가치", 18, False, GRAY), ("이동", 18, False, GRAY)],
   align=PP_ALIGN.CENTER)
tb(s, RX, 8.66, C2W, 1.10,
   [("FDP 지원만으로는 여러 공급사 중 하나입니다. 마이크론과 키옥시아도 이미 지원합니다. 차별화는 워크로드를 표준 정책으로 바꿔 주는 소프트웨어와 현장 엔지니어링에서 생깁니다.", 18.75, False, GRAY)])
footer(s, "출처: FDP 플랫폼 전략 · 구글 Captive와 FDP 팩트체크 · FDP 오픈소스 생태계 리서치 (2026년 8월)", 4)
notes(s, "풀커스텀은 고객당 부채를 쌓는 구조라 지속 불가능합니다. 왜 이번에는 다른가: 앞 세대 배치 표준(Streams·ZNS)은 호스트 소프트웨어 부담으로 실패했지만, FDP는 힌트를 주지 않아도 동작하는 하위 호환 설계라 호스트 복잡도가 구조적으로 낮고, 수요자(구글·메타)가 설계했으며, 리눅스 커널 6.16이 블록 write streams로 수용해 메인라인 진입을 시작했습니다. 남은 공백은 응용·AI 계층(LMCache 등) — 오픈소스 조직 투자의 표적입니다. 효과가 워크로드 의존적이라는 한계는 곧 고객 협업이 전략의 핵심 부품이라는 증거입니다. NVIDIA의 CMX·SCADA는 GPU 직결 스토리지라는 다른 응용 분야로, 호스트 경로 표준인 FDP와 구분해 병행합니다.")

# =========================================================
# S4. 실행과 결론
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "실행과 결론",
       "공급자 우위인 지금, 물량 계약을 공동 플랫폼 계약으로 격상합니다",
       "장기 물량과 워크로드 협력을 한 계약으로 묶고, 생태계 조직을 강화합니다. 목표는 세 번째 비중 전환입니다.")

LW4 = 11.00
card(s, MX, 2.95, LW4, 3.00,
     "제휴 패키지 · 문을 여는 법",
     "물량을 주고, 워크로드를 받습니다",
     "공급 부족 국면에 장기 물량과 가격 확약을 주고, 워크로드 정보와 공동 로드맵을 계약된 권리로 받습니다. 마이크론과 앤트로픽이 2026년 6월 같은 구조의 계약을 맺었고, 삼성도 앤트로픽 투자에 이미 참여하고 있습니다. 1순위는 FDP를 함께 설계해 온 구글입니다.",
     tint=True)
card(s, MX, 6.15, LW4, 2.70,
     "조직 · 두 트랙과 결정 요청",
     "제품은 기존 조직으로, 생태계는 새 조직으로 갑니다",
     "FDP 제품 개발은 데이터센터 전담 조직의 강점을 그대로 씁니다. 오픈소스 기여자와 메인테이너, 고객 현장 상주 엔지니어는 새로 키워야 하며, 이 증원이 이번 보고의 결정 요청입니다.",
     tint=False)

RX4, RW4 = 12.12, 7.09
tb(s, RX4, 2.95, RW4, 0.36, [("세 번째 비중 전환", 20.25, True, BLUE)])
phases = [
    ("① 낸드를 SSD로", "2005년부터 · 완제품으로 올라섰습니다", False),
    ("② SSD를 서버 SSD로", "2017년부터 · 점유율 38.2% 1위", False),
    ("③ 서버 SSD를 FDP SSD로", "지금 · 표준과 생태계로 지킵니다", True),
]
PY, PH, PGAP = 3.48, 0.98, 0.34
for i, (t1, t2, hot) in enumerate(phases):
    y = PY + i * (PH + PGAP)
    if hot:
        rect(s, RX4, y, RW4, PH, fill=BLUE)
        tb(s, RX4 + 0.36, y + 0.12, RW4 - 0.72, 0.38, [(t1, 19.5, True, WHITE)])
        tb(s, RX4 + 0.36, y + 0.53, RW4 - 0.72, 0.34, [(t2, 18, False, WHITE)])
    else:
        rect(s, RX4, y, RW4, PH, fill=WHITE, line=LINE, line_w=0.75)
        tb(s, RX4 + 0.36, y + 0.12, RW4 - 0.72, 0.38, [(t1, 19.5, True, INK)])
        tb(s, RX4 + 0.36, y + 0.53, RW4 - 0.72, 0.34, [(t2, 18, False, GRAY)])
    if i < 2:
        rect(s, RX4 + RW4 / 2 - 0.13, y + PH + 0.05, 0.26, 0.24,
             fill=BLUE_T2, shape=MSO_SHAPE.DOWN_ARROW)
stat(s, RX4, 7.55, RW4, "2027년 이전",
     "실행 창 · 경쟁사가 생태계 자리를 선점하기 전", num_color=INK)

band(s, 9.02, 1.28, "결론",
     "혼자 잘하는 회사에서 고객과 함께 잘하는 회사로, 지금 심어야 다음 하강기에 부드럽게 내려갑니다",
     main_size=24.0)
footer(s, "출처: 삼성 SSD 전략적 방향성 보고서 v1.1 · 마이크론 앤트로픽 계약 공시 (2026년 8월)", 5)
notes(s, "워크로드 정보는 전략의 원료인데 고객은 쉽게 내주지 않습니다. 공급 부족인 지금이 물량과 맞바꿀 유일한 계절이며, 체결 데드라인은 다운턴 조기경보와 연동해 관리합니다. 마이크론 앤트로픽 계약은 구조의 선례이고(워크로드 공유 깊이는 비공개), 고객 유형에 따라 구글형(직접 계약)과 AI 네이티브형(3자 구조)으로 설계가 갈립니다. 수익화 원칙: 소프트웨어는 기본 탑재 무과금, 회수는 하드웨어 점유율과 가격 결정력. 결정 요청 두 가지와 실행 과제 6건(KV Cache FDP 업스트림 1호 과제 포함)은 보고서 부록 D를 참조합니다.")

prs.save(os.path.abspath(OUT))
print(f"생성 완료: {os.path.abspath(OUT)} ({len(prs.slides._sldIdLst)}장)")
