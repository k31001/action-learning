# -*- coding: utf-8 -*-
"""지난 20년 메모리 다운턴 복기 덱 (요약 1장 + 다운턴별 보충 5장).

디자인 시스템: ssd-strategy.pptx / storyline-overview.pptx 승계
  20 x 11.25 in 캔버스 / Arial 단일 폰트 / Samsung Blue #1428A0 단일 액센트
  헤더(조직명·문서등급·킥커·33pt 액션 타이틀·21pt 리드·헤어라인) / 푸터(출처·페이지)
  틴트 카드 #F4F6FC(무테) · 아웃라인 카드 흰색+#D9D9D9 0.75pt · 다크 블루 정리 밴드

실행: .venv/bin/python outputs/presentation/scripts/generate_downturn_review_pptx.py
  (사전에 generate_downturn_assets.py 로 차트 PNG 생성 필요)
출력: outputs/presentation/downturn-review.pptx
콘텐츠 소스: wiki/downturn/downturn-history.md
             sources/raw-notes/memory-downturn-history-research-2026-08-22.md
             sources/raw-notes/nand-market-history-research-2026-08-22.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLUE = RGBColor(0x14, 0x28, 0xA0)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xD9, 0xD9, 0xD9)
TINT = RGBColor(0xF4, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

MX = 0.79
CW = 18.42
GRADE = "[문서등급 표기]"
TOTAL = 6

BASE = os.path.dirname(__file__)
ASSETS = os.path.join(BASE, "..", "assets")
OUT = os.path.join(BASE, "..", "downturn-review.pptx")

prs = Presentation()
prs.slide_width = Emu(18288000)   # 20.00 in
prs.slide_height = Emu(10287000)  # 11.25 in
BLANK = prs.slide_layouts[6]


def _font(run, size, bold, color):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
       wrap=True, spacing=1.15):
    """paras: [(text, size, bold, color)] 항목당 문단 1개."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, size, bold, color) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = t
        _font(r, size, bold, color)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def header(slide, kicker, title, lead):
    tb(slide, MX, 0.46, 8.0, 0.34, [("삼성전자 메모리사업부", 18, True, BLUE)])
    tb(slide, 13.21, 0.46, 6.0, 0.34, [(GRADE, 18, False, GRAY)], align=PP_ALIGN.RIGHT)
    tb(slide, MX, 0.93, CW, 0.34, [(kicker, 18, True, BLUE)])
    tb(slide, MX, 1.33, CW, 0.62, [(title, 33, True, INK)])
    tb(slide, MX, 2.04, CW, 0.44, [(lead, 21, False, GRAY)])
    rect(slide, MX, 2.62, CW, 0.014, fill=LINE)


def footer(slide, source, no):
    tb(slide, MX, 10.49, 15.6, 0.34, [(source, 15, False, GRAY)])
    tb(slide, 17.55, 10.49, 1.66, 0.34, [(f"{no:02d} / {TOTAL:02d}", 18, False, GRAY)],
       align=PP_ALIGN.RIGHT)


def pic(slide, name, x, y, w):
    path = os.path.join(ASSETS, name)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


SRC_MAIN = ("출처: wiki/downturn/downturn-history.md · memory-downturn·nand-market-history-"
            "research-2026-08-22.md (기관 혼합 시계열) · * 역산·추정")

# ================= S1 요약 =================
s = prs.slides.add_slide(BLANK)
header(s, "메모리 다운턴 복기 | 2006-2025 · 5건",
       "지난 20년 다운턴 5건: 공급발은 길었고, 수요발은 짧지만 갈수록 빠르고 깊어졌다",
       "DRAM+NAND 합산. NAND는 침투기(DT12·16)엔 완충, DT19부터 완전 동조: 전체 산업이 한 몸으로 떨어진다")

pic(s, "downturn_timeline.png", MX, 2.82, 10.85)
pic(s, "downturn_scatter.png", 11.95, 2.82, 7.26)

# ---- 비교 표 ----
TY = 6.90
COLS = [(MX, 1.95), (2.74, 1.85), (4.59, 4.60), (9.19, 1.85), (11.04, 4.10), (15.14, 4.07)]
HEADS = ["다운턴", "발원", "원인", "낙폭(전체·D·N)", "핵심 대응", "결과"]
for (cx, cw), htxt in zip(COLS, HEADS):
    tb(s, cx, TY, cw - 0.18, 0.28, [(htxt, 15, True, GRAY)])
rect(s, MX, TY + 0.32, CW, 0.014, fill=LINE)

ROWS = [
    ("DT08", "'07-'09 · 9분기*", ("공급발+수요충격", GRAY),
     "6강 캐파 경쟁 가격 붕괴 + 금융위기 수요 급정지",
     "-26%*", "D -34 · N -14",
     "삼성 무감산 버티기, 40nm 최초, 위기 후 투자 배증",
     "Qimonda 파산, 업계 손실 $7B, 삼성 2010 최대 이익"),
    ("DT12", "'10-'12 · 9분기*", ("공급발(침식)", GRAY),
     "대만·엘피다 증산 + PC 부진, 2년 저가 지속",
     "-19%*", "D -33 · N -2",
     "Line-16 12조, 20nm·모바일 전환 선행, HDD 매각",
     "Elpida 파산, 6강에서 3강으로, Micron 저점 인수"),
    ("DT16", "'15-'16 · 6분기*", ("수요발(침식)", BLUE),
     "PC 출하 감소·스마트폰 둔화, 20nm 공급 증가",
     "+2%*", "D -11 · N +14",
     "증설 대신 3D NAND 전환, Micron Inotera $4.1B",
     "파산 0, 2017 슈퍼사이클 직결, 중국 진입 결정"),
    ("DT19", "'18-'19 · 5분기*", ("수요발(급락)", BLUE),
     "서버·모바일 재고 조정, 미중 분쟁 증폭",
     "-34%*", "D -37.6 · N -27",
     "Micron·SK 감산, 삼성 무감산·CapEx 유지·완주",
     "5분기 최단 V자, 1Q20 DRAM 44.1%, HBM 축소 씨앗"),
    ("DT23", "'22-'23 · 6분기*", ("수요발(침식→급락)", BLUE),
     "팬데믹 특수 소멸·금리, 재고 인지 실패 6개월",
     "-45%*", "D -45 · N -45",
     "무감산 선언 후 감산 선회(23-04), SK는 HBM 집중",
     "DS 최대 적자 -14.9조, 사상 최속, HBM 주도권 SK로"),
]
RH = 0.52
ry = TY + 0.40
for name, era, (origin, ocolor), cause, sc1, sc2, action, result in ROWS:
    tb(s, COLS[0][0], ry + 0.02, COLS[0][1] - 0.18, RH,
       [(name, 18, True, INK), (era, 13.5, False, GRAY)], spacing=1.02)
    tb(s, COLS[1][0], ry + 0.02, COLS[1][1] - 0.18, RH, [(origin, 15, True, ocolor)],
       spacing=1.02)
    tb(s, COLS[2][0], ry + 0.02, COLS[2][1] - 0.22, RH, [(cause, 15, False, INK)],
       spacing=1.02)
    tb(s, COLS[3][0], ry + 0.02, COLS[3][1] - 0.18, RH,
       [(sc1, 16.5, True, INK), (sc2, 13.5, False, GRAY)], spacing=1.02)
    tb(s, COLS[4][0], ry + 0.02, COLS[4][1] - 0.22, RH, [(action, 15, False, INK)],
       spacing=1.02)
    tb(s, COLS[5][0], ry + 0.02, COLS[5][1] - 0.10, RH, [(result, 15, False, INK)],
       spacing=1.02)
    ry += RH
    if name != "DT23":
        rect(s, MX, ry - 0.05, CW, 0.01, fill=LINE)

# ---- 인사이트 밴드 ----
BY = ry + 0.06
rect(s, MX, BY, CW, 0.56, fill=BLUE)
band_items = [
    "① 발원이 기간 결정: 공급 9, 수요 5-6분기",
    "② 낙폭 확대: -26 → -34 → -45% (전체)",
    "③ NAND: 침투기 완충 → DT19부터 동조",
    "④ 직전 성공이 다음 함정: 2019 무감산·HBM",
]
bw = CW / 4
for i, item in enumerate(band_items):
    tb(s, MX + 0.30 + i * bw, BY, bw - 0.45, 0.56, [(item, 15, True, WHITE)],
       anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
footer(s, SRC_MAIN, 1)

# ================= S2-S6 다운턴별 보충 =================


def kpi_tile(slide, x, y, w, h, label, value, vsize=27, vcolor=INK):
    rect(slide, x, y, w, h, fill=TINT)
    tb(slide, x + 0.24, y + 0.14, w - 0.48, 0.26, [(label, 15, False, GRAY)])
    tb(slide, x + 0.24, y + 0.42, w - 0.48, h - 0.5, [(value, vsize, True, vcolor)])


def outline_card(slide, x, y, w, h, label):
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, line_w=0.75)
    tb(slide, x + 0.36, y + 0.24, w - 0.72, 0.34, [(label, 18.75, True, BLUE)])


def bullets(slide, x, y, w, h, items, size=17.25, spacing=1.12):
    tb(slide, x, y, w, h, [("·  " + t, size, False, INK) for t in items],
       spacing=spacing)


def detail_slide(no, kicker, title, lead, kpis, cause, resp_l, resp_r, outcome,
                 lesson, mini, field, source):
    sl = prs.slides.add_slide(BLANK)
    header(sl, kicker, title, lead)
    # 좌: KPI 2x2  (kpis = [지속, 낙폭, 최악분기 라벨, 최악분기 값, 발원])
    tw, th, gap = 2.67, 1.06, 0.16
    dur, depth, wq, wpct, origin_v = kpis
    tiles = [("하강 지속기간", dur, 27, INK),
             ("연매출 낙폭 (DRAM)", depth, 27, INK),
             (f"최악 분기 ({wq}, QoQ)", wpct, 27, INK),
             ("발원·형태", origin_v, 18, BLUE)]
    for i, (lab, val, vs, vc) in enumerate(tiles):
        cx = MX + (i % 2) * (tw + gap)
        cy = 2.82 + (i // 2) * (th + gap)
        kpi_tile(sl, cx, cy, tw, th, lab, val, vsize=vs, vcolor=vc)
    # 좌: 미니 차트 + 현장 수치
    pic(sl, mini, MX, 5.24, 5.50)
    rect(sl, MX, 8.32, 5.50, 2.03, fill=TINT)
    tb(sl, MX + 0.28, 8.56, 5.50 - 0.56, 0.30, [("현장 수치", 15.75, True, BLUE)])
    tb(sl, MX + 0.28, 8.94, 5.50 - 0.56, 1.3, [(field, 16.5, False, INK)], spacing=1.18)
    # 우: 카드 3단
    RX, RW = 6.55, 12.66
    outline_card(sl, RX, 2.82, RW, 2.28, "원인·전개")
    bullets(sl, RX + 0.36, 3.42, RW - 0.72, 1.55, cause)
    outline_card(sl, RX, 5.24, RW, 2.94, "기업 대응")
    subw = (RW - 3 * 0.36) / 2
    lx = RX + 0.36
    rx2 = RX + 0.36 + subw + 0.36
    tb(sl, lx, 5.82, subw, 0.30, [(resp_l[0], 16.5, True, INK)])
    bullets(sl, lx, 6.20, subw, 1.85, resp_l[1:], size=16.5, spacing=1.12)
    tb(sl, rx2, 5.82, subw, 0.30, [(resp_r[0], 16.5, True, GRAY)])
    bullets(sl, rx2, 6.20, subw, 1.85, resp_r[1:], size=16.5, spacing=1.12)
    rect(sl, RX + 0.36 + subw + 0.18, 5.82, 0.011, 2.10, fill=LINE)
    outline_card(sl, RX, 8.32, RW, 2.03, "결과·교훈")
    bullets(sl, RX + 0.36, 8.90, RW - 0.72, 0.85, outcome)
    tb(sl, RX + 0.36, 9.78, RW - 0.72, 0.45, [(lesson, 17.25, True, INK)])
    footer(sl, source, no)


detail_slide(
    2, "메모리 다운턴 복기 | Appendix A · DT08 (2007-2009)",
    "DT08: 6강 캐파 경쟁에 금융위기가 겹쳐 9분기, 업계의 이익이 소멸했다",
    "공급발로 시작해 수요 충격으로 끝난 복합 다운턴. 버틴 1위와 소진된 5위가 갈렸다",
    ["~9분기*", "-34%*", "4Q08", "-36%", "공급발+수요충격"],
    ["6강 체제의 12인치 증설 경쟁: 가격 2007년 -85%, 2008년 -58% 원가 이하 붕괴",
     "2008 Q4 금융위기 수요 급정지: 분기 매출 -36%, 단일 분기 사상 최대 낙폭",
     "업계 전원 적자: 2008년 DRAM 업계 합산 순손실 $7B",
     "NAND도 동반 하강: 2008 매출 -14%($13.9B→$12.0B), DRAM보다는 얕음"],
    ["버틴·확장한 쪽 (삼성)",
     "무감산 버티기 + 40nm DDR3 세계 최초 양산(생산성 +60%)",
     "2009-01 DS/DMC 2부문 통합, 임원 연봉 -20%",
     "퇴출 확인 후 2010 메모리 투자 5.5조 → 9조 배증"],
    ["밀린·퇴출된 쪽",
     "Qimonda: 누적손실 $30억, 정부 지원에도 2009-01 파산",
     "하이닉스: 채권단 관리, 2009년 투자 ~1조* 동결",
     "Micron: FY09 CapEx -81%, 사실상 투자 중단"],
    ["퇴출 직후 현물가 급등: 공급자 1개 퇴출이 가격 균형을 바꾸는 과점 실증",
     "삼성 2009 영업이익 10.9조, 2010년 당시 사상 최대 실적 수확"],
    "교훈: 공급발은 원가·현금의 체력전. 회복은 구조 변화(퇴출)와 함께 온다",
    "downturn_mini_dt08.png",
    "512Mb DDR2 $6.8 → $0.5 (-93%)\n2010 반도체 투자 12.7조(메모리 9조)",
    "출처: wiki/downturn/downturn-history.md · dram-chicken-game-history-2026-08-05.md · "
    "samsung-downturn-actions-2007-2023-2026-08-07.md · * 역산·추정")

detail_slide(
    3, "메모리 다운턴 복기 | Appendix B · DT12 (2010-2012)",
    "DT12: 급락 없는 2년의 저가가 엘피다를 소진시키고 3강 과점을 완성했다",
    "공급발 침식형. 수요의 축이 PC에서 모바일로 넘어가는 심판대를 겸했다",
    ["~9분기*", "-33%*", "4Q10", "-20%", "공급발·침식형"],
    ["회복 1년 만의 재하강: 대만·엘피다 증산 + PC 부진 + 유럽 재정위기",
     "DDR3 2Gb 가격 2011년 한 해 -85%, 2011년 산업 매출 -25%",
     "낙차 대신 지속: 저가 2년이 재무 체력을 소진시키는 침식형",
     "NAND 가격도 원가 이하(31¢/GB): Toshiba 30% 감산(2012-07) 후 +20%"],
    ["버틴·확장한 쪽 (삼성)",
     "Line-16 12조 착공·가동, 30nm → 20nm 세계 최초 연속",
     "모바일 전환 선행: LPDDR3 세계 최초, Austin 로직 전환",
     "HDD 매각($1.375B) 경량화. 엘피다 입찰은 불참"],
    ["밀린·퇴출된 쪽",
     "Elpida: PC → 모바일 전환 실패 + 엔고, 부채 4,480억 엔 파산(2012-02)",
     "대만 진영(파워칩·프로모스·난야) 범용 DRAM 축차 퇴장",
     "하이닉스: SK그룹 인수(2012)로 채권단 관리 종결"],
    ["6강에서 3강으로. 이후 모든 사이클의 구조적 전제가 된 과점 완성",
     "Micron, 엘피다를 ~$2.5B에 저점 인수(2013): 모바일 스케일 점프"],
    "교훈: 침식형은 체력전이자 전환 심판대. 세대·수요 전환 완주가 생사를 갈랐다",
    "downturn_mini_dt12.png",
    "DDR3 2Gb 2011년 -85%\nNAND 매출은 보합(-2%*): 모바일 침투가 완충",
    "출처: wiki/downturn/downturn-history.md · dram-chicken-game-history-2026-08-05.md · "
    "nand-market-history-research-2026-08-22.md · * 역산·추정")

detail_slide(
    4, "메모리 다운턴 복기 | Appendix C · DT16 (2015-2016)",
    "DT16: 3강 과점의 첫 시험은 얕게 끝났지만, 중국 진입이 이 창에서 결정됐다",
    "수요발 침식형, 5건 중 최소 낙폭. 규율이 작동했고 유산은 시장 밖에서 왔다",
    ["~6분기*", "-11%*", "4Q15", "-9.1%", "수요발·침식형"],
    ["PC 출하 감소 + 스마트폰 성장 둔화, 20nm 전환발 공급 증가",
     "PC DRAM 계약가 2016년 -34%, 4GB 모듈 저점 $12.5",
     "2H15 재고발 가격 인하 사이클, 저점은 1H16",
     "NAND는 비동행: 3D 전환 제약으로 +14%*, 전체 메모리는 +2%* 유지"],
    ["3강의 규율",
     "증설 대신 3D NAND 전환 투자로 이동: DRAM 공급 억제",
     "파산·구제 없음: 가격 규율로 흡수한 첫 다운턴",
     "Micron, Inotera 잔여 지분 $4.1B 저점 인수(2016-12)"],
    ["시장 밖의 결정 (중국)",
     "칭화유니, Micron $23B 인수 시도(2015-07) 미국 반대로 무산",
     "중국, 자국 팹 건설로 선회: YMTC·CXMT 설립(2016)",
     "오늘의 CXMT 저가 잠식 리스크가 여기서 출발"],
    ["2016 H2 반등, 2017-18 슈퍼사이클(+77%, +38%)로 직결",
     "공급 규율의 성공 경험이 3강 절제 균형의 학습 사례로 축적"],
    "교훈: 과점은 얕은 수요 조정을 흡수한다. 최대 유산은 진입자의 결정이었다",
    "downturn_mini_dt16.png",
    "PC DRAM 계약가 -34% (2016, YoY)\nDRAM 공급 억제 = 2017 부족의 씨앗",
    "출처: wiki/downturn/downturn-history.md · memory-downturn·nand-market-history-"
    "research-2026-08-22.md · * 역산·추정")

detail_slide(
    5, "메모리 다운턴 복기 | Appendix D · DT19 (2018-2019)",
    "DT19: 재고 조정만으로 -37.6%, 5분기 최단. 무감산 1위가 점유율로 보상받았다",
    "수요발 급락형. 공급 규율은 유지됐고 회복은 재고 소진과 함께 왔다",
    ["~5분기*", "-37.6%", "4Q18", "-18.3%", "수요발·급락형"],
    ["슈퍼사이클 중 과발주한 데이터센터·스마트폰 고객의 재고 조정",
     "미중 분쟁·Huawei 제재가 수요 불확실성 증폭",
     "DRAM 고정가 연간 -50% 이상, 삼성 반도체 영업이익 -69%",
     "NAND 동반 -27%*(4Q18 -16.8% 동시 진입), 정전·감산으로 2H19 먼저 회복"],
    ["삼성: 무감산 + 완주",
     "인위적 감산 없음 + CapEx 22.6조(-5%) 규모 유지",
     "세대 전환 완주: 1z DRAM·6세대 V낸드·EUV 준비·HBM2E 발표",
     "비전 2030 선언. 동시에 HBM팀 축소(니치 판단)"],
    ["경쟁사: 명시적 감산",
     "Micron 최초 공식 감산(2019-03, 웨이퍼 -5%)",
     "SK하이닉스(2019-07): NAND -15%, D램 캐파 축소·CIS 전환",
     "Kioxia/WD 요카이치 정전(2019-06): 비자발 공급 축소"],
    ["4Q19 바닥, 2020 COVID 클라우드 수요로 V자 회복",
     "삼성 1Q20 DRAM 점유율 44.1%: 무감산·투자 유지의 보상"],
    "교훈: 재고발은 스스로 끝난다. 그러나 이 성공 공식이 2023 오판의 근거가 됐다",
    "downturn_mini_dt19.png",
    "산업 -37.6%: 당시 기준 최대 낙폭\n삼성 반도체 OP 44.6 → 14.0조 (-69%)",
    "출처: wiki/downturn/downturn-history.md · samsung-2019-downturn-2017-2019-actions-"
    "2026-08-16.md · nand-market-history-research-2026-08-22.md · * 역산·추정")

detail_slide(
    6, "메모리 다운턴 복기 | Appendix E · DT23 (2022-2023)",
    "DT23: 인지 실패 6개월이 사상 최속·최심을 만들고, 회복의 과실은 배분이 갈랐다",
    "수요발이 침식형으로 시작해 급락형으로 끝났다. 원인 오판이 대응을 무효화한 사례",
    ["~6분기*", "-45%*", "4Q22", "-32.5%", "수요발·침식→급락"],
    ["팬데믹 특수 소멸 + 금리 인상, 고객 재고 대조정",
     "'일시적' 해석 6개월: DS 재고 16.5조 → 29.1조(+76.6%) 방치",
     "3Q22 -28.9%, 4Q22 -32.5% 연속: 2008년 이후 최대 분기 낙폭",
     "NAND도 -45%($67.1B→$36.7B) 동일 낙폭: D·N 완전 동조화"],
    ["삼성: 뒤늦은 선회",
     "무감산 선언(22-10) → 재확인(23-01) → 감산 공식화(23-04-07)",
     "그 사이 1Q23 DS -4.58조, 연간 -14.88조 사상 최대 적자",
     "CapEx 48.4조·R&D 28.34조 사상 최대 유지(HBM·DDR5 재배치)"],
    ["경쟁사: 급축 + 집중",
     "SK하이닉스 CapEx -56.2%: HBM 선택·집중",
     "Micron CapEx -42%, 3사 전원 감산 국면",
     "소모전 교범 불발: 전원 체력 보유 과점에서 퇴출 없음"],
    ["감산 + AI/HBM 수요로 2Q23 반등(+20.4%), 2024년 +75%",
     "HBM 주도권은 SK로: 2019년 축소 결정의 대가 지불"],
    "교훈: 원인·형태 오판은 대응 전체를 무효화한다. 총량 유지보다 배분이 승부",
    "downturn_mini_dt23.png",
    "분기 정점 → 저점 -62%* (3분기): 사상 최속\n2Q22 $25.6B* → 1Q23 $9.7B*",
    "출처: wiki/downturn/downturn-history.md · samsung-downturn-actions-2007-2023-"
    "2026-08-07.md · nand-market-history-research-2026-08-22.md · * 역산·추정")

prs.save(OUT)
print("saved:", os.path.normpath(OUT))
