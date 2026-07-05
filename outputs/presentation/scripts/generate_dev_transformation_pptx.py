"""
개발실 체질 전환 발표자료 — 3슬라이드 · IBM Carbon Design System · 고밀도

빌드 원천 — `outputs/presentation/dev-transformation-outline.md`
  슬라이드 1  WHY   — LTA→SCA 사건 + 신문섭 북극성 명제 + 계약 3단 진화 + Micron 4요소 + 타임라인
  슬라이드 2  WHAT  — As-Is→To-Be 7행 + Palantir FDE 벤치마크 + 리스크↔이점
  슬라이드 3  HOW   — 4대 축 + 3-Phase 로드맵 + KPI

설계 원칙: 여백 최소화 + 각진(sharp) 카드 + 얇은 헤어라인 그리드. 콘텐츠가 슬라이드를 꽉 채우도록 확장.
타이포그래피 스케일(점진적 위계, pt): title 23 · section 14 · cardhead 13 · body 12 · sub 11 · caption 10 · num 22 · data 13.5
디자인: IBM Carbon — Blue-60 #0f62fe(단일 액센트) · Yellow-30 · Red-60 · Green-50 ·
        Gray-100/70/20/10 · Blue-10 · IBM Plex Sans(KR) + IBM Plex Mono(숫자).
        팔레트 토큰은 `generate_carbon_pptx.py`와 동일.

기존 `generate_pptx.py` helper 재사용하되 generate_pptx.THEME 를 Carbon 값으로 in-place 치환.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import generate_pptx as G
from generate_pptx import (
    add_text, add_rect, add_line, add_footer, remove_existing_slides,
    SLIDE_W, SLIDE_H, TEMPLATE,
)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
OUTPUT = os.path.join(ROOT, 'presentation', 'dev-org-transformation.pptx')

TOTAL = 3

# ============================================================
# IBM Carbon 팔레트 — generate_pptx.THEME in-place 치환
# ============================================================
CARBON = {
    'samsung_blue':  RGBColor(0x0F, 0x62, 0xFE),  # Blue-60 (단일 액센트)
    'deep_navy':     RGBColor(0x16, 0x16, 0x16),  # Gray-100
    'amber':         RGBColor(0xF1, 0xC2, 0x1B),  # Yellow-30
    'amber_light':   RGBColor(0xFC, 0xF4, 0xD6),  # Yellow-10 (warning bg / SO WHAT)
    'dark_text':     RGBColor(0x16, 0x16, 0x16),  # Gray-100
    'gray_caption':  RGBColor(0x52, 0x52, 0x52),  # Gray-70
    'light_gray':    RGBColor(0xE0, 0xE0, 0xE0),  # Gray-20 (border-subtle)
    'soft_blue_bg':  RGBColor(0xED, 0xF5, 0xFF),  # Blue-10 (info bg)
    'soft_white_bg': RGBColor(0xF4, 0xF4, 0xF4),  # Gray-10 (layer-01)
    'white':         RGBColor(0xFF, 0xFF, 0xFF),
    'red_alert':     RGBColor(0xDA, 0x1E, 0x28),  # Red-60
    'green_pos':     RGBColor(0x24, 0xA1, 0x48),  # Green-50
}
G.THEME.update(CARBON)
THEME = G.THEME
PURPLE   = RGBColor(0x8A, 0x3F, 0xFC)  # Purple-60
RED_10   = RGBColor(0xFF, 0xF1, 0xF1)
GREEN_10 = RGBColor(0xDE, 0xFB, 0xE6)
GRAY_20  = RGBColor(0xE0, 0xE0, 0xE0)
BLUE_40  = RGBColor(0x78, 0xA9, 0xFF)  # Blue-40 (dark surface 위 링크색)

FONT_KO = 'IBM Plex Sans KR'   # 한글 (미설치 시 시스템 한글 폰트로 대체)
FONT_MO = 'IBM Plex Mono'      # 숫자·코드

# ── 타이포그래피 스케일 (점진적 위계, 2단계 상향) ──────────────
T_TITLE   = 23    # 슬라이드 H1
T_KICKER  = 12    # eyebrow (mono)
T_SECTION = 14    # 섹션 제목
T_CARD    = 13    # 카드/행 강조 제목
T_BODY    = 12    # 본문 primary
T_SUB     = 11    # 본문 secondary
T_CAP     = 10    # 미세 라벨·태그
T_NUM     = 22    # 키넘버 (mono)
T_DATA    = 13.5  # 데이터 값 강조

MARGIN = Inches(0.32)
CW = SLIDE_W - MARGIN * 2


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])


def carbon_header(slide, kicker, title):
    add_rect(slide, MARGIN, Inches(0.26), Inches(0.06), Inches(0.72), fill=THEME['samsung_blue'])
    add_text(slide, MARGIN + Inches(0.2), Inches(0.24), Inches(11), Inches(0.26),
             kicker, font=FONT_MO, size=T_KICKER, bold=True, color=THEME['samsung_blue'])
    add_text(slide, MARGIN + Inches(0.2), Inches(0.47), Inches(12), Inches(0.52),
             title, font=FONT_KO, size=T_TITLE, bold=True, color=THEME['dark_text'])
    add_line(slide, MARGIN, Inches(1.1), SLIDE_W - MARGIN, Inches(1.1),
             color=THEME['light_gray'], width=Emu(9525))


def sect(slide, x, y, w, text):
    add_text(slide, x, y, w, Inches(0.3), text, font=FONT_KO, size=T_SECTION, bold=True, color=THEME['dark_text'])


def sowhat(slide, y, text, label='SO WHAT'):
    add_rect(slide, MARGIN, y, CW, Inches(0.5), fill=THEME['amber_light'])
    add_rect(slide, MARGIN, y, Inches(0.06), Inches(0.5), fill=THEME['amber'])
    add_text(slide, MARGIN + Inches(0.2), y + Inches(0.13), Inches(1.5), Inches(0.26),
             label, font=FONT_MO, size=T_SUB, bold=True, color=THEME['dark_text'])
    add_text(slide, MARGIN + Inches(1.6), y + Inches(0.08), CW - Inches(1.8), Inches(0.36),
             text, font=FONT_KO, size=T_BODY, bold=True, color=THEME['dark_text'], valign='middle')


# ============================================================
# 슬라이드 1: WHY
# ============================================================

def build_slide_1(prs):
    slide = _blank(prs)
    carbon_header(slide, 'WHY NOW', '왜 지금 — LTA가 SCA가 되었다')

    # 북극성 명제 배너
    qb_y = Inches(1.2)
    add_rect(slide, MARGIN, qb_y, CW, Inches(0.78), fill=THEME['soft_blue_bg'])
    add_rect(slide, MARGIN, qb_y, Inches(0.06), Inches(0.78), fill=THEME['samsung_blue'])
    add_text(slide, MARGIN + Inches(0.22), qb_y + Inches(0.1), CW - Inches(0.44), Inches(0.32),
             '“고객의 아키텍처 안으로 들어가 수요를 함께 설계하는 기업이 이긴다.” — 칩을 많이 파는 기업이 아니라.',
             font=FONT_KO, size=T_CARD, bold=True, color=THEME['dark_text'])
    add_text(slide, MARGIN + Inches(0.22), qb_y + Inches(0.47), CW - Inches(0.44), Inches(0.24),
             '신문섭, Bain & Company (2026-06-18 인터뷰)', font=FONT_KO, size=T_SUB, color=THEME['gray_caption'])

    colY = Inches(2.16)
    lw = Inches(6.15)
    lx = MARGIN
    rx = MARGIN + lw + Inches(0.22)
    rw = CW - lw - Inches(0.22)

    # 좌: 계약 3단 진화
    sect(slide, lx, colY, lw, '계약 구조의 3단 진화 — 요구 역량의 이동')
    stages = [
        ('1  SPOT / 분기', '가격이 유일한 변수', '원가·수율·납기', THEME['gray_caption']),
        ('2  LTA + 선급금', '물량·가격 3~5년 락인 · 선급금 10~30% · 2027 비트 고정가', '캐파 계획·공급 신뢰성', THEME['samsung_blue']),
        ('3  SCA 전략적 계약', 'LTA 위에 공동설계·운영통합·자본연계 적층', '워크로드 이해·공동설계·선제 제안', THEME['amber']),
    ]
    sy = colY + Inches(0.36)
    for t, d, cap, c in stages:
        h = Inches(0.94)
        add_rect(slide, lx, sy, lw, h, fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        add_rect(slide, lx, sy, Inches(0.06), h, fill=c)
        add_text(slide, lx + Inches(0.2), sy + Inches(0.09), lw - Inches(0.34), Inches(0.28),
                 t, font=FONT_KO, size=T_CARD, bold=True, color=c)
        add_text(slide, lx + Inches(0.2), sy + Inches(0.4), lw - Inches(0.34), Inches(0.28),
                 d, font=FONT_KO, size=T_BODY, color=THEME['dark_text'], line_spacing=1.0)
        add_text(slide, lx + Inches(0.2), sy + Inches(0.68), lw - Inches(0.34), Inches(0.24),
                 '요구 역량  ' + cap, font=FONT_KO, size=T_SUB, bold=True, color=THEME['gray_caption'])
        sy = sy + h + Inches(0.09)

    # 우: Micron–Anthropic 4요소
    sect(slide, rx, colY, rw, 'Micron ↔ Anthropic SCA (2026-06-22)')
    comps = [
        ('다년 공급', 'HBM·DRAM·SSD 전 포트폴리오', 'LTA에도'),
        ('공동 최적화', 'Claude 워크로드 맞춤 공동설계', 'SCA 신규'),
        ('운영 통합', 'Micron 전사 Claude 배치', 'SCA 신규'),
        ('자본 연계', 'Series H 전략 투자 ($965B)', 'SCA 신규'),
    ]
    cy = colY + Inches(0.36)
    for name, desc, tag in comps:
        h = Inches(0.52)
        isnew = tag == 'SCA 신규'
        add_rect(slide, rx, cy, rw, h, fill=THEME['white'], line=THEME['light_gray'])
        add_text(slide, rx + Inches(0.16), cy, Inches(1.6), h, name, font=FONT_KO, size=T_BODY, bold=True,
                 color=THEME['amber'] if isnew else THEME['gray_caption'], valign='middle')
        add_text(slide, rx + Inches(1.82), cy, rw - Inches(2.9), h, desc, font=FONT_KO, size=T_SUB,
                 color=THEME['dark_text'], valign='middle')
        add_rect(slide, rx + rw - Inches(1.02), cy + Inches(0.11), Inches(0.9), Inches(0.3),
                 fill=THEME['amber_light'] if isnew else THEME['soft_white_bg'])
        add_text(slide, rx + rw - Inches(1.02), cy + Inches(0.15), Inches(0.9), Inches(0.22), tag,
                 font=FONT_KO, size=T_CAP, bold=True, color=THEME['dark_text'] if isnew else THEME['gray_caption'], align='center')
        cy = cy + h + Inches(0.06)

    # 키넘버 (gray-100 다크 스트립, mono)
    nums = [('16건', 'Micron SCA'), ('~$100B', '최소 계약매출'), ('$22B', '예치금·금융약정')]
    nw = (rw - Inches(0.16)) / 3
    ny = cy + Inches(0.06)
    for i, (v, l) in enumerate(nums):
        x = rx + (nw + Inches(0.08)) * i
        add_rect(slide, x, ny, nw, Inches(0.74), fill=THEME['dark_text'])
        add_text(slide, x + Inches(0.14), ny + Inches(0.08), nw - Inches(0.28), Inches(0.4),
                 v, font=FONT_MO, size=T_NUM, bold=True, color=THEME['white'])
        add_text(slide, x + Inches(0.14), ny + Inches(0.5), nw - Inches(0.28), Inches(0.2),
                 l, font=FONT_KO, size=T_CAP, color=GRAY_20)

    # 하단 타임라인 (전폭)
    tl_y = Inches(5.78)
    add_text(slide, MARGIN, tl_y - Inches(0.28), Inches(11), Inches(0.26),
             '체질 전환을 요구하는 사건의 누적 — 일회성이 아니라 구조 전환', font=FONT_KO, size=T_SUB, bold=True, color=THEME['dark_text'])
    events = [
        ('2025-06', 'SK hynix 커스텀 HBM 인증', False),
        ('2025-10', 'OpenAI Stargate LOI', False),
        ('2025~26', 'LTA 선급금 10~30% 체제화', False),
        ('2026-05', 'Anthropic Series H · 3사', False),
        ('2026-06', 'Micron–Anthropic SCA', True),
    ]
    ew = (CW - Inches(0.32)) / 5
    for i, (d, t, hot) in enumerate(events):
        x = MARGIN + (ew + Inches(0.08)) * i
        add_rect(slide, x, tl_y, ew, Inches(0.56), fill=THEME['dark_text'] if hot else THEME['soft_white_bg'],
                 line=None if hot else THEME['light_gray'])
        add_text(slide, x + Inches(0.12), tl_y + Inches(0.07), ew - Inches(0.24), Inches(0.2),
                 d, font=FONT_MO, size=T_SUB, bold=True, color=THEME['amber'] if hot else THEME['samsung_blue'])
        add_text(slide, x + Inches(0.12), tl_y + Inches(0.29), ew - Inches(0.24), Inches(0.24),
                 t, font=FONT_KO, size=T_CAP, color=THEME['white'] if hot else THEME['dark_text'], line_spacing=0.95)

    sowhat(slide, Inches(6.5),
           '고객이 사는 것은 칩이 아니라 “내 워크로드를 이해하고 아키텍처를 함께 최적화하는 파트너”다.')
    add_footer(slide, 1, TOTAL, 'WHY · LTA → SCA')
    return slide


# ============================================================
# 슬라이드 2: WHAT
# ============================================================

def build_slide_2(prs):
    slide = _blank(prs)
    carbon_header(slide, 'WHAT', '무엇을 — 수주 이행자에서 기술 파트너로')

    colY = Inches(1.24)
    lw = Inches(6.7)
    lx = MARGIN
    rx = MARGIN + lw + Inches(0.22)
    rw = CW - lw - Inches(0.22)

    # 좌: As-Is → To-Be 7행
    sect(slide, lx, colY, lw, '개발실 역할의 재정의 — As-Is → To-Be')
    hy = colY + Inches(0.36)
    add_rect(slide, lx, hy, lw, Inches(0.32), fill=THEME['dark_text'])
    add_text(slide, lx + Inches(0.14), hy + Inches(0.05), Inches(1.5), Inches(0.22), '차원', font=FONT_KO, size=T_SUB, bold=True, color=THEME['white'])
    add_text(slide, lx + Inches(1.7), hy + Inches(0.05), Inches(2.4), Inches(0.22), 'As-Is · 수주 이행자', font=FONT_KO, size=T_SUB, bold=True, color=GRAY_20)
    add_text(slide, lx + Inches(4.2), hy + Inches(0.05), Inches(2.4), Inches(0.22), 'To-Be · 기술 파트너', font=FONT_KO, size=T_SUB, bold=True, color=BLUE_40)
    rows = [
        ('요구사항', '확정 스펙 수령', '요구를 공동 정의'),
        ('제안', 'RFQ 응답', '로드맵 기반 선제 제안'),
        ('기술 방향', '표준·고객 추종', '유리한 기술 요소 드라이브'),
        ('고객 접점', '영업 뒤 후방 지원', '엔지니어가 고객 옆 상주'),
        ('정보 흐름', '요구 내려오면 착수', '워크로드 상시 센싱'),
        ('모델링 범위', '메모리 디바이스 단품', '랙·DC 전체 시스템 모델'),
        ('성공 지표', '품질·납기·수율(QCD)', 'QCD + 공동설계·채택률'),
    ]
    ry = hy + Inches(0.32)
    for i, (dim, a, b) in enumerate(rows):
        h = Inches(0.47)
        hi = (dim == '모델링 범위')  # 계단 시각화로 이어지는 행 강조
        add_rect(slide, lx, ry, lw, h, fill=THEME['soft_blue_bg'] if hi else (THEME['soft_white_bg'] if i % 2 == 0 else THEME['white']), line=THEME['light_gray'])
        add_text(slide, lx + Inches(0.14), ry, Inches(1.55), h, dim, font=FONT_KO, size=T_BODY, bold=True, color=THEME['samsung_blue'] if hi else THEME['dark_text'], valign='middle')
        add_text(slide, lx + Inches(1.7), ry, Inches(2.45), h, a, font=FONT_KO, size=T_SUB, color=THEME['gray_caption'], valign='middle')
        add_text(slide, lx + Inches(4.05), ry, Inches(0.28), h, '→', font=FONT_MO, size=T_BODY, bold=True, color=THEME['samsung_blue'], valign='middle')
        add_text(slide, lx + Inches(4.2), ry, Inches(2.4), h, b, font=FONT_KO, size=T_SUB, bold=True, color=THEME['dark_text'], valign='middle')
        ry = ry + h

    # 좌하: 모델링 범위 확대 — 계단형 (디바이스 → 서버 → 랙 → DC)
    scope_title_y = ry + Inches(0.08)
    add_text(slide, lx, scope_title_y, lw, Inches(0.24),
             '↑ 모델링 범위 확대 — 디바이스 ⊂ 서버 ⊂ 랙 ⊂ 데이터센터', font=FONT_KO, size=T_SUB, bold=True, color=THEME['dark_text'])
    base_y = Inches(6.4)          # 막대 하단 기준선
    top_min = scope_title_y + Inches(0.3)  # 가장 높은 막대의 상단
    max_h = base_y - top_min
    scope = [
        ('디바이스', '데이터시트', '현재', 0.42, THEME['gray_caption']),
        ('서버', 'v0.1', 'P1', 0.62, BLUE_40),
        ('랙', 'v1.0', 'P2', 0.8, THEME['samsung_blue']),
        ('데이터센터', 'v2.0', 'P3', 1.0, RGBColor(0x00, 0x43, 0xCE)),
    ]
    sw = (lw - Inches(0.18)) / 4
    for i, (name, ver, ph, lvl, c) in enumerate(scope):
        x = lx + (sw + Inches(0.06)) * i
        bh = Emu(int(max_h * lvl))
        y = base_y - bh
        add_rect(slide, x, y, sw, bh, fill=THEME['white'], line=c, line_width=Emu(12700))
        add_rect(slide, x, y, sw, Inches(0.04), fill=c)  # 상단 액센트
        add_text(slide, x + Inches(0.1), y + Inches(0.07), sw - Inches(0.2), Inches(0.24),
                 name, font=FONT_KO, size=T_SUB, bold=True, color=c)
        add_text(slide, x + Inches(0.1), y + bh - Inches(0.28), sw - Inches(0.2), Inches(0.22),
                 f'{ver} · {ph}', font=FONT_MO, size=T_CAP, bold=True, color=THEME['gray_caption'])

    # 우상: FDE 벤치마크
    sect(slide, rx, colY, rw, '벤치마크 · Palantir FDE')
    fy = colY + Inches(0.36)
    fh = Inches(2.95)
    add_rect(slide, rx, fy, rw, fh, fill=THEME['white'], line=THEME['light_gray'])
    add_rect(slide, rx, fy, Inches(0.06), fh, fill=THEME['samsung_blue'])
    add_text(slide, rx + Inches(0.2), fy + Inches(0.11), rw - Inches(0.38), Inches(0.5),
             '“Delta” — 고객사에 상주하는 엔지니어. Anthropic·OpenAI가 GTM 전략으로 채택(Palantir 640% 동력).',
             font=FONT_KO, size=T_SUB, color=THEME['gray_caption'], line_spacing=1.0)
    maps = [
        ('고객사 상주', 'Co-Design Pod'),
        ('한 고객, 많은 능력', '파일럿 집중→확대'),
        ('말한 요구 vs 실제', '요구 공동 정의'),
        ('gravel → paved', '재사용 설계 플랫폼'),
        ('성과로 평가', '개정 KPI'),
    ]
    my = fy + Inches(0.72)
    for f, d in maps:
        add_text(slide, rx + Inches(0.2), my, Inches(2.7), Inches(0.26), '· ' + f, font=FONT_KO, size=T_BODY, color=THEME['gray_caption'])
        add_text(slide, rx + Inches(2.86), my, Inches(0.26), Inches(0.26), '→', font=FONT_MO, size=T_BODY, bold=True, color=THEME['samsung_blue'])
        add_text(slide, rx + Inches(3.14), my, Inches(2.1), Inches(0.26), d, font=FONT_KO, size=T_BODY, bold=True, color=THEME['dark_text'])
        my = my + Inches(0.32)
    add_text(slide, rx + Inches(0.2), my + Inches(0.02), rw - Inches(0.38), Inches(0.44),
             '메모리 변형: FDE(상주) + 시스템 아키텍트·모델링(성능·파워 정량화) — 모델을 무기로 들고 들어간다.',
             font=FONT_KO, size=T_SUB, italic=True, color=THEME['samsung_blue'], line_spacing=1.0)

    # 우하: 리스크 ↔ 이점
    ry3 = fy + fh + Inches(0.12)
    half = (rw - Inches(0.12)) / 2
    bh = Inches(1.78)
    add_rect(slide, rx, ry3, half, bh, fill=RED_10, line=THEME['light_gray'])
    add_rect(slide, rx, ry3, half, Inches(0.3), fill=THEME['red_alert'])
    add_text(slide, rx + Inches(0.13), ry3 + Inches(0.05), half - Inches(0.2), Inches(0.22), '안 하면 · 리스크', font=FONT_KO, size=T_SUB, bold=True, color=THEME['white'])
    for i, t in enumerate(['SCA 수주 배제', '커스텀 전환기 고착', '2nd source·가격력 상실', '미래 기술 선점 실패']):
        add_text(slide, rx + Inches(0.13), ry3 + Inches(0.38) + Inches(0.34) * i, half - Inches(0.24), Inches(0.3),
                 '· ' + t, font=FONT_KO, size=T_SUB, color=THEME['dark_text'])
    bx = rx + half + Inches(0.12)
    add_rect(slide, bx, ry3, half, bh, fill=GREEN_10, line=THEME['light_gray'])
    add_rect(slide, bx, ry3, half, Inches(0.3), fill=THEME['green_pos'])
    add_text(slide, bx + Inches(0.13), ry3 + Inches(0.05), half - Inches(0.2), Inches(0.22), '하면 · 이점', font=FONT_KO, size=T_SUB, bold=True, color=THEME['white'])
    for i, t in enumerate(['지속 매출(락인+선급금)', '수익률 프리미엄', '미래 기술 선점', '자본 연계 옵션']):
        add_text(slide, bx + Inches(0.13), ry3 + Inches(0.38) + Inches(0.34) * i, half - Inches(0.24), Inches(0.3),
                 '· ' + t, font=FONT_KO, size=T_SUB, color=THEME['dark_text'])

    sowhat(slide, Inches(6.5),
           '리스크는 비가역·이점은 복리 → 조기 전환의 기대값이 압도적. “아키텍처 안으로”는 검증된 조직 형태(FDE)다.')
    add_footer(slide, 2, TOTAL, 'WHAT · 역할 전환 + FDE')
    return slide


# ============================================================
# 슬라이드 3: HOW
# ============================================================

def build_slide_3(prs):
    slide = _blank(prs)
    carbon_header(slide, 'HOW', '어떻게 — 4대 축과 3-Phase 실행')

    # 4대 축
    ly = Inches(1.24)
    sect(slide, MARGIN, ly, Inches(12), '전환 전략 — 4대 축')
    axes = [
        ('축 1 · 기술', THEME['samsung_blue'], ['워크로드 랩 신설', '시스템 성능·파워 모델', '커스텀 플랫폼·임베디드 SW']),
        ('축 2 · 문화', THEME['amber'], ['"정답 구현"→"가설 제안"', '실패 허용 예산', '트레이드오프 토론 표준화']),
        ('축 3 · 조직', THEME['green_pos'], ['Co-Design Pod (=FDE)', '시스템 아키텍트·모델링 조직', '워크로드 인텔리전스']),
        ('축 4 · 일하는 방식', PURPLE, ['로드맵 교차 리뷰 (분기)', '선행 시제품(PoA) 사이클', 'AI 도구 내재화 (RS-7)']),
    ]
    cw = (CW - Inches(0.24)) / 4
    cy = ly + Inches(0.38)
    for i, (name, c, items) in enumerate(axes):
        x = MARGIN + (cw + Inches(0.08)) * i
        h = Inches(1.64)
        add_rect(slide, x, cy, cw, h, fill=THEME['white'], line=THEME['light_gray'])
        add_rect(slide, x, cy, cw, Inches(0.36), fill=c)
        add_text(slide, x + Inches(0.13), cy + Inches(0.06), cw - Inches(0.26), Inches(0.26),
                 name, font=FONT_KO, size=T_CARD, bold=True, color=THEME['white'])
        iy = cy + Inches(0.48)
        for it in items:
            add_text(slide, x + Inches(0.13), iy, cw - Inches(0.24), Inches(0.38),
                     '· ' + it, font=FONT_KO, size=T_BODY, color=THEME['dark_text'], line_spacing=0.98)
            iy = iy + Inches(0.39)

    # 3-Phase 로드맵
    py = Inches(3.42)
    sect(slide, MARGIN, py, Inches(8), '3-Phase 로드맵')
    phases = [
        ('Phase 1 · 90일 — 증명', THEME['samsung_blue'], ['Co-Design Pod 1호 발족 (FDE 벤치마크)', '워크로드 랩 + 시스템 모델 v0.1', '선제 제안 1호 · KPI 개정안']),
        ('Phase 2 · 1년 — 제도화', THEME['amber'], ['Pod 3~5개 · 아키텍트·모델링 조직', '모델 v1.0(랙) · PoA 연 4회', '★ SCA형 계약 1건 이상 수주']),
        ('Phase 3 · 3년 — 표준화', THEME['green_pos'], ['커스텀 플랫폼 (리드타임 −50%)', '모델 v2.0(DC·고객 공용)', '"전략적 인프라 파트너" IR 공시']),
    ]
    pw = (CW - Inches(0.16)) / 3
    pcy = py + Inches(0.38)
    for i, (t, c, items) in enumerate(phases):
        x = MARGIN + (pw + Inches(0.08)) * i
        h = Inches(1.42)
        hot = i == 1
        add_rect(slide, x, pcy, pw, h, fill=THEME['amber_light'] if hot else THEME['soft_white_bg'], line=THEME['light_gray'])
        add_rect(slide, x, pcy, Inches(0.06), h, fill=c)
        add_text(slide, x + Inches(0.18), pcy + Inches(0.09), pw - Inches(0.32), Inches(0.26),
                 t, font=FONT_KO, size=T_CARD, bold=True, color=c)
        iy = pcy + Inches(0.44)
        for it in items:
            add_text(slide, x + Inches(0.18), iy, pw - Inches(0.34), Inches(0.32),
                     '· ' + it, font=FONT_KO, size=T_SUB, color=THEME['dark_text'], line_spacing=0.98)
            iy = iy + Inches(0.32)

    # KPI 스트립
    ky = Inches(5.2)
    sect(slide, MARGIN, ky, Inches(8), '성공 지표 (현재 → 1년 → 3년)')
    kpis = [
        ('선제 제안', '~0→12→40건'),
        ('로드맵 채택률', '—→20→35%'),
        ('고객 교류 시간', '낮음→10→25%'),
        ('커스텀 매출', '낮음→추적→30%+'),
        ('시스템 모델', '디바이스→랙→DC'),
        ('아키텍트 인력', '없음→조직→Pod당1+'),
    ]
    kw = (CW - Inches(0.4)) / 6
    kcy = ky + Inches(0.38)
    for i, (l, v) in enumerate(kpis):
        x = MARGIN + (kw + Inches(0.08)) * i
        add_rect(slide, x, kcy, kw, Inches(0.72), fill=THEME['soft_white_bg'], line=THEME['light_gray'])
        add_text(slide, x + Inches(0.11), kcy + Inches(0.1), kw - Inches(0.22), Inches(0.22),
                 l, font=FONT_KO, size=T_CAP, color=THEME['gray_caption'])
        add_text(slide, x + Inches(0.11), kcy + Inches(0.36), kw - Inches(0.22), Inches(0.28),
                 v, font=FONT_KO, size=T_BODY, bold=True, color=THEME['samsung_blue'])

    sowhat(slide, Inches(6.5),
           '90일 안에 파일럿 Pod·선제 제안으로 증명 → 1년 안에 SCA형 계약 1건으로 제도화의 근거를 만든다.',
           label='NEXT STEP')
    add_footer(slide, 3, TOTAL, 'HOW · 4축 + 로드맵')
    return slide


# ============================================================
# Main
# ============================================================

def main():
    print(f'Loading template (theme only): {TEMPLATE}')
    prs = Presentation(TEMPLATE)
    remove_existing_slides(prs)
    print('Palette: IBM Carbon · Type scale 23/14/13/12/11/10 + num 22 · 3 dense slides')

    for i, builder in enumerate([build_slide_1, build_slide_2, build_slide_3], 1):
        print(f'  Building slide {i}: {builder.__name__}...')
        builder(prs)

    print(f'Saving to: {OUTPUT}')
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f'Done. Output size: {size_kb:.1f} KB · {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
