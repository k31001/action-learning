"""
팀 얼라인먼트용 3매 발표자료 — "왜 시나리오 플래닝으로 풀어야 하는가"

디자인 컨셉: Black & White 베이스 + Blue / Orange 포인트.
- 텍스트 박스의 풀배경 색 채우기 자제 (좌측 컬러 바, chip 같은 작은 영역에서만)
- 카드 = 흰 배경 + 얇은 회색 테두리 (#E5E7EB)
- 강조 = 좌측 3px 컬러 바, 작은 chip, bold 텍스트로
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# 재사용 — drawing primitives 만 (THEME 은 새로 정의)
from generate_pptx import (
    add_text, add_rect, add_line,
    FONT_KO, FONT_EN, SLIDE_W, SLIDE_H,
)

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
TEMPLATE = os.path.join(ROOT, 'presentation', 'template.pptx')
OUTPUT = os.path.join(ROOT, 'presentation', 'team-alignment-scenario-planning.pptx')

TOTAL = 3
SECTION = '시나리오 플래닝 도입 제안'


# ============================================================
# Tone-down 컬러 시스템 (B&W + Blue·Orange 포인트)
# ============================================================
T = {
    # 베이스
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
    'ink':         RGBColor(0x11, 0x11, 0x11),  # 거의 검정, 제목·강조
    'text':        RGBColor(0x2D, 0x2D, 0x2D),  # 본문
    'caption':     RGBColor(0x6B, 0x72, 0x80),  # 회색 캡션
    'divider':     RGBColor(0xE5, 0xE7, 0xEB),  # 라인·테두리
    'soft_bg':     RGBColor(0xFA, 0xFA, 0xFA),  # 살짝의 음영 (드물게)

    # 포인트
    'blue':        RGBColor(0x25, 0x63, 0xEB),  # 메인 블루
    'blue_deep':   RGBColor(0x1E, 0x40, 0xAF),  # 더 짙은 블루 (제목)
    'orange':      RGBColor(0xEA, 0x58, 0x0C),  # 메인 오렌지

    # 상태 (드물게)
    'red':         RGBColor(0xB9, 0x1C, 0x1C),
}

BORDER_W = Emu(6350)         # 1px @ 96dpi 느낌
BORDER_W_THICK = Emu(12700)  # 2px


# ============================================================
# 공통 컴포넌트 — 톤다운 카드
# ============================================================

def card(slide, x, y, w, h, *, border=None, fill=None, accent_left=None, accent_top=None):
    """카드 = 흰 배경 + 얇은 회색 테두리. fill 옵션은 거의 흰 (soft_bg) 만 권장."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else T['white']
    shape.line.color.rgb = border if border is not None else T['divider']
    shape.line.width = BORDER_W if border is None else BORDER_W_THICK
    shape.shadow.inherit = False
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
    # 좌측 액센트 바
    if accent_left is not None:
        bar_w = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_left
        bar.line.fill.background()
        bar.shadow.inherit = False
    # 상단 액센트 바
    if accent_top is not None:
        bar_h = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_top
        bar.line.fill.background()
        bar.shadow.inherit = False
    return shape


def chip(slide, x, y, w, h, text, *, fg, bg=None, bordered=True):
    """작은 라벨 chip. bg 없으면 흰 + 컬러 테두리."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if bg is None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = T['white']
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    if bordered:
        shape.line.color.rgb = fg
        shape.line.width = BORDER_W
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_EN
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = fg
    return shape


def page_title(slide, kicker, title, subtitle=None):
    """페이지 헤더 — 작은 kicker(블루) + 큰 타이틀(잉크) + 캡션 부제."""
    # 좌측 풀높이 가는 액센트 (블루 1px)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(0.04), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = T['blue']
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide, Inches(0.5), Inches(0.42), Inches(12.83), Inches(0.28),
             kicker, font=FONT_EN, size=10, bold=True, color=T['blue'])
    add_text(slide, Inches(0.5), Inches(0.68), Inches(12.83), Inches(0.55),
             title, font=FONT_KO, size=23, bold=True, color=T['ink'])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.22), Inches(12.83), Inches(0.32),
                 subtitle, font=FONT_KO, size=10.5, color=T['caption'])
    add_line(slide, Inches(0.5), Inches(1.62), Inches(12.83), Inches(1.62),
             color=T['divider'], width=BORDER_W)


def footer_line(slide, page_no):
    """페이지 하단 푸터 — 분리선 + 좌측 라벨 + 페이지 번호."""
    add_line(slide, Inches(0.5), Inches(7.18), Inches(12.83), Inches(7.18),
             color=T['divider'], width=BORDER_W)
    add_text(slide, Inches(0.5), Inches(7.25), Inches(8), Inches(0.22),
             f'삼성전자  |  {SECTION}',
             font=FONT_KO, size=8.5, color=T['caption'])
    add_text(slide, Inches(11.8), Inches(7.25), Inches(1.0), Inches(0.22),
             f'{page_no} / {TOTAL}',
             font=FONT_EN, size=8.5, color=T['caption'], align='right')


def bullet_block(slide, x, y, w, h, items, *, size=10, color=None, gap=0.04):
    """• 글머리 다중 라인."""
    color = color or T['text']
    text = '\n'.join(f'•  {it}' for it in items)
    add_text(slide, x, y, w, h, text, size=size, color=color, line_spacing=1.32)


def banner(slide, y, label, body, *, accent=None):
    """하단 띠 — 흰 배경 + 좌측 컬러 바 + 라벨 + 본문 (풀배경 채움 금지)."""
    accent = accent or T['orange']
    H = Inches(0.36)
    # 카드(흰배경 + 얇은 회색 테두리)
    card(slide, Inches(0.5), y, Inches(12.83), H, accent_left=accent)
    add_text(slide, Inches(0.72), y, Inches(1.6), H,
             label, font=FONT_EN, size=9, bold=True, color=accent,
             valign='middle')
    add_text(slide, Inches(2.25), y, Inches(10.4), H,
             body, font=FONT_KO, size=10, bold=True, color=T['ink'],
             valign='middle')


# ============================================================
# SLIDE 1 — WHY
# ============================================================

def build_slide_1_why(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))

    page_title(
        slide,
        kicker='PART 1  ·  WHY',
        title='왜 시나리오 플래닝인가 — 단순 예측으로는 풀 수 없는 문제',
        subtitle='호황 정점·점유율 추락·다운턴 사이클이 동시에 진행 중인 지금, 우리는 '
                 '"하나의 미래" 가 아니라 "여러 가능한 미래" 를 동시에 준비해야 한다',
    )

    # =========================================
    # 좌측 — 포컬 이슈 + KPI + 본질 불확실성
    # =========================================
    LX = Inches(0.5)
    LW = Inches(6.0)

    # Focal Issue 카드 (흰 배경 + 두꺼운 테두리 + 좌측 오렌지 바)
    fi_y = Inches(1.85)
    fi_h = Inches(2.0)
    card(slide, LX, fi_y, LW, fi_h, border=T['ink'], accent_left=T['orange'])
    add_text(slide, LX + Inches(0.3), fi_y + Inches(0.22), LW - Inches(0.55), Inches(0.25),
             'FOCAL ISSUE  ·  우리가 풀어야 할 단 하나의 질문',
             font=FONT_EN, size=10, bold=True, color=T['orange'])
    add_text(slide, LX + Inches(0.3), fi_y + Inches(0.6), LW - Inches(0.55), Inches(1.0),
             'AI 메모리 시대에 삼성전자 메모리사업부가\n'
             '2030~2035년에도 글로벌 리더십을 유지하기 위해\n'
             '어떤 전략적 결정을',
             font=FONT_KO, size=14, color=T['text'], line_spacing=1.30)
    add_text(slide, LX + Inches(0.3), fi_y + Inches(1.55), LW - Inches(0.55), Inches(0.4),
             '지금 내려야 하는가?',
             font=FONT_KO, size=18, bold=True, color=T['ink'])

    # KPI 3개 (흰 배경 + 상단 컬러 바)
    KY = Inches(3.97)
    KH = Inches(0.85)
    kw = (LW - Inches(0.2)) / 3
    kpis = [
        ('호황 정점', '2026 Q1 매출\n사상 최대 $50.4B', T['orange']),
        ('회복 궤도', 'HBM 점유율\n17% → 35% 반등', T['blue']),
        ('미상 리스크', '다음 다운턴\n시점·규모 미지', T['ink']),
    ]
    for i, (label, body, accent) in enumerate(kpis):
        kx = LX + kw * i + Inches(0.1) * i
        card(slide, kx, KY, kw, KH, accent_top=accent)
        add_text(slide, kx + Inches(0.14), KY + Inches(0.16),
                 kw - Inches(0.28), Inches(0.22),
                 label, font=FONT_KO, size=9, bold=True, color=accent)
        add_text(slide, kx + Inches(0.14), KY + Inches(0.38),
                 kw - Inches(0.28), Inches(0.48),
                 body, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 line_spacing=1.20)

    # 3대 본질 불확실성
    add_text(slide, LX, Inches(5.0), LW, Inches(0.3),
             '그리고 우리를 둘러싼 3대 본질적 불확실성',
             font=FONT_KO, size=11, bold=True, color=T['ink'])

    UY = Inches(5.34)
    UH = Inches(1.62)
    uw = (LW - Inches(0.16)) / 3
    uncertainties = [
        ('DF1', 'AI 수요', '슈퍼사이클\nvs\n거품 붕괴',
         '빅테크 CapEx $725B 지속 vs ROI 미실현 95%', T['blue']),
        ('DF2', '미중 지정학', '디커플링\nvs\n관리된 공존',
         'MATCH 법안 vs H20 재허용 정책 진동', T['blue']),
        ('DF3', '메모리 아키텍처', 'HBM 지속\nvs\n패러다임 전환',
         '3D DRAM·PIM·CXL 부상 시점', T['blue']),
    ]
    for i, (df, name, poles, evi, accent) in enumerate(uncertainties):
        ux = LX + uw * i + Inches(0.08) * i
        card(slide, ux, UY, uw, UH, accent_top=accent)
        chip(slide, ux + Inches(0.12), UY + Inches(0.16),
             Inches(0.42), Inches(0.22), df, fg=accent)
        add_text(slide, ux + Inches(0.6), UY + Inches(0.16), uw - Inches(0.7), Inches(0.22),
                 name, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 valign='middle')
        add_text(slide, ux + Inches(0.1), UY + Inches(0.5), uw - Inches(0.2), Inches(0.66),
                 poles, font=FONT_KO, size=10, bold=True, color=T['text'],
                 align='center', line_spacing=1.10)
        add_text(slide, ux + Inches(0.1), UY + Inches(1.2), uw - Inches(0.2), Inches(0.4),
                 evi, font=FONT_KO, size=8, color=T['caption'],
                 align='center', line_spacing=1.20)

    # =========================================
    # 우측 — 단선 예측 vs 시나리오 플래닝 비교 + 검증 사례
    # =========================================
    RX = Inches(6.83)
    RW = Inches(6.0)

    add_text(slide, RX, Inches(1.85), RW, Inches(0.3),
             '단선 예측  vs  시나리오 플래닝',
             font=FONT_KO, size=13, bold=True, color=T['ink'])
    add_text(slide, RX, Inches(2.13), RW, Inches(0.25),
             '같은 정보로도 결과는 근본적으로 다르다',
             font=FONT_KO, size=10, color=T['caption'])

    CY = Inches(2.5)
    CH = Inches(2.5)
    cw = (RW - Inches(0.16)) / 2

    # 단선 예측 카드 (좌측 회색 바)
    card(slide, RX, CY, cw, CH, accent_left=T['caption'])
    add_text(slide, RX + Inches(0.22), CY + Inches(0.14),
             cw - Inches(0.34), Inches(0.28),
             '✕  단선 예측 (Single Forecast)',
             font=FONT_KO, size=11, bold=True, color=T['ink'])
    add_line(slide, RX + Inches(0.22), CY + Inches(0.46),
             RX + cw - Inches(0.22), CY + Inches(0.46),
             color=T['divider'], width=BORDER_W)
    bullet_block(
        slide, RX + Inches(0.22), CY + Inches(0.54),
        cw - Inches(0.44), Inches(1.95),
        [
            '"가장 그럴 듯한 미래" 하나만 가정',
            '예측이 빗나가면 대응 시점 늦음',
            '호황 가정 → 다운턴 시 적자 직격',
            '협상 가정 → 디커플링 시 시안 팹 손실',
            '도전 가설(HBM 패러다임 전환) 누락',
            '결정 = "맞히기" 게임  →  높은 실패 비용',
        ],
        size=10, color=T['text'],
    )

    # 시나리오 플래닝 카드 (좌측 블루 바, 테두리 강조)
    rx2 = RX + cw + Inches(0.16)
    card(slide, rx2, CY, cw, CH, accent_left=T['blue'], border=T['blue'])
    add_text(slide, rx2 + Inches(0.22), CY + Inches(0.14),
             cw - Inches(0.34), Inches(0.28),
             '✓  시나리오 플래닝 (Shell Method)',
             font=FONT_KO, size=11, bold=True, color=T['blue_deep'])
    add_line(slide, rx2 + Inches(0.22), CY + Inches(0.46),
             rx2 + cw - Inches(0.22), CY + Inches(0.46),
             color=T['divider'], width=BORDER_W)
    bullet_block(
        slide, rx2 + Inches(0.22), CY + Inches(0.54),
        cw - Inches(0.44), Inches(1.95),
        [
            '4~5개 "가능한 미래" 를 동시에 가정',
            '시나리오마다 함의·기회·위협 사전 도출',
            '어떤 미래가 와도 작동하는 Robust 전략',
            'Main Bet + Side Bet 보험 구조',
            'EWI 로 분기점 자동 감지 → 30일 내 실행',
            '결정 = "준비" 게임  →  높은 적응 속도',
        ],
        size=10, color=T['text'],
    )

    # 검증 사례 — 시나리오 플래닝 실제 활용 기업 4개 (삼성 제거 + Disney 교체)
    EY = Inches(5.2)
    EH = Inches(1.78)
    card(slide, RX, EY, RW, EH, accent_top=T['orange'])
    add_text(slide, RX + Inches(0.22), EY + Inches(0.18), RW - Inches(0.44), Inches(0.25),
             'PROOF  ·  시나리오 플래닝을 실제 도입·활용한 기업',
             font=FONT_EN, size=10, bold=True, color=T['orange'])
    add_text(slide, RX + Inches(0.22), EY + Inches(0.43), RW - Inches(0.44), Inches(0.22),
             '50년간 글로벌 대기업이 표준 의사결정 도구로 채택해 온 방법론',
             font=FONT_KO, size=9, color=T['caption'])

    cases = [
        ('Shell',         '1970s~', '오일쇼크 사전 시나리오 → 단독 흑자 유지\n시나리오 플래닝의 원조 도입 사례'),
        ('ExxonMobil',    '1990s~', '연 단위 시나리오 갱신 → 장기 CapEx 의사결정\n10년 단위 에너지 전망 표준화'),
        ('Daimler-Benz',  '1980s~', '글로벌 자동차 시장 시나리오 4종 운영\n모빌리티 R&D 포트폴리오 결정 도구'),
        ('UPS',           '1990s~', '20년 물류 시나리오 4종 정기 갱신\n전 세계 거점 투자 의사결정에 활용'),
    ]
    case_w = (RW - Inches(0.4)) / 4
    for i, (org, year, lesson) in enumerate(cases):
        ex_x = RX + Inches(0.22) + case_w * i
        add_text(slide, ex_x, EY + Inches(0.78),
                 case_w - Inches(0.15), Inches(0.28),
                 org, font=FONT_EN, size=12, bold=True, color=T['blue_deep'])
        add_text(slide, ex_x, EY + Inches(1.05),
                 case_w - Inches(0.15), Inches(0.2),
                 year, font=FONT_EN, size=9, bold=True, color=T['orange'])
        add_text(slide, ex_x, EY + Inches(1.27),
                 case_w - Inches(0.15), Inches(0.48),
                 lesson, font=FONT_KO, size=8.5, color=T['text'],
                 line_spacing=1.25)

    footer_line(slide, 1)


# ============================================================
# SLIDE 2 — HOW
# ============================================================

def build_slide_2_how(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))

    page_title(
        slide,
        kicker='PART 2  ·  HOW',
        title='8단계 워크플로우 — 우리 메모리 사업에 이렇게 적용한다',
        subtitle='Shell 시나리오 플래닝의 표준 8단계 프로세스를 메모리사업부 포컬 이슈에 매핑하면 '
                 '각 단계마다 명확한 산출물이 나온다 (위키에서 이미 1차 작업 완료)',
    )

    # 상단 8단계 흐름도 — 흰 배경 + 상단 액센트 색만 (블루/오렌지) + 단계 번호 chip
    FY = Inches(1.85)
    FH = Inches(1.1)
    FX = Inches(0.5)
    FW = Inches(12.33)

    steps = [
        ('1', 'Focal\nIssue', '단 하나의 질문', T['blue']),
        ('2', 'STEEP\n요인', 'S/T/E/E/P 50개', T['blue']),
        ('3', 'Impact ×\nUncertainty', '상위 20개', T['blue']),
        ('4', 'Driving\nForces', '2개 + 보조 1개', T['blue']),
        ('5', '2×2\n매트릭스', '4 + 와일드카드', T['orange']),
        ('6', '시나리오\n내러티브', '각 함의·분기점', T['orange']),
        ('7', 'Main / Side\nBet', '보험 구조', T['orange']),
        ('8', 'Robust +\nEWI', '실행 + 모니터', T['orange']),
    ]
    box_w = (FW - Inches(0.16) * 7) / 8
    for i, (num, name, desc, accent) in enumerate(steps):
        bx = FX + (box_w + Inches(0.16)) * i
        card(slide, bx, FY, box_w, FH, accent_top=accent)
        # STEP N chip (얇은 컬러 텍스트만)
        add_text(slide, bx + Inches(0.08), FY + Inches(0.14),
                 box_w - Inches(0.16), Inches(0.2),
                 f'STEP {num}', font=FONT_EN, size=9, bold=True,
                 color=accent, align='center')
        # 이름
        add_text(slide, bx + Inches(0.05), FY + Inches(0.4),
                 box_w - Inches(0.1), Inches(0.46),
                 name, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 align='center', line_spacing=1.10)
        # 설명
        add_text(slide, bx + Inches(0.05), FY + Inches(0.86),
                 box_w - Inches(0.1), Inches(0.22),
                 desc, font=FONT_KO, size=8, color=T['caption'],
                 align='center')
        # 화살표 (얇은 회색)
        if i < len(steps) - 1:
            ax1 = bx + box_w + Inches(0.02)
            ax2 = bx + box_w + Inches(0.14)
            ay = FY + Inches(0.55)
            add_line(slide, ax1, ay, ax2, ay,
                     color=T['divider'], width=BORDER_W)

    # 중단 — 좌측 매핑 카드 / 우측 2x2 매트릭스
    MY = Inches(3.15)
    MH = Inches(3.45)

    # 매핑 카드 (흰 배경)
    MX = Inches(0.5)
    MW = Inches(7.55)
    card(slide, MX, MY, MW, MH)
    add_text(slide, MX + Inches(0.25), MY + Inches(0.18), MW - Inches(0.5), Inches(0.3),
             '각 단계  →  우리 메모리 사업의 산출물 (위키에서 1차 완료)',
             font=FONT_KO, size=11, bold=True, color=T['ink'])
    add_line(slide, MX + Inches(0.25), MY + Inches(0.55),
             MX + MW - Inches(0.25), MY + Inches(0.55),
             color=T['divider'], width=BORDER_W)

    mappings = [
        ('Step 1', 'Focal Issue',       '2030~2035 글로벌 리더십 유지 = 지금 무엇을 결정?'),
        ('Step 2', 'STEEP 50개',         'Social·Tech·Environment·Economy·Political 환경 요인'),
        ('Step 3', 'I×U 매트릭스',       '50 → 20개 핵심 불확실성 선별 (5점 척도)'),
        ('Step 4', 'Driving Forces',    'DF1 (AI 수요) × DF2 (미중 지정학)  +  DF3 (메모리 아키텍처)'),
        ('Step 5', '2×2 매트릭스',       '시나리오 A·B·C·D  +  와일드카드 E'),
        ('Step 6', '시나리오 내러티브',   '각 시나리오별 5년 후 풍경 + 분기점 + 모니터링 지표'),
        ('Step 7', 'Main / Side Bet',   'Main = B 르네상스 (확률 35%) / Side = A·C·D·E 보험'),
        ('Step 8', 'Robust + EWI',      'RS-1~RS-9  +  21개 자동 트리거  +  D1~D17 결정 묶음'),
    ]
    ROW_Y0 = MY + Inches(0.68)
    ROW_H = Inches(0.34)
    for i, (step, label, outcome) in enumerate(mappings):
        ry = ROW_Y0 + ROW_H * i
        add_text(slide, MX + Inches(0.25), ry, Inches(0.7), ROW_H,
                 step, font=FONT_EN, size=9, bold=True, color=T['orange'],
                 valign='middle')
        add_text(slide, MX + Inches(0.95), ry, Inches(1.65), ROW_H,
                 label, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 valign='middle')
        add_text(slide, MX + Inches(2.55), ry, Inches(0.3), ROW_H,
                 '→', font=FONT_EN, size=11, bold=True, color=T['blue'],
                 valign='middle')
        add_text(slide, MX + Inches(2.85), ry, MW - Inches(3.1), ROW_H,
                 outcome, font=FONT_KO, size=9.5, color=T['text'],
                 valign='middle')

    # 우측 — 2x2 매트릭스 미니 (흰 배경 + 컬러 코드만)
    QX = Inches(8.3)
    QW = Inches(4.53)
    card(slide, QX, MY, QW, MH)
    add_text(slide, QX + Inches(0.2), MY + Inches(0.18), QW - Inches(0.4), Inches(0.3),
             '핵심 산출물 미리보기  ·  2×2 시나리오 매트릭스',
             font=FONT_KO, size=11, bold=True, color=T['ink'])
    add_text(slide, QX + Inches(0.2), MY + Inches(0.45), QW - Inches(0.4), Inches(0.25),
             'DF1 (AI 수요) × DF2 (미중 지정학)  =  4 + 1개 시나리오',
             font=FONT_KO, size=9, color=T['caption'])

    # 매트릭스 그리기
    GX = QX + Inches(0.4)
    GY = MY + Inches(0.95)
    GW = QW - Inches(0.8)
    GH = Inches(2.15)
    half_w = Emu(int(GW / 2))
    half_h = Emu(int(GH / 2))
    # 4분면 — 흰 배경 + 컬러 보더로 시나리오 구분
    quadrants = [
        # col, row, name, prob, border
        (1, 0, 'A · 황금 요새',    '26%', T['blue']),
        (1, 1, 'B · AI 르네상스 ⭐', '35%', T['orange']),
        (0, 0, 'C · 기술 냉전',     '10%', T['caption']),
        (0, 1, 'D · 조용한 재편',   '23%', T['caption']),
    ]
    for col, row, name, prob, accent in quadrants:
        qx = GX + half_w * col
        qy = GY + half_h * row
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, half_w, half_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = T['white']
        rect.line.color.rgb = accent
        rect.line.width = BORDER_W
        rect.shadow.inherit = False
        add_text(slide, qx + Inches(0.06), qy + Inches(0.16),
                 half_w - Inches(0.12), Inches(0.28),
                 name, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 align='center')
        add_text(slide, qx + Inches(0.06), qy + Inches(0.48),
                 half_w - Inches(0.12), Inches(0.3),
                 prob, font=FONT_EN, size=12, bold=True, color=accent,
                 align='center')

    # 축 라벨
    add_text(slide, GX, GY - Inches(0.3), GW, Inches(0.25),
             '↑  미중 디커플링 심화', font=FONT_KO, size=8, bold=True,
             color=T['caption'], align='center')
    add_text(slide, GX, GY + GH + Inches(0.04), GW, Inches(0.25),
             '↓  관리된 공존', font=FONT_KO, size=8, bold=True,
             color=T['caption'], align='center')
    add_text(slide, GX - Inches(0.32), GY, Inches(0.3), GH,
             '←\nA I\n버\n블', font=FONT_KO, size=8, bold=True,
             color=T['caption'], align='center', valign='middle',
             line_spacing=1.0)
    add_text(slide, GX + GW + Inches(0.02), GY, Inches(0.3), GH,
             '슈\n퍼\n사\n이\n클\n→', font=FONT_KO, size=8, bold=True,
             color=T['caption'], align='center', valign='middle',
             line_spacing=1.0)

    # 와일드카드 주석
    add_text(slide, QX + Inches(0.2), MY + MH - Inches(0.35),
             QW - Inches(0.4), Inches(0.25),
             '※ 와일드카드 E "패러다임 전환" (6%) — DF3 Pole B 실현',
             font=FONT_KO, size=8.5, italic=True, color=T['caption'])

    # 하단 SO WHAT 띠 (좌측 오렌지 바, 흰 배경)
    banner(slide, Inches(6.7), 'SO WHAT',
           '8단계 모두 우리 케이스에 1차 매핑 완료 — '
           '팀 합의만 떨어지면 본격 산출물 정제 + 결정 묶음 처리로 진입 가능',
           accent=T['orange'])

    footer_line(slide, 2)


# ============================================================
# SLIDE 3 — WHAT
# ============================================================

def build_slide_3_what(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))

    page_title(
        slide,
        kicker='PART 3  ·  WHAT',
        title='무엇이 도출되는가 — 5개 시나리오 + 두 가지 실행 레이어',
        subtitle='시나리오 플래닝의 최종 산출물은 "Main Bet + Side Bet" 베팅 구조와, '
                 '이를 실행으로 전환하는 "Robust 전략" 과 "EWI" 두 레이어다',
    )

    # ============================================
    # 상단 — 5개 시나리오 가로 1줄 (그대로 유지)
    # ============================================
    SY = Inches(1.85)
    SH = Inches(1.55)
    SX = Inches(0.5)
    SW = Inches(12.33)
    sw = (SW - Inches(0.16) * 4) / 5

    scenarios = [
        ('B', 'AI 르네상스',   '35%', 'MAIN BET',
         '동서 양쪽 시장 + HBM 리더십 회복 → 최대 수혜', T['orange']),
        ('A', '황금 요새',     '26%', 'SIDE BET',
         '디커플링 + AI 지속 → 서방 HBM 독점 공급자', T['blue']),
        ('D', '조용한 재편',   '23%', 'SIDE BET',
         'AI 거품 + 공존 → 비용·차별화 경쟁 (2022~23형)', T['caption']),
        ('C', '기술 냉전',     '10%', 'SIDE BET',
         '이중 충격 → 시안 팹 + 수요 동시 직면', T['caption']),
        ('E', '패러다임 전환', '6%',  'WILDCARD',
         '3D DRAM·PIM·CXL 부상 → R&D 포트폴리오 재검토', T['ink']),
    ]
    for i, (code, name, prob, tag, desc, accent) in enumerate(scenarios):
        sx = SX + (sw + Inches(0.16)) * i
        card(slide, sx, SY, sw, SH, accent_top=accent)
        add_text(slide, sx + Inches(0.16), SY + Inches(0.16),
                 Inches(0.7), Inches(0.6),
                 code, font=FONT_EN, size=28, bold=True, color=accent)
        chip(slide, sx + sw - Inches(1.0), SY + Inches(0.16),
             Inches(0.86), Inches(0.22), tag, fg=accent)
        add_text(slide, sx + Inches(0.16), SY + Inches(0.72),
                 sw - Inches(0.32), Inches(0.28),
                 name, font=FONT_KO, size=11, bold=True, color=T['ink'])
        add_text(slide, sx + Inches(0.16), SY + Inches(0.98),
                 sw - Inches(0.32), Inches(0.22),
                 prob, font=FONT_EN, size=11, bold=True, color=accent)
        add_text(slide, sx + Inches(0.16), SY + Inches(1.21),
                 sw - Inches(0.32), Inches(0.34),
                 desc, font=FONT_KO, size=8.5, color=T['text'],
                 line_spacing=1.22)

    # ============================================
    # 중단 — 2컬럼: Robust 전략 (좌) / EWI 트리거 (우)
    #            각 4개 예시, 헤더에 정의 + 미니 흐름도
    # ============================================
    EY = Inches(3.55)
    EH = Inches(3.05)
    EX = Inches(0.5)
    EW = Inches(12.33)
    cw = (EW - Inches(0.2)) / 2

    # ┌────────────────────────────────────────
    # │ COLUMN 1 — ROBUST 전략
    # └────────────────────────────────────────
    c1x = EX
    card(slide, c1x, EY, cw, EH, accent_top=T['blue'])

    # 헤더 — 큰 타이틀 + 정의
    add_text(slide, c1x + Inches(0.28), EY + Inches(0.2),
             Inches(3.5), Inches(0.3),
             'ROBUST 전략', font=FONT_EN, size=14, bold=True, color=T['blue_deep'])
    add_text(slide, c1x + Inches(2.05), EY + Inches(0.24),
             cw - Inches(2.5), Inches(0.28),
             '어떤 미래가 와도 + 가치를 만드는 전략',
             font=FONT_KO, size=10.5, color=T['ink'])

    # 정의 — "5개 시나리오 모두에서 가치 검증" 미니 시각화
    DY = EY + Inches(0.6)
    add_text(slide, c1x + Inches(0.28), DY,
             Inches(2.4), Inches(0.22),
             '5개 시나리오 가치 매트릭스',
             font=FONT_KO, size=8.5, color=T['caption'])
    # 5개 ✓ 박스
    chk_x0 = c1x + Inches(2.75)
    for i, lbl in enumerate(['B', 'A', 'D', 'C', 'E']):
        cx = chk_x0 + Inches(0.4) * i
        chip(slide, cx, DY - Inches(0.02), Inches(0.32), Inches(0.24),
             lbl, fg=T['blue'])
    add_text(slide, c1x + cw - Inches(1.4), DY,
             Inches(1.15), Inches(0.22),
             '모두 +가치 ✓', font=FONT_EN, size=9, bold=True,
             color=T['blue'], align='right')

    add_text(slide, c1x + Inches(0.28), DY + Inches(0.26),
             cw - Inches(0.56), Inches(0.18),
             '예시 4 / 총 9개',
             font=FONT_EN, size=8.5, color=T['caption'])
    add_line(slide, c1x + Inches(0.28), DY + Inches(0.5),
             c1x + cw - Inches(0.28), DY + Inches(0.5),
             color=T['divider'], width=BORDER_W)

    # 4개 RS 행
    rs_examples = [
        ('RS-1', '옵션형 캐파',
         '필요할 때 켜고, 안 필요할 때 끈다',
         '캐파 30%를 트리거 모듈로 분리 → 호황 풀가동·다운턴 자동 감산',
         '특히 D · C'),
        ('RS-2', '바벨 포트폴리오',
         '양 끝에 베팅, 중간은 버린다',
         'HBM·고부가 + 범용 양 끝 / 어정쩡한 가운데 제품군은 정리',
         '특히 A · B'),
        ('RS-3', '고객 록인',
         '떠나면 비싸지게 만든다',
         'CMX·SCADA·FDP 소프트웨어 + 다년 LTA + 사이트 전담 엔지니어',
         '특히 B · D'),
        ('RS-5', '재무 규율',
         '호황엔 절제, 다운턴엔 사수',
         '호황 capex 절제 + 다운턴 하한 4조 원 / 년 + 순현금 30조 원 버퍼',
         '특히 C · D'),
    ]
    body_y0 = DY + Inches(0.6)
    body_h = EY + EH - body_y0 - Inches(0.18)
    rs_h = (body_h - Inches(0.04) * 3) / 4
    for i, (code, name, principle, body, applies) in enumerate(rs_examples):
        ry = body_y0 + (rs_h + Inches(0.04)) * i
        # 좌측 컬러 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     c1x + Inches(0.28), ry, Inches(0.05), rs_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = T['blue']
        bar.line.fill.background()
        bar.shadow.inherit = False
        # 줄 1: 코드 + 이름 + 우측 적용 시나리오
        add_text(slide, c1x + Inches(0.42), ry + Inches(0.02),
                 Inches(0.6), Inches(0.22),
                 code, font=FONT_EN, size=10, bold=True, color=T['blue'])
        add_text(slide, c1x + Inches(1.05), ry + Inches(0.02),
                 cw - Inches(2.6), Inches(0.22),
                 name, font=FONT_KO, size=11, bold=True, color=T['ink'])
        add_text(slide, c1x + cw - Inches(1.4), ry + Inches(0.04),
                 Inches(1.15), Inches(0.2),
                 applies, font=FONT_KO, size=8.5, bold=True,
                 color=T['orange'], align='right')
        # 줄 2: 원리 (이탤릭 회색)
        add_text(slide, c1x + Inches(0.42), ry + Inches(0.25),
                 cw - Inches(0.7), Inches(0.2),
                 f'"{principle}"',
                 font=FONT_KO, size=9, italic=True, color=T['caption'])
        # 줄 3: 메모리 적용
        add_text(slide, c1x + Inches(0.42), ry + Inches(0.45),
                 cw - Inches(0.7), Inches(0.26),
                 body, font=FONT_KO, size=9, color=T['text'],
                 line_spacing=1.25)

    # ┌────────────────────────────────────────
    # │ COLUMN 2 — EWI 트리거
    # └────────────────────────────────────────
    c2x = EX + cw + Inches(0.2)
    card(slide, c2x, EY, cw, EH, accent_top=T['orange'])

    # 헤더
    add_text(slide, c2x + Inches(0.28), EY + Inches(0.2),
             Inches(1.6), Inches(0.3),
             'EWI', font=FONT_EN, size=14, bold=True, color=T['orange'])
    add_text(slide, c2x + Inches(0.95), EY + Inches(0.24),
             cw - Inches(1.2), Inches(0.28),
             'Early Warning Indicator  ·  분기점이 보이기 전에 30일 빨리 감지',
             font=FONT_KO, size=10.5, color=T['ink'])

    # 작동 흐름 미니 다이어그램
    DY2 = EY + Inches(0.6)
    add_text(slide, c2x + Inches(0.28), DY2,
             Inches(2.4), Inches(0.22),
             '작동 흐름  ·  예측이 아니라 자동 감지',
             font=FONT_KO, size=8.5, color=T['caption'])

    flow_steps = ['지표 모니터링', '임계치 도달', '경보 발령', '30일 내 행동']
    flow_x0 = c2x + Inches(2.5)
    fcell_w = (cw - Inches(2.7)) / 4
    for i, step in enumerate(flow_steps):
        fx = flow_x0 + fcell_w * i
        # 박스
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      fx, DY2 - Inches(0.02),
                                      fcell_w - Inches(0.06), Inches(0.24))
        rect.fill.solid()
        rect.fill.fore_color.rgb = T['white']
        rect.line.color.rgb = T['orange']
        rect.line.width = BORDER_W
        rect.shadow.inherit = False
        add_text(slide, fx, DY2 - Inches(0.02),
                 fcell_w - Inches(0.06), Inches(0.24),
                 step, font=FONT_KO, size=7.5, bold=True, color=T['ink'],
                 align='center', valign='middle')
        if i < 3:
            ax = fx + fcell_w - Inches(0.06)
            ay = DY2 + Inches(0.1)
            add_text(slide, ax - Inches(0.04), ay - Inches(0.08),
                     Inches(0.1), Inches(0.16),
                     '›', font=FONT_EN, size=10, bold=True,
                     color=T['orange'], align='center')

    add_text(slide, c2x + Inches(0.28), DY2 + Inches(0.26),
             cw - Inches(0.56), Inches(0.18),
             '예시 4 / 총 21개 (수요·거시·경쟁·정책 + 수요 변곡 선행)',
             font=FONT_EN, size=8.5, color=T['caption'])
    add_line(slide, c2x + Inches(0.28), DY2 + Inches(0.5),
             c2x + cw - Inches(0.28), DY2 + Inches(0.5),
             color=T['divider'], width=BORDER_W)

    # 4개 EWI 행
    ewi_examples = [
        ('수요', 'HBM 현물가  6개월  −30%',
         '시나리오 C 경보  →  HBM 신규 설비 집행 중단·소재 비축'),
        ('거시', '빅테크 AI CapEx YoY  <10%  (3분기 연속)',
         '시나리오 D 경보  →  RS-1 캐파 즉시 조절·다운턴 EWI 발동'),
        ('경쟁', 'CXMT  HBM3  양산 진척 확인',
         '중국 추격 가속  →  RS-3 고객 록인 강화·기술 격차 재투자'),
        ('정책', '시안 팹  라이선스 갱신 불발',
         '디커플링 가속  →  텍사스·평택 NAND 대체 라인 즉시 가동'),
    ]
    ewi_y0 = DY2 + Inches(0.6)
    ewi_body_h = EY + EH - ewi_y0 - Inches(0.18)
    ewi_h = (ewi_body_h - Inches(0.04) * 3) / 4
    for i, (sig_type, trigger, action) in enumerate(ewi_examples):
        ry = ewi_y0 + (ewi_h + Inches(0.04)) * i
        # 좌측 컬러 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     c2x + Inches(0.28), ry, Inches(0.05), ewi_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = T['orange']
        bar.line.fill.background()
        bar.shadow.inherit = False
        # 줄 1: 신호 유형 chip + 지표/임계치
        chip(slide, c2x + Inches(0.42), ry + Inches(0.02),
             Inches(0.55), Inches(0.22), sig_type, fg=T['orange'])
        add_text(slide, c2x + Inches(1.05), ry + Inches(0.02),
                 cw - Inches(1.25), Inches(0.22),
                 trigger, font=FONT_KO, size=10, bold=True, color=T['ink'],
                 valign='middle')
        # 줄 2: 자동 행동 (한 단계 들여서)
        add_text(slide, c2x + Inches(0.42), ry + Inches(0.28),
                 cw - Inches(0.7), Inches(0.45),
                 action, font=FONT_KO, size=9, color=T['text'],
                 line_spacing=1.30)

    # 하단 제안 띠
    banner(slide, Inches(6.75), '제안',
           '본 방법론으로 풀어가는 데 팀 합의 — 다음 단계는 STEEP·DF·시나리오 산출물 팀 리뷰와 보강',
           accent=T['orange'])

    footer_line(slide, 3)


# ============================================================
# Layout helper
# ============================================================

def _blank_layout(prs):
    layouts = prs.slide_layouts
    for idx in (6, 5, 0):
        if idx < len(layouts):
            return layouts[idx]
    return layouts[0]


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation(TEMPLATE)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 템플릿에 남아 있는 슬라이드 모두 제거 (XML rels 까지 깨끗하게)
    sldIdLst = prs.slides._sldIdLst
    rels = prs.part.rels
    for sld_id in list(sldIdLst):
        rId = sld_id.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        sldIdLst.remove(sld_id)
        if rId and rId in rels:
            prs.part.drop_rel(rId)

    print('Building 3-slide team alignment deck...')
    build_slide_1_why(prs)
    print('  Slide 1 (WHY) done')
    build_slide_2_how(prs)
    print('  Slide 2 (HOW) done')
    build_slide_3_what(prs)
    print('  Slide 3 (WHAT) done')

    print(f'Saving to: {OUTPUT}')
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f'Done. {len(prs.slides)} slides · {size_kb:.1f} KB')


if __name__ == '__main__':
    main()
