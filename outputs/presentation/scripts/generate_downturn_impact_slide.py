#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SP-2 다운턴 -- 원인·파급·문제 정의 1장 슬라이드 (v5).

단일 소스: wiki/downturn/samsung-impact.md §5
출력:     outputs/presentation/downturn-scenario-impact.pptx
재생성:   .venv/bin/python outputs/presentation/scripts/generate_downturn_impact_slide.py

v5 (2026-08-25): 수요 막대의 기준 구성을 2025 실측(TrendForce·Yole:
  DRAM $1,657억 중 HBM $340억·NAND $697억 -> 지수 14/56/30)으로 교체.
  공급 두 케이스는 수요 막대 대신 공급 전망 차트(3사 CAPEX 합산 $54->64->95B,
  CXMT 캐파 11->15%E·매출 8->14%E)로. 발표자 노트에 상세 해설 추가.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 토큰 ────────────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x14, 0x28, 0xA0)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
G_MID  = RGBColor(0x8A, 0x8A, 0x8E)
G_LINE = RGBColor(0xD9, 0xD9, 0xD9)
G_BG   = RGBColor(0xF6, 0xF6, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
T_DEM  = RGBColor(0xE9, 0xED, 0xF8)
T_SUP  = RGBColor(0xF3, 0xF0, 0xE9)
T_SHF  = RGBColor(0xEF, 0xEC, 0xF4)
# 제품 3색 (블루 램프) + 옅은 버전(현재 막대)
P_FULL = [RGBColor(0x14, 0x28, 0xA0), RGBColor(0x5A, 0x73, 0xCC), RGBColor(0xAD, 0xB9, 0xE4)]
P_LITE = [RGBColor(0xC4, 0xCA, 0xE7), RGBColor(0xD8, 0xDE, 0xF1), RGBColor(0xEA, 0xED, 0xF8)]
S_CUR  = RGBColor(0xC9, 0xC9, 0xCC)   # 공급 차트 현재/과거 막대
RED    = RGBColor(0xD9, 0x30, 0x25)   # 실측된 전조 (리스크 의미 전달 전용)
GREEN  = RGBColor(0x00, 0xA6, 0x51)   # 미발동 신호 (의미 전달 전용)
S_FUT  = RGBColor(0x6B, 0x6B, 0x70)   # 공급 차트 전망 막대

FONT = 'SamsungOneKorean'

def _font(run, size, bold=False, color=INK, name=FONT):
    f = run.font
    f.name, f.size, f.bold = name, Pt(size), bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set('typeface', name)

def txt(slide, x, y, w, h, parts, *, size=10, align=PP_ALIGN.LEFT,
        leading=1.12, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(parts, str):
        parts = [[(parts, {})]]
    elif parts and isinstance(parts[0], tuple):
        parts = [parts]
    for i, line in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if isinstance(line, str):
            line = [(line, {})]
        for text, kw in line:
            r = p.add_run()
            r.text = text
            _font(r, kw.get('size', size), kw.get('bold', False), kw.get('color', INK))
    return tb

def box(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, rnd=False,
        radius=0.055, dash=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if rnd:
        s.adjustments[0] = radius
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
        if dash:
            s.line.dash_style = dash
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def hline(slide, x, y, w, color=G_LINE, weight=0.5):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False

def vline(slide, x, y, h, color=G_LINE, weight=0.5):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False

# ── 데이터 (전부 위키 기인용) ────────────────────────────────────────────────
# 2025 글로벌 메모리 매출 실측: DRAM $1,657억(그중 HBM $340억) + NAND $697억
# = $2,354억 (TrendForce·Yole, wiki/concepts/{hbm-market,memory-market-overview}.md)
# -> 지수 100: HBM 14 · DRAM(HBM 제외) 56 · NAND 30
PRODUCTS = ['HBM', '기타 DRAM', 'NAND·SSD']
BASE_MIX = [14, 56, 30]

SCEN = [
    dict(cause='투자 자금이 끊기는 경우', group=0, chart='demand',
         head='「AI 인프라 자금줄 경색,\n데이터센터 발주 일제히 보류」',
         mk=('전조', 'Meta FCF -91% · Amazon FCF 적자 전환, 4사 CapEx ~$750B의 조달 의존'),
         mix=[7, 42, 24], newseg=0, note=None,
         prob='팔 곳 잃은 HBM4·서버 캐파:\nCIS 전환(공용 80%)인가,\ntake-or-pay 방어인가'),
    dict(cause='필요량이 줄어드는 경우', group=0, chart='demand',
         head='「AI는 호황인데 메모리는\n제자리… 효율화의 역설」',
         mk=('동인', 'KV 캐시 오프로드(H100 동시 사용자 10배) · CXL 풀링 · 모델 경량화'),
         mix=[11, 45, 35], newseg=0, note=None,
         prob='밀려난 수요는 SSD로 온다:\n서버 고객을 묶어둘\n표준·SW 고리의 확보'),
    dict(cause='CAPEX가 몰리는 경우', group=1, chart='capex',
         head='「신규 팹 동시 가동에 공급\n과잉… 치킨게임 재점화 우려」',
         mk=('공급', 'Micron 아이다호 · SK하이닉스 용인 · 국내 신규 팹, 2028~29 동시 가동(착공 확정)'),
         note='27E·28Eᵉ = 리서치 기반 추정',
         prob='범용·NAND 감산 결단의\n30일 규율: 2023년 6개월\n지연의 재발 차단'),
    dict(cause='후발이 파고드는 경우', group=1, chart='cxmt',
         head='「중국산 메모리 범용 시장\n잠식… 가격 하단이 사라졌다」',
         mk=('공급', 'CXMT 캐파 점유 11%→15%E(2028) · HBM3 월 6만장, YMTC Xtacking 본딩'),
         note='하단 점유의 구조적 상승',
         prob='DDR4·LPDDR4 하단 철수\n시점과 전환처(CIS·차량),\nHBM 인증 장벽 사수'),
    dict(cause='제품 정의가 바뀌는 경우', group=2, chart='demand',
         head='「HBM 시대 저무나…\n주문은 차세대 메모리로」',
         mk=('동인', '3D DRAM(4F² VCT) · zHBM · CXL 채택 개시로 표준 HBM 주문 이동'),
         mix=[6, 50, 35], newseg=8, note='수요가 차세대로 이동',
         prob='표준 HBM에서 zHBM·4F²\n3D DRAM으로의 전환기,\n별동대·R&D 하한 사수'),
]

# 공급 차트 데이터 (위키 기인용)
# DRAM 3사(삼성·SK·Micron) + NAND 6사(산업: 삼성·키옥시아·SK/솔리다임·Micron·샌디스크·YMTC)
# 부문별 CAPEX $B (TrendForce·각사 IR, memory-capex-history.md). 2027E·28E는 과잉 유지 가정(시나리오).
# 27E: DRAM 58.2(장비 매출 2027 +7.8% 프록시) + NAND 23.0(Kioxia ¥470B 고원·YMTC 2배)
# 28E: 27E 유지 가정. 근거: sources/articles/memory-capex-outlook-2027-2028-2026-08-26.md
CAPEX3 = [('2025', 46.5, 21.1, False), ('26E', 54.0, 22.2, False),
          ('27E', 58.2, 23.0, True), ('28E', 58.2, 23.0, True)]
CXMT   = [('캐파', 11, 15), ('매출', 8, 14)]            # 점유율 % (현재, 전망E)

GROUPS = [
    (0, 2, T_DEM, [('수요발', True), (' · 수요가 꺼진다', False)]),
    (2, 2, T_SUP, [('공급발', True), (' · 공급이 넘친다', False)]),
    (4, 1, T_SHF, [('전환발', True), (' · 판이 바뀐다', False)]),
]

# ── 슬라이드 ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

MX, MW = 0.47, 12.393

txt(slide, MX, 0.26, MW, 0.38, [[
    ('다운턴의 다섯 갈래:  ', {'bold': True, 'color': BLUE, 'size': 20}),
    ('꺼지는 수요가 둘, 넘치는 공급이 둘, 바뀌는 판이 하나다', {'bold': True, 'size': 20}),
]])
hline(slide, MX, 0.72, MW, BLUE, 1.4)

RAIL_W = 1.10
col_g = 0.12
col_w = (MW - RAIL_W - 4 * col_g) / 5
col_x = [MX + RAIL_W + i * (col_w + col_g) for i in range(5)]

GB_Y, GB_H = 0.86, 0.32
for st, n, tint, label in GROUPS:
    gx = col_x[st]
    gw = n * col_w + (n - 1) * col_g
    box(slide, gx, GB_Y, gw, GB_H, fill=tint, rnd=True, radius=0.13)
    txt(slide, gx, GB_Y + 0.065, gw, 0.2, [[(t, {'bold': b, 'size': 11}) for t, b in label]],
        align=PP_ALIGN.CENTER, leading=1.0)

CH_Y = 1.30
for i, sc in enumerate(SCEN):
    txt(slide, col_x[i], CH_Y, col_w, 0.22, [[(sc['cause'], {'bold': True, 'size': 12.5})]], leading=1.0)

for i in range(1, 5):
    vline(slide, col_x[i] - col_g / 2, 1.28, 4.94)

HL_Y, HL_H = 1.62, 0.64
txt(slide, MX, HL_Y + 0.04, RAIL_W - 0.14, 0.5,
    [[('그날의', {'bold': True, 'color': G_MID, 'size': 10})],
     [('헤드라인', {'bold': True, 'color': G_MID, 'size': 10})],
     [('(가상)', {'color': G_MID, 'size': 10})]],
    align=PP_ALIGN.RIGHT, leading=1.12)
for i, sc in enumerate(SCEN):
    x = col_x[i]
    box(slide, x, HL_Y, col_w, HL_H, fill=WHITE, line=G_LINE, line_w=0.75)
    box(slide, x, HL_Y, col_w, 0.032, fill=INK)
    txt(slide, x + 0.10, HL_Y + 0.09, col_w - 0.20, 0.5,
        [[(ln, {'bold': True, 'size': 10.5})] for ln in sc['head'].split('\n')], leading=1.18)

MK_Y = 2.40
txt(slide, MX, MK_Y + 0.02, RAIL_W - 0.14, 0.4,
    [[('무엇이', {'bold': True, 'color': G_MID, 'size': 10})],
     [('움직이나', {'bold': True, 'color': G_MID, 'size': 10})]],
    align=PP_ALIGN.RIGHT, leading=1.12)
for i, sc in enumerate(SCEN):
    pre, body = sc['mk']
    txt(slide, col_x[i], MK_Y, col_w, 0.62, [[
        (pre + '  ', {'bold': True, 'color': BLUE, 'size': 10}),
        (body, {'size': 10, 'color': INK})]], leading=1.16)

# ── 파급 차트 밴드 ──────────────────────────────────────────────────────────
BAR_LY = 3.00
# 좌측 레일: 수요 막대 범례
txt(slide, MX, BAR_LY, RAIL_W - 0.14, 0.18,
    [[('수요 막대', {'bold': True, 'color': G_MID, 'size': 10})]], align=PP_ALIGN.RIGHT, leading=1.0)
lg_items = list(zip(PRODUCTS, P_FULL)) + [('차세대 이동분', None)]
for li, (name, c) in enumerate(lg_items):
    ly = BAR_LY + 0.26 + li * 0.215
    if c is not None:
        box(slide, MX + 0.02, ly + 0.015, 0.11, 0.11, fill=c)
    else:
        box(slide, MX + 0.02, ly + 0.015, 0.11, 0.11, fill=WHITE, line=G_MID, line_w=0.8, dash=2)
    txt(slide, MX + 0.19, ly, RAIL_W - 0.20, 0.16, [[(name, {'color': GRAY, 'size': 10})]],
        leading=1.0, wrap=False)
txt(slide, MX, BAR_LY + 1.30, RAIL_W - 0.14, 0.34,
    [[('2025 구성', {'color': G_MID, 'size': 10})], [('= 100 기준', {'color': G_MID, 'size': 10})]],
    align=PP_ALIGN.RIGHT, leading=1.15)

# 컬럼별 차트 소제목
CHART_TITLE = {'demand': '제품 수요 (2025=100)', 'capex': '메모리 CAPEX ($B)', 'cxmt': 'CXMT 점유율 (%)'}
for i, sc in enumerate(SCEN):
    txt(slide, col_x[i], BAR_LY, col_w, 0.16,
        [[(CHART_TITLE[sc['chart']], {'bold': True, 'color': G_MID, 'size': 10})]],
        align=PP_ALIGN.CENTER, leading=1.0)

BASE_Y = 4.78
DSCALE = 0.0128    # 수요 지수 1 -> in

for i, sc in enumerate(SCEN):
    x = col_x[i]
    hline(slide, x + 0.10, BASE_Y, col_w - 0.20, G_LINE, 0.9)

    if sc['chart'] == 'demand':
        bw, bgap = 0.40, 0.14
        pair_x = x + (col_w - (2 * bw + bgap)) / 2
        yy = BASE_Y
        for v, c in zip(BASE_MIX, P_FULL):
            h = v * DSCALE
            box(slide, pair_x, yy - h, bw, h, fill=c, line=WHITE, line_w=0.75)
            yy -= h
        txt(slide, pair_x - 0.1, BASE_Y - 100 * DSCALE - 0.21, bw + 0.2, 0.16,
            [[('100', {'bold': True, 'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)
        sx = pair_x + bw + bgap
        yy = BASE_Y
        tot = sum(sc['mix'])
        for v, c in zip(sc['mix'], P_FULL):
            h = v * DSCALE
            box(slide, sx, yy - h, bw, h, fill=c, line=WHITE, line_w=0.75)
            yy -= h
        if sc['newseg']:
            h = sc['newseg'] * DSCALE
            box(slide, sx, yy - h, bw, h, fill=WHITE, line=G_MID, line_w=0.9, dash=2)
            yy -= h
            tot += sc['newseg']
        txt(slide, sx - 0.1, yy - 0.21, bw + 0.2, 0.16,
            [[(str(tot), {'bold': True, 'color': BLUE, 'size': 10.5})]], align=PP_ALIGN.CENTER, leading=1.0)
        txt(slide, pair_x - 0.08, BASE_Y + 0.05, bw + 0.16, 0.14,
            [[('2025', {'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)
        txt(slide, sx - 0.08, BASE_Y + 0.05, bw + 0.16, 0.14,
            [[('시나리오', {'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)

    elif sc['chart'] == 'capex':
        # DRAM 3사 + NAND 6사 부문별 CAPEX 누적 (2027E·28E = 과잉 유지 가정)
        # 미니 범례
        txt(slide, x, BAR_LY + 0.175, col_w, 0.14, [[
            ('■ ', {'color': S_FUT, 'size': 10}), ('DRAM 3사  ', {'color': GRAY, 'size': 10}),
            ('■ ', {'color': S_CUR, 'size': 10}), ('NAND 6사', {'color': GRAY, 'size': 10})]],
            align=PP_ALIGN.CENTER, leading=1.0)
        bw, bgap = 0.30, 0.135
        total_w = 4 * bw + 3 * bgap
        bx0 = x + (col_w - total_w) / 2
        CSCALE = 1.10 / 82
        for bi, (label, dram, nand, assumed) in enumerate(CAPEX3):
            bx = bx0 + bi * (bw + bgap)
            hd, hn = dram * CSCALE, nand * CSCALE
            if assumed:
                box(slide, bx, BASE_Y - hd, bw, hd, fill=RGBColor(0xA9, 0xA9, 0xAE), line=WHITE, line_w=0.75)
                box(slide, bx, BASE_Y - hd - hn, bw, hn, fill=RGBColor(0xE2, 0xE2, 0xE5), line=G_MID, line_w=0.7, dash=2)
            else:
                box(slide, bx, BASE_Y - hd, bw, hd, fill=S_FUT, line=WHITE, line_w=0.75)
                box(slide, bx, BASE_Y - hd - hn, bw, hn, fill=S_CUR, line=WHITE, line_w=0.75)
            tot = round(dram + nand)
            txt(slide, bx - 0.12, BASE_Y - hd - hn - 0.20, bw + 0.24, 0.16,
                [[(f'{tot}ᵉ' if assumed else str(tot),
                   {'bold': True, 'color': (G_MID if assumed else INK), 'size': 10})]],
                align=PP_ALIGN.CENTER, leading=1.0)
            txt(slide, bx - 0.12, BASE_Y + 0.05, bw + 0.24, 0.14,
                [[(label, {'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)

    elif sc['chart'] == 'cxmt':
        # CXMT 점유: 캐파 11 -> 15E(2028), 매출 8 -> 14E(2027)
        bw, bgap, ggap = 0.34, 0.10, 0.36
        total_w = 4 * bw + 2 * bgap + ggap
        bx0 = x + (col_w - total_w) / 2
        XSCALE = 1.15 / 15
        for gi, (label, cur, fut) in enumerate(CXMT):
            gx0 = bx0 + gi * (2 * bw + bgap + ggap)
            for bi, (v, col, lab) in enumerate([(cur, S_CUR, str(cur)), (fut, S_FUT, f'{fut}E')]):
                bx = gx0 + bi * (bw + bgap)
                h = v * XSCALE
                box(slide, bx, BASE_Y - h, bw, h, fill=col)
                txt(slide, bx - 0.12, BASE_Y - h - 0.21, bw + 0.24, 0.16,
                    [[(lab, {'bold': True, 'color': (INK if bi else G_MID), 'size': 10})]],
                    align=PP_ALIGN.CENTER, leading=1.0)
            txt(slide, gx0 - 0.06, BASE_Y + 0.05, 2 * bw + bgap + 0.12, 0.14,
                [[(label, {'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)

    if sc['note']:
        txt(slide, x, BASE_Y + 0.235, col_w, 0.15, [[(sc['note'], {'bold': True, 'color': GRAY, 'size': 10})]],
            align=PP_ALIGN.CENTER, leading=1.0)

# 풀어야 할 문제
PB_Y = 5.34
txt(slide, MX, PB_Y + 0.05, RAIL_W - 0.14, 0.4,
    [[('풀어야 할', {'bold': True, 'color': G_MID, 'size': 10})],
     [('문제', {'bold': True, 'color': G_MID, 'size': 10})]],
    align=PP_ALIGN.RIGHT, leading=1.12)
for i, sc in enumerate(SCEN):
    x = col_x[i]
    box(slide, x, PB_Y + 0.02, 0.035, 0.52, fill=BLUE)
    txt(slide, x + 0.12, PB_Y + 0.03, col_w - 0.12, 0.55,
        [[(ln, {'bold': True, 'size': 10.5})] for ln in sc['prob'].split('\n')], leading=1.18)

BY = 6.18
box(slide, MX, BY, MW, 0.34, fill=G_BG, rnd=True, radius=0.10)
txt(slide, MX + 0.16, BY + 0.075, MW - 0.32, 0.2, [[
    ('종합  ', {'bold': True, 'color': BLUE, 'size': 10.5}),
    ('수요발은 수요 막대가 줄고, 공급발은 수요가 그대로인 채 공급 곡선이 커진다.  ', {'size': 10.5}),
    ('원인이 다르면 풀어야 할 문제도 다르다.', {'bold': True, 'size': 10.5}),
]], leading=1.0)

txt(slide, MX, BY + 0.46, MW - 1.65, 0.34, [[
    ('수요 막대 기준: 2025 글로벌 메모리 매출 실측 구성(DRAM $1,657억 중 HBM $340억 · NAND $697억, TrendForce·Yole)=100, 시나리오 막대는 방향 가정 적용 · '
     '공급 차트: DRAM 3사·NAND 6사 부문별 CAPEX(TrendForce·각사 IR, 27Eᵉ=장비 매출 +7.8% 프록시+Kioxia ¥470B 고원·28Eᵉ=유지 가정)·CXMT 점유(TrendForce·FT) · 헤드라인은 가상 예시 · 대응 전략은 별도 장 · 기준일 2026-08',
     {'color': GRAY, 'size': 9})]], leading=1.25)
txt(slide, MX + MW - 1.5, BY + 0.46, 1.5, 0.18, [[('[문서등급 표기]', {'color': GRAY, 'size': 9})]],
    align=PP_ALIGN.RIGHT, leading=1.0)

# ── 발표자 노트 ─────────────────────────────────────────────────────────────
NOTES = """[슬라이드 위치] SP-2 다운턴 시나리오 플래닝 트랙의 영향 진단 장. 앞 장들이 "다운턴이 오는가"(EWI)와 "다섯 시나리오가 무엇인가"를 다뤘다면, 이 장은 원인 기준으로 다섯 갈래를 한눈에 정리하고 각 갈래가 남기는 "풀어야 할 문제"까지 연결한다. 대응 전략(DP-1~7 대비, DR-1~6 대응)은 다음 장에서 다룬다.

[왜 원인 중심인가] 다운턴 기간은 원인이 정한다는 패턴이 지난 20년 복기(DT08·DT12·DT16·DT19·DT23)에서 확인됐다. 수요발은 짧고 깊게(단기전), 공급발은 길게(장기전) 간다. 그래서 속도·기간 축은 이 장에서 제외했고, 원인만 맞히면 싸움의 성격이 따라온다. 시나리오 코드(DT-A 등)와 확률 수치도 발표 층에서는 제외했다. 확률은 추정 정밀도가 낮아 정성 판단만 남겼고, 정량 추정 근거는 wiki/downturn/scenario-matrix.md에 유지되어 있다.

[다섯 갈래 해설]
1. 투자 자금이 끊기는 경우(수요발·조달): 최종 수요가 아니라 돈이 먼저 끊기는 경로. 이미 전조가 실측됐다. Meta의 분기 FCF가 -91%(약 $7.8억까지 축소), Amazon은 TTM FCF가 적자로 전환했는데 4사 합산 CapEx는 ~$750B(+82% YoY)로 계속 늘고 있다. 지출과 현금창출의 방향이 갈라진 상태는 지속 기간에 한계가 있고, 신용 이벤트(네오클라우드·DC SPV 디폴트) 하나로 발주 동결이 시작될 수 있다. 이때 가장 먼저 비는 곳은 가장 많이 투자한 고부가(HBM·서버) 캐파다.
2. 필요량이 줄어드는 경우(수요발·원단위): AI는 성장하는데 작업당 메모리 소요가 꺾이는 경로. NVIDIA Dynamo 계열 KV 캐시 오프로드는 같은 H100으로 동시 사용자를 10배로 늘렸다. 같은 하드웨어로 10배를 처리한다는 것은 10배를 사지 않아도 된다는 뜻이기도 하다. CXL 풀링·모델 경량화도 같은 방향. 가격 폭락 없이 성장률만 꺾여서 위기감이 생기지 않는 것이 최대 위험이며, NAND·SSD는 오프로드 수혜로 오히려 수요가 는다(막대에서 유일하게 커지는 세그먼트).
3. CAPEX가 몰리는 경우(공급발·기존 진영): 수요는 멀쩡한데 호황기에 결정된 투자가 한꺼번에 캐파로 도착하는 경로. 차트는 부문별 CAPEX로, DRAM 3사(삼성·SK하이닉스·Micron)가 2025 $46.5B에서 2026E $54.0B, NAND 6사(삼성·키옥시아·SK/솔리다임·Micron·샌디스크·YMTC 산업 합계)가 $21.1B에서 $22.2B로 올라선다(TrendForce·각사 IR). 2027E·28E 막대는 리서치 기반 추정(ᵉ)이다: 2027E 합산 ≈$81B는 DRAM 3사 ≈$58B(반도체 장비 시장 전망의 DRAM 장비 매출 2027 +7.8%를 프록시로 적용)와 NAND 6사 ≈$23B(Kioxia FY26~28 연평균 ¥4,700억 고원 유지 계획 +66% vs FY25, Kioxia-SanDisk 2026 합산 +40%, YMTC 캐파 2026~27 약 2배)로 구성했고, 2028E는 공표치 부재로 27E 수준 유지 가정이다(TrendForce 2026-07-30: 신규 캐파의 실질 산출 기여가 2028년 본격화 = 투자 고원 지속과 정합). 실측·계획과 구분되도록 옅은 색·점선·ᵉ로 표기했다. 산출 방법과 한계는 sources/articles/memory-capex-outlook-2027-2028-2026-08-26.md 참조. 투자가 캐파로 바뀌는 리드타임 2~3년을 감안하면 2028~29에 Micron 아이다호, SK하이닉스 용인 1기, 국내 신규 팹이 동시에 가동된다(모두 착공이 끝난 확정 사실). 여기에 3강 절제가 무너지면(점유율 목표 선언, 조기 램프업) 급락으로 전환된다. 다섯 갈래 중 감산이 직접 효과를 내는 유일한 경우다.
4. 후발이 파고드는 경우(공급발·신흥): 보조금 기반 후발이 하단 가격대를 영구히 끌어내리는 경로. CXMT는 글로벌 DRAM 캐파 점유 11%에서 2028년 15%로, 매출 점유는 8%(2025 Q3)에서 14%(2027E)로 오르는 전망이고 HBM3도 월 6만 장 양산에 들어간다. NAND에서는 YMTC가 Xtacking 하이브리드 본딩으로 같은 일을 한다. 손실이 나도 퇴출되지 않는 공급자라서 치킨게임(소모전)의 승리 조건 자체가 없고, 다른 갈래와 달리 회복이 없다. 지금도 진행 중인 유일한 갈래다.
5. 제품 정의가 바뀌는 경우(전환발): 수요가 사라지는 게 아니라 다른 제품군으로 이동하는 경로. 3D DRAM(4F² VCT 셀), zHBM(커스텀 적층), CXL 채택이 시작되면 표준 HBM 주문이 옮겨간다. 막대 꼭대기의 점선 세그먼트가 그 이동분이다. 감산도 계약 방어도 무의미하고, 미리 그 자리에 가 있는 것(별동대)만 유효하다.

[수요 막대 읽는 법] 왼쪽 옅은 막대는 2025 글로벌 메모리 매출의 실측 구성을 100으로 지수화한 것이다. DRAM $1,657억(그중 HBM $340억) + NAND $697억 = $2,354억(TrendForce·Yole) 이므로 HBM 14, HBM 제외 DRAM 56, NAND·SSD 30. 오른쪽 진한 막대는 각 시나리오의 방향 가정(wiki/downturn/samsung-impact.md §5.3)을 이 실측 구성에 적용한 시나리오 값이다. 기준은 실측, 변화 폭은 시나리오 가정이라는 점을 구분해서 설명할 것. 서버/범용 DRAM 분리 실측치는 공개 소스에 없어 DRAM은 한 세그먼트로 묶었다.

[공급 차트 읽는 법] CAPEX가 몰리는 경우는 수요 막대가 무의미하다(수요는 그대로). 대신 공급의 선행 지표인 부문별 CAPEX(DRAM 3사 + NAND 6사 누적)를 보여준다. 2025~26E는 실측·계획, 2027E·28E는 시나리오 가정이라는 구분을 반드시 언급할 것. 후발의 경우도 마찬가지로 CXMT의 캐파·매출 점유 전망을 보여준다. 두 차트 모두 "수요가 아니라 공급 쪽 곡선이 커진다"는 것이 메시지다.

[풀어야 할 문제] 전략 자체가 아니라 전략이 풀어야 할 문제를 정의한 것이다(전략은 다음 장). 1번: HBM4 캐파는 +50% 증설이 잡혀 있는데 조달이 끊기면 팔 곳이 없다. DRAM-CIS 설비 공용률 80%인 전환 옵션과 take-or-pay·NTB 계약 방어선 중 무엇을 어디까지 쓸 것인가. 2번: 필요량이 줄어드는 경로에서 HBM을 떠난 수요는 KV 캐시 오프로드를 타고 SSD로 온다. 문제는 그 수요가 우리 eSSD(1위 38.2%)에 머물게 할 고리다. 하드웨어 스펙만으로는 교체가 쉬우니, 서버 고객의 스토리지 스택에 깊이 박히는 표준·소프트웨어 층(예: FDP 같은 데이터 배치 표준의 주도권)이 잠금 고리의 후보가 된다. 원단위 추적 지표(DX-8) 확보도 병행 과제. 3번: 2022-10 무감산 선언에서 2023-04 감산 공식화까지 6개월이 걸렸고 그 사이 DS 적자 -4.58조가 쌓였다. 같은 지연을 막을 30일 결정 규율이 문제. 4번: DDR4·LPDDR4 하단은 CXMT 침투로 돌아오지 않는다. 철수 시점과 라인 전환처(CIS·차량용), 그리고 HBM 인증 장벽 유지가 문제. 5번: 2019년 HBM팀 해체가 점유 40%에서 17%로의 추락으로 이어진 전례가 있다. 전환기에 별동대와 R&D 하한을 배분 논리에서 지켜내는 것이 문제.

[문제와 대응 방향의 연결] 다섯 문제는 다음 장(대응 전략)의 후보 방향과 자연스럽게 이어진다. 1·4번(캐파의 전환처, 하단 철수 후 전환)은 라인을 다른 제품으로 돌릴 수 있는 몸을 갖추는 방향을, 2번(SSD로 오는 수요를 묶어둘 표준·SW 고리)은 서버 고객 락인을 만드는 스토리지 소프트웨어 층의 주도권을, 5번(전환기 별동대·R&D 하한)은 zHBM 등 차세대 메모리 연구의 보존·집중을 가리킨다. 발표 시 결론을 못박지 말고 '문제가 가리키는 방향'으로만 열어 둘 것.

[출처] wiki/downturn/samsung-impact.md §5(단일 소스), scenario-matrix.md, downturn-history.md, key-drivers.md, memory-capex-outlook-2027-2028(2027~28 CAPEX 추정). 수치 원출처: hyperscaler-q2-2026-actuals(FCF·CapEx), kv-cache-ssd-offload(10배), memory-capex-history(3사 CAPEX), apple-cxmt-china-dram·TrendForce(CXMT), enterprise-ssd-market-1q26(38.2%), fab-toolset-commonality(공용률 80%), samsung-downturn-actions(2023 감산 지연). 헤드라인은 상황 이해를 돕는 가상 예시 문장으로 실제 보도가 아니다."""
slide.notes_slide.notes_text_frame.text = NOTES


# ═════════════════════════════════════════════════════════════════════════════
# 근접도 슬라이드 공용 빌더 — 시나리오별 "발현 근접도 평가" (슬라이드 2~6)
# 단일 소스: wiki/downturn/samsung-impact.md §5.1c 이하
# ═════════════════════════════════════════════════════════════════════════════
EST_FILL = RGBColor(0xE2, 0xE2, 0xE5)

def _numfmt(v):
    s = f'{v:,.1f}'
    return s[:-2] if s.endswith('.0') else s

def draw_bar_series(sl, px, pw, series, *, annot_color=BLUE):
    """세로 막대. series: [(라벨, 값, 주석, 추정여부)]. 추정 막대는 옅은 색·점선·ᵉ."""
    A_BASE = 5.30
    vmax = max(v for _, v, _, _ in series)
    scale = 3.30 / vmax
    n = len(series)
    bw, bgap = 0.62, 0.24
    need = n * bw + (n - 1) * bgap
    if need > pw - 0.2:
        bw = (pw - 0.2 - (n - 1) * bgap) / n
        need = n * bw + (n - 1) * bgap
    bx0 = px + (pw - need) / 2
    hline(sl, px + 0.1, A_BASE, pw - 0.2, G_LINE, 0.9)
    for bi, (label, v, annot, est) in enumerate(series):
        bx = bx0 + bi * (bw + bgap)
        h = v * scale
        if est:
            box(sl, bx, A_BASE - h, bw, h, fill=EST_FILL, line=G_MID, line_w=0.8, dash=2)
        else:
            box(sl, bx, A_BASE - h, bw, h, fill=S_FUT)
        txt(sl, bx - 0.15, A_BASE - h - 0.22, bw + 0.3, 0.16,
            [[(_numfmt(v) + ('ᵉ' if est else ''), {'bold': True, 'color': (G_MID if est else INK), 'size': 10.5})]],
            align=PP_ALIGN.CENTER, leading=1.0)
        if annot:
            txt(sl, bx - 0.15, A_BASE - h - 0.40, bw + 0.3, 0.15,
                [[(annot, {'color': annot_color, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)
        txt(sl, bx - 0.15, A_BASE + 0.05, bw + 0.3, 0.14,
            [[(label, {'color': G_MID, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)

def proximity_slide(spec):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    txt(sl, MX, 0.26, MW, 0.38, [[
        (spec['title_blue'], {'bold': True, 'color': BLUE, 'size': 20}),
        (spec['title_rest'], {'bold': True, 'size': 20})]])
    hline(sl, MX, 0.72, MW, BLUE, 1.4)
    txt(sl, MX, 0.84, MW, 0.2, [[(spec['subtitle'], {'color': GRAY, 'size': 10.5})]], leading=1.0)

    PA_X, PA_W = MX, 3.92
    PB_X, PB_W = MX + 4.10, 3.92
    PC_X, PC_W = MX + 8.20, MW - 8.20
    P_Y = 1.22
    for px, pw, title in ((PA_X, PA_W, spec['panelA']), (PB_X, PB_W, spec['panelB']),
                          (PC_X, PC_W, spec['panelC'])):
        box(sl, px, P_Y, pw, 0.045, fill=BLUE)
        head, rest = title.split('  ', 1)
        txt(sl, px, P_Y + 0.12, pw, 0.2, [[
            (head + '   ', {'bold': True, 'color': BLUE, 'size': 12.5}),
            (rest, {'bold': True, 'size': 12.5})]], leading=1.0)

    spec['chart'](sl, PA_X, PA_W)
    txt(sl, PA_X, 5.58, PA_W, 0.34,
        [[(c, {'size': 10})] for c in spec['chart_caption']], leading=1.2)

    for ti, (label, val_parts, sub, vsize) in enumerate(spec['tiles']):
        ty = 1.52 + ti * 1.42
        box(sl, PB_X, ty, PB_W, 1.26, fill=G_BG, line=G_LINE, line_w=0.5, rnd=True, radius=0.06)
        txt(sl, PB_X + 0.16, ty + 0.12, PB_W - 0.32, 0.17, [[(label, {'bold': True, 'color': GRAY, 'size': 10.5})]], leading=1.0)
        txt(sl, PB_X + 0.16, ty + 0.36, PB_W - 0.32, 0.34,
            [[(t, {'bold': True, 'color': c, 'size': vsize}) for t, c in val_parts]], leading=1.0)
        txt(sl, PB_X + 0.16, ty + 0.82, PB_W - 0.32, 0.32, [[(sub, {'color': GRAY, 'size': 10})]], leading=1.15)

    def crow(y, color, bold_t, rest_t):
        box(sl, PC_X + 0.02, y + 0.028, 0.10, 0.10, fill=color)
        txt(sl, PC_X + 0.22, y, PC_W - 0.24, 0.32,
            [[(bold_t, {'bold': True, 'size': 10}), (rest_t, {'size': 10, 'color': GRAY})]], leading=1.12)

    y = 1.52
    txt(sl, PC_X, y, PC_W, 0.17, [[(f"실측된 전조 {len(spec['red_rows'])}", {'bold': True, 'color': RED, 'size': 10.5})]], leading=1.0)
    y += 0.22
    for b, r in spec['red_rows']:
        crow(y, RED, b, r); y += 0.25
    y += 0.13
    txt(sl, PC_X, y, PC_W, 0.17, [[(f"미발동 발화 신호 {len(spec['green_rows'])}", {'bold': True, 'color': GREEN, 'size': 10.5})]], leading=1.0)
    y += 0.22
    for b, r in spec['green_rows']:
        crow(y, GREEN, b, r); y += 0.25
    y += 0.13
    txt(sl, PC_X, y, PC_W, 0.17, [[('완충', {'bold': True, 'color': GRAY, 'size': 10.5})]], leading=1.0)
    y += 0.22
    txt(sl, PC_X + 0.02, y, PC_W - 0.04, 0.5,
        [[(spec['buffer'][0], {'size': 10})],
         [(spec['buffer'][1], {'bold': True, 'size': 10})]], leading=1.25)
    wy = y + 0.59
    box(sl, PC_X, wy, PC_W, 1.10, fill=G_BG, rnd=True, radius=0.06)
    txt(sl, PC_X + 0.14, wy + 0.12, PC_W - 0.28, 0.9,
        [[('감시선  ', {'bold': True, 'color': BLUE, 'size': 10.5}), (spec['watch'][0], {'size': 10.5})]] +
        [[(w, {'size': 10.5})] for w in spec['watch'][1:]], leading=1.25)

    box(sl, MX, 6.18, MW, 0.34, fill=G_BG, rnd=True, radius=0.10)
    txt(sl, MX + 0.16, 6.255, MW - 0.32, 0.2, [[
        ('종합  ', {'bold': True, 'color': BLUE, 'size': 10.5}),
        (spec['verdict_plain'], {'size': 10.5}),
        (spec['verdict_bold'], {'bold': True, 'size': 10.5})]], leading=1.0)

    txt(sl, MX, 6.64, MW - 1.65, 0.34, [[(spec['footer'], {'color': GRAY, 'size': 9})]], leading=1.25)
    txt(sl, MX + MW - 1.5, 6.64, 1.5, 0.18, [[('[문서등급 표기]', {'color': GRAY, 'size': 9})]],
        align=PP_ALIGN.RIGHT, leading=1.0)
    sl.notes_slide.notes_text_frame.text = spec['notes']
    return sl


# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — "투자 자금이 끊기는 경우" 발현 근접도 (wiki §5.1c)
# ═════════════════════════════════════════════════════════════════════════════
NOTES2 = """[슬라이드 위치] 앞 장(다섯 갈래)의 1번 시나리오(투자 자금이 끊기는 경우)에 대한 심화 장. 질문은 하나다: 이 시나리오는 지금 얼마나 가까이 와 있는가. 구성은 지출(왼쪽), 창출(가운데), 판정(오른쪽)의 3단이다.

[왼쪽: 지출 패널] 4사(Alphabet·Amazon·Microsoft·Meta) 합산 AI CapEx는 2024 약 $205B, 2025 약 $410B(+100%), 2026E 약 $750B(+82%)로 2년 새 3.7배가 됐다. 2026E는 Q2 실적 발표 후 네 곳 전원이 가이던스를 올린 결과이고 삭감 발표는 0건이다(Amazon ~$220B, Google $195~205B, Meta $130~145B, MS ~$190B + FY27 $255~260B 상향). 2027E $1,100B+는 JPMorgan 전망으로 점선ᵉ 표기했다. 참고로 4사는 토지·전력·셸 같은 장기 자산은 조기 확약하고 칩은 필요 몇 달 전에 사는 디리스킹 관행을 공통 채택 중이다.

[가운데: 창출 패널] 지출을 받쳐줄 이익·현금의 실측이 핵심이다. 첫째, AI CapEx가 영업현금흐름에서 차지하는 비율이 2024년 76%에서 2026E 94%로 올라 자기잠식 임계점에 접근했다. 내부 현금만으로는 지출을 유지하기 어려워지는 지점이며, 그 다음 수단이 부채·SPV 조달이라서 이 시나리오(조달 경색)의 노출이 커진다. 둘째, Meta의 분기 FCF가 -91%($784M)로 급감했고 Amazon은 TTM FCF가 약 -$7.6B로 적자 전환해 다이버전스가 개별 회사가 아니라 패턴이 될 조짐이다. 셋째, Bain 프레임에 따르면 2030년 지출 궤도를 수익성 있게 회수하려면 연 $2조의 신규 AI 매출이 필요한데 자금 갭이 연 $800B다. 지출의 지속 가능성 자체가 외부 조달에 걸려 있다는 뜻이다.

[오른쪽: 판정 보드] 실측된 전조 3건(빨강)과 미발동 발화 신호 4건(초록)을 구분하는 것이 이 장의 핵심 규율이다. 전조: FCF 다이버전스 2사, CapEx/현금흐름 94%, 공급자(SK하이닉스) 영업이익률 사상 최대(~76%)라는 후기순환 신호. 미발동: CapEx 삭감 발표 0건(감별 트리거는 YoY -25%), GPU 임대가 firming(H100 ~$2.69·H200 ~$4.38, 트리거는 6개월 -35% 급락), DC 착공·발주 취소 미관측(파이프라인 55.9GW 유지), Google Cloud 백로그 $514B로 오히려 +$50B 증가. 완충 요인으로 take-or-pay·NTB 계약 락인이 발화 후에도 초기 6~12개월의 매출을 흡수한다. 그래서 같은 발화라도 커버리지 만기가 몰려 있으면 절벽, 분산되어 있으면 완만한 하강이 된다.

[판정을 말로 하면] "가능성이 오르는 국면이지, 발현 국면이 아니다." 확률 숫자는 표기하지 않되, 방향은 분명하다: 괴리(전조)는 계속 벌어지고 있고, 발화는 아직 시작되지 않았다. 다음 변곡은 FCF 반전이 실제 발주 동결로 번지는 순간이며, 빅테크 FCF·조달 스프레드와 GPU 임대가가 그 최선행 감시선이다. 질문이 나오면: 이 감시선은 대시보드 EWI로 이미 배선되어 있고, 트리거 도달 시 30일 내 대응 규율과 연결된다.

[출처] wiki/downturn/samsung-impact.md §5.1c(단일 소스). 수치 원출처: ai-capex(연도별 CapEx), hyperscaler-q2-2026-actuals(가이던스 상향·FCF·임대가), hyperscaler-q2-2026-capex(백로그·SK OPM), ai-compute-economics-gap(Bain $2조·$800B), visualizations 하이퍼스케일러 구성(76%→94%). 2027E는 JPMorgan 전망 인용."""

SPEC_A = dict(
    title_blue='끊기는 경우의 근접도:  ',
    title_rest='전조는 실측되었지만, 발화 신호는 아직 없다',
    subtitle='1번 시나리오(투자 자금이 끊기는 경우)의 발현 가능성 평가: 지출 속도와 창출 속도의 괴리, 그리고 발화 전 체크리스트',
    panelA='지출  AI CapEx, 4사 합산 ($B)',
    panelB='창출  이익·현금이 따라가지 못한다',
    panelC='판정 보드  발화 전 체크리스트',
    chart=lambda sl, px, pw: draw_bar_series(sl, px, pw, [
        ('2024', 205, '', False), ('2025', 410, '+100%', False),
        ('26E', 750, '+82%', False), ('27E', 1100, 'JPM 전망', True)]),
    chart_caption=['Q2 실적 후 4사 전원 상향, 삭감 발표 0건.',
                   '장기 자산(토지·전력·셸)은 조기 확약 관행 지속'],
    tiles=[
        ('CapEx / 영업현금흐름', [('76%', G_MID), ('  →  ', G_MID), ('94%', BLUE)],
         "2024 → 2026E, '자기잠식' 임계점 근접", 19),
        ('하이퍼스케일러 FCF 반전', [('Meta -91%', BLUE), (' · ', G_MID), ('Amazon 적자', BLUE)],
         '분기 $784M로 급감 · TTM ~-$7.6B, 두 번째 반전', 16),
        ('필요 매출 vs 자금 갭 (Bain)', [('$2조/년', BLUE), ('  vs 갭 ', G_MID), ('$800B/년', BLUE)],
         '2030 수익성 회수에 필요한 신규 매출과 부족분', 19),
    ],
    red_rows=[
        ('FCF 다이버전스 2사', ': Meta·Amazon 반전'),
        ('CapEx/현금흐름 94%', ': 자기잠식 임계'),
        ('공급자 이익 정점', ': SK OPM ~76% 사상 최대'),
    ],
    green_rows=[
        ('CapEx 삭감 0건', ': 트리거 YoY -25%'),
        ('GPU 임대가 firming', ': H100 $2.69·H200 $4.38'),
        ('DC 취소 미관측', ': 파이프라인 55.9GW 유지'),
        ('Cloud 백로그 견조', ': Google $514B(+$50B QoQ)'),
    ],
    buffer=['take-or-pay·NTB 락인이 발화 후에도 초기 6~12개월을 흡수.',
            '그래서 관건은 커버리지의 만기 구조다'],
    watch=['다음 변곡 = FCF 반전이',
           '발주 동결로 번지는 순간.',
           '빅테크 FCF·조달 스프레드와 GPU',
           '임대가(6개월 -35%)가 최선행 감시선'],
    verdict_plain='지출은 2년 새 3.7배가 됐는데 내부 현금 여력은 임계에 닿았다. 전조 3건 실측, 발화 신호 0건.  ',
    verdict_bold='가능성이 오르는 국면이지, 아직 발현 국면은 아니다.',
    footer=('출처: wiki/downturn/samsung-impact.md §5.1c · CapEx: 각사 가이던스 종합(2026E ~$750B)·27Eᵉ JPMorgan 전망 · FCF·백로그·임대가: Q2 2026 실적 보도 종합 · '
            'CapEx/현금흐름·자금 갭: SemiAnalysis 계열·Bain · 기준일 2026-08'),
    notes=NOTES2,
)
# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 3 — "필요량이 줄어드는 경우" 발현 근접도 (wiki §5.1d)
# ═════════════════════════════════════════════════════════════════════════════
NOTES3 = """[슬라이드 위치] 2번 시나리오(필요량이 줄어드는 경우)의 근접도 심화 장. 이 시나리오는 "전조 없음"이 아니라 "전조 과잉"이 특징이다. 효율화는 모델(KV 캐시 90% 절감), 시스템(CXL 서버 -25%, H100 10배 서빙), 가격(토큰 비용 280분의 1) 세 계층 모두에서 이미 실측·상용화됐다. 그런데도 메모리 시장은 사상 최대 호황이고 전망은 상향 중이다. 이 모순이 시나리오의 본질이다: 효율화는 수요를 죽이는 힘이 아니라 총량 폭증과 매 분기 경주하는 힘이며, 현재까지는 매번 지고 있다.

[왼쪽: 수요 패널] HBM 비트수요 증가율은 2025 +130% → 2026E +70% → 2027E +50~60%로 3년 연속 감속 전망(TrendForce)이다. 증가율의 1차 미분이 계속 마이너스라는 것이 실측된 방향이고, 레벨은 여전히 고성장이며 절대 시장 전망은 오히려 상향 중($551.6B→$889.3B, 2027 $1.28조)이라는 것이 반대 방향이다. 이 두 사실이 공존하는 것이 지금 국면이다.

[가운데: 절감 패널] 토큰당 추론 비용은 GPT-3.5급 고정 성능 기준 2022-11 $20.00 → 2024-10 $0.07로 18개월에 280분의 1(Stanford AI Index), 성능 고정 비용은 연 중앙값 50배 하락에 2024년 이후 연 200배로 가속(Epoch AI). KV 캐시는 DeepSeek-V4-Pro가 100만 토큰 컨텍스트에서 이전 대비 10%만 소요(90% 절감, arXiv 2606.19348). CXL은 Meta가 재활용 DDR4를 수백만 대 규모로 실배치해 일부 추론 서버를 최대 25% 절감 — 논문이 아니라 운영 인프라다.

[오른쪽: 판정 보드] 전조 4건: 원단위 수직 낙하, 모델 절감 상용화, 시스템 실배치, 가속기당 HBM 탑재 정체 조짐(Rubin 288GB 동결, Rubin Ultra 1TB→최저 256GB 구성 검토 — 단 명시 사유는 공급 부족이므로 수요발 효율화로 오독 금지, 관측 오염 주의). 미발동 4건: DX-8(메모리 매출 성장률 < AI 지출 성장률 2분기 연속) 미발동 — 오히려 전망 상향, DX-1 CapEx 감속 없음(+82%·삭감 0건), DX-2 임대가 하락 없음(H100 $2.69 firming), 토큰 총소비 증가율 둔화 없음(Google 월 3,200조 토큰, YoY 7배). 완충: Jevons 상쇄가 완벽 작동 중 — 에이전트(인간의 5배 토큰)·코딩(토큰 비중 11%→50%+)이 절감분 이상을 되사는 중.

[판정을 말로 하면] 승부는 두 힘의 순차다. 원단위 하락은 연 50배 속도의 확정 추세라 되돌아오지 않는다. 토큰 폭증은 에이전트 채택 곡선에 달린 가변 추세다. 가변 쪽이 꺾이는 순간이 T=0이고 그 시점에 가격 폭락 같은 극적 신호는 없다. 그래서 감시 대상은 효율화가 아니라 상쇄력의 증가율이며, DX-8 순방향과 토큰 총소비 증가율을 한 쌍으로 판독한다. 이 시나리오만의 상방: KV 오프로드·계층화는 HBM 소요를 깎는 동시에 SSD 수요를 늘린다 — eSSD 1위(38.2%) 삼성에게 유일하게 헤지가 내장된 다운턴이다.

[출처] wiki/downturn/samsung-impact.md §5.1d(단일 소스). 수치 원출처: dt-b-demand-efficiency-signals-2026-08-28(토큰 비용·KV 캐시·CXL·비트수요·토큰 소비), kv-cache-ssd-offload-ecosystem(10배 서빙), semianalysis-vera-rubin(Rubin 탑재), hyperscaler-q2-2026-actuals(CapEx·임대가), enterprise-ssd-market-1q26(38.2%)."""

SPEC_B = dict(
    title_blue='줄어드는 경우의 근접도:  ',
    title_rest='원단위 폭락은 실측되었지만, 총량 폭증이 아직 삼키고 있다',
    subtitle='2번 시나리오(필요량이 줄어드는 경우)의 발현 가능성 평가: 원단위 폭락과 총량 폭증의 경주, 그리고 발화 전 체크리스트',
    panelA='수요  HBM 비트수요 증가율 (%, YoY)',
    panelB='절감  원단위가 수직 낙하한다',
    panelC='판정 보드  발화 전 체크리스트',
    chart=lambda sl, px, pw: draw_bar_series(sl, px, pw, [
        ('2025', 130, '', False), ('26E', 70, '', False), ('27E', 55, 'TrendForce', True)]),
    chart_caption=['증가율의 1차 미분이 3년 연속 마이너스.',
                   '레벨은 여전히 고성장, 절대 전망은 오히려 상향'],
    tiles=[
        ('토큰당 추론 비용 (성능 고정)', [('280분의 1', BLUE)],
         '18개월 새 $20 → $0.07, 연 중앙값 50배 하락', 19),
        ('KV 캐시 원단위', [('90% 절감', BLUE)],
         'DeepSeek-V4-Pro, 100만 토큰 컨텍스트 실측', 19),
        ('CXL 실배치', [('서버 -25%', BLUE)],
         'Meta 재활용 DDR4 수백만 대, 실운용 진입', 19),
    ],
    red_rows=[
        ('원단위 수직 낙하', ': 18개월 280분의 1'),
        ('모델 절감 상용화', ': KV 캐시 90% 절감 도달'),
        ('시스템 실배치', ': H100 10배 서빙·CXL -25%'),
        ('탑재량 정체 조짐', ': Rubin Ultra 하향 검토'),
    ],
    green_rows=[
        ('DX-8 괴리 미발동', ': 시장 전망 $889B로 상향'),
        ('CapEx 감속 없음', ': 4사 +82%, 삭감 0건'),
        ('임대가 하락 없음', ': H100 $2.69 firming'),
        ('토큰 증가율 견조', ': Google YoY 7배'),
    ],
    buffer=['Jevons 상쇄 작동 중: 에이전트·코딩이 절감분 이상을 되산다.',
            '관측 대상은 효율화가 아니라 상쇄력의 증가율이다'],
    watch=['DX-8 순방향(메모리 매출',
           'vs AI 지출 성장률 괴리)과',
           '토큰 총소비 증가율이 한 쌍.',
           '탑재 하향은 사유 병기 판독'],
    verdict_plain='전조는 다섯 갈래 중 가장 두껍게 실측됐지만 총량 폭증이 매 분기 이기는 중. 발화 신호 0건.  ',
    verdict_bold='조용한 국면이지만, 상쇄가 꺾이면 경보 없이 발현된다.',
    footer=('출처: wiki/downturn/samsung-impact.md §5.1d · 토큰 비용: Stanford AI Index·Epoch AI · KV 캐시: arXiv 2606.19348 · '
            'CXL: Meta 실배치 보도 · 비트수요: TrendForce · 기준일 2026-08'),
    notes=NOTES3,
)

# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 4 — "CAPEX가 몰리는 경우" 발현 근접도 (wiki §5.1e)
# ═════════════════════════════════════════════════════════════════════════════
NOTES4 = """[슬라이드 위치] 3번 시나리오(CAPEX가 몰리는 경우)의 근접도 심화 장. DT-A 판과의 결정적 차이는 예고 시간이다 — DX-4(투입-출하 갭·증설 공시)는 선행 12개월로 전 지표 중 가장 길고, 장비 리드타임 12~24개월 때문에 지금 보이는 발주·착공이 곧 2028년의 공급이다. 다섯 시나리오 중 유일하게 "미래의 공급을 오늘 관측할 수 있는" 시나리오이며, 그 관측 결과가 이 슬라이드다.

[왼쪽: 공급 패널] DRAM 3사(삼성·SK·Micron) + NAND 6사 합산 CAPEX는 2025 $67.6B(46.5+21.1) → 2026E $76.2B(+13%, 삼성 20.0·SK 20.5·Micron 13.5) → 2027E ≈$81Bᵉ(리서치 추정: DRAM 장비 매출 +7.8% 프록시 + Kioxia ¥470B 고원·YMTC 2배) → 2028E ≈$81Bᵉ(유지 가정). 투자→캐파 리드타임 2~3년이므로 산출 도래 창은 2027H2~2028에 집중된다. 2026년 8월 한 달 사이에만 SK 54조(용인 Y2+M17) 이사회 승인, 삼성 P5 조기 본공사, Kioxia-SanDisk $31B+ 일본 투자 보도가 겹쳤고, SEMI는 300mm 장비 지출이 2026 $133B → 2027 $151B로 연속 사상 최대를 경신한다고 본다.

[가운데: 장전 패널] 장비 리드타임 12~24개월 — 오늘 발주분은 취소 불가 상태로 2028년경 도착한다. 2027년 NAND는 공급 증가율이 수요를 추월해 2H27부터 타이트가 풀린다(TrendForce, 중국 비트 점유 ~19%) — 6사 축이 DT-C의 첫 실험장이 되고, DRAM은 2027년에도 sufficiency -1~-2% 적자다. 공급사 재고는 2~4주로 역대 최저(2024말 13~17주) — 현재는 과잉의 반대 극단이나, 바로 이 숫자가 확정 일괄 발주를 정당화하는 명분이 된다(2차 자료 기반, 업계 추정 병기).

[오른쪽: 판정 보드] 전조 4건: CAPEX 재가속(67.6→76.2→81ᵉ), 신규 팹 2027H2~2028 창 동시 확정(용인 Y1 클린룸 27-02·M15X 26-11 양산·삼성 P4/P5·Micron ID1/ID2·Kioxia K2), 램프 앞당기기 경쟁 개시(용인 3개월·P5 6개월 조기 — 절제 이탈의 회색 전단계), 후발 무규율 증설(CXMT 2027 42만 wpm 목표·상하이 팹 장비 반입). 미발동 4건: DX-4 투입-출하 갭 미발동(DRAM 수급 적자 지속), 절제 이탈 공식화 미관측(Micron CEO "discipline" 4회, 주요 고객 수요의 50~2/3만 충족 발언), DX-6 현물-계약 스프레드 미발동(계약가 Q2 +58~63% 상승 지속), 재고일수 반전 없음(2~4주 최저). 완충: take-or-pay SCA/LTA가 신규 캐파를 사전 판매로 묶고(Micron SCA 16건·최소 계약 매출 ~$100B·예치금 $22B), HBM의 웨이퍼 3~4배 소비가 범용 캐파를 잠식하며, 반독점 계류가 치킨게임 유인을 낮춘다.

[판정을 말로 하면] "장전은 끝났고 방아쇠는 당겨지지 않았다." 과거 패턴: 2017~18 CAPEX 정점의 청구서는 2019년에, 2021~22 정점의 청구서는 2023년 최악 다운턴으로 돌아왔다 — 정점에서 과잉까지 1~2년, 착공에서 산출까지 2~3년. 이번 정점(2026~27E)의 만기는 2028년이다. 순서에 주목: NAND 2H27 완화가 DRAM 2028의 예고편이며, 가동률 50%의 Kioxia와 캐파 2배의 YMTC가 완화 국면에서 보이는 행동이 3강의 2028년을 미리 보여준다. 그리고 이 시나리오의 최대 함정은 도착 후가 아니라 지금이다 — 재고 최저·수요 미충족이라는 현재의 숫자가 확정 일괄 발주의 명분이 되므로, DP-2(Shell 선행·장비 단계 반입) 규율의 시험대는 2028년이 아니라 이번 발주 사이클이다.

[출처] wiki/downturn/samsung-impact.md §5.1e(단일 소스). 수치 원출처: memory-capex-history·memory-capex-outlook-2027-2028(CAPEX), dt-c-capex-oversupply-signals-2026-08-28(팹 일정·리드타임·수급·재고), bloomberg-micron-ceo-virginia(discipline), micron-q3-fy26(SCA $100B), price-trends(계약가)."""

SPEC_C = dict(
    title_blue='몰리는 경우의 근접도:  ',
    title_rest='2028에 도착할 캐파는 장전됐고, 발화 신호는 전부 미발동이다',
    subtitle='3번 시나리오(CAPEX가 몰리는 경우)의 발현 가능성 평가: 취소 불가 파이프라인의 실측, 그리고 발화 전 체크리스트',
    panelA='공급  9사 CAPEX 합산 ($B)',
    panelB='장전  파이프라인은 취소 불가 상태다',
    panelC='판정 보드  발화 전 체크리스트',
    chart=lambda sl, px, pw: draw_bar_series(sl, px, pw, [
        ('2025', 67.6, '', False), ('26E', 76.2, '+13%', False),
        ('27E', 81, '리서치 추정', True), ('28E', 81, '유지 가정', True)]),
    chart_caption=['DRAM 3사 + NAND 6사. 투자→캐파 리드타임 2~3년,',
                   '산출 도래 창은 2027H2~2028에 집중'],
    tiles=[
        ('장비 리드타임', [('12~24개월', BLUE)],
         '오늘 발주분은 취소 불가 상태로 2028년경 도착', 19),
        ('2027 NAND 수급', [('공급 > 수요', BLUE), ('  역전', G_MID)],
         '2H27 완화 개시, 6사 축이 첫 실험장', 18),
        ('공급사 재고', [('2~4주', BLUE), ('  역대 최저', G_MID)],
         '이 숫자가 증설 가속의 명분이 된다', 18),
    ],
    red_rows=[
        ('CAPEX 재가속', ': 67.6 → 76.2 → 81ᵉ $B'),
        ('신규 팹 동시 확정', ': 용인·P4/P5·ID1/2·K2'),
        ('램프 앞당기기', ': 용인 3개월·P5 6개월 조기'),
        ('후발 무규율 증설', ': CXMT 42만 wpm 목표'),
    ],
    green_rows=[
        ('투입-출하 갭 (DX-4)', ': DRAM 수급 -1~-2% 적자'),
        ('절제 이탈 미관측', ': "discipline" 공언 유지'),
        ('현물-계약 (DX-6)', ': 계약가 +58~63% 상승 중'),
        ('재고 반전 없음', ': 2~4주 역대 최저 지속'),
    ],
    buffer=['take-or-pay SCA 사전 판매 + HBM 웨이퍼 3~4배 잠식이 완충.',
            '함정은 도착 후가 아니라 지금의 확정 일괄 발주다'],
    watch=['DX-4 투입-출하 갭의 부호',
           '전환(최장 12개월 선행)과',
           '절제 이탈 발언, 그리고 NAND',
           '2H27 완화가 2028의 예고편'],
    verdict_plain='구조적 전제는 전부 실측됐고 2028 캐파 도래는 확정에 가깝다. 전조 4건 실측, 발화 신호 0건.  ',
    verdict_bold='장전은 끝났고, 방아쇠는 당겨지지 않았다.',
    footer=('출처: wiki/downturn/samsung-impact.md §5.1e · CAPEX: TrendForce·각사 IR, 27~28Eᵉ 리서치 추정 · 팹 일정·리드타임: 2026-08 업계 보도 종합 · '
            '수급·재고: TrendForce·업계 추정 · 기준일 2026-08'),
    notes=NOTES4,
)

# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 5 — "후발이 파고드는 경우" 발현 근접도 (wiki §5.1f)
# ═════════════════════════════════════════════════════════════════════════════
NOTES5 = """[슬라이드 위치] 4번 시나리오(후발이 파고드는 경우)의 근접도 심화 장. 이 시나리오는 방아쇠를 기다리지 않는다 — DT-A~C가 "언제 오는가"의 문제라면 DT-D는 "얼마나 진행됐는가"의 문제다.

[왼쪽: 물량 패널] CXMT 캐파는 2020년 3만 → 2026년 8월 현재 30만 wpm(허페이 2+베이징 1)으로 6년 만에 10배가 됐고, 연말 35만이면 마이크론(37.5만)과 물량이 동급이 된다. 2026년 7월 STAR 상장으로 $8.6B를 조달(상장일 +466%)해 상하이·베이징 제2팹까지 실현되면 60만 wpm 이상, 2030년 점유 30%가 자체 공언 목표다. 캐파 점유는 11%→15%E(2028, FT), 매출 점유는 Q4 2025 7.67%로 세계 4위(Omdia). 보조금 기반이라 손실이 퇴출로 이어지지 않는 구조에 이제 시장 자본까지 얹혔다.

[가운데: 문턱 패널] 2026년 들어 질적 문턱 두 개가 넘어갔다. 하나는 수율 — 17nm DDR5 수율 90% 돌파 보도로 삼성 동세대(92~93%)와의 격차가 오차 범위 수준으로 좁혀졌다(공식 확인 없는 업계 소식통 보도임을 반드시 병기). 다른 하나는 인증 — 애플이 중국 내수용 기기 검증에 착수했고, HP·Asus·Acer가 인증을 완료하고 노트북에 넣기 시작했다. "중국 내수 저가 전용"이라는 방어선이 서방 OEM으로 뚫린 것이다. 모바일에서는 중국 내수 LPDDR 30%+ 점유·LPDDR6로 Xiaomi 플래그십 SoC 공급 경쟁 진입, 서버에서는 Tencent 다년 계약(¥200억+ 보도)·Alibaba Cloud 고객 확보. NAND 축은 완료형이다 — YMTC가 Q2 2026 출하 점유 14%로 마이크론·키옥시아를 제치고 3위, 2027년 말 1위 공언. DRAM에서 2~3년 뒤 벌어질 일의 리허설을 NAND에서 지금 보고 있다.

[오른쪽: 판정 보드] 전조 4건: 캐파 30만 wpm(6년 10배·IPO $8.6B), DDR5 수율 90% 보도, 서방 인증 개방(HP·Asus·Acer), NAND 침식 완료형(YMTC 3위). 미발동 4건: 애플-CXMT 인증 정식 승인 미결(DX-5 리트머스 — 상원 초당파 서한 vs 9월 미중 정상회담 후 허용 관측이 맞서는 정치전), 허페이 본디드 DRAM 파일럿의 양산 전환(EUV 우회) 미확정, HBM3 양산 공식화 미발동(샘플링 단계·캐파 20% 배정 계획만), 그리고 핵심 역설 — 레거시 가격 앵커 하락 미발동: 정반대로 DDR4 현물 +172% 급등·DDR4>DDR5 역전. 3사가 HBM·DDR5로 웨이퍼를 옮기며 DDR4 공급이 수요보다 빨리 소멸했기 때문이고, CXMT조차 DDR4를 정리하고 위로 올라가는 중이다. 침식의 '구조'(공급자가 사라지지 않고 상향 침투)는 완성되고 있는데 '증상'(가격 하락)은 슈퍼사이클이 가리고 있다 — 임계 기반 경보만 보면 영원히 안 울리는 침식형의 전형이다.

[판정을 말로 하면] 완충은 실재하지만 유통기한이 있다. EUV 차단으로 sub-10nm이 막혀 있고 HBM4 이후 베이스다이에 필요한 선단 로직이 없다 — 그러나 본디드 DRAM 파일럿과 YMTC-CXMT 하이브리드 본딩 제휴가 그 천장을 시험 중이다. 감시선은 셋: 9월 미중 정상회담 전후 애플 건의 결말(EWI cxmt_apple_qualification), 캐파 점유 15% 도달·HBM3 양산 공식화(DX-5), 그리고 공급 부족이 풀리는 순간 레거시 가격이 어디에 착지하는가. 특히 마지막이 진짜 판정 순간이다 — 다음 다운턴이 오면 낮아진 앵커가 한꺼번에 드러난다.

[출처] wiki/downturn/samsung-impact.md §5.1f(단일 소스). 수치 원출처: hybrid-bonding-structures-china-limits(캐파·IPO·EUV 천장), dt-d-china-entrant-signals-2026-08-28(수율·OEM·점유율·가격·YMTC), apple-cxmt-china-dram(애플 건·캐파 점유 전망). 수율 90%는 공식 미확인 보도임을 발표 시 명시할 것."""

SPEC_D = dict(
    title_blue='파고드는 경우의 근접도:  ',
    title_rest='침식의 구조는 완성 단계, 가격 앵커만 슈퍼사이클이 붙들고 있다',
    subtitle='4번 시나리오(후발이 파고드는 경우)의 발현 가능성 평가: 물량·수율·인증의 실측, 그리고 발화 전 체크리스트',
    panelA='물량  CXMT DRAM 캐파 (만 wpm)',
    panelB='문턱  질의 장벽이 넘어가고 있다',
    panelC='판정 보드  발화 전 체크리스트',
    chart=lambda sl, px, pw: draw_bar_series(sl, px, pw, [
        ('2020', 3, '', False), ('2023', 12, '', False), ('2024', 17, '', False),
        ('26.8', 30, '10배', False), ('26말', 35, '', True), ('장기', 60, '계획', True)]),
    chart_caption=['연말 35만 wpm이면 마이크론(37.5만)과 동급.',
                   'STAR IPO $8.6B가 상하이·베이징 2팹 실탄'],
    tiles=[
        ('DDR5 수율 (보도, 공식 미확인)', [('>90%', BLUE)],
         '삼성 동세대 92~93%와 오차 범위 수준', 19),
        ('서방 OEM 인증', [('HP·Asus·Acer', BLUE)],
         '2026 중반 인증 완료, 노트북 채택 개시', 16),
        ('NAND 선행 실현', [('YMTC 14%', BLUE), ('  3위', G_MID)],
         'Q2 2026 출하, 마이크론·키옥시아 추월', 19),
    ],
    red_rows=[
        ('캐파 30만 wpm', ': 6년 10배, IPO $8.6B'),
        ('DDR5 수율 90% 보도', ': 세대 격차 소멸 수준'),
        ('서방 인증 개방', ': HP·Asus·Acer 채택 개시'),
        ('NAND 침식 완료형', ': YMTC 14% 글로벌 3위'),
    ],
    green_rows=[
        ('애플 인증 미결', ': DX-5, 9월 회담이 변수'),
        ('HBM3 양산 미공식', ': 샘플링, 캐파 20% 배정설'),
        ('EUV 우회 미확정', ': 본디드 DRAM 파일럿 단계'),
        ('가격 앵커 미하락', ': DDR4 현물 +172% 역설'),
    ],
    buffer=['AI 공급 부족이 침식을 가리는 중. sub-10nm·HBM4 경로는 봉쇄.',
            '낮아진 앵커는 다운턴 도착과 동시에 드러난다'],
    watch=['애플-CXMT 건의 결말(9월',
           '미중 회담 전후)과 HBM3 양산',
           '공식화, 공급 부족 해소 뒤의',
           '레거시 가격 착지점이 판정 순간'],
    verdict_plain='물량·수율·인증·NAND 전조 4건 실측, 발화(가격 하락)만 슈퍼사이클이 가리고 있다.  ',
    verdict_bold='유일하게 이미 진행 중인 갈래, 다운턴과 동시에 표면화된다.',
    footer=('출처: wiki/downturn/samsung-impact.md §5.1f · 캐파·IPO: SemiAnalysis 종합 · 수율·OEM: 2026-08 업계 보도(공식 미확인 표기) · '
            '점유율: Omdia·FT · 기준일 2026-08'),
    notes=NOTES5,
)

# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 6 — "제품 정의가 바뀌는 경우" 발현 근접도 (wiki §5.1g)
# ═════════════════════════════════════════════════════════════════════════════
def draw_ladder(sl, px, pw):
    """제품 정의권 이동 사다리: 5단 계단. 실측 = 진한 채움, 로드맵·목업 = 점선."""
    steps = [
        ('HBM3E', '표준', True),
        ('HBM4', '로직화', True),
        ('HBM4E', '커스텀E', False),
        ('NVHBM', '고객 정의', True),
        ('zHBM', '3D 적층', False),
    ]
    A_BASE = 5.30
    sw = (pw - 0.3) / 5
    sh0, dh = 0.52, 0.52
    for i, (name, sub, solid) in enumerate(steps):
        sx = px + 0.15 + i * sw
        h = sh0 + i * dh
        if solid:
            box(sl, sx + 0.03, A_BASE - h, sw - 0.06, h, fill=S_FUT)
            ncol, scol = WHITE, RGBColor(0xE3, 0xE3, 0xE6)
        else:
            box(sl, sx + 0.03, A_BASE - h, sw - 0.06, h, fill=EST_FILL, line=G_MID, line_w=0.8, dash=2)
            ncol, scol = INK, GRAY
        txt(sl, sx + 0.03, A_BASE - h + 0.07, sw - 0.06, 0.16,
            [[(name, {'bold': True, 'color': ncol, 'size': 10})]], align=PP_ALIGN.CENTER, leading=1.0)
        txt(sl, sx, A_BASE + 0.05, sw, 0.14,
            [[(sub, {'color': G_MID, 'size': 9})]], align=PP_ALIGN.CENTER, leading=1.0)
    hline(sl, px + 0.1, A_BASE, pw - 0.2, G_LINE, 0.9)
    txt(sl, px + 0.15, A_BASE - 3.05, pw - 0.3, 0.16,
        [[('정의권: 메모리사 표준 → 고객 설계', {'color': GRAY, 'size': 10})]], leading=1.0)

NOTES6 = """[슬라이드 위치] 5번 시나리오(제품 정의가 바뀌는 경우)의 근접도 심화 장. 이 시나리오는 다른 넷과 성격이 다르다 — 시장 전체는 성장하는데 특정 자산군(2D DRAM 캐파·표준 HBM 라인)만 가치가 급락하는 경우로, 감산·계약·가격 규율이라는 다운턴 매뉴얼이 통째로 작동하지 않는다. 그래서 근접도 평가의 축도 가격·재고가 아니라 "제품 정의권이 누구에게 있는가"다.

[왼쪽: 구조 패널] 제품 정의권 이동의 5단 사다리. ① HBM3/3E 표준 커머디티(실측 — 현재 매출의 사실상 전부) ② HBM4 base die 로직 파운드리화(실측 완료 — 3사 전원, 삼성 SF4·SK/Micron TSMC N12, ISSCC 2026) ③ HBM4E 커스텀 본격화(설계 2026 중반 완료, 2027 양산, 2027 HBM 수요 ~40% 전망 — 점선) ④ NVHBM(2026-08-26 NVIDIA 공식 발표 — 메모리 컨트롤러를 스택 안으로, BW +30%·전력 -15% vs 표준 HBM4E, 첫 파트너가 메모리사가 아니라 Amazon Trainium4라는 점이 상징적. GPU 적용은 Feynman 2028) ⑤ zHBM(Hot Chips 2026 스펙 — 전력 -70%·BW 2.3배, 목업 단계·양산 미정 — 점선). 진한 칸은 실측, 점선 칸은 로드맵·목업이라는 구분을 반드시 언급할 것.

[가운데: 판 갈이 패널] 실측된 것은 '구조의 이동', 미발동인 것은 '돈의 이동'이다. 커스텀 HBM 매출 비중은 여전히 ~0%이고 발화 트리거(30%, 경계 15%)까지 전 구간이 비어 있다. 오히려 현세대 커스텀 HBM4 물량은 "거의 다 캔슬"(영업 수장, 2026-08-03) — 커머디티가 해마다 세대교체되는 속도전 구간에서는 2~3년 걸리는 커스텀이 나올 때쯤 이미 뒤처지기 때문이다. 이 내부 논쟁(영업 "커스텀 퇴조" vs 상품기획 "전 고객 zHBM")은 층위 차이다: 현세대 커스텀은 퇴조가 맞고, 차세대(HBM4E부터, zHBM·cHBM)는 커스텀이 로드맵의 중심이다. 내부 1차 자료: 7월 미주 하이퍼스케일러 4사 미팅에서 아젠다 5~6개를 들고 갔는데 "다 필요 없고 zHBM", 예전의 "표준 아니면 안 쓴다"가 "일단 써볼게"로 바뀌었다는 것이 상품기획팀장의 증언이다.

[오른쪽: 판정 보드] 전조 4건: base die 로직화 완료, 고객 주도 규격 등장(NVHBM·Marvell 공동 정의·Broadcom 커스텀 설계 70%+), 고객 요구 전면화(4사 미팅), 조직 선행 재편(삼성 HBM4부터 표준·커스텀 이원화 + 커스텀 엔지니어 250명 증원). 미발동 4건: 커스텀 매출 30% 돌파(현재 ~0%), 총시장 vs 자사 특정 제품군 성장률 괴리(표준 포함 전 제품 쇼티지·계약가 상승 지속), 차세대 양산 채택 공시(zHBM 목업·NVHBM 2028), CXL 물량 구조 재편(풀링은 구글 연구진 회의론과 논쟁 중, 어태치율 실측 없음 — 확장 용도 프로덕션 진입은 Azure·Meta로 실측). 완충: 세대교체 속도가 커스텀 개발을 앞지르는 구간 + 삼성은 메모리+파운드리(SF4 내재화) 동시 보유 3사 유일 — 전환이 와도 수혜측에 설 여지.

[판정을 말로 하면] "잠복 국면"이다. 단, 이 시나리오의 특성상 발화 신호가 켜진 뒤에는 대응 수단이 없다 — 감산도 계약도 진부화를 막지 못하므로, 유일한 응수는 미리 그 자리에 있는 것(DP-5 차세대 별동대, DP-2 옵션형 캐파)이다. 감시선의 목적도 "언제 대응할까"가 아니라 "별동대 배분을 언제 늘릴까"를 정하는 데 있다: 분기 커스텀 매출 비중(경계 15%/트리거 30%), NVHBM 추가 파트너·zHBM 양산 공시 카운트, 총시장 vs 자사 표준 HBM 성장률 괴리(DR-1 감별 1순위). 3D DRAM(4F² VCT)은 ISSCC 2026에서 코어 면적 17.0%→2.7% 실증, 양산은 2029~30으로 이 트랙 시계의 끝자락이다.

[출처] wiki/downturn/samsung-impact.md §5.1g(단일 소스). 수치 원출처: dt-e-product-redefinition-signals-2026-08-28(NVHBM·zHBM·HBM4E·생태계·CXL), semianalysis-isscc-2026(base die 로직화·3D DRAM), logic-paradigm-shifts-3d-dram(조직 재편·커스텀 전망), choi-jangseok·lee-changsoo 인터뷰(내부 1차), lta-to-sca-transition(커스텀 비중 EWI)."""

SPEC_E = dict(
    title_blue='바뀌는 경우의 근접도:  ',
    title_rest='규격 정의권은 고객으로 넘어가는데, 매출은 아직 표준에 있다',
    subtitle='5번 시나리오(제품 정의가 바뀌는 경우)의 발현 가능성 평가: 정의권 이동의 실측, 그리고 발화 전 체크리스트',
    panelA='구조  제품 정의권의 이동 사다리',
    panelB='판 갈이  구조는 움직였고, 돈은 아직이다',
    panelC='판정 보드  발화 전 체크리스트',
    chart=draw_ladder,
    chart_caption=['진한 칸 = 실측(①②④ 발표), 점선 = 로드맵·목업.',
                   'NVHBM 첫 파트너는 메모리사가 아니라 Amazon'],
    tiles=[
        ('커스텀 HBM 매출 비중', [('~0%', BLUE), ('  vs 트리거 ', G_MID), ('30%', G_MID)],
         '로드맵(27~28 과반 전망)과 매출의 괴리가 스트레스', 19),
        ('NVHBM (고객 정의 base die)', [('BW +30%', BLUE), (' · ', G_MID), ('전력 -15%', BLUE)],
         '08-26 NVIDIA 발표, 첫 파트너 Amazon Trainium4', 16),
        ('zHBM (3D 수직 적층)', [('전력 -70%', BLUE), (' · ', G_MID), ('BW 2.3배', BLUE)],
         'Hot Chips 2026 스펙, 목업 단계·양산 미정', 16),
    ],
    red_rows=[
        ('base die 로직화 완료', ': HBM4 3사 전원 SF4·N12'),
        ('고객 주도 규격 등장', ': NVHBM·Marvell 공동 정의'),
        ('고객 요구 전면화', ': 4사 미팅 "다 필요 없고 zHBM"'),
        ('조직 선행 재편', ': 커스텀 이원화·250명 증원'),
    ],
    green_rows=[
        ('커스텀 매출 30%', ': 현재 ~0%, 경계 15%'),
        ('성장률 괴리 미관측', ': 표준 포함 전 제품 쇼티지'),
        ('양산 채택 공시 없음', ': zHBM 목업·NVHBM 2028'),
        ('CXL 재편 미발동', ': 풀링 회의론, 실측 없음'),
    ],
    buffer=['연간 세대교체가 2~3년 커스텀을 앞질러 현세대 커스텀은 퇴조.',
            '판정 축은 가격·재고가 아니라 정의권의 소재다'],
    watch=['분기 커스텀 매출 비중',
           '(경계 15% / 트리거 30%)과',
           'NVHBM 파트너·zHBM 양산 공시가',
           '별동대 증배분의 타이머다'],
    verdict_plain='정의권 이동은 공식화됐지만 돈은 아직 표준에 있다. 전조 4건 실측, 발화 신호 0건.  ',
    verdict_bold='잠복 국면이며, 미리 가 있는 것만 유효한 갈래다.',
    footer=('출처: wiki/downturn/samsung-impact.md §5.1g · NVHBM: NVIDIA 발표(2026-08-26) · zHBM: Hot Chips 2026 · '
            '커스텀 비중·HBM4E: TrendForce·사내 인터뷰 · 기준일 2026-08'),
    notes=NOTES6,
)

# 슬라이드 2~6 생성 (시나리오 순서 고정: A → B → C → D → E)
for _spec in (SPEC_A, SPEC_B, SPEC_C, SPEC_D, SPEC_E):
    proximity_slide(_spec)

OUT = 'outputs/presentation/downturn-scenario-impact.pptx'
prs.save(OUT)
print('saved:', OUT)
