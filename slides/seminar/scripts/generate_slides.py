"""
Claude Code 세미나 — 스킬 & 커넥터 슬라이드 (각 1매) 생성기.

산출물:
  - slides/seminar/claude-code-skills.pptx       (1슬라이드)
  - slides/seminar/claude-code-connectors.pptx   (1슬라이드)

설계 원칙:
  - 두 슬라이드는 동일한 4-카드 그리드(2×2) 레이아웃.
  - 액센트 색만 다르게 — 스킬은 copper, 커넥터는 slate.
  - 모노그램(원형 + 단일 글자)으로 아이콘 표현 — 이모지 폰트 의존성 회피.
  - 하단 1줄에 출처 URL 명시.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
OUT_DIR = os.path.join(ROOT, 'slides', 'seminar')

# ----- Palette (Anthropic-feel) -----
INK         = RGBColor(0x26, 0x20, 0x1C)   # 본문 텍스트
INK_SOFT    = RGBColor(0x4A, 0x42, 0x3A)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
GRAY_LINE   = RGBColor(0xD9, 0xD3, 0xC8)
CREAM       = RGBColor(0xF4, 0xF0, 0xE8)   # 카드 배경
PAGE_BG     = RGBColor(0xFB, 0xF8, 0xF2)   # 페이지 배경
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

COPPER      = RGBColor(0xCC, 0x78, 0x5C)   # 스킬 액센트
COPPER_SOFT = RGBColor(0xF4, 0xDC, 0xD0)
SLATE       = RGBColor(0x4A, 0x55, 0x68)   # 커넥터 액센트
SLATE_SOFT  = RGBColor(0xDD, 0xE2, 0xEA)

AMBER       = RGBColor(0xD9, 0x77, 0x06)   # 보안 경고 배너
AMBER_SOFT  = RGBColor(0xFE, 0xF3, 0xC7)

FONT_KO = '맑은 고딕'
FONT_EN = 'Calibri'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# Drawing primitives
# ============================================================

def _set_run(run, text, *, font=FONT_KO, size=12, bold=False, color=None, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=FONT_KO, size=12, bold=False,
             color=None, align='left', valign='top', wrap=True, line_spacing=None,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP,
        'middle': MSO_ANCHOR.MIDDLE,
        'bottom': MSO_ANCHOR.BOTTOM,
    }[valign]
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, font=font, size=size, bold=bold, italic=italic,
                 color=color or INK)
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_width=None,
             rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    shape.shadow.inherit = False
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
    if rounded:
        # Tighten corner radius
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=None, width=None):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color or GRAY_LINE
    if width:
        line.line.width = width
    return line


def add_circle(slide, x, y, d, *, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_monogram(slide, cx, cy, *, letter, accent_color, size=Inches(0.6)):
    """원형 + 가운데 단일 글자 (모노그램 아이콘)."""
    add_circle(slide, cx, cy, size, fill=accent_color)
    add_text(slide,
             cx, cy, size, size,
             letter,
             font=FONT_EN, size=20, bold=True, color=WHITE,
             align='center', valign='middle')


def add_page_background(slide, color=PAGE_BG):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)


# ============================================================
# Composable layout pieces
# ============================================================

def add_header(slide, title, subtitle, accent_color):
    # Accent left bar
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.08), Inches(0.95),
             fill=accent_color)
    # Title
    add_text(slide, Inches(0.78), Inches(0.42), Inches(12), Inches(0.55),
             title, font=FONT_KO, size=24, bold=True, color=INK)
    # Subtitle
    add_text(slide, Inches(0.78), Inches(1.02), Inches(12), Inches(0.4),
             subtitle, font=FONT_KO, size=12, color=GRAY)
    # Bottom rule
    add_line(slide, Inches(0.5), Inches(1.55), Inches(12.83), Inches(1.55),
             color=GRAY_LINE, width=Emu(6350))


def add_footer(slide, *, sources, page_label):
    add_line(slide, Inches(0.5), Inches(7.0), Inches(12.83), Inches(7.0),
             color=GRAY_LINE, width=Emu(6350))
    add_text(slide, Inches(0.5), Inches(7.08), Inches(11), Inches(0.32),
             '출처  ·  ' + sources,
             font=FONT_EN, size=9, color=GRAY)
    add_text(slide, Inches(11.83), Inches(7.08), Inches(1.0), Inches(0.32),
             page_label, font=FONT_EN, size=9, color=GRAY, align='right')


def add_card(slide, x, y, w, h, *, monogram, accent_color, headline, body):
    """4-카드 한 셀."""
    # Card background
    add_rect(slide, x, y, w, h, fill=CREAM, rounded=True)
    # Top accent strip
    strip_h = Inches(0.06)
    add_rect(slide, x + Inches(0.15), y + Inches(0.18),
             Inches(0.45), strip_h, fill=accent_color)

    # Monogram (top-left)
    mono_d = Inches(0.55)
    add_monogram(slide, x + Inches(0.32), y + Inches(0.4),
                 letter=monogram, accent_color=accent_color, size=mono_d)

    # Headline
    add_text(slide,
             x + Inches(1.05), y + Inches(0.43),
             w - Inches(1.2), Inches(0.45),
             headline,
             font=FONT_KO, size=14, bold=True, color=INK,
             valign='middle')

    # Body
    add_text(slide,
             x + Inches(0.35), y + Inches(1.05),
             w - Inches(0.7), h - Inches(1.2),
             body,
             font=FONT_KO, size=10.5, color=INK_SOFT,
             line_spacing=1.25, valign='top')


def add_callout_banner(slide, x, y, w, h, *, label, body,
                       accent=AMBER, soft=AMBER_SOFT):
    """하단 강조 배너 (보안 경고 등)."""
    add_rect(slide, x, y, w, h, fill=soft)
    add_rect(slide, x, y, Inches(0.1), h, fill=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.05),
             Inches(2.0), h - Inches(0.1),
             label, font=FONT_EN, size=10, bold=True, color=accent,
             valign='middle')
    add_text(slide, x + Inches(2.4), y + Inches(0.05),
             w - Inches(2.6), h - Inches(0.1),
             body, font=FONT_KO, size=10.5, bold=True, color=INK,
             valign='middle')


# ============================================================
# Slide builders
# ============================================================

def build_skills_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_page_background(slide)
    add_header(
        slide,
        title='Claude Code 기본 제공 스킬 (Built-in Skills)',
        subtitle='SKILL.md로 패키징된 작업 지침 — 동적 로드로 Claude의 능력을 확장',
        accent_color=COPPER,
    )

    # 4-card grid (2×2)
    grid_top   = Inches(1.85)
    card_w     = Inches(6.20)
    card_h     = Inches(2.45)
    gutter     = Inches(0.20)
    left_x     = Inches(0.50)
    right_x    = left_x + card_w + gutter
    row2_y     = grid_top + card_h + gutter

    add_card(
        slide, left_x, grid_top, card_w, card_h,
        monogram='S', accent_color=COPPER,
        headline='SKILL.md 기반 동적 로드 지침',
        body=(
            '디렉토리 + SKILL.md (YAML frontmatter + Markdown) 한 묶음.\n'
            'description 매칭으로 Claude가 자동 호출 — 또는 /스킬-이름으로 명시 호출.\n'
            '본문은 호출 시점에만 컨텍스트에 로드 → 평소 토큰 비용 거의 0.'
        ),
    )

    add_card(
        slide, right_x, grid_top, card_w, card_h,
        monogram='A', accent_color=COPPER,
        headline='Anthropic 공식 번들 — docx · pdf · pptx · xlsx',
        body=(
            'github.com/anthropics/skills 공식 레포 공개.\n'
            'Claude.ai · Claude Code · Anthropic API 어디서든 동일 동작.\n'
            'Claude의 문서 능력을 뒤에서 구동하는 동일 스킬셋.\n'
            '설치: /plugin marketplace add anthropics/skills'
        ),
    )

    add_card(
        slide, left_x, row2_y, card_w, card_h,
        monogram='/', accent_color=COPPER,
        headline='Claude Code 내장 슬래시 스킬 6종',
        body=(
            '/simplify  — 최근 변경 코드를 3개 리뷰 에이전트로 점검·수정\n'
            '/batch     — 대규모 변경을 5~30단위 분해 후 worktree 병렬 실행\n'
            '/debug     — 디버그 로그 분석 / /loop — 프롬프트 반복 실행\n'
            '/claude-api — Anthropic SDK 코드 마이그레이션\n'
            '/fewer-permission-prompts — 권한 프롬프트 줄이기'
        ),
    )

    add_card(
        slide, right_x, row2_y, card_w, card_h,
        monogram='?', accent_color=COPPER,
        headline='발견 + 호출 방식',
        body=(
            '/skills 로 사용 가능한 전체 목록 확인.\n'
            '자동 발동 = description 매칭 / 명시 = /이름 (인자 가능: /fix-issue 123).\n'
            'disable-model-invocation: true 로 자동 발동 차단.\n'
            '저장 위치: enterprise / ~/.claude/skills/ (개인) / .claude/skills/ (프로젝트) / 플러그인'
        ),
    )

    add_footer(
        slide,
        sources=('code.claude.com/docs/en/skills   ·   '
                 'code.claude.com/docs/en/commands   ·   '
                 'github.com/anthropics/skills'),
        page_label='Claude Code 세미나 · 1/2',
    )


def build_connectors_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(slide)
    add_header(
        slide,
        title='Claude Code 커넥터 (Connectors via MCP)',
        subtitle='Model Context Protocol — 외부 도구·DB·API를 Claude의 도구로 노출',
        accent_color=SLATE,
    )

    grid_top   = Inches(1.85)
    card_w     = Inches(6.20)
    card_h     = Inches(2.10)
    gutter     = Inches(0.20)
    left_x     = Inches(0.50)
    right_x    = left_x + card_w + gutter
    row2_y     = grid_top + card_h + gutter

    add_card(
        slide, left_x, grid_top, card_w, card_h,
        monogram='M', accent_color=SLATE,
        headline='MCP 서버 ↔ Claude Code 어댑터',
        body=(
            '외부 시스템(GitHub · Slack · DB 등)이 MCP 서버를 노출하면\n'
            'Claude Code가 그 도구·리소스·프롬프트를 세션에서 직접 호출.\n'
            '"다른 도구에서 채팅창으로 데이터를 복사·붙여넣고 있다면 커넥트할 시점."'
        ),
    )

    add_card(
        slide, right_x, grid_top, card_w, card_h,
        monogram='D', accent_color=SLATE,
        headline='디렉토리 커넥터 — claude.com/connectors',
        body=(
            '카테고리: Productivity · Communication · Data · Sales & Marketing ·\n'
            'Code · Design · Financial · Health & Wellness 등.\n'
            '대표: GitHub · Slack · Notion · Asana · Atlassian · Airtable ·\n'
            'Stripe · Sentry · Linear · Figma · Google Drive.'
        ),
    )

    add_card(
        slide, left_x, row2_y, card_w, card_h,
        monogram='+', accent_color=SLATE,
        headline='커스텀 / 자체 MCP 서버',
        body=(
            '디렉토리 외 통합은 modelcontextprotocol/servers (오픈소스 모음)\n'
            '또는 MCP SDK로 직접 구축. 로컬 도구·내부 시스템 연결에 활용.'
        ),
    )

    add_card(
        slide, right_x, row2_y, card_w, card_h,
        monogram='>', accent_color=SLATE,
        headline='추가 명령 + 관리',
        body=(
            'HTTP(권장): claude mcp add --transport http notion https://mcp.notion.com/mcp\n'
            'SSE(deprecated) / stdio(로컬, --env KEY=VAL  --  명령).\n'
            '관리: claude mcp list / get / remove · 세션 내 /mcp (OAuth · 상태).\n'
            '스코프: local / project (.mcp.json 공유) / user.'
        ),
    )

    # Security banner
    add_callout_banner(
        slide,
        Inches(0.5), Inches(6.45),
        Inches(12.33), Inches(0.45),
        label='SECURITY',
        body=('신뢰할 수 있는 MCP 서버만 연결할 것 — 외부 콘텐츠를 가져오는 서버는 '
              'prompt injection 위험에 특히 유의.'),
        accent=AMBER, soft=AMBER_SOFT,
    )

    add_footer(
        slide,
        sources=('code.claude.com/docs/en/mcp   ·   '
                 'claude.com/connectors   ·   '
                 'modelcontextprotocol.io'),
        page_label='Claude Code 세미나 · 2/2',
    )


# ============================================================
# Build
# ============================================================

def _new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def build_skills_pptx(path):
    prs = _new_prs()
    build_skills_slide(prs)
    prs.save(path)
    return path


def build_connectors_pptx(path):
    prs = _new_prs()
    build_connectors_slide(prs)
    prs.save(path)
    return path


def build_combined_pptx(path):
    """두 슬라이드를 한 덱으로 합친 버전 (옵션)."""
    prs = _new_prs()
    build_skills_slide(prs)
    build_connectors_slide(prs)
    prs.save(path)
    return path


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    p1 = build_skills_pptx(os.path.join(OUT_DIR, 'claude-code-skills.pptx'))
    p2 = build_connectors_pptx(os.path.join(OUT_DIR, 'claude-code-connectors.pptx'))
    p3 = build_combined_pptx(os.path.join(OUT_DIR, 'claude-code-skills-and-connectors.pptx'))
    for p in (p1, p2, p3):
        size_kb = os.path.getsize(p) / 1024
        print(f'  ✓ {os.path.relpath(p, ROOT)}  ({size_kb:.1f} KB)')


if __name__ == '__main__':
    main()
